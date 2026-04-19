"""Generate TOMOSHIKA ロジック整理 deck as PPTX (36 slides).

Matches the content of TOMOSHIKA ロジック整理.html one-to-one.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ========== COLORS ==========
C = {
    'paper':   RGBColor(0xFB, 0xF8, 0xF1),
    'ink':     RGBColor(0x00, 0x00, 0x00),
    'ink2':    RGBColor(0x14, 0x14, 0x14),
    'muted':   RGBColor(0x3A, 0x3A, 0x3A),
    'line':    RGBColor(0xC8, 0xC3, 0xB6),
    'line2':   RGBColor(0xDC, 0xD6, 0xC6),
    'dark':    RGBColor(0x1A, 0x1A, 0x1A),
    'white':   RGBColor(0xFF, 0xFF, 0xFF),
    'paper_d': RGBColor(0xF5, 0xEF, 0xE2),
    'paper_d2':RGBColor(0xE3, 0xDC, 0xC9),
    'accent':  RGBColor(0xB9, 0x5A, 0x1F),  # orange/terracotta
    'accent2': RGBColor(0x1C, 0x4D, 0x9B),  # blue
    'ok':      RGBColor(0x2F, 0x74, 0x3A),
    'warn':    RGBColor(0xB4, 0x6A, 0x0C),
    'chip':    RGBColor(0xEF, 0xEA, 0xDD),
    'th':      RGBColor(0xED, 0xE6, 0xD2),
    'callout_bg': RGBColor(0xE8, 0xEE, 0xF8),
    'confirm_bg': RGBColor(0xFF, 0xF3, 0xD6),
    'gold':    RGBColor(0xD6, 0xA0, 0x3A),
    'gold_l':  RGBColor(0xE8, 0xC4, 0x6E),
    'cream':   RGBColor(0xE3, 0xDC, 0xC9),
}

FONT = 'Noto Sans JP'
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s


def rrect(slide, x, y, w, h, color, line=None, adj=0.04):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    if hasattr(s, 'adjustments') and len(s.adjustments) > 0:
        s.adjustments[0] = adj
    return s


def tbox(slide, x, y, w, h, text, size=14, color=None, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT,
         letter_spacing=None):
    """Plain text box (no fill, no border)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color
    return tb


def tbox_rich(slide, x, y, w, h, runs, size=14, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=None):
    """runs: list of (text, bold, color) tuples; \n in text creates new paragraph."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    first = True
    for run_spec in runs:
        if len(run_spec) == 4:
            txt, bold, col, sz = run_spec
        elif len(run_spec) == 3:
            txt, bold, col = run_spec; sz = size
        else:
            txt, bold = run_spec; col = None; sz = size
        parts = txt.split('\n')
        for idx, part in enumerate(parts):
            if idx > 0 or not first:
                if idx > 0:
                    p = tf.add_paragraph()
                    p.alignment = align
                    if line_spacing:
                        p.line_spacing = line_spacing
                else:
                    p = tf.paragraphs[-1]
            else:
                p = tf.paragraphs[0]
                p.alignment = align
                if line_spacing:
                    p.line_spacing = line_spacing
            if part == '':
                continue
            r = p.add_run()
            r.text = part
            r.font.name = font
            r.font.size = Pt(sz)
            r.font.bold = bold
            if col is not None:
                r.font.color.rgb = col
            first = False
    return tb


def kicker(slide, text, color=None):
    if color is None:
        color = C['ink2']
    # dot
    dot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.55), Inches(0.12), Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb = C['accent']; dot.line.fill.background()
    # text
    tb = slide.shapes.add_textbox(Inches(0.85), Inches(0.45), Inches(10), Inches(0.35))
    tf = tb.text_frame; tf.margin_left = Inches(0); tf.margin_top = Inches(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = color


def page_footer(slide, num):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(4), Inches(0.3))
    tf = tb.text_frame; tf.margin_left = Inches(0); tf.margin_top = Inches(0)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = 'TOMOSHIKA'
    r.font.name = FONT; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = C['muted']
    tb2 = slide.shapes.add_textbox(Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3))
    tf2 = tb2.text_frame; tf2.margin_left = Inches(0); tf2.margin_top = Inches(0)
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = f'{num:02d}'
    r2.font.name = FONT; r2.font.size = Pt(10); r2.font.bold = True; r2.font.color.rgb = C['muted']


def bg(slide, color=None):
    if color is None:
        color = C['paper']
    rect(slide, 0, 0, 13.333, 7.5, color)


def title_h2(slide, text, y=0.95, size=28):
    tbox(slide, 0.65, y, 12, 0.8, text, size=size, color=C['ink'], bold=True)


def mktable(slide, x, y, w, h, data, col_widths=None, font_size=12,
            header=True, header_bg=None, header_color=None, bold_first_col=False):
    """Create table; data = list of rows (list of strings or list of (text, bold) tuples)."""
    if header_bg is None:
        header_bg = C['th']
    if header_color is None:
        header_color = C['ink']
    rows = len(data); cols = len(data[0])
    tab_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = tab_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for r_idx, row in enumerate(data):
        for c_idx, cell_val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.TOP
            # header/body colors
            if header and r_idx == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_bg
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = C['white']
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
            # Parse runs
            if isinstance(cell_val, str):
                runs = [(cell_val, (header and r_idx == 0) or (bold_first_col and c_idx == 0))]
            else:
                runs = cell_val  # list of (text, bold) or (text, bold, color)
            p = tf.paragraphs[0]
            for ri, rn in enumerate(runs):
                if isinstance(rn, tuple):
                    if len(rn) == 3:
                        txt, bold, col = rn
                    else:
                        txt, bold = rn; col = None
                else:
                    txt = rn; bold = False; col = None
                if '\n' in txt:
                    parts = txt.split('\n')
                    for pi, pt in enumerate(parts):
                        if pi > 0:
                            p = tf.add_paragraph()
                        if pt == '': continue
                        r = p.add_run(); r.text = pt
                        r.font.name = FONT
                        r.font.size = Pt(font_size)
                        r.font.bold = bold
                        r.font.color.rgb = col or (header_color if (header and r_idx == 0) else C['ink'])
                else:
                    r = p.add_run(); r.text = txt
                    r.font.name = FONT
                    r.font.size = Pt(font_size)
                    r.font.bold = bold
                    r.font.color.rgb = col or (header_color if (header and r_idx == 0) else C['ink'])
    return tab_shape


# ========== SLIDE BUILDERS ==========

def slide_01_cover():
    s = blank()
    bg(s, C['dark'])
    # Accent dot + kicker
    dot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(4.2), Inches(0.14), Inches(0.14))
    dot.fill.solid(); dot.fill.fore_color.rgb = C['gold_l']; dot.line.fill.background()
    tbox(s, 1.15, 4.12, 10, 0.35, '三者会議資料 / 2026-04-18',
         size=14, color=C['paper_d2'], bold=True)
    # Title
    tbox(s, 0.9, 4.55, 12, 1.8,
         'TOMOSHIKA\nロジック整理',
         size=54, color=C['white'], bold=True)
    # Subtitle
    tbox(s, 0.9, 6.3, 12, 0.7,
         '学生が適性検査を受けてから、結果画面に推薦企業が表示されるまでの流れを\n4つのフェーズに分けて整理したもの',
         size=14, color=C['paper_d2'])
    # Meta
    tbox(s, 0.9, 6.98, 12, 0.35,
         '中西製作所 様      MIL 様      エヌイチ',
         size=12, color=C['paper_d2'], bold=True)
    # Brand corner
    tbox(s, 0.9, 0.4, 6, 0.3, 'TOMOSHIKA / 三者会議',
         size=10, color=RGBColor(0xA3, 0x9D, 0x8C), bold=True)


def slide_02_overview():
    s = blank(); bg(s); kicker(s, 'OVERVIEW'); title_h2(s, '全体の流れ')
    # Flow nodes
    labels = [
        ('PHASE 1', '適性検査を受ける', '【MIL様】', '学生が51問に回答する'),
        ('PHASE 2 →', 'データを送る', '【MIL様 → Dify】', 'MIL様が検査結果をDifyに送る'),
        ('PHASE 3', '企業のマッチング', '【Dify】', '受け取った点数で企業リストと照合する'),
        ('PHASE 4 →', '結果を返す', '【Dify → MIL様】', 'マッチング結果をMIL様に返す'),
    ]
    node_w = 2.9; gap = 0.17
    x0 = 0.65
    for i, (n, title, owner, desc) in enumerate(labels):
        x = x0 + i * (node_w + gap)
        rrect(s, x, 1.85, node_w, 1.75, C['white'], line=C['line'])
        tbox(s, x + 0.15, 1.95, node_w - 0.3, 0.3, n, size=10, color=C['ink2'], bold=True)
        tbox(s, x + 0.15, 2.25, node_w - 0.3, 0.45, title, size=15, color=C['ink'], bold=True)
        tbox(s, x + 0.15, 2.72, node_w - 0.3, 0.3, owner, size=11, color=C['accent'], bold=True)
        tbox(s, x + 0.15, 3.03, node_w - 0.3, 0.55, desc, size=11, color=C['ink2'])
    # 三社の役割
    tbox(s, 0.65, 3.85, 8, 0.45, '三社の役割', size=18, color=C['ink'], bold=True)
    data = [
        ['会社', '何をするか'],
        [[('MIL様', True)], '学生に適性検査を受けてもらう画面を作る。検査結果をDifyに送り、返ってきた結果を学生に見せる'],
        [[('エヌイチ', True)], '中西製作所様のDifyアカウント上に「検査結果を受け取って企業をマッチングするしくみ」を構築・運用する'],
        [[('中西製作所様', True)], 'Difyアカウントの持ち主。どの企業を推薦するか、どういう条件で中西を出すかの最終決定権を持つ'],
    ]
    mktable(s, 0.65, 4.4, 12, 2.4, data, col_widths=[2.2, 9.8], font_size=12)
    page_footer(s, 2)


def slide_03_dify():
    s = blank(); bg(s); kicker(s, 'OVERVIEW'); title_h2(s, 'Difyについて')
    # bullets
    items = [
        [('Dify', True, C['ink']), ('は、プログラムの処理を組み立てて動かすためのツール（クラウドサービス）', False, C['ink'])],
        [('このDifyのアカウントは', False, C['ink']), ('中西製作所様のもの', True, C['accent'])],
        [('エヌイチがこのアカウント上でしくみ（ワークフロー）を作っている', False, C['ink'])],
        [('フェーズ2・4で必要になる接続先の情報や認証キーも、中西製作所様のアカウントから発行される', False, C['ink'])],
    ]
    y = 1.85
    for runs in items:
        # bullet dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), Inches(y + 0.14), Inches(0.08), Inches(0.08))
        dot.fill.solid(); dot.fill.fore_color.rgb = C['ink']; dot.line.fill.background()
        tbox_rich(s, 1.0, y, 11.5, 0.5, runs, size=16)
        y += 0.55
    # ASCII structure
    tbox(s, 0.65, 4.3, 10, 0.45, 'アカウント構造のイメージ', size=18, color=C['ink'], bold=True)
    rrect(s, 0.65, 4.85, 12, 2.1, C['white'], line=C['line'])
    ascii_box = (
        '中西製作所様 Difyアカウント（オーナー）\n'
        '  ┌──────────────────────────────────────────┐\n'
        '  │  ワークフロー：適性検査 → 企業マッチング     │\n'
        '  │  （エヌイチが構築・運用）                     │\n'
        '  └──────────────────────────────────────────┘\n'
        '                  ↑ API接続\n'
        '         MIL様のサーバー（学生に画面を出す）'
    )
    tbox(s, 0.85, 4.95, 11.8, 2.0, ascii_box,
         size=11, color=C['ink'], font='Courier New')
    page_footer(s, 3)


def divider(num_str, title, subtitle, section_label, page):
    s = blank(); bg(s, C['dark'])
    dot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(0.55), Inches(0.14), Inches(0.14))
    dot.fill.solid(); dot.fill.fore_color.rgb = C['gold_l']; dot.line.fill.background()
    tbox(s, 1.15, 0.47, 10, 0.35, section_label,
         size=13, color=C['cream'], bold=True)
    # Number
    tbox(s, 0.9, 1.5, 6, 2.5, num_str, size=140, color=C['gold'], bold=True)
    # Title
    tbox(s, 0.9, 4.0, 12, 1.2, title, size=52, color=C['white'], bold=True)
    # Subtitle
    tbox_rich(s, 0.9, 5.3, 12, 1.8,
              [(subtitle, False, C['cream'], 15)],
              size=15, line_spacing=1.4)


def slide_04_div_p1():
    divider('01', '適性検査を受ける',
            '学生が51問の質問に答えると、その回答から11個の数値（点数）が算出される。\n'
            'この点数がフェーズ3の企業マッチングの材料になる。',
            'PHASE 01 / MIL 様', 4)


def slide_05_q1_q24():
    s = blank(); bg(s); kicker(s, 'PHASE 01 / 質問の中身')
    title_h2(s, 'メインの質問: Q1〜Q24（志向性を測る24問）', size=24)
    tbox_rich(s, 0.65, 1.6, 12.5, 0.4,
              [('この24問で', False, C['ink']),
               ('「この学生がどんな仕事に向いているか」を6タイプで数値化', True, C['ink']),
               ('する。', False, C['ink'])],
              size=13)
    tbox(s, 0.65, 2.0, 12.5, 0.65,
         '24問は6つのタイプに4問ずつ割り振られていて、各問は「1点（あてはまらない）〜 5点（あてはまる）」で答える。\n'
         '同じタイプの4問の合計がそのタイプの点数になる。',
         size=13, color=C['ink'])
    # ASCII example
    rrect(s, 0.65, 2.75, 12, 0.75, C['white'], line=C['line'])
    tbox(s, 0.85, 2.83, 11.8, 0.7,
         '例: R型（現実型）の4問にそれぞれ 5点, 4点, 3点, 3点 と答えた場合\n'
         '→ R型スコア = 5 + 4 + 3 + 3 = 15点（20点満点中）',
         size=12, color=C['ink'], font='Courier New')
    # Table
    tbox(s, 0.65, 3.6, 12, 0.4, 'これを6タイプ分計算する', size=16, color=C['ink'], bold=True)
    data = [
        ['タイプ', '対応する設問', '算出方法', '点数の範囲'],
        [[('R型', True), ('（現実型）', False)], 'Q1〜Q4', '4問の合計', '4〜20点'],
        [[('I型', True), ('（研究型）', False)], 'Q5〜Q8', '4問の合計', '4〜20点'],
        [[('A型', True), ('（芸術型）', False)], 'Q9〜Q12', '4問の合計', '4〜20点'],
        [[('S型', True), ('（社会型）', False)], 'Q13〜Q16', '4問の合計', '4〜20点'],
        [[('E型', True), ('（企業型）', False)], 'Q17〜Q20', '4問の合計', '4〜20点'],
        [[('C型', True), ('（慣習型）', False)], 'Q21〜Q24', '4問の合計', '4〜20点'],
    ]
    mktable(s, 0.65, 4.05, 12, 2.5, data, col_widths=[2.2, 3.0, 3.5, 3.3], font_size=11)
    tbox_rich(s, 0.65, 6.6, 12.5, 0.4,
              [('点数が高いほど、そのタイプの志向が強い', True, C['ink']),
               ('ということ。', False, C['ink'])], size=13)
    page_footer(s, 5)


def slide_06_6types():
    s = blank(); bg(s); kicker(s, 'PHASE 01 / 質問の中身')
    title_h2(s, '6タイプそれぞれの意味')
    data = [
        ['タイプ', 'どんな人か', '向いている仕事'],
        [[('R型', True), ('（現実型）', False)], '手を動かすのが好き。技術や機械に興味がある', 'エンジニア、製造、建設'],
        [[('I型', True), ('（研究型）', False)], '分析や調べものが好き。データや理論に興味がある', '研究開発、コンサル、分析'],
        [[('A型', True), ('（芸術型）', False)], '創造や表現が好き。自由な発想を大事にする', 'デザイナー、企画、クリエイティブ'],
        [[('S型', True), ('（社会型）', False)], '人の役に立つのが好き。対人支援に興味がある', '教育、福祉、人事'],
        [[('E型', True), ('（企業型）', False)], 'リーダーシップを発揮したい。経営や組織運営に興味がある', '営業、マネジメント、起業'],
        [[('C型', True), ('（慣習型）', False)], '正確さや管理が好き。堅実に物事を進める', '経理、法務、事務、品質管理'],
    ]
    mktable(s, 0.65, 1.9, 12, 4.8, data, col_widths=[2.5, 5.5, 4.0], font_size=13)
    page_footer(s, 6)


def slide_07_q25_q51():
    s = blank(); bg(s); kicker(s, 'PHASE 01 / 質問の中身')
    title_h2(s, '特殊な質問: Q25〜Q51')
    data = [
        ['設問', '何を測るか', '回答形式', 'どう使うか'],
        ['Q25〜Q36', '行動特性（ポータブルスキル）', '各問1〜5点',
         [('このうち', False), ('Q30（主体性）とQ31（変革意欲）', True), ('だけをフェーズ3で使う', False)]],
        ['Q37〜Q48', '価値観（キャリア・アンカー）', '各問1〜5点', '今回の判定では使わない'],
        [[('Q49', True)], '嘘回答の傾向（L尺度）', '1〜5点', '5点の場合、6タイプの点数を割り引く（フェーズ3のStep 0）'],
        [[('Q50', True)], 'BtoB企業が好きか', '1〜5点', '3点以上→BtoB企業リスト、3点未満→BtoC企業リストを出す'],
        [[('Q51', True)], 'メーカー/インフラ企業が好きか', '1〜5点', 'Q50と合わせて中西製作所様との相性を判定する'],
    ]
    mktable(s, 0.65, 1.9, 12, 4.8, data, col_widths=[1.6, 3.3, 2.0, 5.1], font_size=12)
    page_footer(s, 7)


def slide_08_result_rule():
    s = blank(); bg(s); kicker(s, 'PHASE 01 / 回答結果の決まり方')
    title_h2(s, '回答結果の決まり方')
    tbox_rich(s, 0.65, 1.55, 12, 0.4,
              [('6タイプの点数を高い順に並べて、', False, C['ink']),
               ('1位と2位', True, C['ink']),
               ('を取る。これが「この学生はどんなタイプか」の結論。', False, C['ink'])],
              size=14)
    tbox(s, 0.65, 2.05, 12, 0.4, '具体例①: ある学生の6タイプの点数', size=16, color=C['ink'], bold=True)
    # Big nums row 1
    types = [('R型', '8', False, None), ('I型', '6', False, None),
             ('A型', '7', False, None), ('S型（1位）', '18', True, C['accent']),
             ('E型（2位）', '10', True, C['accent2']), ('C型', '5', False, None)]
    draw_bignums(s, 0.65, 2.55, 12, types)
    bullets1 = [
        [('一番高い点数は ', False, C['ink']), ('S型の18点', True, C['ink']), (' → これが ', False, C['ink']), ('1位タイプ', True, C['ink'])],
        [('二番目に高いのは ', False, C['ink']), ('E型の10点', True, C['ink']), (' → これが ', False, C['ink']), ('2位タイプ', True, C['ink'])],
        [('この学生は', False, C['ink']), ('「人の役に立ちたい（S型）」が一番強く、「リーダーシップ（E型）」も持っている', True, C['ink'])],
    ]
    y = 3.55
    for runs in bullets1:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.1), Inches(0.07), Inches(0.07))
        dot.fill.solid(); dot.fill.fore_color.rgb = C['ink']; dot.line.fill.background()
        tbox_rich(s, 1.0, y, 11.5, 0.35, runs, size=12); y += 0.32
    # Another example
    tbox(s, 0.65, 4.65, 12, 0.4, 'もう1つの例', size=16, color=C['ink'], bold=True)
    types2 = [('R型', '10', False, None), ('I型', '8', False, None),
              ('A型（2位）', '16', True, C['accent2']), ('S型', '12', False, None),
              ('E型（1位）', '19', True, C['accent']), ('C型', '7', False, None)]
    draw_bignums(s, 0.65, 5.1, 12, types2)
    bullets2 = [
        [('1位: ', False, C['ink']), ('E型（19点）', True, C['ink']), (' → 2位: ', False, C['ink']), ('A型（16点）', True, C['ink'])],
        [('この学生は', False, C['ink']), ('「経営・リーダーシップ（E型）」が一番で、「創造力（A型）」もある', True, C['ink'])],
    ]
    y = 6.15
    for runs in bullets2:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.1), Inches(0.07), Inches(0.07))
        dot.fill.solid(); dot.fill.fore_color.rgb = C['ink']; dot.line.fill.background()
        tbox_rich(s, 1.0, y, 11.5, 0.35, runs, size=12); y += 0.32
    page_footer(s, 8)


def draw_bignums(slide, x0, y0, total_w, types):
    n = len(types); gap = 0.12
    w = (total_w - gap * (n - 1)) / n
    for i, (t, v, hi, hi_color) in enumerate(types):
        x = x0 + i * (w + gap)
        line = hi_color if hi else C['line']
        bg_col = C['white']
        if hi and hi_color:
            # light tint
            bg_col = RGBColor(
                min(255, hi_color[0] // 8 + 240),
                min(255, hi_color[1] // 8 + 240),
                min(255, hi_color[2] // 8 + 240),
            )
        rrect(slide, x, y0, w, 0.95, bg_col, line=line)
        tbox(slide, x, y0 + 0.07, w, 0.3, t, size=10, color=C['ink2'], bold=True, align=PP_ALIGN.CENTER)
        tbox(slide, x, y0 + 0.35, w, 0.55, v, size=24, color=C['ink'], bold=True, align=PP_ALIGN.CENTER, font='Courier New')


def slide_09_points():
    s = blank(); bg(s); kicker(s, 'PHASE 01 / 回答結果の決まり方')
    title_h2(s, 'ポイント')
    cards = [
        ('1位タイプ', 'メインの推薦企業を決める基準', '1位タイプがメインの推薦企業を決める基準になる'),
        ('2位タイプ', '補足の推薦企業', '「あなたにはこういう面もありますよ」という補足の推薦企業になる'),
        ('同点の場合', 'R → I → A → S → E → C', '1位と2位が同点の場合は R → I → A → S → E → C の順で先にくる方が1位'),
    ]
    cw = 3.9; gap = 0.2
    for i, (label, head, sub) in enumerate(cards):
        x = 0.65 + i * (cw + gap)
        rrect(s, x, 2.0, cw, 3.5, C['white'], line=C['line'])
        tbox(s, x + 0.25, 2.2, cw - 0.5, 0.35, label.upper(), size=10, color=C['ink2'], bold=True)
        tbox(s, x + 0.25, 2.6, cw - 0.5, 1.4, head, size=18, color=C['ink'], bold=True)
        tbox(s, x + 0.25, 4.0, cw - 0.5, 1.5, sub, size=12, color=C['ink'])
    page_footer(s, 9)


def slide_10_11values():
    s = blank(); bg(s); kicker(s, 'PHASE 01 / アウトプット')
    title_h2(s, 'フェーズ1でできあがるもの（＝フェーズ2で送る11個の数値）', size=22)
    data = [
        ['項目', '値の例', '何の点数か', 'フェーズ3のどこで使うか'],
        ['R', '8', 'R型の合計点（4〜20点）', '1位・2位タイプの判定に使う'],
        ['I', '6', 'I型の合計点', '同上'],
        ['A', '7', 'A型の合計点', '同上'],
        ['S', '18', 'S型の合計点', '同上'],
        ['E', '10', 'E型の合計点', '同上'],
        ['C', '5', 'C型の合計点', '同上'],
        ['Q30', '4', '主体性の点数（1〜5点）', '中西製作所様向きかの判定に使う'],
        ['Q31', '4', '変革意欲の点数（1〜5点）', '同上'],
        ['Q49', '3', '嘘回答検知の点数（1〜5点）', '5点のときだけ6タイプの点数を補正する'],
        ['Q50', '4', 'BtoB志向の点数（1〜5点）', 'BtoB企業 or BtoC企業の切り替えに使う'],
        ['Q51', '4', 'メーカー志向の点数（1〜5点）', '中西製作所様との相性判定に使う'],
    ]
    mktable(s, 0.65, 1.8, 12, 5.1, data, col_widths=[1.1, 1.2, 4.5, 5.2], font_size=11)
    page_footer(s, 10)


def slide_11_p1_confirm():
    s = blank(); bg(s); kicker(s, 'PHASE 01 / 現在の設計')
    title_h2(s, '現在の設計 と 確認事項')
    # Callout
    rrect(s, 0.65, 2.0, 12, 2.6, C['callout_bg'], line=C['accent2'])
    rect(s, 0.65, 2.0, 0.08, 2.6, C['accent2'])
    tbox(s, 0.9, 2.15, 11.5, 0.4, 'ℹ️ 現在の設計', size=16, color=C['ink'], bold=True)
    items = [
        [('6タイプの合計点の算出', True, C['ink']),
         (': MIL様側で合計してからDifyに送る設計になっている（Dify側は合計済みの点数を受け取る）', False, C['ink'])],
        [('嘘回答補正（Q49=5のとき点数を×0.8する処理）', True, C['ink']),
         (': Dify側で自動的に処理する設計になっている（MIL様は素の点数を送ればOK）', False, C['ink'])],
    ]
    y = 2.65
    for runs in items:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y + 0.12), Inches(0.08), Inches(0.08))
        dot.fill.solid(); dot.fill.fore_color.rgb = C['ink']; dot.line.fill.background()
        tbox_rich(s, 1.2, y, 11.2, 0.85, runs, size=13, line_spacing=1.3); y += 0.85
    # Confirm box
    rrect(s, 0.65, 4.9, 12, 0.9, C['confirm_bg'], line=C['warn'])
    tbox_rich(s, 0.9, 5.05, 11.5, 0.7,
              [('→ 会議で上記の認識に相違がないか確認する', True, RGBColor(0x6B, 0x3A, 0x00), 15)],
              size=15)
    page_footer(s, 11)


def slide_12_div_p2():
    divider('02', '検査結果をDifyに送る',
            'フェーズ1で算出された11個の数値を、MIL様のサーバーから中西製作所様のDifyに送る。\n'
            '送ると自動的にフェーズ3の判定処理が走り、結果が返ってくる（フェーズ4）。',
            'PHASE 02 / MIL 様 → Dify', 12)


def slide_13_one_trip():
    s = blank(); bg(s); kicker(s, 'PHASE 02 / 通信の構造')
    title_h2(s, 'フェーズ2（送る）とフェーズ4（返ってくる）は1回のやり取り', size=22)
    tbox(s, 0.65, 1.6, 12, 0.6,
         '電話でいうと、MIL様が電話をかけて質問し、Difyが電話口で答えを返す ── これが1回で完結する。',
         size=14, color=C['ink'])
    # ASCII
    rrect(s, 0.65, 2.5, 12, 3.5, C['white'], line=C['line'])
    tbox(s, 0.85, 2.65, 11.8, 3.3,
         'MIL様のサーバー                     中西製作所様のDifyアカウント\n'
         '                                    （エヌイチが構築・運用）\n\n'
         '  ① 11個の数値を送る ─────→    ② 受け取って判定処理を実行\n'
         '                                     （フェーズ3の処理が動く）\n\n'
         '  ④ 結果を受け取る   ←─────    ③ 判定結果を返す\n'
         '                                     （フェーズ4）',
         size=13, color=C['ink'], font='Courier New')
    page_footer(s, 13)


def slide_14_required():
    s = blank(); bg(s); kicker(s, 'PHASE 02 / 準備')
    title_h2(s, '送るために必要なもの')
    tbox_rich(s, 0.65, 1.55, 12.5, 0.75,
              [('MIL様がDifyにデータを送るには、以下の', False, C['ink']),
               ('2つの情報', True, C['ink']),
               ('が必要。どちらも', False, C['ink']),
               ('中西製作所様のDifyアカウントから発行されるもの', True, C['ink']),
               ('で、エヌイチが代行してMIL様に共有する。', False, C['ink'])],
              size=12, line_spacing=1.35)
    data = [
        ['必要なもの', '何か', '誰が発行するか', '誰に渡すか'],
        [[('エンドポイントURL', True), ('（接続先）', False)], 'データの送り先アドレス', '中西製作所様のDifyアカウント', 'MIL様に共有'],
        [[('APIキー', True), ('（認証キー）', False)], '「この人は使っていいですよ」という合鍵', '中西製作所様のDifyアカウント', 'MIL様に共有'],
    ]
    mktable(s, 0.65, 2.4, 12, 1.5, data, col_widths=[3.0, 4.0, 2.8, 2.2], font_size=11)
    tbox(s, 0.65, 4.05, 12, 0.4, '準備の流れ', size=16, color=C['ink'], bold=True)
    steps = [
        'エヌイチが中西製作所様のDifyアカウント上でしくみを完成させる',
        'Difyの管理画面で「公開」すると、エンドポイントURLとAPIキーが自動発行される',
        'この2つをエヌイチからMIL様に共有する',
        'MIL様はこの2つを自社サーバーに設定して、学生が検査を完了するたびにデータを送れるようにする',
    ]
    y = 4.55
    for i, t in enumerate(steps, 1):
        # numbered bullet
        oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y + 0.04), Inches(0.28), Inches(0.28))
        oval.fill.solid(); oval.fill.fore_color.rgb = C['ink']; oval.line.fill.background()
        tbox(s, 0.7, y + 0.02, 0.28, 0.3, str(i), size=11, color=C['white'], bold=True, align=PP_ALIGN.CENTER)
        tbox(s, 1.1, y, 11.8, 0.4, t, size=12, color=C['ink'])
        y += 0.38
    # Callout
    rrect(s, 0.65, 6.15, 12, 0.75, C['callout_bg'], line=C['accent2'])
    rect(s, 0.65, 6.15, 0.08, 0.75, C['accent2'])
    tbox_rich(s, 0.9, 6.22, 11.5, 0.6,
              [('MIL様が何かを発行してエヌイチや中西製作所様に渡す必要はありません。', True, C['ink'], 13),
               ('\n中西製作所様のDifyから出た情報を、MIL様が受け取って使うだけです。', False, C['ink'], 12)],
              line_spacing=1.3)
    page_footer(s, 14)


def slide_15_send_data():
    s = blank(); bg(s); kicker(s, 'PHASE 02 / データ')
    title_h2(s, '送るデータの中身')
    tbox(s, 0.65, 1.6, 12, 0.4, 'フェーズ1でできた11個の数値をまとめて送る。', size=14, color=C['ink'])
    rrect(s, 0.65, 2.15, 12, 4.2, C['white'], line=C['line'])
    tbox(s, 0.9, 2.3, 11.5, 4.0,
         '送信データの中身（例）:\n\n'
         '  R   = 8        ← R型の合計点\n'
         '  I   = 6        ← I型の合計点\n'
         '  A   = 7        ← A型の合計点\n'
         '  S   = 18       ← S型の合計点\n'
         '  E   = 10       ← E型の合計点\n'
         '  C   = 5        ← C型の合計点\n'
         '  Q30 = 4        ← 主体性の点数\n'
         '  Q31 = 4        ← 変革意欲の点数\n'
         '  Q49 = 3        ← 嘘回答検知の点数\n'
         '  Q50 = 4        ← BtoB志向の点数\n'
         '  Q51 = 4        ← メーカー志向の点数',
         size=12, color=C['ink'], font='Courier New')
    tbox(s, 0.65, 6.5, 12, 0.4, '全部で11個、すべて数値。', size=14, color=C['ink'], bold=True)
    page_footer(s, 15)


def slide_16_p2_confirm():
    s = blank(); bg(s); kicker(s, 'PHASE 02 / 現在の設計')
    title_h2(s, '現在の設計 と 確認・相談事項')
    rrect(s, 0.65, 2.0, 12, 2.7, C['callout_bg'], line=C['accent2'])
    rect(s, 0.65, 2.0, 0.08, 2.7, C['accent2'])
    tbox(s, 0.9, 2.15, 11.5, 0.4, 'ℹ️ 現在の設計', size=16, color=C['ink'], bold=True)
    items = [
        [('データの項目名', True, C['ink']),
         (': R, I, A, S, E, C, Q30, Q31, Q49, Q50, Q51 でDify側の受け口を作っている', False, C['ink'])],
        [('APIキー', True, C['ink']),
         (': まだ発行していない。Difyの有料プラン（Professional: 月額約9,000円）に切り替えた後に発行・共有する予定', False, C['ink'])],
    ]
    y = 2.7
    for runs in items:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y + 0.12), Inches(0.08), Inches(0.08))
        dot.fill.solid(); dot.fill.fore_color.rgb = C['ink']; dot.line.fill.background()
        tbox_rich(s, 1.2, y, 11.2, 0.95, runs, size=13, line_spacing=1.3); y += 0.95
    # Two confirm boxes
    rrect(s, 0.65, 5.0, 5.85, 1.0, C['confirm_bg'], line=C['warn'])
    tbox_rich(s, 0.85, 5.12, 5.5, 0.8,
              [('→ 確認\n', True, RGBColor(0x6B, 0x3A, 0x00), 14),
               ('項目名がMIL様側の実装と合っているか', False, C['ink2'], 12)],
              line_spacing=1.3)
    rrect(s, 6.8, 5.0, 5.85, 1.0, C['confirm_bg'], line=C['warn'])
    tbox_rich(s, 7.0, 5.12, 5.5, 0.8,
              [('→ 相談\n', True, RGBColor(0x6B, 0x3A, 0x00), 14),
               ('APIキーの共有方法・タイミング', False, C['ink2'], 12)],
              line_spacing=1.3)
    page_footer(s, 16)


def slide_17_div_p3():
    divider('03', '企業のマッチング',
            'MIL様から受け取った11個の数値をもとに、「この学生にどの企業を推薦するか」を自動で判定する。',
            'PHASE 03 / Dify / エヌイチ構築', 17)


def slide_18_judge_inputs():
    s = blank(); bg(s); kicker(s, 'PHASE 03 / 判定に使うもの')
    title_h2(s, '判定に使うもの')
    data = [
        ['何を', '説明'],
        [[('11個の数値', True), ('（フェーズ2で受け取ったもの）', False)],
         '学生の志向性タイプ・行動特性・業界志向'],
        [[('企業リスト', True), ('（Difyの中にあらかじめ登録してあるもの）', False)],
         '7タイプ × BtoB/BtoC = 14パターンの企業名と推薦理由。津丸様提供のExcel「紹介企業」シートに基づき4/10に登録済み'],
    ]
    mktable(s, 0.65, 2.0, 12, 2.4, data, col_widths=[4.0, 8.0], font_size=13)
    page_footer(s, 18)


def slide_19_logic_overview():
    s = blank(); bg(s); kicker(s, 'PHASE 03 / 判定のしくみ')
    title_h2(s, '判定のしくみ（津丸様との確認に基づく確定ロジック）', size=22)
    tbox_rich(s, 0.65, 1.55, 12.5, 0.5,
              [('判定は', False, C['ink']), ('4段階', True, C['ink']),
               ('あり、', False, C['ink']), ('上から順番にチェックしていく', True, C['ink']),
               ('。どこかのステップに当てはまったら、そこで結果が確定する。下のステップは見ない。', False, C['ink'])],
              size=13, line_spacing=1.3)
    steps = [
        ('gray', '0', '前処理: 嘘回答の補正', ''),
        ('orange', '1', '最優先: 「志向が定まっていない」学生を検出', '→ 当てはまったら → トモキャリに誘導して終了'),
        ('green', '2', '「主体性・変革意欲が高い」かつ「BtoB志向」の学生を検出', '→ 当てはまったら → 中西製作所様を推薦して終了'),
        ('green', '3', '「中西と相性のいいタイプ＋業界志向」の学生を検出', '→ 当てはまったら → 中西製作所様を推薦して終了'),
        ('blue', '4', 'タイプに合った企業を推薦', ''),
    ]
    col_map = {'gray': RGBColor(0x3A, 0x3A, 0x3A), 'orange': C['accent'],
               'green': C['ok'], 'blue': C['accent2']}
    y = 2.3
    for (color, num, main, sub) in steps:
        rrect(s, 0.65, y, 12, 0.8, C['white'], line=C['line'])
        rrect(s, 0.85, y + 0.12, 0.55, 0.55, col_map[color])
        tbox(s, 0.85, y + 0.13, 0.55, 0.55, num, size=18, color=C['white'], bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tbox(s, 1.55, y + 0.1, 11, 0.3, main, size=13, color=C['ink'], bold=True)
        if sub:
            tbox(s, 1.55, y + 0.4, 11, 0.3, sub, size=11, color=C['ink2'])
        y += 0.9
    page_footer(s, 19)


def slide_20_step0():
    s = blank(); bg(s); kicker(s, 'PHASE 03 / Step 0')
    title_h2(s, 'Step 0: 嘘回答の補正（前処理）', size=24)
    tbox(s, 0.65, 1.6, 12, 0.35, '目的', size=15, color=C['ink'], bold=True)
    tbox(s, 0.65, 1.95, 12, 0.4, '自分をよく見せようとして回答を盛っている可能性がある学生の点数を割り引く。',
         size=13, color=C['ink'])
    tbox(s, 0.65, 2.4, 12, 0.35, 'しくみ', size=15, color=C['ink'], bold=True)
    items = [
        [('Q49（嘘回答検知）の点数が ', False, C['ink']), ('5点（満点）のときだけ', True, C['ink']),
         ('、6タイプ（R, I, A, S, E, C）の点数すべてに', False, C['ink']), ('×0.8', True, C['ink']),
         ('をかける', False, C['ink'])],
        [('Q49が5点以外のときは、何もしない', False, C['ink'])],
    ]
    y = 2.8
    for runs in items:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.12), Inches(0.08), Inches(0.08))
        dot.fill.solid(); dot.fill.fore_color.rgb = C['ink']; dot.line.fill.background()
        tbox_rich(s, 1.0, y, 11.8, 0.4, runs, size=12); y += 0.38
    tbox(s, 0.65, 3.7, 12, 0.35, '例', size=15, color=C['ink'], bold=True)
    data = [
        ['学生', 'Q49の点数', '補正前のE型', '補正前のR型', '→ 補正後のE型', '→ 補正後のR型'],
        ['Xさん', [('5点', True)], '20点', '15点', [('16点', True), ('（×0.8）', False)], [('12点', True), ('（×0.8）', False)]],
        ['Yさん', '3点', '20点', '15点', '20点（そのまま）', '15点（そのまま）'],
    ]
    mktable(s, 0.65, 4.1, 12, 1.5, data, col_widths=[1.5, 1.8, 2.0, 2.0, 2.4, 2.3], font_size=11)
    tbox(s, 0.65, 5.75, 12, 0.4, 'この補正後の点数を使って、以降のStep 1〜4を判定する。',
         size=13, color=C['ink'])
    page_footer(s, 20)


def slide_21_step1():
    s = blank(); bg(s); kicker(s, 'PHASE 03 / Step 1')
    title_h2(s, 'Step 1: 志向が定まっていない学生を検出する（最優先）', size=22)
    tbox(s, 0.65, 1.55, 12, 0.35, '目的', size=15, color=C['ink'], bold=True)
    tbox_rich(s, 0.65, 1.9, 12.5, 0.7,
              [('6タイプの点数にほとんど差がない学生 ＝ どの仕事に向いているか判断できない学生を検出して、'
                '企業を推薦せずに', False, C['ink']),
               ('キャリア相談（トモキャリ）に誘導する', True, C['ink']),
               ('。', False, C['ink'])], size=13, line_spacing=1.35)
    tbox(s, 0.65, 2.7, 12, 0.35, '条件', size=15, color=C['ink'], bold=True)
    tbox_rich(s, 0.65, 3.05, 12.5, 0.4,
              [('6タイプの中で一番高い点数と一番低い点数の差が ', False, C['ink']),
               ('1点以下', True, C['ink'])], size=13)
    tbox(s, 0.65, 3.6, 12, 0.35, '例', size=15, color=C['ink'], bold=True)
    data = [
        ['学生', 'R', 'I', 'A', 'S', 'E', 'C', '最高−最低', '結果'],
        ['Aさん', '12', '12', '11', '12', '12', '11',
         [('12−11 = ', False), ('1点', True)],
         [('→ ', False), ('トモキャリ誘導', True), ('（差がなさすぎて判定できない）', False)]],
        ['Bさん', '15', '12', '8', '18', '10', '6',
         [('18−6 = ', False), ('12点', True)],
         '→ 次のStep 2へ（はっきり差がある）'],
    ]
    mktable(s, 0.65, 4.0, 12, 1.4, data,
            col_widths=[1.2, 0.6, 0.6, 0.6, 0.6, 0.6, 0.6, 2.0, 5.2], font_size=11)
    tbox(s, 0.65, 5.55, 12, 0.5,
         '→ Aさんのように全タイプの点数が横並びの場合は、企業を推薦しても意味がないのでトモキャリに案内する。',
         size=13, color=C['ink'])
    page_footer(s, 21)


def slide_22_step2():
    s = blank(); bg(s); kicker(s, 'PHASE 03 / Step 2')
    title_h2(s, 'Step 2: 主体性・変革意欲が高く、BtoB志向の学生を検出する', size=22)
    tbox(s, 0.65, 1.55, 12, 0.35, '目的', size=14, color=C['ink'], bold=True)
    tbox_rich(s, 0.65, 1.9, 12.5, 0.4,
              [('中西製作所様の社風に合う学生を検出して、', False, C['ink']),
               ('中西製作所様を推薦する', True, C['ink']), ('。', False, C['ink'])],
              size=12)
    tbox(s, 0.65, 2.3, 12, 0.35, '背景（津丸様との確認に基づく）', size=14, color=C['ink'], bold=True)
    tbox_rich(s, 0.65, 2.65, 12.5, 0.7,
              [('中西製作所様は「安定した基盤がありながらベンチャー精神で変革を推進する」社風。'
                'この社風に合う人材は、', False, C['ink']),
               ('主体性（Q30）と変革意欲（Q31）が両方高い人', True, C['ink']),
               ('。ただし中西製作所様はBtoB企業なので、', False, C['ink']),
               ('BtoC志向が強い学生には中西を案内しない', True, C['ink']),
               ('。', False, C['ink'])],
              size=11, line_spacing=1.3)
    tbox_rich(s, 0.65, 3.4, 12.5, 0.35,
              [('条件（3つを', False, C['ink']), ('すべて', True, C['ink']),
               ('満たす）', False, C['ink'])],
              size=14)
    cond_items = [
        [('Q30（主体性）が ', False), ('4点以上', True), ('（5点満点中）', False)],
        [('Q31（変革意欲）が ', False), ('4点以上', True), ('（5点満点中）', False)],
        [('Q50（BtoB志向）が ', False), ('3点以上', True), ('（＝ BtoB寄り）', False)],
    ]
    y = 3.8
    for i, runs in enumerate(cond_items, 1):
        tbox(s, 0.75, y, 0.4, 0.35, f'{i}.', size=12, color=C['ink'], bold=True)
        tbox_rich(s, 1.15, y, 11.5, 0.35, runs, size=12); y += 0.32
    data = [
        ['学生', 'Q30', 'Q31', 'Q50', '結果'],
        ['Aさん', '4', '5', '4', [('3つとも満たす → ', False), ('中西製作所様を推薦', True)]],
        ['Bさん', '5', '3', '4', 'Q31が4未満 → 次のStep 3へ'],
        ['Cさん', '4', '4', '2', 'Q50が3未満（BtoC寄り）→ 次のStep 3へ'],
        ['Dさん', '3', '4', '4', 'Q30が4未満 → 次のStep 3へ'],
    ]
    mktable(s, 0.65, 4.85, 12, 1.8, data,
            col_widths=[1.0, 0.8, 0.8, 0.8, 8.6], font_size=10)
    tbox(s, 0.65, 6.75, 12, 0.4,
         '→ Cさんは主体性・変革意欲は高いが、BtoC志向が強いため中西製作所様（BtoB企業）には案内せず、Step 4の通常判定でBtoC企業を推薦する。',
         size=10, color=C['ink2'])
    page_footer(s, 22)


def slide_23_step3():
    s = blank(); bg(s); kicker(s, 'PHASE 03 / Step 3')
    title_h2(s, 'Step 3: 中西と相性のいいタイプ × 業界志向の学生を検出する', size=22)
    tbox_rich(s, 0.65, 1.55, 12.5, 0.65,
              [('目的: Step 2に当てはまらなかったけど、', False, C['ink']),
               ('タイプ的にも業界志向的にも中西製作所様と合っている学生', True, C['ink']),
               ('を拾って、', False, C['ink']),
               ('中西製作所様を推薦する', True, C['ink']),
               ('。以下の3つの条件を', False, C['ink']),
               ('すべて', True, C['ink']),
               ('満たす必要がある。', False, C['ink'])],
              size=11, line_spacing=1.3)
    # 3 cards
    cards = [
        ('条件 ①', 'タイプの相性',
         '1位タイプか2位タイプのどちらかが E型・I型・A型 のいずれかであること。\n'
         '・E型（経営・リーダーシップ志向）\n・I型（研究・分析志向）\n・A型（創造・企画志向）\n\n'
         '1位でなくても、2位に入っていればOK。'),
        ('条件 ②', '業界志向の相性',
         'Q50（BtoB志向）とQ51（メーカー志向）の平均が3点以上であること。\n\n'
         '→ BtoB企業やメーカーに興味がある ＝ 中西製作所様の業界と合っている。'),
        ('条件 ③', 'BtoB志向',
         'Q50（BtoB志向）が3点以上であること。\n\n'
         '→ 中西製作所様はBtoB企業なので、BtoC寄りの学生には案内しない。'),
    ]
    cw = 3.9; gap = 0.2
    for i, (lbl, head, body) in enumerate(cards):
        x = 0.65 + i * (cw + gap)
        rrect(s, x, 2.35, cw, 2.55, C['white'], line=C['line'])
        tbox(s, x + 0.2, 2.5, cw - 0.4, 0.3, lbl, size=10, color=C['ink2'], bold=True)
        tbox(s, x + 0.2, 2.8, cw - 0.4, 0.45, head, size=14, color=C['ink'], bold=True)
        tbox(s, x + 0.2, 3.3, cw - 0.4, 1.6, body, size=9, color=C['ink'])
    tbox(s, 0.65, 5.0, 12, 0.35, '例', size=14, color=C['ink'], bold=True)
    data = [
        ['学生', '1位', '2位', 'Q50', 'Q51', '平均', '①', '②', '③', '結果'],
        ['Aさん', 'E型(19)', 'A型(16)', '4', '4', '4.0', '✅', '✅', '✅',
         [('→ ', False), ('中西製作所様を推薦', True)]],
        ['Bさん', 'S型(18)', 'I型(15)', '4', '3', '3.5', '✅', '✅', '✅',
         [('→ ', False), ('中西製作所様を推薦', True)]],
        ['Cさん', 'S型(18)', 'R型(15)', '4', '4', '4.0', '❌', '✅', '✅',
         '→ Step 4へ（タイプが対象外）'],
        ['Dさん', 'E型(19)', 'A型(16)', '2', '4', '3.0', '✅', '✅', '❌',
         '→ Step 4へ（BtoC寄り）'],
    ]
    mktable(s, 0.65, 5.4, 12, 1.6, data,
            col_widths=[0.9, 1.1, 1.1, 0.55, 0.55, 0.65, 0.5, 0.5, 0.5, 5.65], font_size=9)
    tbox(s, 0.65, 7.05, 12, 0.35,
         '→ Dさんはタイプ・業界志向の平均は合うが、Q50が3未満でBtoC寄りなので中西製作所様は案内せず、Step 4（通常判定）に進む。',
         size=9, color=C['ink2'])
    page_footer(s, 23)


def slide_24_step4():
    s = blank(); bg(s); kicker(s, 'PHASE 03 / Step 4')
    title_h2(s, 'Step 4: タイプに合った企業を推薦する（通常判定）', size=22)
    tbox(s, 0.65, 1.55, 12, 0.35, '目的', size=14, color=C['ink'], bold=True)
    tbox_rich(s, 0.65, 1.9, 12.5, 0.4,
              [('Step 1〜3のどれにも当てはまらなかった学生に対し、', False, C['ink']),
               ('その学生のタイプに合った企業を推薦する', True, C['ink']),
               ('。', False, C['ink'])], size=12)
    tbox(s, 0.65, 2.3, 12, 0.35, 'しくみ', size=14, color=C['ink'], bold=True)
    tbox(s, 0.65, 2.65, 12, 0.4, '以下の2つを出す。', size=12, color=C['ink'])
    bullets = [
        [('メインの推薦企業', True, C['ink']),
         (': 1位タイプに対応する企業リストから出す', False, C['ink'])],
        [('補足の推薦企業', True, C['ink']),
         (': 2位タイプに対応する企業リストから出す（「あなたにはこんな面もありますよ」という補足）', False, C['ink'])],
    ]
    y = 3.05
    for i, runs in enumerate(bullets, 1):
        tbox(s, 0.75, y, 0.3, 0.3, f'{i}.', size=12, color=C['ink'], bold=True)
        tbox_rich(s, 1.1, y, 11.5, 0.4, runs, size=12); y += 0.35
    tbox_rich(s, 0.65, 3.8, 12.5, 0.35,
              [('さらに、', False, C['ink']),
               ('Q50（BtoB志向）の点数', True, C['ink']),
               ('で、BtoB企業リストとBtoC企業リストを切り替える。', False, C['ink'])],
              size=12)
    b = [
        [('Q50が ', False), ('3点以上', True), (' → BtoB企業リストを使う', False)],
        [('Q50が ', False), ('3点未満', True), (' → BtoC企業リストを使う', False)],
    ]
    y = 4.25
    for runs in b:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y + 0.12), Inches(0.07), Inches(0.07))
        dot.fill.solid(); dot.fill.fore_color.rgb = C['ink']; dot.line.fill.background()
        tbox_rich(s, 1.0, y, 11.5, 0.35, [(r[0], r[1], C['ink']) for r in runs], size=12); y += 0.32
    tbox(s, 0.65, 5.0, 12, 0.4, '具体例: 1位 = S型、2位 = R型、Q50 = 4点 の学生の場合',
         size=14, color=C['ink'], bold=True)
    rrect(s, 0.65, 5.45, 12, 1.5, C['white'], line=C['line'])
    tbox(s, 0.85, 5.55, 11.8, 1.4,
         '1位 S型 → Q50が4点（3以上なのでBtoB） → 企業リストの「S型 × BtoB」を見る\n'
         '  → メイン推薦: ベネッセ, リクルートマネジメント, パソナ, ワークスアプリ, LITALICO\n\n'
         '2位 R型 → Q50が4点（3以上なのでBtoB） → 企業リストの「R型 × BtoB」を見る\n'
         '  → 補足推薦: 三菱電機, 村田製作所, 安川電機, 栗田工業, THK',
         size=11, color=C['ink'], font='Courier New')
    page_footer(s, 24)


def slide_25_company_list():
    s = blank(); bg(s); kicker(s, 'PHASE 03 / 企業リスト')
    title_h2(s, '企業リスト（Difyの中に登録済み）', size=24)
    tbox(s, 0.65, 1.55, 12, 0.7,
         '津丸様から提供いただいたExcelシート「紹介企業」にもとづいて、以下のリストをDifyの中に登録している（4/10登録済み）。'
         '各企業に対する推薦理由テキストも登録済み。',
         size=11, color=C['ink2'])
    data = [
        ['タイプ', 'BtoB企業（Q50が3点以上のとき）', 'BtoC企業（Q50が3点未満のとき）'],
        [[('中西型', True), ('（Step 2,3で該当した場合）', False)],
         '中西製作所, キーエンス, 平田機工, マクニカ, ダイフク',
         '※中西該当かつBtoC寄りの場合は中西を案内せず、Step 4で通常判定に回す'],
        [[('R型', True), ('（現実型）', False)],
         '三菱電機, 村田製作所, 安川電機, 栗田工業, THK',
         'トヨタ, 任天堂, スノーピーク, ヤマハ, マキタ'],
        [[('I型', True), ('（研究型）', False)],
         '野村総研, 島津製作所, 中外製薬, オムロン, ユーザーベース',
         'Google, 武田薬品, コーセー, エーザイ, メルカリ'],
        [[('A型', True), ('（芸術型）', False)],
         '電通, 博報堂, チームラボ, Sansan, 乃村工藝社',
         'サンリオ, オリエンタルランド, 資生堂, 集英社, 無印良品'],
        [[('S型', True), ('（社会型）', False)],
         'ベネッセ, リクルートマネジメント, パソナ, ワークスアプリ, LITALICO',
         '日本赤十字, 損保ジャパン, 星野リゾート, Zoff, クラシコム'],
        [[('E型', True), ('（企業型）', False)],
         '三菱商事, 伊藤忠, 光通信, サイバーエージェント, レバレジーズ',
         'ソフトバンク, 楽天, ファーストリテイリング, ZOZO'],
        [[('C型', True), ('（慣習型）', False)],
         '日本生命, 東京海上, 三菱UFJ, 富士通, オービック',
         'JR東日本, ANA, 日本郵便, 公務員, ゆうちょ銀行'],
    ]
    mktable(s, 0.65, 2.3, 12, 4.7, data, col_widths=[2.5, 4.75, 4.75], font_size=10)
    page_footer(s, 25)


def slide_26_matching_summary():
    s = blank(); bg(s); kicker(s, 'PHASE 03 / まとめ')
    title_h2(s, '判定結果と企業リストの組み合わせまとめ', size=22)
    data = [
        ['判定結果', '企業リストのどこを使うか', '出力例'],
        [[('Step 2 or 3 で', False), ('中西該当', True), ('（Q50≧3のみ）', False)],
         '「中西型 × BtoB」の行を使う', '中西製作所, キーエンス, 平田機工…'],
        [[('Step 4 の', False), ('メイン推薦', True)],
         '1位タイプの行。Q50でBtoB/BtoC切替', '1位がE型 × BtoB → 三菱商事, 伊藤忠…'],
        [[('Step 4 の', False), ('補足推薦', True)],
         '2位タイプの行。Q50でBtoB/BtoC切替', '2位がA型 × BtoB → 電通, 博報堂…'],
        [[('Step 1 で', False), ('未分化', True)],
         '企業リストは使わない', 'トモキャリ誘導（企業推薦なし）'],
    ]
    mktable(s, 0.65, 1.9, 12, 2.0, data, col_widths=[3.5, 4.5, 4.0], font_size=11)
    rrect(s, 0.65, 4.1, 12, 1.9, C['callout_bg'], line=C['accent2'])
    rect(s, 0.65, 4.1, 0.08, 1.9, C['accent2'])
    tbox(s, 0.9, 4.2, 11.5, 0.4, 'ℹ️ 現在の設計', size=14, color=C['ink'], bold=True)
    items = [
        [('企業リスト', True, C['ink']),
         (': 4/10に津丸様提供のExcelに基づいて登録済み。推薦理由テキストも全タイプ分登録済み', False, C['ink'])],
        [('推薦理由テキストの作成', True, C['ink']),
         (': エヌイチ側でExcelの内容をもとに作成・登録した。中西製作所様側での最終チェックはまだ', False, C['ink'])],
    ]
    y = 4.65
    for runs in items:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y + 0.12), Inches(0.07), Inches(0.07))
        dot.fill.solid(); dot.fill.fore_color.rgb = C['ink']; dot.line.fill.background()
        tbox_rich(s, 1.2, y, 11.2, 0.7, runs, size=11, line_spacing=1.3); y += 0.65
    rrect(s, 0.65, 6.15, 5.85, 0.8, C['confirm_bg'], line=C['warn'])
    tbox_rich(s, 0.85, 6.25, 5.5, 0.6,
              [('→ 確認\n', True, RGBColor(0x6B, 0x3A, 0x00), 12),
               ('企業リストの内容に変更・追加がないか', False, C['ink2'], 11)])
    rrect(s, 6.8, 6.15, 5.85, 0.8, C['confirm_bg'], line=C['warn'])
    tbox_rich(s, 7.0, 6.25, 5.5, 0.6,
              [('→ 相談\n', True, RGBColor(0x6B, 0x3A, 0x00), 12),
               ('推薦理由テキストの最終チェックを誰がいつやるか', False, C['ink2'], 11)])
    page_footer(s, 26)


def slide_27_div_p4():
    divider('04', '結果をMIL様に返す',
            'フェーズ3の判定が終わると、その結果がMIL様に自動的に返される。\n'
            'フェーズ2でMIL様がデータを送ったのと同じ通信の「返事」として返ってくるので、MIL様側で特別な設定は不要。\n'
            'APIキーも不要（フェーズ2で送信したときに認証は済んでいるため）。',
            'PHASE 04 / Dify → MIL 様', 27)


def slide_28_return_data():
    s = blank(); bg(s); kicker(s, 'PHASE 04 / 返却データ')
    title_h2(s, '返ってくるデータの中身')
    tbox(s, 0.65, 1.55, 12, 0.45,
         '判定結果として以下の8つの情報が返ってくる。MIL様はこれをもとに学生の結果画面を表示する。',
         size=12, color=C['ink'])
    data = [
        ['項目', '何が入るか', 'MIL様での使い方'],
        [[('判定カテゴリ', True)], '「中西該当」「タイプ別マッチング」「トモキャリ誘導」のいずれか', 'どの結果画面を出すかの判別に使う'],
        [[('1位タイプ', True)], 'R〜Cのいずれか', '「あなたは〇〇タイプです」として表示'],
        [[('2位タイプ', True)], 'R〜Cのいずれか', '「〇〇な面もあります」として補足表示'],
        [[('メインの推薦企業', True)], '企業名のリスト（最大5社）', '結果画面にメインとして企業名を表示'],
        [[('メインの推薦理由', True)], '各企業への推薦理由テキスト', '企業名の下に表示'],
        [[('補足の推薦企業', True)], '企業名のリスト(通常判定時のみ)', '「こんな企業も」として表示'],
        [[('補足の推薦理由', True)], '各企業への推薦理由テキスト', '同上'],
        [[('メッセージ', True)], '学生への一言メッセージ', '結果画面の冒頭 or 末尾に表示'],
    ]
    mktable(s, 0.65, 2.1, 12, 4.8, data, col_widths=[2.4, 5.0, 4.6], font_size=11)
    page_footer(s, 28)


def slide_29_return_diff():
    s = blank(); bg(s); kicker(s, 'PHASE 04 / 返却データ')
    title_h2(s, '判定カテゴリごとに返ってくる内容の違い', size=22)
    data = [
        ['', '中西該当（Step 2 or 3）', 'タイプ別マッチング（Step 4）', 'トモキャリ誘導（Step 1）'],
        [[('メインの推薦企業', True)], '中西製作所様＋同系統5社', '1位タイプに合う企業', 'なし（企業推薦しない）'],
        [[('補足の推薦企業', True)], 'なし（中西への誘導を優先）', '2位タイプに合う企業', 'なし'],
        [[('メッセージ', True)], '中西製作所様向けメッセージ', 'タイプに合わせたメッセージ', 'トモキャリへの誘導文'],
    ]
    mktable(s, 0.65, 2.0, 12, 3.0, data, col_widths=[2.6, 3.2, 3.2, 3.0], font_size=12)
    page_footer(s, 29)


def slide_30_example1():
    s = blank(); bg(s); kicker(s, 'PHASE 04 / 具体例')
    title_h2(s, '具体例: Step 2で中西該当になった学生の場合', size=22)
    rrect(s, 0.65, 2.0, 12, 4.7, C['white'], line=C['line'])
    tbox(s, 0.85, 2.1, 11.8, 4.5,
         '返ってくるデータ:\n\n'
         '  判定カテゴリ    = 中西該当\n'
         '  1位タイプ       = S型（社会型）\n'
         '  2位タイプ       = R型（現実型）\n'
         '  メインの推薦企業 = 中西製作所, キーエンス, 平田機工, マクニカ, ダイフク\n'
         '  メインの推薦理由 = （各社の「なぜあなたに合うか」テキスト）\n'
         '  補足の推薦企業  = なし\n'
         '  補足の推薦理由  = なし\n'
         '  メッセージ     = 「あなたの主体性と変革意欲は、安定基盤×ベンチャー志向の\n'
         '                   企業で大きく花開きます。」',
         size=12, color=C['ink'], font='Courier New')
    page_footer(s, 30)


def slide_31_example2():
    s = blank(); bg(s); kicker(s, 'PHASE 04 / 具体例')
    title_h2(s, '具体例: Step 4で通常判定になった学生の場合', size=22)
    rrect(s, 0.65, 1.75, 12, 2.9, C['white'], line=C['line'])
    tbox(s, 0.85, 1.85, 11.8, 2.8,
         '返ってくるデータ:\n\n'
         '  判定カテゴリ    = タイプ別マッチング\n'
         '  1位タイプ       = E型（企業型）\n'
         '  2位タイプ       = A型（芸術型）\n'
         '  メインの推薦企業 = 三菱商事, 伊藤忠, 光通信, サイバーエージェント, レバレジーズ\n'
         '  メインの推薦理由 = （各社の理由テキスト）\n'
         '  補足の推薦企業  = 電通, 博報堂, チームラボ, Sansan, 乃村工藝社\n'
         '  補足の推薦理由  = （各社の理由テキスト）\n'
         '  メッセージ     = 「あなたはE型の特性が最も強く、A型の側面も備わっています。」',
         size=11, color=C['ink'], font='Courier New')
    rrect(s, 0.65, 4.85, 12, 1.2, C['callout_bg'], line=C['accent2'])
    rect(s, 0.65, 4.85, 0.08, 1.2, C['accent2'])
    tbox(s, 0.9, 4.92, 11.5, 0.3, 'ℹ️ 現在の設計', size=12, color=C['ink'], bold=True)
    items = [
        [('データの形式', True, C['ink']),
         (': エヌイチ側で上記の構成で作っている。MIL様側の実装に合わせて調整可能（上馬場様からの回答待ち）', False, C['ink'])],
        [('企業名の形式', True, C['ink']),
         (': 現在はカンマ区切りの文字列で返している', False, C['ink'])],
    ]
    y = 5.25
    for runs in items:
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y + 0.1), Inches(0.07), Inches(0.07))
        dot.fill.solid(); dot.fill.fore_color.rgb = C['ink']; dot.line.fill.background()
        tbox_rich(s, 1.2, y, 11.2, 0.35, runs, size=10); y += 0.34
    rrect(s, 0.65, 6.2, 5.85, 0.8, C['confirm_bg'], line=C['warn'])
    tbox_rich(s, 0.85, 6.3, 5.5, 0.6,
              [('→ 確認\n', True, RGBColor(0x6B, 0x3A, 0x00), 12),
               ('データの項目名・形式がMIL様の実装に合うか', False, C['ink2'], 11)])
    rrect(s, 6.8, 6.2, 5.85, 0.8, C['confirm_bg'], line=C['warn'])
    tbox_rich(s, 7.0, 6.3, 5.5, 0.6,
              [('→ 相談\n', True, RGBColor(0x6B, 0x3A, 0x00), 12),
               ('企業名を配列とカンマ区切りのどちらで返すか', False, C['ink2'], 11)])
    page_footer(s, 31)


def slide_32_div_agenda():
    divider('※', '会議で扱う事項まとめ',
            '共有 ／ 確認 ／ 相談 の3種類に分けて整理する。',
            'AGENDA', 32)


def slide_33_share():
    s = blank(); bg(s); kicker(s, 'AGENDA / 共有')
    title_h2(s, '共有（エヌイチから三社に伝える事項）', size=24)
    data = [
        ['#', '内容'],
        ['1', '6タイプの合計点はMIL様側で計算し、Difyには合計済みの点数を送る設計になっている'],
        ['2', '嘘回答補正（Q49=5のとき×0.8）はDify側で処理する設計になっている'],
        ['3', '企業リスト・推薦理由テキストは4/10に津丸様のExcelに基づいて登録済み'],
        ['4', '中西該当の場合でもQ50が3点未満（BtoC寄り）の学生には中西製作所様を案内しない。通常判定（Step 4）に回してBtoC企業を推薦する'],
        ['5', 'APIキーとエンドポイントURLは中西製作所様のDifyアカウントから発行される。Dify有料プラン切替後にMIL様へ共有する'],
        ['6', 'フェーズ2（送信）とフェーズ4（返却）は1回の通信の行きと帰り。フェーズ4でMIL様側の追加設定は不要'],
    ]
    mktable(s, 0.65, 1.9, 12, 4.5, data, col_widths=[0.8, 11.2], font_size=12)
    page_footer(s, 33)


def slide_34_confirm():
    s = blank(); bg(s); kicker(s, 'AGENDA / 確認')
    title_h2(s, '確認（三社間で認識を合わせる事項）', size=24)
    data = [
        ['#', '内容', '確認相手'],
        ['7', 'データの項目名（R, I, A, S, E, C, Q30, Q31, Q49, Q50, Q51）がMIL様の実装と合っているか', 'MIL様'],
        ['8', '各点数の範囲（6タイプは4〜20、その他は1〜5）に認識のずれがないか', 'MIL様・中西製作所様'],
        ['9', '返却データの項目名・形式がMIL様の画面実装に合うか', 'MIL様'],
        ['10', '企業リストの内容に変更・追加がないか', '中西製作所様'],
        ['11', '偏差値化（強調フラグ）は初期運用では保留し、1,000件以降で実装する方針でよいか', '中西製作所様'],
    ]
    mktable(s, 0.65, 1.9, 12, 4.0, data, col_widths=[0.8, 8.6, 2.6], font_size=12)
    page_footer(s, 34)


def slide_35_discuss():
    s = blank(); bg(s); kicker(s, 'AGENDA / 相談')
    title_h2(s, '相談（会議で議論・決定する事項）', size=24)
    data = [
        ['#', '内容', '論点'],
        ['12', 'APIキーの共有方法・タイミング', 'メール・Slackなど。セキュリティ上の取り扱い'],
        ['13', '企業名の返却形式', '配列とカンマ区切り文字列のどちらがMIL様にとって扱いやすいか'],
        ['14', '推薦理由テキストの最終チェック', '誰がいつやるか。現状はエヌイチが作成・登録した内容のまま'],
    ]
    mktable(s, 0.65, 1.9, 12, 2.8, data, col_widths=[0.8, 4.4, 6.8], font_size=12)
    page_footer(s, 35)


def slide_36_end():
    s = blank(); bg(s, C['dark'])
    dot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(4.2), Inches(0.14), Inches(0.14))
    dot.fill.solid(); dot.fill.fore_color.rgb = C['gold_l']; dot.line.fill.background()
    tbox(s, 1.15, 4.12, 10, 0.35, 'END / 2026-04-18', size=13, color=C['paper_d2'], bold=True)
    tbox(s, 0.9, 4.55, 12, 2.2,
         'ありがとう\nございました',
         size=68, color=C['white'], bold=True)
    tbox(s, 0.9, 6.7, 12, 0.4,
         'ご質問・ご指摘は随時お願いいたします。',
         size=15, color=C['paper_d2'])


# ========== BUILD ==========

slide_01_cover()
slide_02_overview()
slide_03_dify()
slide_04_div_p1()
slide_05_q1_q24()
slide_06_6types()
slide_07_q25_q51()
slide_08_result_rule()
slide_09_points()
slide_10_11values()
slide_11_p1_confirm()
slide_12_div_p2()
slide_13_one_trip()
slide_14_required()
slide_15_send_data()
slide_16_p2_confirm()
slide_17_div_p3()
slide_18_judge_inputs()
slide_19_logic_overview()
slide_20_step0()
slide_21_step1()
slide_22_step2()
slide_23_step3()
slide_24_step4()
slide_25_company_list()
slide_26_matching_summary()
slide_27_div_p4()
slide_28_return_data()
slide_29_return_diff()
slide_30_example1()
slide_31_example2()
slide_32_div_agenda()
slide_33_share()
slide_34_confirm()
slide_35_discuss()
slide_36_end()

import os
out = '/Users/kyouyuu/claude/clients/nakanishi/TOMOSHIKA_ロジック整理.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Size: {os.path.getsize(out):,} bytes')
print(f'Slides: {len(prs.slides)}')
