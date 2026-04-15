#!/usr/bin/env python3
"""Generate TSUBOMIYA proposal PPTX — 27 slides."""

import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'tools'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
TEAL = RGBColor(0x42, 0xA4, 0xAF)
SLATE = RGBColor(0x4A, 0x6C, 0x8C)
CORAL = RGBColor(0xC9, 0x54, 0x54)
SAGE = RGBColor(0x3A, 0x8F, 0x63)
NAVY = RGBColor(0x1C, 0x35, 0x57)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF9, 0xFA, 0xFB)
HL_TEAL = RGBColor(0xEA, 0xF6, 0xF7)
HL_GREEN = RGBColor(0xF0, 0xFD, 0xF4)
HL_RED = RGBColor(0xFE, 0xF2, 0xF2)
HL_YELLOW = RGBColor(0xFE, 0xF3, 0xC7)
HL_BLUE = RGBColor(0xEF, 0xF6, 0xFF)
DIV = RGBColor(0xE5, 0xE7, 0xEB)
TITLE_C = RGBColor(0x11, 0x18, 0x27)
BODY_C = RGBColor(0x1F, 0x29, 0x37)
SUB_C = RGBColor(0x4B, 0x55, 0x63)
LABEL_C = RGBColor(0x6B, 0x72, 0x80)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_HEAD = RGBColor(0x1E, 0x40, 0x5F)
ISSUE_BAR = RGBColor(0xEF, 0x44, 0x44)
OK_BAR = RGBColor(0x16, 0xA3, 0x4A)
OK_TEXT = RGBColor(0x16, 0x6D, 0x37)
WARN_BAR = RGBColor(0xD9, 0x77, 0x06)
INFO_BAR = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_ACC = RGBColor(0xE0, 0xF2, 0xF4)

FONT = 'Noto Sans JP'
SW, SH = Inches(13.33), Inches(7.50)
Mv = 0.50  # margin
CW = 12.33

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH


# ── Utility ──

def ns():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    return slide


def rc(slide, x, y, w, h, color, rounded=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    if rounded and hasattr(s, 'adjustments') and len(s.adjustments) > 0:
        s.adjustments[0] = 0.03
    return s


def bk(slide, x, y, w, h, txt, size,
       bg=None, tc=None, bold=False, align=PP_ALIGN.LEFT,
       anchor=MSO_ANCHOR.TOP, bar=None, rounded=True,
       ml=0.12, mr=0.08, mt=0.06, mb=0.06):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    if bg:
        s.fill.solid(); s.fill.fore_color.rgb = bg
    else:
        s.fill.background()
    s.line.fill.background()
    if rounded and hasattr(s, 'adjustments') and len(s.adjustments) > 0:
        s.adjustments[0] = 0.02
    tf = s.text_frame
    tf.word_wrap = True; tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(ml); tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(mb)
    for i, line in enumerate(txt.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(1); p.space_after = Pt(1)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = tc or BODY_C; r.font.name = FONT
    if bar:
        rc(slide, x, y, 0.05, h, bar)
    return s


def bd(slide, cx, cy, r, txt, bg_c=None, tc=None):
    s = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(r * 2), Inches(r * 2))
    s.fill.solid(); s.fill.fore_color.rgb = bg_c or TEAL; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = txt
    run.font.size = Pt(max(8, int(r * 26))); run.font.bold = True
    run.font.color.rgb = tc or WHITE; run.font.name = FONT


def tb(slide, title, subtitle=None):
    rc(slide, 0, 0, 13.33, 0.06, TEAL)
    bk(slide, Mv, 0.35, 8.5, 0.55, title, 27,
       tc=TITLE_C, bold=True, rounded=False, ml=0.08)
    bk(slide, 10.50, 0.45, 2.60, 0.35, '株式会社エヌイチ', 13,
       tc=SUB_C, bold=True, align=PP_ALIGN.RIGHT, rounded=False)
    if subtitle:
        bk(slide, Mv, 0.95, 10.0, 0.35, subtitle, 13,
           tc=SUB_C, rounded=False, ml=0.08)
        rc(slide, Mv, 1.38, CW, 0.025, TEAL)
        return 1.55
    else:
        rc(slide, Mv, 0.98, CW, 0.025, TEAL)
        return 1.15


def sec_div(num, title, desc):
    slide = ns()
    rc(slide, 0, 0, 0.12, 7.50, TEAL)
    bk(slide, 1.00, 1.50, 3.00, 2.50, num, 96,
       tc=DIV, bold=True, rounded=False, ml=0.04)
    bk(slide, 1.00, 3.80, 10.0, 0.80, title, 36,
       tc=TITLE_C, bold=True, rounded=False, ml=0.04)
    bk(slide, 1.00, 4.70, 10.0, 0.50, desc, 16,
       tc=SUB_C, rounded=False, ml=0.04)
    rc(slide, 1.00, 5.40, 4.00, 0.04, TEAL)


def sh(slide, x, y, w, txt, size=13):
    rc(slide, x, y + 0.02, 0.05, 0.26, TEAL)
    bk(slide, x + 0.15, y, w - 0.15, 0.30, txt, size,
       tc=TITLE_C, bold=True, rounded=False)


# =============================================
# SLIDE 1: Cover
# =============================================
slide = ns()
rc(slide, 0, 0, 0.15, 7.50, TEAL)
rc(slide, 0.15, 0, 13.18, 0.65, DARK_HEAD)
bk(slide, 0.50, 0.12, 5.0, 0.42, '株式会社エヌイチ', 14,
   tc=WHITE, rounded=False)
bk(slide, 1.5, 2.10, 10.33, 1.10, '生成AI支援ご提案書', 48,
   tc=BLACK, bold=True, align=PP_ALIGN.CENTER, rounded=False)
bk(slide, 1.5, 3.40, 10.33, 0.55, '株式会社TSUBOMIYA 様向け', 24,
   tc=SUB_C, align=PP_ALIGN.CENTER, rounded=False)
rc(slide, 5.17, 4.15, 3.0, 0.035, TEAL)
bk(slide, 0.67, 6.50, 2.5, 0.30, '2026年4月15日', 13,
   tc=SUB_C, rounded=False)
bk(slide, 9.00, 6.45, 3.80, 0.35, '株式会社エヌイチ', 15,
   tc=BODY_C, bold=True, align=PP_ALIGN.RIGHT, rounded=False)


# =============================================
# SLIDE 2: TOC
# =============================================
slide = ns()
y = tb(slide, '目次')
sections = [
    ('01', 'お打合せ振り返り', '4/15 初回お打合せの振り返り'),
    ('02', '現状', '貴社の事業・AI活用・マーケティングの現状'),
    ('03', '課題', '成長を阻む構造要因の整理'),
    ('04', 'ゴール', '3ヶ月後・6ヶ月後の目指す姿'),
    ('05', '提案詳細', 'Kawaru Coach のご提案'),
]
for i, (num, title, desc) in enumerate(sections):
    ry = y + 0.15 + i * 1.05
    rc(slide, Mv, ry, CW, 0.90, CARD, rounded=True)
    rc(slide, Mv, ry, 0.06, 0.90, TEAL)
    bk(slide, Mv + 0.25, ry + 0.10, 0.65, 0.35, num, 18,
       bg=TEAL, tc=WHITE, bold=True, align=PP_ALIGN.CENTER,
       anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
    bk(slide, Mv + 1.10, ry + 0.10, 4.00, 0.35, title, 18,
       tc=TITLE_C, bold=True, rounded=False, ml=0.08)
    bk(slide, Mv + 1.10, ry + 0.50, 10.00, 0.30, desc, 12,
       tc=SUB_C, rounded=False, ml=0.08)


# =============================================
# SLIDE 3: Section 01 divider
# =============================================
sec_div('01', 'お打合せ振り返り', '4/15 初回お打合せの振り返り')


# =============================================
# SLIDE 4: Recap
# =============================================
slide = ns()
y = tb(slide, '4/15 初回お打合せ概要', '振り返り')

# Left panel: meeting info
lw = 5.80
rc(slide, Mv, y, lw, 5.60, WHITE, rounded=True)
rc(slide, Mv, y, 0.06, 5.60, TEAL)
sh(slide, Mv + 0.20, y + 0.15, lw - 0.40, '基本情報', 14)
info = [
    ('日時', '2026年4月15日（火）'),
    ('形式', 'オンライン（約30分）'),
    ('先方', '松尾雄太様（代表取締役）'),
    ('自社', '大川龍之介'),
]
iy = y + 0.55
for label, val in info:
    bk(slide, Mv + 0.20, iy, 0.85, 0.30, label, 10,
       bg=LIGHT_ACC, tc=TEAL, bold=True,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
    bk(slide, Mv + 1.15, iy, lw - 1.55, 0.30, val, 12,
       bg=CARD, tc=BODY_C, rounded=False, ml=0.10)
    iy += 0.40

iy += 0.15
sh(slide, Mv + 0.20, iy, lw - 0.40, 'お伺いした内容', 14)
iy += 0.40
points = [
    'リラクゼーションサロン3店舗経営、2025年12月法人化',
    'フランチャイズ展開を準備中（FC募集段階）',
    'AI顧問サービスを複数社で比較検討中',
    '「やらないとやばい」という危機感をお持ち',
    'LP・広告の内製化、本部業務の効率化にご関心',
]
for p in points:
    bk(slide, Mv + 0.20, iy, lw - 0.40, 0.30, f'●  {p}', 11,
       bg=CARD, tc=SUB_C)
    iy += 0.36

# Right panel: next actions
rx = Mv + lw + 0.25
rw = CW - lw - 0.25
rc(slide, rx, y, rw, 5.60, WHITE, rounded=True)
rc(slide, rx, y, 0.06, 5.60, SAGE)
sh(slide, rx + 0.20, y + 0.15, rw - 0.40, 'ネクストアクション', 14)

bk(slide, rx + 0.20, y + 0.60, rw - 0.40, 1.00,
   'エヌイチ\n→ 提案資料のお送り（本資料）', 13,
   bg=HL_GREEN, tc=BODY_C, bar=OK_BAR)

bk(slide, rx + 0.20, y + 1.80, rw - 0.40, 1.00,
   '松尾様\n→ 4月中にご判断、5月スタートをご希望', 13,
   bg=HL_BLUE, tc=BODY_C, bar=INFO_BAR)

bk(slide, rx + 0.20, y + 3.20, rw - 0.40, 1.60,
   '松尾様のお言葉\n\n「AIにお金をかけとかないと、\n知らないが故に20万30万\n50万払っちゃうなって思って」', 12,
   bg=HL_TEAL, tc=BODY_C, bar=TEAL)


# =============================================
# SLIDE 5: Section 02 divider
# =============================================
sec_div('02', '現状', '貴社の事業・AI活用・マーケティングの現状')


# =============================================
# SLIDE 6: Current State 1 - Business Overview
# =============================================
slide = ns()
y = tb(slide, '3店舗展開×FC募集フェーズ', '現状①　事業概要')

# Left: business card
lw = 6.50
rc(slide, Mv, y, lw, 5.60, CARD, rounded=True)
rc(slide, Mv, y, 0.06, 5.60, TEAL)

info2 = [
    ('事業内容', 'リラクゼーションサロン「利率アップ」'),
    ('店舗数', '直営3店舗'),
    ('法人化', '2025年12月（株式会社TSUBOMIYA）'),
    ('本部体制', '松尾様＋2〜3名'),
    ('FC展開', '募集中（加盟店はまだゼロ）'),
    ('その他', 'セラピスト向けコミュニティ運営中'),
]
ty = y + 0.20
for label, val in info2:
    bk(slide, 0.70, ty, 1.50, 0.35, label, 11,
       bg=LIGHT_ACC, tc=TEAL, bold=True,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ml=0.05)
    bk(slide, 2.30, ty, 4.50, 0.35, val, 13,
       tc=BODY_C, rounded=False, ml=0.10)
    ty += 0.48

ty += 0.20
bk(slide, 0.70, ty, 6.10, 0.80,
   '「フランチャイズの加盟数が伸びたら嬉しい。\nそれが僕のセンターです」', 13,
   bg=HL_TEAL, tc=BODY_C, bar=TEAL,
   anchor=MSO_ANCHOR.MIDDLE)

# Right: visual structure
rx = 7.30
rw = 5.53
rc(slide, rx, y, rw, 5.60, WHITE, rounded=True)
rc(slide, rx, y, 0.06, 5.60, SLATE)

bk(slide, rx + 0.20, y + 0.15, rw - 0.40, 0.35, '事業構造', 14,
   tc=TITLE_C, bold=True, rounded=False, ml=0.08)

boxes = [
    ('直営3店舗', 'サロン「利率アップ」\n大阪エリア', TEAL),
    ('FC展開（準備中）', 'フランチャイズ加盟店募集中\n現在加盟ゼロ', SLATE),
    ('コミュニティ', 'セラピスト向け\n先生×受講生マッチング構想', NAVY),
]
by = y + 0.70
for title, desc, color in boxes:
    rc(slide, rx + 0.20, by, rw - 0.40, 1.40, CARD, rounded=True)
    rc(slide, rx + 0.20, by, 0.06, 1.40, color)
    bk(slide, rx + 0.40, by + 0.12, rw - 0.80, 0.35, title, 14,
       tc=BODY_C, bold=True, rounded=False, ml=0.08)
    bk(slide, rx + 0.40, by + 0.52, rw - 0.80, 0.70, desc, 11,
       tc=SUB_C, rounded=False, ml=0.08)
    by += 1.55


# =============================================
# SLIDE 7: Current State 2 - AI Usage
# =============================================
slide = ns()
y = tb(slide, 'AI活用は「手探り」の段階', '現状②　AI活用状況')

# Table
tools = [
    ('ChatGPT / Gemini', 'テキスト生成（簡単な質問・調べもの）', '時々', CARD),
    ('Claude Code', 'LP制作を試行（完成度80%、モバイル崩れ）', '試行段階', HL_YELLOW),
    ('Dify / n8n 等', '未使用', '—', HL_RED),
    ('業務自動化ワークフロー', '未経験', '—', HL_RED),
]

# Header
rc(slide, Mv, y, CW, 0.40, TEAL)
bk(slide, Mv + 0.10, y, 3.00, 0.40, 'ツール', 12,
   tc=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)
bk(slide, Mv + 3.10, y, 6.00, 0.40, '活用状況', 12,
   tc=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)
bk(slide, Mv + 9.10, y, 3.23, 0.40, '頻度', 12,
   tc=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)

ty = y + 0.48
for tool, usage, freq, bg_c in tools:
    rc(slide, Mv, ty, CW, 0.50, bg_c, rounded=True)
    bk(slide, Mv + 0.10, ty, 3.00, 0.50, tool, 13,
       tc=BODY_C, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)
    bk(slide, Mv + 3.10, ty, 6.00, 0.50, usage, 12,
       tc=SUB_C, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)
    bk(slide, Mv + 9.10, ty, 3.23, 0.50, freq, 12,
       tc=SUB_C, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)
    ty += 0.58

ty += 0.30
# Quotes
bk(slide, Mv, ty, CW, 0.55,
   '「かなり疎いんですよ」—— 松尾様', 14,
   bg=HL_TEAL, tc=BODY_C, bar=TEAL, anchor=MSO_ANCHOR.MIDDLE)

ty += 0.70
bk(slide, Mv, ty, CW, 0.55,
   '「やらないとやばいなっていう感じがひしひしと伝わってくる」—— 松尾様', 14,
   bg=HL_TEAL, tc=BODY_C, bar=TEAL, anchor=MSO_ANCHOR.MIDDLE)

ty += 0.70
bk(slide, Mv, ty, CW, 0.55,
   'Claude CodeでLPを自作する行動力はある。正しいガイドがあれば加速できるポテンシャル。', 13,
   bg=HL_GREEN, tc=OK_TEXT, bar=OK_BAR, anchor=MSO_ANCHOR.MIDDLE)


# =============================================
# SLIDE 8: Current State 3 - Marketing
# =============================================
slide = ns()
y = tb(slide, 'LP・広告は外注に依存', '現状③　マーケティング体制')

items = [
    ('LP制作', '外注（1本あたり10〜50万円）', ISSUE_BAR),
    ('広告運用', '代行会社に委託（「任せきり」状態）', ISSUE_BAR),
    ('SNS・クリエイティブ', '自社で作りたいが、コストがかかるなら控えている', WARN_BAR),
    ('チラシ等', '「無料で作れるなら作る。お金かかるなら作る人減っちゃう」', WARN_BAR),
]

ty = y + 0.10
for label, desc, bar_c in items:
    rc(slide, Mv, ty, CW, 0.80, CARD, rounded=True)
    rc(slide, Mv, ty, 0.06, 0.80, bar_c)
    bk(slide, Mv + 0.25, ty + 0.08, 2.80, 0.30, label, 13,
       bg=LIGHT_ACC, tc=TEAL, bold=True,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ml=0.05)
    bk(slide, Mv + 3.20, ty + 0.08, CW - 3.40, 0.60, desc, 13,
       tc=BODY_C, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.10)
    ty += 0.92

ty += 0.20
bk(slide, Mv, ty, CW, 0.70,
   '「AIにお金をかけとかないと、知らないが故に20万30万50万払っちゃうなって思って」—— 松尾様', 14,
   bg=HL_TEAL, tc=BODY_C, bar=TEAL, anchor=MSO_ANCHOR.MIDDLE)

ty += 0.85
bk(slide, Mv, ty, CW, 0.55,
   '内製化の手段がないため、結果的に高コスト構造になっている【推察】', 13,
   bg=HL_YELLOW, tc=WARN_BAR, bar=WARN_BAR, anchor=MSO_ANCHOR.MIDDLE)


# =============================================
# SLIDE 9: Section 03 divider
# =============================================
sec_div('03', '課題', '成長を阻む構造要因の整理')


# =============================================
# SLIDE 10: Issue Map (Logic Tree)
# =============================================
slide = ns()
y = tb(slide, '成長を阻む4つの構造要因', '課題①　全体マップ')
content_h = 5.70
n_branches = 4
branch_h = content_h / n_branches

root_x, root_w = 0.40, 2.00
rc(slide, root_x, y, root_w, content_h, HL_RED, rounded=True)
rc(slide, root_x, y, 0.06, content_h, ISSUE_BAR)
bk(slide, root_x + 0.12, y + 0.10, root_w - 0.22, content_h - 0.20,
   'FC展開と\n事業成長を\n支える\n「本部機能」が\nAI活用で\n強化できて\nいない', 12,
   tc=BODY_C, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, rounded=False)

mid_x = root_x + root_w + 0.30
mid_w = 2.50
leaf_x = mid_x + mid_w + 0.30
leaf_w = 13.33 - leaf_x - 0.25

branches = [
    ('A. クリエイティブ・\nマーケが外注依存', [
        ('LP制作に10〜50万の外注コスト', True),
        ('広告運用を代行に完全委託', True),
        ('制作スピードが外注に依存', False),
    ]),
    ('B. AI活用の\n知識・スキル不足', [
        ('AI全般に「かなり疎い」と自認', True),
        ('Claude CodeでLP制作は80%止まり', True),
        ('適切なツール選定ができない', False),
    ]),
    ('C. FC展開の\n基盤が未整備', [
        ('FC加盟店ゼロ、募集段階', True),
        ('オペレーション標準化が課題', True),
        ('マニュアル・ナレッジ未整備', False),
    ]),
    ('D. 少人数本部の\nリソース制約', [
        ('本部スタッフは松尾様＋2〜3名', True),
        ('「誰でもできる作業」が時間を圧迫', True),
        ('コア業務に集中できていない', False),
    ]),
]

for i, (branch_label, leaves) in enumerate(branches):
    by = y + i * branch_h + 0.06
    bh = branch_h - 0.12
    conn_y = by + bh / 2 - 0.02
    rc(slide, root_x + root_w, conn_y, 0.30, 0.04, LABEL_C)

    rc(slide, mid_x, by, mid_w, bh, HL_BLUE, rounded=True)
    rc(slide, mid_x, by, 0.06, bh, INFO_BAR)
    bk(slide, mid_x + 0.12, by + 0.04, mid_w - 0.22, bh - 0.08,
       branch_label, 11, tc=BODY_C, bold=True,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, rounded=False)

    n_leaves = len(leaves)
    leaf_item_h = bh / n_leaves
    for j, (leaf_text, is_fact) in enumerate(leaves):
        ly = by + j * leaf_item_h + 0.03
        lh = leaf_item_h - 0.06
        conn_ly = ly + lh / 2 - 0.02
        rc(slide, mid_x + mid_w, conn_ly, 0.30, 0.04, LABEL_C)

        label_text = '事実' if is_fact else '推察'
        label_bg = LIGHT_ACC if is_fact else HL_YELLOW
        label_tc = TEAL if is_fact else WARN_BAR
        bar_c = TEAL if is_fact else WARN_BAR

        rc(slide, leaf_x, ly, leaf_w, lh, CARD, rounded=True)
        rc(slide, leaf_x, ly, 0.05, lh, bar_c)
        bk(slide, leaf_x + 0.10, ly + 0.02, 0.55, lh - 0.04,
           label_text, 8, bg=label_bg, tc=label_tc, bold=True,
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ml=0.02)
        bk(slide, leaf_x + 0.70, ly, leaf_w - 0.80, lh,
           leaf_text, 10, tc=BODY_C, rounded=False, ml=0.05,
           anchor=MSO_ANCHOR.MIDDLE)


# =============================================
# SLIDE 11: Issue 2 - Cost structure
# =============================================
slide = ns()
y = tb(slide, '外注コストが成長投資を圧迫', '課題②　コスト構造')

# Big badge
bd(slide, 11.50, y + 0.80, 0.70, 'LP1本\n10〜50\n万円', bg_c=CORAL)

# Cost table
rc(slide, Mv, y + 0.10, 8.50, 0.40, TEAL)
bk(slide, Mv + 0.10, y + 0.10, 2.80, 0.40, '外注項目', 12,
   tc=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)
bk(slide, Mv + 3.00, y + 0.10, 2.80, 0.40, '想定コスト', 12,
   tc=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)
bk(slide, Mv + 5.90, y + 0.10, 2.60, 0.40, '年間試算', 12,
   tc=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)

rows = [
    ('LP制作', '10〜50万円/本', '40〜200万円（4本）'),
    ('広告運用代行', '月額手数料（運用費の20%）', '年間数十万円〜'),
    ('バナー・SNS素材', '数万円/点', '都度発生'),
]
ty = y + 0.58
for item, cost, annual in rows:
    rc(slide, Mv, ty, 8.50, 0.48, CARD, rounded=True)
    bk(slide, Mv + 0.10, ty, 2.80, 0.48, item, 13,
       tc=BODY_C, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)
    bk(slide, Mv + 3.00, ty, 2.80, 0.48, cost, 12,
       tc=SUB_C, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)
    bk(slide, Mv + 5.90, ty, 2.60, 0.48, annual, 12,
       tc=CORAL, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)
    ty += 0.56

ty += 0.30
bk(slide, Mv, ty, CW, 0.65,
   '「AIにお金をかけとかないと、知らないが故に20万30万50万払っちゃうなって思って」\n—— 松尾様', 14,
   bg=HL_TEAL, tc=BODY_C, bar=TEAL, anchor=MSO_ANCHOR.MIDDLE)

ty += 0.80
bk(slide, Mv, ty, CW, 0.55,
   'LP1本を外注するコストで、AI顧問6ヶ月分の大半を賄える → 内製化の投資対効果は極めて高い', 14,
   bg=HL_GREEN, tc=OK_TEXT, bold=True, bar=OK_BAR, anchor=MSO_ANCHOR.MIDDLE)


# =============================================
# SLIDE 12: Issue 3 - FC standardization
# =============================================
slide = ns()
y = tb(slide, 'FC展開に「再現性の基盤」が不可欠', '課題③　FC標準化')

cards = [
    ('オペレーション標準化', '属人的。マニュアル未整備',
     'AIで業務マニュアルを\n体系化・自動更新', TEAL),
    ('集客の再現性', 'LP・広告を外注依存',
     '自社でLP制作・広告\nクリエイティブを量産', SLATE),
    ('本部の管理効率', '少人数で手作業中心',
     '売上集計・報告・連絡を\n自動化', NAVY),
]

cw = 3.70
gap = 0.57
sx = Mv + (CW - (3 * cw + 2 * gap)) / 2

for i, (title, before, after, color) in enumerate(cards):
    cx = sx + i * (cw + gap)
    # Card
    rc(slide, cx, y + 0.10, cw, 5.40, CARD, rounded=True)
    rc(slide, cx, y + 0.10, 0.06, 5.40, color)
    # Title
    bk(slide, cx + 0.15, y + 0.25, cw - 0.30, 0.45, title, 15,
       tc=BODY_C, bold=True, align=PP_ALIGN.CENTER,
       anchor=MSO_ANCHOR.MIDDLE, rounded=False)
    # Before
    bk(slide, cx + 0.15, y + 0.85, cw - 0.30, 0.30, '現状', 10,
       bg=CORAL, tc=WHITE, bold=True, align=PP_ALIGN.CENTER,
       anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
    bk(slide, cx + 0.15, y + 1.25, cw - 0.30, 0.80, before, 12,
       bg=HL_RED, tc=BODY_C, bar=ISSUE_BAR, anchor=MSO_ANCHOR.MIDDLE)
    # Arrow
    bk(slide, cx + cw / 2 - 0.20, y + 2.15, 0.40, 0.30, '▼', 16,
       tc=TEAL, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, rounded=False)
    # After
    bk(slide, cx + 0.15, y + 2.55, cw - 0.30, 0.30, 'AI活用後', 10,
       bg=SAGE, tc=WHITE, bold=True, align=PP_ALIGN.CENTER,
       anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
    bk(slide, cx + 0.15, y + 2.95, cw - 0.30, 0.80, after, 12,
       bg=HL_GREEN, tc=BODY_C, bar=OK_BAR, anchor=MSO_ANCHOR.MIDDLE)

# Bottom note
bk(slide, Mv, y + 4.20, CW, 0.55,
   'AI活用スキルを身につけることが、FC展開の基盤づくりに直結する', 14,
   bg=HL_TEAL, tc=TEAL, bold=True, bar=TEAL, anchor=MSO_ANCHOR.MIDDLE)


# =============================================
# SLIDE 13: Section 04 divider
# =============================================
sec_div('04', 'ゴール', '3ヶ月後・6ヶ月後の目指す姿')


# =============================================
# SLIDE 14: Goal 1 - Phase 1 LP
# =============================================
slide = ns()
y = tb(slide, 'Phase 1：LPを自力で作れる', '3ヶ月後の姿')

# Before/After columns
col_w = 5.80
gap = 0.73

# Before
bx = Mv
rc(slide, bx, y + 0.10, col_w, 5.40, CARD, rounded=True)
rc(slide, bx, y + 0.10, 0.06, 5.40, CORAL)
bk(slide, bx + 0.20, y + 0.25, 1.40, 0.35, 'Before', 13,
   bg=CORAL, tc=WHITE, bold=True, align=PP_ALIGN.CENTER,
   anchor=MSO_ANCHOR.MIDDLE, ml=0.03)

before_items = [
    ('LP制作', '外注で10〜50万円、納期2週間'),
    ('LP更新', '修正のたびに外注依頼・追加費用'),
    ('FC募集LP', '未着手 or 外注検討中'),
    ('バナー・SNS', '外注 or 作成を控えている'),
]
by = y + 0.80
for label, desc in before_items:
    bk(slide, bx + 0.20, by, col_w - 0.40, 0.30, label, 11,
       bg=HL_RED, tc=CORAL, bold=True, anchor=MSO_ANCHOR.MIDDLE, ml=0.08)
    bk(slide, bx + 0.20, by + 0.35, col_w - 0.40, 0.40, desc, 12,
       tc=SUB_C, rounded=False, ml=0.12, anchor=MSO_ANCHOR.MIDDLE)
    by += 0.90

# Arrow
bk(slide, Mv + col_w + 0.10, y + 2.50, 0.53, 0.50, '→', 28,
   tc=TEAL, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, rounded=False)

# After
ax = Mv + col_w + gap
rc(slide, ax, y + 0.10, col_w, 5.40, CARD, rounded=True)
rc(slide, ax, y + 0.10, 0.06, 5.40, SAGE)
bk(slide, ax + 0.20, y + 0.25, 1.40, 0.35, 'After', 13,
   bg=SAGE, tc=WHITE, bold=True, align=PP_ALIGN.CENTER,
   anchor=MSO_ANCHOR.MIDDLE, ml=0.03)

after_items = [
    ('LP制作', 'Claude Codeで自作、数時間で公開'),
    ('LP更新', '自分で即修正、A/Bテストも即座に'),
    ('FC募集LP', 'FC募集用LPを自作・公開済み'),
    ('バナー・SNS', 'AIで素材を量産、PDCAを自分で回す'),
]
ay = y + 0.80
for label, desc in after_items:
    bk(slide, ax + 0.20, ay, col_w - 0.40, 0.30, label, 11,
       bg=HL_GREEN, tc=OK_TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE, ml=0.08)
    bk(slide, ax + 0.20, ay + 0.35, col_w - 0.40, 0.40, desc, 12,
       tc=SUB_C, rounded=False, ml=0.12, anchor=MSO_ANCHOR.MIDDLE)
    ay += 0.90


# =============================================
# SLIDE 15: Goal 2 - Phase 2 Automation
# =============================================
slide = ns()
y = tb(slide, 'Phase 2：本部業務を自動化できる', '6ヶ月後の姿')

col_w = 5.80
gap = 0.73

# Before
bx = Mv
rc(slide, bx, y + 0.10, col_w, 5.40, CARD, rounded=True)
rc(slide, bx, y + 0.10, 0.06, 5.40, CORAL)
bk(slide, bx + 0.20, y + 0.25, 1.40, 0.35, 'Before', 13,
   bg=CORAL, tc=WHITE, bold=True, align=PP_ALIGN.CENTER,
   anchor=MSO_ANCHOR.MIDDLE, ml=0.03)

before2 = [
    ('売上集計', '3店舗の数字を手動集計'),
    ('口コミ対応', '1件ずつ手打ちで返信'),
    ('リマインド', 'スタッフが手動連絡'),
    ('SNS投稿', '毎回30分かけて文面作成'),
]
by = y + 0.80
for label, desc in before2:
    bk(slide, bx + 0.20, by, col_w - 0.40, 0.30, label, 11,
       bg=HL_RED, tc=CORAL, bold=True, anchor=MSO_ANCHOR.MIDDLE, ml=0.08)
    bk(slide, bx + 0.20, by + 0.35, col_w - 0.40, 0.40, desc, 12,
       tc=SUB_C, rounded=False, ml=0.12, anchor=MSO_ANCHOR.MIDDLE)
    by += 0.90

bk(slide, Mv + col_w + 0.10, y + 2.50, 0.53, 0.50, '→', 28,
   tc=TEAL, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, rounded=False)

ax = Mv + col_w + gap
rc(slide, ax, y + 0.10, col_w, 5.40, CARD, rounded=True)
rc(slide, ax, y + 0.10, 0.06, 5.40, SAGE)
bk(slide, ax + 0.20, y + 0.25, 1.40, 0.35, 'After', 13,
   bg=SAGE, tc=WHITE, bold=True, align=PP_ALIGN.CENTER,
   anchor=MSO_ANCHOR.MIDDLE, ml=0.03)

after2 = [
    ('売上集計', '毎日自動集計→レポートが届く'),
    ('口コミ対応', 'AIが返信案を自動生成→確認して投稿'),
    ('リマインド', '前日に自動でリマインド送信'),
    ('SNS投稿', '写真を入れるだけで投稿文が自動生成'),
]
ay = y + 0.80
for label, desc in after2:
    bk(slide, ax + 0.20, ay, col_w - 0.40, 0.30, label, 11,
       bg=HL_GREEN, tc=OK_TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE, ml=0.08)
    bk(slide, ax + 0.20, ay + 0.35, col_w - 0.40, 0.40, desc, 12,
       tc=SUB_C, rounded=False, ml=0.12, anchor=MSO_ANCHOR.MIDDLE)
    ay += 0.90


# =============================================
# SLIDE 16: Goal 3 - Automation Examples
# =============================================
slide = ns()
y = tb(slide, '「ここまで自分でできる」自動化の具体例', '御社で実現できること')

examples = [
    ('① LP制作・更新', 'Claude Code',
     '外注10〜50万円→自分で数時間。\nA/Bテストも即座に回せる。'),
    ('② SNS投稿文の自動生成', 'Dify + AI',
     '写真を入れるだけで\nInstagram投稿文が自動生成。'),
    ('③ 口コミ返信の自動下書き', 'n8n + AI',
     'Googleマップの口コミを自動検知\n→AIが返信案を作成。'),
    ('④ 予約リマインド自動送信', 'n8n + LINE',
     '予約前日に自動でリマインド。\nノーショー率の削減に直結。'),
    ('⑤ 売上日報の自動集計', 'n8n + スプレッドシート',
     '3店舗の売上を毎日自動集計。\nレポートがLINEに届く。'),
    ('⑥ FC向けマニュアル整備', 'Claude Code',
     'オペレーション情報をAIが整理\n→標準化マニュアルとして出力。'),
]

cw = 3.80
ch = 2.35
gap_x = 0.47
gap_y = 0.30
sx = Mv + (CW - (3 * cw + 2 * gap_x)) / 2

for i, (title, tool, desc) in enumerate(examples):
    col = i % 3
    row = i // 3
    cx = sx + col * (cw + gap_x)
    cy = y + 0.10 + row * (ch + gap_y)

    colors = [TEAL, SLATE, SAGE, TEAL, SLATE, NAVY]
    color = colors[i]

    rc(slide, cx, cy, cw, ch, CARD, rounded=True)
    rc(slide, cx, cy, 0.06, ch, color)
    bk(slide, cx + 0.20, cy + 0.12, cw - 0.35, 0.35, title, 13,
       tc=BODY_C, bold=True, rounded=False, ml=0.08)
    bk(slide, cx + 0.20, cy + 0.48, cw - 0.35, 0.28, tool, 10,
       bg=LIGHT_ACC, tc=TEAL, bold=True,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
    bk(slide, cx + 0.20, cy + 0.85, cw - 0.35, 1.30, desc, 12,
       tc=SUB_C, rounded=False, ml=0.08)


# =============================================
# SLIDE 17: Section 05 divider
# =============================================
sec_div('05', '提案詳細', 'Kawaru Coach のご提案')


# =============================================
# SLIDE 18: Comparison - Why Coach
# =============================================
slide = ns()
y = tb(slide, '「自分でできる」が最大のROI', '提案根拠　アプローチ比較')

cols = [
    ('① 外注し続ける', CORAL, [
        ('初期コスト', '低い（都度払い）'),
        ('スキル蓄積', '✕ なし'),
        ('長期コスト', '✕ 永続的に外注費'),
        ('LP制作', '外注依存が続く'),
        ('業務自動化', '対象外'),
        ('FC展開貢献', '低い'),
    ]),
    ('② ツールだけ導入', WARN_BAR, [
        ('初期コスト', '低い（月額数千円）'),
        ('スキル蓄積', '△ 自力で限界あり'),
        ('長期コスト', '△ 活用しきれず放置リスク'),
        ('LP制作', '80%で止まる'),
        ('業務自動化', '何から始めるかわからない'),
        ('FC展開貢献', '限定的'),
    ]),
    ('③ コーチ型\n（Kawaru Coach）', TEAL, [
        ('初期コスト', '中（月額20万円）'),
        ('スキル蓄積', '◎ 体系的に習得'),
        ('長期コスト', '◎ 6ヶ月後に自走'),
        ('LP制作', '100%を自力で作れる'),
        ('業務自動化', 'ロードマップ通りに実現'),
        ('FC展開貢献', '標準化・内製化に直結'),
    ]),
]

col_w = 3.60
gap_c = 0.57
sx = Mv + (CW - (3 * col_w + 2 * gap_c)) / 2

for i, (col_title, color, items) in enumerate(cols):
    cx = sx + i * (col_w + gap_c)
    # Header
    bk(slide, cx, y + 0.05, col_w, 0.55, col_title, 14,
       bg=color, tc=WHITE, bold=True, align=PP_ALIGN.CENTER,
       anchor=MSO_ANCHOR.MIDDLE, ml=0.05)
    # Items
    iy = y + 0.70
    for label, val in items:
        bk(slide, cx, iy, col_w, 0.25, label, 9,
           bg=LIGHT_ACC if i == 2 else CARD,
           tc=TEAL if i == 2 else LABEL_C,
           bold=True, anchor=MSO_ANCHOR.MIDDLE, ml=0.08)
        bk(slide, cx, iy + 0.28, col_w, 0.40, val, 12,
           bg=CARD, tc=BODY_C if i == 2 else SUB_C,
           bold=(i == 2), anchor=MSO_ANCHOR.MIDDLE, ml=0.10)
        iy += 0.73


# =============================================
# SLIDE 19: Kawaru Coach - 4 Strengths
# =============================================
slide = ns()
y = tb(slide, 'Kawaru Coach — 選ばれる4つの理由', 'サービス紹介')

strengths = [
    ('5,000名超の\nAIコンサルタント', '業種・規模・フェーズに\n応じて最適な\n担当者を選定。\n「知識が浅い」リスクなし。', TEAL),
    ('経営・事業・組織の\n3視点で伴走', 'ツールの使い方だけでなく\n「事業にどう効くか」を\n設計。局所改善に\nとどまらない。', SLATE),
    ('ゴールは「自走」\n顧問依存にしない', '契約終了後も\n自分で回せる状態を\n作ることが使命。\nマニュアル整備まで支援。', SAGE),
    ('月額制×最低3ヶ月\n解約リスクは我々が取る', '成果が出なければ\n継続理由がない設計。\n常に成果にコミット。', NAVY),
]

cw = 2.75
gap_s = 0.39
sx = Mv + (CW - (4 * cw + 3 * gap_s)) / 2

for i, (title, desc, color) in enumerate(strengths):
    cx = sx + i * (cw + gap_s)
    rc(slide, cx, y + 0.05, cw, 3.40, CARD, rounded=True)
    rc(slide, cx, y + 0.05, 0.06, 3.40, color)
    bd(slide, cx + cw - 0.35, y + 0.35, 0.22, str(i + 1), bg_c=color)
    bk(slide, cx + 0.15, y + 0.20, cw - 0.60, 0.70, title, 13,
       tc=BODY_C, bold=True, rounded=False, ml=0.05,
       anchor=MSO_ANCHOR.MIDDLE)
    bk(slide, cx + 0.15, y + 1.00, cw - 0.30, 2.30, desc, 11,
       tc=SUB_C, rounded=False, ml=0.05)

# BIRDY case study
case_y = y + 3.65
rc(slide, Mv, case_y, CW, 1.80, HL_TEAL, rounded=True)
rc(slide, Mv, case_y, 0.06, 1.80, TEAL)
bk(slide, Mv + 0.20, case_y + 0.10, 2.00, 0.30, '導入事例', 11,
   bg=TEAL, tc=WHITE, bold=True, align=PP_ALIGN.CENTER,
   anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
bk(slide, Mv + 2.40, case_y + 0.10, 4.00, 0.30, '株式会社BIRDY様', 13,
   tc=BODY_C, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)
bk(slide, Mv + 0.20, case_y + 0.50, CW - 0.40, 0.50,
   '「ITリテラシーが低い会社でも、ちゃんと導入・運用まで伴走してくれます」', 15,
   tc=BODY_C, bold=True, rounded=False, ml=0.10, anchor=MSO_ANCHOR.MIDDLE)
bk(slide, Mv + 0.20, case_y + 1.05, CW - 0.40, 0.60,
   '代表自身がAI未経験 → バックオフィス効率化・代表業務の仕組み化を実現。経営資源を成長領域にシフト。', 12,
   tc=SUB_C, rounded=False, ml=0.10)


# =============================================
# SLIDE 20: 3 Phases
# =============================================
slide = ns()
y = tb(slide, '支援の進め方（3フェーズ）', '支援プロセス')

phases = [
    ('フェーズ 01', '「何をすべきか」を\n言語化する', [
        '経営目標・業務課題・組織体制を整理',
        'AIで取り組むべき優先課題を明文化',
        '成果物：ロードマップ・優先課題リスト',
        '初期費用10万円',
    ], TEAL),
    ('フェーズ 02', 'AI顧問が隣で動き\n実装を前に進める', [
        '月2回のMTGで進捗確認・実装支援',
        'MTG外もチャットで相談対応',
        'Phase 1: LP制作 → Phase 2: 業務自動化',
        '月額20万円',
    ], SLATE),
    ('フェーズ 03', '成功体験を定着し\n自走化する', [
        '新しい自動化を自力で設計・構築',
        'ナレッジ蓄積・マニュアル整備',
        'FC展開に向けた業務基盤の完成',
        '契約終了後も自分で回せる状態に',
    ], SAGE),
]

cw = 3.50
total_w = 3 * cw + 2 * 0.70
sx = Mv + (CW - total_w) / 2

for i, (label, title, bullets, color) in enumerate(phases):
    cx = sx + i * (cw + 0.70)
    bk(slide, cx + cw / 2 - 0.60, y + 0.05, 1.20, 0.32, label, 11,
       bg=color, tc=WHITE, bold=True, align=PP_ALIGN.CENTER,
       anchor=MSO_ANCHOR.MIDDLE)
    ct = y + 0.48
    ch = 5.20
    rc(slide, cx, ct, cw, ch, WHITE, rounded=True)
    rc(slide, cx, ct, 0.06, ch, color)
    bk(slide, cx + 0.15, ct + 0.12, cw - 0.30, 0.70, title, 14,
       bg=HL_TEAL, tc=BODY_C, bold=True,
       anchor=MSO_ANCHOR.MIDDLE, bar=color)
    by = ct + 0.95
    for b in bullets:
        bk(slide, cx + 0.15, by, cw - 0.30, 0.35, f'●  {b}', 11,
           bg=CARD, tc=SUB_C)
        by += 0.42
    if i < 2:
        ax = cx + cw + 0.15
        bk(slide, ax, y + 3.00, 0.40, 0.40, '→', 22,
           tc=TEAL, bold=True, align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE, rounded=False)


# =============================================
# SLIDE 21: 6-month Roadmap
# =============================================
slide = ns()
y = tb(slide, '6ヶ月ロードマップ', '実施スケジュール')

# Phase 1
bk(slide, Mv, y + 0.05, 2.00, 0.35, 'Phase 1', 13,
   bg=TEAL, tc=WHITE, bold=True, align=PP_ALIGN.CENTER,
   anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
bk(slide, Mv + 2.10, y + 0.05, 4.00, 0.35, 'クリエイティブ内製化（1〜3ヶ月目）', 13,
   tc=BODY_C, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)

p1_months = [
    ('1ヶ月目', 'AI基礎×LP制作入門', 'Claude Codeの基本操作を習得。簡易LPを1本制作'),
    ('2ヶ月目', 'LP実践×広告クリエイティブ', 'FC募集用LPを実戦で制作。バナー・SNS素材も作成'),
    ('3ヶ月目', 'LP改善×運用自走', 'A/Bテスト・改善サイクルを自分で回せる状態に'),
]

# Header
hdr_y = y + 0.50
rc(slide, Mv, hdr_y, CW, 0.32, TEAL)
bk(slide, Mv + 0.10, hdr_y, 1.40, 0.32, '月', 10,
   tc=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.05)
bk(slide, Mv + 1.60, hdr_y, 3.60, 0.32, 'テーマ', 10,
   tc=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.05)
bk(slide, Mv + 5.30, hdr_y, 7.03, 0.32, '到達ゴール', 10,
   tc=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.05)

ry = hdr_y + 0.38
for month, theme, goal in p1_months:
    rc(slide, Mv, ry, CW, 0.45, CARD, rounded=True)
    bk(slide, Mv + 0.10, ry, 1.40, 0.45, month, 12,
       tc=TEAL, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.05)
    bk(slide, Mv + 1.60, ry, 3.60, 0.45, theme, 12,
       tc=BODY_C, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.05)
    bk(slide, Mv + 5.30, ry, 7.03, 0.45, goal, 11,
       tc=SUB_C, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.05)
    ry += 0.52

# Phase 2
ry += 0.25
bk(slide, Mv, ry, 2.00, 0.35, 'Phase 2', 13,
   bg=NAVY, tc=WHITE, bold=True, align=PP_ALIGN.CENTER,
   anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
bk(slide, Mv + 2.10, ry, 4.00, 0.35, '業務自動化（4〜6ヶ月目）', 13,
   tc=BODY_C, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.08)

ry += 0.45
rc(slide, Mv, ry, CW, 0.32, NAVY)
bk(slide, Mv + 0.10, ry, 1.40, 0.32, '月', 10,
   tc=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.05)
bk(slide, Mv + 1.60, ry, 3.60, 0.32, 'テーマ', 10,
   tc=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.05)
bk(slide, Mv + 5.30, ry, 7.03, 0.32, '到達ゴール', 10,
   tc=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.05)

p2_months = [
    ('4ヶ月目', '業務棚卸し×自動化入門', 'Dify/n8nの基本操作。最初の1ワークフロー稼働'),
    ('5ヶ月目', '複数ワークフロー構築', '口コミ返信・売上集計・リマインド等を自動化'),
    ('6ヶ月目', '自走化×FC基盤整備', '新しい自動化を自力で設計。FCマニュアル整備'),
]

ry += 0.38
for month, theme, goal in p2_months:
    rc(slide, Mv, ry, CW, 0.45, CARD, rounded=True)
    bk(slide, Mv + 0.10, ry, 1.40, 0.45, month, 12,
       tc=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.05)
    bk(slide, Mv + 1.60, ry, 3.60, 0.45, theme, 12,
       tc=BODY_C, bold=True, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.05)
    bk(slide, Mv + 5.30, ry, 7.03, 0.45, goal, 11,
       tc=SUB_C, anchor=MSO_ANCHOR.MIDDLE, rounded=False, ml=0.05)
    ry += 0.52


# =============================================
# SLIDE 22: Pricing
# =============================================
slide = ns()
y = tb(slide, '費用・契約条件', '費用')

lw = 6.20
# Price table
rc(slide, Mv, y, lw, 4.50, WHITE, rounded=True)
rc(slide, Mv, y, 0.06, 4.50, TEAL)
sh(slide, Mv + 0.20, y + 0.12, lw - 0.40, '料金プラン', 14)

price_items = [
    ('プラン', 'Kawaru Coach ライトプラン'),
    ('月額', '200,000円（税別）'),
    ('MTG頻度', '月2回（1回60分）'),
    ('MTG外サポート', 'チャット相談対応'),
    ('初期費用', '100,000円（言語化フェーズ）'),
    ('最低契約期間', '3ヶ月'),
    ('推奨契約期間', '6ヶ月（Phase 1 + Phase 2 完走）'),
    ('お支払い', '月ごと'),
]

ty = y + 0.55
for label, val in price_items:
    bk(slide, Mv + 0.20, ty, 1.80, 0.32, label, 10,
       bg=LIGHT_ACC, tc=TEAL, bold=True,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
    bk(slide, Mv + 2.10, ty, lw - 2.30, 0.32, val, 13,
       tc=BODY_C, rounded=False, ml=0.10, anchor=MSO_ANCHOR.MIDDLE)
    ty += 0.42

# Right: ROI
rx = Mv + lw + 0.25
rw = CW - lw - 0.25
rc(slide, rx, y, rw, 4.50, WHITE, rounded=True)
rc(slide, rx, y, 0.06, 4.50, SAGE)
sh(slide, rx + 0.20, y + 0.12, rw - 0.40, '投資対効果', 14)

bk(slide, rx + 0.20, y + 0.55, rw - 0.40, 0.80,
   '6ヶ月間の投資総額\n初期10万＋月額20万×6ヶ月\n= 130万円（税別）', 14,
   bg=HL_TEAL, tc=BODY_C, bold=True, bar=TEAL, anchor=MSO_ANCHOR.MIDDLE)

roi_items = [
    'LP外注費の削減：年間40〜200万円',
    '6ヶ月後に自走→以降のコンサル費用ゼロ',
    '本部業務の自動化による工数削減効果',
    'FC展開の基盤構築→事業成長を加速',
]
ry = y + 1.55
for item in roi_items:
    bk(slide, rx + 0.20, ry, rw - 0.40, 0.35, f'●  {item}', 11,
       bg=CARD, tc=SUB_C)
    ry += 0.42

ry += 0.15
bk(slide, rx + 0.20, ry, rw - 0.40, 0.80,
   'LP1〜2本分の外注費で\n6ヶ月間AIの専属コーチがつく', 14,
   bg=HL_GREEN, tc=OK_TEXT, bold=True, bar=OK_BAR,
   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# =============================================
# SLIDE 23: Next Actions
# =============================================
slide = ns()
y = tb(slide, 'ネクストアクション')

actions = [
    ('STEP 1 — ご検討・ご質問（〜4月末）',
     '本資料の内容をご確認いただき、ご不明点・ご質問があればお気軽にご連絡ください。\nオンラインでの追加打ち合わせも可能です。'),
    ('STEP 2 — 初回ヒアリング（5月第1週）',
     'ご契約後、AI活用ロードマップ策定のための詳細ヒアリングを実施。\n御社の業務・課題を深掘りし、6ヶ月のゴールを一緒に設定します。'),
    ('STEP 3 — 伴走スタート（5月中旬〜）',
     '専属コンサルタントのアサイン完了後、月2回のMTGで伴走開始。\nPhase 1のLP制作からスタートします。'),
]

n = len(actions)
avail_h = 7.30 - y - 0.10
rg = 0.18
rh = (avail_h - rg * (n - 1)) / n
step_colors = [TEAL, SLATE, NAVY]

for i, (title, desc) in enumerate(actions):
    ry = y + 0.10 + i * (rh + rg)
    rc(slide, Mv + 0.80, ry, CW - 0.80, rh, WHITE, rounded=True)
    rc(slide, Mv + 0.80, ry, 0.06, rh, step_colors[i])
    bd(slide, Mv + 0.35, ry + rh / 2, 0.26, str(i + 1), bg_c=step_colors[i])
    if i < n - 1:
        rc(slide, Mv + 0.33, ry + rh, 0.04, rg, step_colors[i])
    bk(slide, Mv + 1.10, ry + 0.10, 8.00, 0.34, title, 14,
       bg=HL_BLUE, tc=BODY_C, bold=True,
       anchor=MSO_ANCHOR.MIDDLE, bar=step_colors[i])
    bk(slide, Mv + 1.10, ry + 0.52, CW - 1.50, rh - 0.60, desc, 12,
       tc=SUB_C, rounded=False, ml=0.10)

# Contact
bk(slide, Mv, 6.85, CW, 0.35,
   'お問い合わせ：大川龍之介  Mail: r.okawa@n1-inc.co.jp  TEL: 080-2646-2420', 11,
   bg=CARD, tc=LABEL_C, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# =============================================
# SLIDE 24: Closing
# =============================================
slide = ns()
rc(slide, 0, 0, 0.15, 7.50, TEAL)
rc(slide, 0.15, 0, 13.18, 0.65, DARK_HEAD)
bk(slide, 0.50, 0.12, 5.0, 0.42, '株式会社エヌイチ', 14,
   tc=WHITE, rounded=False)
bk(slide, 1.20, 1.60, 10.93, 2.20,
   '私たちが提供するのは\n「ツール」や「研修」ではありません。\n\n'
   '貴社の未来を創る「AI人材」と、\n組織全体の「成長」です。',
   22, bg=HL_TEAL, tc=TITLE_C, bold=True,
   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bar=TEAL)
bk(slide, 2.0, 4.10, 9.33, 0.45,
   '貴社のAIパートナーとして、共に未来を創ることをお約束します。',
   14, tc=SUB_C, align=PP_ALIGN.CENTER, rounded=False)
bk(slide, 3.0, 4.80, 7.33, 0.50, '株式会社エヌイチ', 20,
   tc=BODY_C, bold=True, align=PP_ALIGN.CENTER, rounded=False)
bk(slide, 3.0, 5.30, 7.33, 0.35, 'AI（アイ）ある社会に', 14,
   tc=TEAL, align=PP_ALIGN.CENTER, rounded=False)
bk(slide, 3.50, 5.85, 6.33, 0.85,
   '大川 龍之介（Okawa Ryunosuke）\n'
   'Kawaru事業部\n'
   'TEL: 080-2646-2420  Mail: r.okawa@n1-inc.co.jp',
   10, bg=CARD, tc=LABEL_C, align=PP_ALIGN.CENTER,
   anchor=MSO_ANCHOR.MIDDLE)


# =============================================
# SLIDE 25: Ref divider
# =============================================
slide = ns()
rc(slide, 0, 0, 0.12, 7.50, TEAL)
bk(slide, 1.00, 2.80, 10.0, 0.80, '参考資料', 36,
   tc=TITLE_C, bold=True, rounded=False, ml=0.04)
bk(slide, 1.00, 3.70, 10.0, 0.50, 'Appendix', 20,
   tc=LABEL_C, rounded=False, ml=0.04)
rc(slide, 1.00, 4.40, 4.00, 0.04, TEAL)


# =============================================
# SLIDE 26: Company Overview (correct data)
# =============================================
slide = ns()
y = tb(slide, '会社概要')

lw = 6.50
rc(slide, Mv, y, lw, 5.70, CARD, rounded=True)
info3 = [
    ('会社名', '株式会社エヌイチ'),
    ('経営陣', '代表取締役CEO 奥山幸生\n取締役COO 髙橋悠\n取締役CHRO 米澤浩明'),
    ('資本金', '1,000万円'),
    ('設立', '2023年10月'),
    ('所在地', '東京都新宿区西新宿3-7-30\nフロンティアグラン西新宿12F'),
]
ty = y + 0.20
for label, val in info3:
    h = 0.70 if '\n' in val else 0.35
    bk(slide, 0.70, ty, 1.50, h, label, 11,
       bg=LIGHT_ACC, tc=TEAL, bold=True,
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ml=0.05)
    bk(slide, 2.30, ty, 4.50, h, val, 12,
       tc=BODY_C, rounded=False, ml=0.10, anchor=MSO_ANCHOR.MIDDLE)
    ty += h + 0.12

ty += 0.10
bk(slide, 0.70, ty, 6.10, 0.45,
   'Mission：新しい働く『カタチ』を創る', 13,
   bg=HL_TEAL, tc=BODY_C, bold=True, bar=TEAL,
   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
ty += 0.55
bk(slide, 0.70, ty, 6.10, 0.45,
   'Vision：AI（アイ）ある社会に', 13,
   bg=HL_TEAL, tc=TEAL, bold=True, bar=TEAL,
   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Right: Service cards
svcs = [
    ('Kawaru', '業務自動化SaaS', TEAL),
    ('Kawaru Team', 'フルカスタマイズAI研修', SLATE),
    ('Kawaru BPO', 'AI業務効率化代行', SAGE),
    ('Kawaru Coach', 'AI顧問サービス', NAVY),
]
sx, sw = 7.30, 5.53
sy = y
for name, desc, bar_c in svcs:
    rc(slide, sx, sy, sw, 1.20, CARD, rounded=True)
    rc(slide, sx, sy, 0.06, 1.20, bar_c)
    bk(slide, sx + 0.20, sy + 0.15, sw - 0.40, 0.35, name, 15,
       tc=BODY_C, bold=True, rounded=False, ml=0.08)
    bk(slide, sx + 0.20, sy + 0.55, sw - 0.40, 0.40, desc, 12,
       tc=SUB_C, rounded=False, ml=0.08)
    sy += 1.40


# =============================================
# SLIDE 27: Case Studies
# =============================================
slide = ns()
y = tb(slide, '導入実績', 'あらゆる規模の企業様に導入されています')

companies = [
    '株式会社ブロードリンク', '株式会社中西製作所',
    '株式会社ミズカラ', 'エル・ティー・エス リンク',
    'FRUOR株式会社', '株式会社BIRDY',
    '株式会社オンシナジー', 'キリングループ',
]
bw, bh, bgap = 2.85, 0.33, 0.23
for i, name in enumerate(companies):
    bx = Mv + (i % 4) * (bw + bgap)
    by = y + 0.05 + (i // 4) * (bh + 0.10)
    bk(slide, bx, by, bw, bh, name, 10,
       bg=CARD, tc=SUB_C, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

cases = [
    ('キリングループ様', 'AI未経験、活用イメージなし',
     '全員が「業務でAIを使いたい」と回答\n満足度 8.5点'),
    ('株式会社ミズカラ様', '問い合わせ対応に月283時間',
     '月85時間へ削減\n約200時間＝25営業日分の削減'),
    ('株式会社BIRDY様', '業務の属人化、営業効率の壁',
     'バックオフィス効率化\n経営資源を成長領域にシフト'),
]
cw_c, cg_c = 3.85, 0.39
ct = y + 0.90
ch = 4.60
for i, (co, issue, result) in enumerate(cases):
    cx = Mv + i * (cw_c + cg_c)
    rc(slide, cx, ct, cw_c, ch, CARD, rounded=True)
    bd(slide, cx + 0.35, ct + 0.30, 0.20, str(i + 1))
    bk(slide, cx + 0.65, ct + 0.12, cw_c - 0.85, 0.35, co, 14,
       tc=BODY_C, bold=True, rounded=False)
    iy = ct + 0.60
    bk(slide, cx + 0.12, iy, cw_c - 0.24, 0.85,
       f'課題\n{issue}', 11,
       bg=HL_RED, tc=BODY_C, bar=ISSUE_BAR)
    bk(slide, cx + cw_c / 2 - 0.20, iy + 0.90, 0.40, 0.28,
       '▼', 14, tc=TEAL, align=PP_ALIGN.CENTER,
       anchor=MSO_ANCHOR.MIDDLE, rounded=False)
    ry = iy + 1.25
    bk(slide, cx + 0.12, ry, cw_c - 0.24, 1.60,
       f'成果\n{result}', 12,
       bg=HL_GREEN, tc=BODY_C, bold=True, bar=OK_BAR)


# =============================================
# SAVE
# =============================================
out = '/Users/kyouyuu/claude/output/proposal_TSUBOMIYA_20260415.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
