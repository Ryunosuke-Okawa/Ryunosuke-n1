"""
リードインクス株式会社 提案書 — フレッシュデザイン版
slide_builder.py を使わず直接 python-pptx で書く（一時実験）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ── カラーパレット ──────────────────────────────────
TEAL      = RGBColor(0x42, 0xA4, 0xAF)
TEAL_L    = RGBColor(0xE0, 0xF4, 0xF6)
NAVY      = RGBColor(0x1E, 0x40, 0x5F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x1E, 0x29, 0x3B)
GRAY      = RGBColor(0x64, 0x74, 0x8B)
LGRAY     = RGBColor(0xF1, 0xF5, 0xF9)
MGRAY     = RGBColor(0xE2, 0xE8, 0xF0)
RED       = RGBColor(0xEF, 0x44, 0x44)
RED_L     = RGBColor(0xFE, 0xF2, 0xF2)
GREEN     = RGBColor(0x16, 0xA3, 0x4A)
GREEN_L   = RGBColor(0xF0, 0xFD, 0xF4)
BLUE      = RGBColor(0x3B, 0x82, 0xF6)
BLUE_L    = RGBColor(0xEF, 0xF6, 0xFF)
ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)
ORANGE_L  = RGBColor(0xFE, 0xF3, 0xC7)

FONT = 'Noto Sans JP'
W, H = Inches(13.33), Inches(7.50)

# ── ユーティリティ ──────────────────────────────────

def prs_new():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, x, y, w, h, fill, line=None, radius=0):
    if radius:
        s = slide.shapes.add_shape(
            5,  # ROUNDED_RECTANGLE
            Inches(x), Inches(y), Inches(w), Inches(h))
        try: s.adjustments[0] = radius
        except: pass
    else:
        s = slide.shapes.add_shape(
            1,  # RECTANGLE
            Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def txt(slide, x, y, w, h, text, size, color=BLACK, bold=False,
        align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.TOP,
        wrap=True, ml=0.08, mt=0.05):
    tb = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    tf.vertical_anchor = v_anchor
    tf.margin_left  = Inches(ml)
    tf.margin_right = Inches(0.05)
    tf.margin_top   = Inches(mt)
    tf.margin_bottom= Inches(0.03)
    for i, line in enumerate(str(text).split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = FONT
    return tb

# ── 共通ヘッダー ────────────────────────────────────

def header(slide, title, subtitle=None):
    """タイトルバー（高さ1.0"）。コンテンツ開始Y を返す。"""
    rect(slide, 0, 0, 13.33, 1.00, TEAL)
    txt(slide, 0.40, 0.12, 9.5, 0.76, title, 24, WHITE, bold=True,
        v_anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, 10.60, 0.18, 2.50, 0.60, '株式会社エヌイチ', 11, WHITE,
        align=PP_ALIGN.RIGHT, v_anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        rect(slide, 0, 1.00, 13.33, 0.42, TEAL_L)
        txt(slide, 0.40, 1.03, 12.50, 0.36, subtitle, 12, NAVY)
        return 1.52
    return 1.10

# ── カードユーティリティ ────────────────────────────

def card(slide, x, y, w, h, fill=WHITE, border=MGRAY, bar=None, radius=0.02):
    r = rect(slide, x, y, w, h, fill, border, radius)
    if bar:
        rect(slide, x, y, 0.06, h, bar)
    return r

def bullet_rows(slide, items, x, y, w, item_h, size=13,
                bg=LGRAY, bar=None, gap=0.10, text_color=BLACK):
    """items: str のリスト。各行をカードとして縦に並べる"""
    for i, text in enumerate(items):
        iy = y + i * (item_h + gap)
        card(slide, x, iy, w, item_h, bg, MGRAY, bar, 0.02)
        txt(slide, x + (0.18 if bar else 0.12), iy, w - 0.25, item_h,
            text, size, text_color, v_anchor=MSO_ANCHOR.MIDDLE)

def label_rows(slide, pairs, x, y, lw, vw, row_h=0.40, gap=0.08,
               label_bg=TEAL_L, label_color=TEAL):
    """pairs: [(label, value)] の2列表示"""
    for i, (lbl, val) in enumerate(pairs):
        iy = y + i * (row_h + gap)
        rect(slide, x, iy, lw, row_h, label_bg)
        txt(slide, x + 0.06, iy, lw - 0.08, row_h, lbl, 10,
            label_color, bold=True, v_anchor=MSO_ANCHOR.MIDDLE)
        rect(slide, x + lw + 0.05, iy, vw, row_h, WHITE, MGRAY)
        txt(slide, x + lw + 0.12, iy, vw - 0.15, row_h, val, 12,
            BLACK, v_anchor=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════
# スライド生成
# ════════════════════════════════════════════════════

prs = prs_new()

# ────────────────────────────────────────────────────
# 1. 表紙
# ────────────────────────────────────────────────────
s = blank(prs)
rect(s, 0, 0, 13.33, 7.50, NAVY)
rect(s, 0, 0, 0.30, 7.50, TEAL)
rect(s, 0.30, 5.80, 13.03, 0.04, TEAL)
txt(s, 1.20, 5.90, 11.0, 0.50, '株式会社エヌイチ', 14, WHITE, align=PP_ALIGN.RIGHT)
txt(s, 1.20, 2.00, 11.0, 1.20, '生成AI支援 ご提案書', 52, WHITE, bold=True,
    align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
txt(s, 1.20, 3.40, 11.0, 0.60, 'リードインクス株式会社 様向け', 22, TEAL_L,
    align=PP_ALIGN.CENTER)
txt(s, 1.20, 6.55, 11.0, 0.55, '2026年2月27日', 14, GRAY,
    align=PP_ALIGN.RIGHT)

# ────────────────────────────────────────────────────
# 2. 会社概要
# ────────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '会社概要')
# 左: 会社情報テーブル
card(s, 0.40, cy, 6.00, 5.90, WHITE, MGRAY, TEAL, 0.03)
info = [
    ('会社名',   '株式会社エヌイチ'),
    ('経営陣',   'CEO 奥山幸生 / COO 髙橋悠 / CHRO 米澤浩明'),
    ('資本金',   '1,000万円'),
    ('設立',     '2023年10月'),
    ('所在地',   '東京都新宿区西新宿3-7-30'),
    ('事業',     'AI活用DX支援 / AI研修 / AIスクール'),
    ('Mission',  '新しい「働くカタチ」を創る'),
    ('Vision',   'AI（アイ）ある社会に'),
]
label_rows(s, info, 0.60, cy + 0.20, 1.20, 4.40, row_h=0.52, gap=0.07)

# 右: サービスカード
svcs = [
    ('Kawaru',       '業務自動化SaaS',          TEAL,   TEAL_L),
    ('Kawaru Team',  'フルカスタマイズAI研修',   BLUE,   BLUE_L),
    ('Kawaru BPO',   'AI業務効率化代行',         GREEN,  GREEN_L),
    ('Kawaru Coach', 'AI顧問サービス',           ORANGE, ORANGE_L),
]
sx, sy = 6.85, cy
for name, desc, bar_c, bg_c in svcs:
    card(s, sx, sy, 6.10, 1.30, bg_c, MGRAY, bar_c, 0.03)
    txt(s, sx + 0.25, sy + 0.12, 5.60, 0.45, name, 16, BLACK, bold=True)
    txt(s, sx + 0.25, sy + 0.62, 5.60, 0.45, desc, 13, GRAY)
    sy += 1.45

# ────────────────────────────────────────────────────
# 3. 導入実績
# ────────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '導入実績', 'あらゆる規模の企業様に導入されています')
# 企業名バッジ
companies = [
    '株式会社ブロードリンク', '株式会社中西製作所',
    '株式会社ミズカラ', 'エル・ティー・エス リンク',
    'FRUOR株式会社', '株式会社BIRDY',
    '株式会社オンシナジー', 'キリングループ',
]
bw, bh, bgap = 2.93, 0.38, 0.22
for i, name in enumerate(companies):
    bx = 0.40 + (i % 4) * (bw + bgap)
    by = cy + 0.10 + (i // 4) * (bh + 0.12)
    rect(s, bx, by, bw, bh, LGRAY, MGRAY, 0.02)
    txt(s, bx + 0.10, by, bw - 0.10, bh, name, 11, GRAY,
        align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

# 事例カード
cases = [
    ('キリングループ様',    'AI未経験、活用イメージなし',
     '全員が「業務でAIを使いたい」と回答\n満足度 8.5点'),
    ('株式会社ミズカラ様',  '問い合わせ対応に月283時間',
     '月85時間へ削減\n約200時間＝25営業日分を削減'),
    ('株式会社BIRDY様',    '業務の属人化・営業効率の壁',
     'バックオフィス効率化\n経営資源を成長領域にシフト'),
]
cw, cg, ct = 3.92, 0.27, cy + 1.10
ch = 7.30 - ct
for i, (co, issue, result) in enumerate(cases):
    cx = 0.40 + i * (cw + cg)
    card(s, cx, ct, cw, ch, WHITE, MGRAY, TEAL, 0.03)
    txt(s, cx + 0.22, ct + 0.15, cw - 0.35, 0.38, co, 14, BLACK, bold=True)
    # 課題
    rect(s, cx + 0.15, ct + 0.62, cw - 0.30, 1.10, RED_L, None, 0.02)
    rect(s, cx + 0.15, ct + 0.62, 0.06, 1.10, RED)
    txt(s, cx + 0.30, ct + 0.68, cw - 0.50, 0.28, '課題', 10, RED, bold=True)
    txt(s, cx + 0.30, ct + 1.00, cw - 0.50, 0.58, issue, 12, BLACK)
    txt(s, cx + cw/2 - 0.15, ct + 1.80, 0.30, 0.32, '▼', 14, TEAL,
        align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    # 成果
    rect(s, cx + 0.15, ct + 2.22, cw - 0.30, 1.30, GREEN_L, None, 0.02)
    rect(s, cx + 0.15, ct + 2.22, 0.06, 1.30, GREEN)
    txt(s, cx + 0.30, ct + 2.28, cw - 0.50, 0.28, '成果', 10, GREEN, bold=True)
    txt(s, cx + 0.30, ct + 2.60, cw - 0.50, 0.78, result, 12, BLACK, bold=True)

# ────────────────────────────────────────────────────
# 4. 前回振り返り
# ────────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '前回お打合せの振り返り', '2026年2月27日 初回商談')
# 左パネル: 商談概要
card(s, 0.40, cy, 6.00, 5.90, WHITE, MGRAY, TEAL, 0.03)
txt(s, 0.65, cy + 0.15, 5.50, 0.35, '商談概要', 13, TEAL, bold=True)
overview = [
    ('先方',   'リードインクス株式会社 呉振宇氏（本部長）'),
    ('自社',   '大川龍之介・髙橋悠（COO）'),
    ('内容',   'エヌイチのサービス紹介 → 現状・ニーズのヒアリング'),
    ('先方反応', 'Kawaru Team（フルカスタマイズ研修）に高い関心'),
    ('合意',   '研修を軸に進める方向性で合意。社長稟議を経て詳細設計へ'),
]
label_rows(s, overview, 0.58, cy + 0.58, 1.20, 4.65, row_h=0.52, gap=0.08)

# 右パネル: ネクストアクション
card(s, 6.80, cy, 6.10, 5.90, WHITE, MGRAY, BLUE, 0.03)
txt(s, 7.05, cy + 0.15, 5.60, 0.35, '合意したネクストアクション', 13, BLUE, bold=True)
nas = [
    ('呉振宇氏',    '来週木曜日の社長1on1でAI研修提案について相談'),
    ('エヌイチ',    'AI開発BPOパートナーの紹介準備（社内承認後）'),
]
ny = cy + 0.65
for who, na in nas:
    rect(s, 7.00, ny, 1.20, 0.34, BLUE_L)
    txt(s, 7.05, ny, 1.15, 0.34, who, 11, BLUE, bold=True, v_anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 8.25, ny, 4.45, 0.34, LGRAY, MGRAY)
    txt(s, 8.35, ny, 4.30, 0.34, na, 12, BLACK, v_anchor=MSO_ANCHOR.MIDDLE)
    ny += 0.52

# ────────────────────────────────────────────────────
# 5. 現状の整理
# ────────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '御社の現状整理', '初回ヒアリングで確認できた事実ベースの情報')
facts = [
    'ソフトバンクグループ会社として保険会社・保険代理店向けのSaaSプロダクトを開発・提供',
    'エンジニアは両本部合計で50〜60名。多国籍チーム（日本語話者・中国語話者が混在）',
    'AI活用は個人レベルのChatGPT利用が中心。組織としてのベストプラクティスが未確立',
    'ソフトバンクの方針によりChatGPT Enterprise中心の活用が前提',
    '金融規制上、個人情報を含む本番運用は海外ベンダーへの委託不可。AIのPOCは海外と実施中',
    '3月末に親会社（ソフトバンク）へのAIプロダクト企画書提出を予定',
]
fh = (7.30 - cy - 0.20) / len(facts) - 0.10
for i, fact in enumerate(facts):
    fy = cy + 0.10 + i * (fh + 0.10)
    rect(s, 0.40, fy, 0.08, fh, TEAL)
    rect(s, 0.55, fy, 12.38, fh, LGRAY, MGRAY, 0.02)
    txt(s, 0.75, fy, 11.90, fh, fact, 14, BLACK, v_anchor=MSO_ANCHOR.MIDDLE)

# ────────────────────────────────────────────────────
# 6. 課題のロジックツリー
# ────────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '御社が直面している「3つの壁」', '現状の事実から導かれる構造的な課題')

# 根本課題（左）
card(s, 0.40, cy, 2.20, 5.90, RED_L, None, RED, 0.03)
txt(s, 0.55, cy, 1.90, 5.90, '組織として\nAIを活用・開発\nする力がなく\n顧客への\nAI価値提供が\n遅れている',
    14, BLACK, bold=True, align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

branches = [
    ('AI開発の\n方法論・知識不足', BLUE_L, BLUE, [
        ('AIエージェント開発のプロセスが社内に存在しない', True),
        ('ビジネス要件をAIタスクに分解する方法が不明', True),
        ('個人の我流利用にとどまり組織知が蓄積されない', True),
    ]),
    ('組織環境の制約', ORANGE_L, ORANGE, [
        ('ツール利用がChatGPT Enterprise限定（ソフトバンク方針）', True),
        ('多言語チーム（日中混在）による研修実施の障壁', True),
        ('金融規制により本番運用は国内ベンダー限定', True),
    ]),
    ('POCから本番実装への壁', RED_L, RED, [
        ('海外ベンダーは本番運用に使えない（法規制）', True),
        ('自社エンジニアのAI開発スキルが不足', True),
        ('保険業界特有ユースケースの実現方法が不明', False),
    ]),
]

bh_total = 5.90
branch_h = bh_total / 3

mx, mw = 2.85, 2.50
lx = 5.60

for i, (blabel, bbg, bbar, leaves) in enumerate(branches):
    by = cy + i * branch_h
    # コネクター（根本→中）
    rect(s, 2.60, by + branch_h/2 - 0.02, 0.25, 0.04, GRAY)
    # 中間ノード
    card(s, mx, by + 0.12, mw, branch_h - 0.24, bbg, None, bbar, 0.03)
    txt(s, mx + 0.18, by + 0.12, mw - 0.25, branch_h - 0.24,
        blabel, 13, BLACK, bold=True, align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    # コネクター（中→葉）
    rect(s, mx + mw, by + branch_h/2 - 0.02, 0.25, 0.04, GRAY)
    # 葉ノード
    lw = 13.33 - lx - 0.40
    leaf_h = (branch_h - 0.30) / 3
    for j, (leaf_text, is_fact) in enumerate(leaves):
        ly = by + 0.15 + j * (leaf_h + 0.00)
        lbg = LGRAY
        lbar = TEAL if is_fact else ORANGE
        label = '事実' if is_fact else '推察'
        lc = TEAL if is_fact else ORANGE
        card(s, lx, ly, lw, leaf_h - 0.05, lbg, MGRAY, lbar, 0.02)
        rect(s, lx + 0.15, ly + 0.04, 0.60, leaf_h - 0.13, TEAL_L if is_fact else ORANGE_L)
        txt(s, lx + 0.16, ly + 0.04, 0.58, leaf_h - 0.13, label, 9,
            lc, bold=True, align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
        txt(s, lx + 0.85, ly, lw - 0.95, leaf_h - 0.05,
            leaf_text, 12, BLACK, v_anchor=MSO_ANCHOR.MIDDLE)

# ────────────────────────────────────────────────────
# 7. ご提案のゴール
# ────────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'ご提案のゴール')
# メインメッセージ
rect(s, 0.40, cy, 12.53, 1.40, TEAL_L, None, 0.03)
rect(s, 0.40, cy, 0.10, 1.40, TEAL)
txt(s, 0.65, cy + 0.10, 12.10, 1.20,
    '3ヶ月で、エンジニアが「AI開発の型」を持って\n自走できる組織をつくります',
    26, BLACK, bold=True, v_anchor=MSO_ANCHOR.MIDDLE)

# 4つのポイント
points = [
    ('AI開発の方法論',     'AIエージェント開発のプロセスとビジネス要件からタスクを分解する思考法を習得'),
    ('業界特化の実践',     '保険業界のユースケース（AI-OCR・保険金自動判断）を題材にした実践演習'),
    ('環境に合わせた設計', 'ChatGPT Enterprise前提で設計。実務にそのままつながる内容'),
    ('多言語チーム対応',   '日中バイリンガル資料のオプションで、全メンバーが参加可能'),
]
ph = (7.25 - (cy + 1.60)) / len(points) - 0.08
for i, (title, desc) in enumerate(points):
    py = cy + 1.60 + i * (ph + 0.08)
    rect(s, 0.40, py, 0.08, ph, TEAL)
    card(s, 0.55, py, 12.38, ph, WHITE, MGRAY, None, 0.02)
    txt(s, 0.75, py + 0.06, 2.80, ph - 0.12, title, 14, TEAL, bold=True,
        v_anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 3.55, py + ph * 0.15, 0.02, ph * 0.70, MGRAY)
    txt(s, 3.75, py + 0.06, 9.00, ph - 0.12, desc, 13, BLACK,
        v_anchor=MSO_ANCHOR.MIDDLE)

# ────────────────────────────────────────────────────
# 8. なぜKawaru Teamか（比較）
# ────────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'なぜKawaru Teamか', '一般的なAI研修との違い')

cols = [
    ('一般的なAI研修', LGRAY,   GRAY,  BLACK, [
        '既成カリキュラムをそのまま提供',
        '自社業務・業界への応用は自己解決',
        '研修後の実践サポートなし',
        '多言語チームへの対応が困難',
    ]),
    ('Kawaru Team',   TEAL_L,  TEAL,  BLACK, [
        '御社の課題から0ベースでカリキュラム設計',
        '保険業界のユースケースを題材に構成',
        'ChatGPT Enterprise前提で実務直結',
        '日中バイリンガル対応オプションあり',
    ]),
    ('期待できる成果', GREEN_L, GREEN, BLACK, [
        'AI開発の型が組織に定着',
        '海外BPO依存を段階的に脱却',
        'エンジニアが自走してAIを実装できる',
        '3月末の企画書に間に合う投資対効果',
    ]),
]
cw = (12.53 - 0.40) / 3 - 0.15
cx = 0.40
for col_title, bg, bar, tc, items in cols:
    card(s, cx, cy, cw, 5.90, bg, None, bar, 0.03)
    txt(s, cx + 0.18, cy + 0.12, cw - 0.30, 0.55,
        col_title, 17, tc, bold=True, align=PP_ALIGN.CENTER)
    rect(s, cx + 0.15, cy + 0.72, cw - 0.30, 0.03, bar)
    ih = (5.90 - 0.90) / len(items) - 0.10
    for j, item in enumerate(items):
        iy = cy + 0.88 + j * (ih + 0.10)
        card(s, cx + 0.15, iy, cw - 0.30, ih, WHITE, MGRAY, bar, 0.02)
        txt(s, cx + 0.30, iy, cw - 0.50, ih, item, 13, tc,
            v_anchor=MSO_ANCHOR.MIDDLE, wrap=True)
    cx += cw + 0.15

# ────────────────────────────────────────────────────
# 9. 研修の進め方
# ────────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'Kawaru Teamの進め方', 'ヒアリングからフォローまで、すべて伴走します')

steps = [
    ('STEP 1', 'ヒアリング\n課題整理', TEAL, TEAL_L, [
        '現状・ゴールの深掘りヒアリング',
        '学習対象者のスキルレベル確認',
        '研修テーマ・到達目標の設定',
        'ChatGPT Enterprise環境の確認',
    ]),
    ('STEP 2', 'カリキュラム\n設計', BLUE, BLUE_L, [
        '御社専用カリキュラムを0から設計',
        '保険業界ユースケースを題材に',
        '日中バイリンガル資料の作成（オプション）',
        '平日9〜18時での実施スケジュール調整',
    ]),
    ('STEP 3', '研修実施\nフォロー', GREEN, GREEN_L, [
        '1日2〜3時間 × 3〜5日、25名以内',
        '実施中の質疑対応・進捗確認',
        '実施後フォローアップMTG',
        '助成金申請サポート',
    ]),
]
cw = (12.53 - 0.40) / 3 - 0.20
cx = 0.40
for i, (badge, title, bar, bg, items) in enumerate(steps):
    # バッジ
    rect(s, cx + cw/2 - 0.65, cy, 1.30, 0.42, bar)
    txt(s, cx + cw/2 - 0.65, cy, 1.30, 0.42, badge, 13, WHITE, bold=True,
        align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    # カード
    card(s, cx, cy + 0.48, cw, 5.42, bg, None, bar, 0.03)
    txt(s, cx + 0.18, cy + 0.60, cw - 0.30, 0.90,
        title, 18, BLACK, bold=True, align=PP_ALIGN.CENTER)
    rect(s, cx + 0.15, cy + 1.55, cw - 0.30, 0.03, bar)
    ih = (5.42 - 1.62) / len(items) - 0.10
    for j, item in enumerate(items):
        iy = cy + 0.48 + 1.65 + j * (ih + 0.10)
        card(s, cx + 0.15, iy, cw - 0.30, ih, WHITE, MGRAY, bar, 0.02)
        txt(s, cx + 0.30, iy, cw - 0.50, ih, item, 13, BLACK,
            v_anchor=MSO_ANCHOR.MIDDLE)
    # 矢印
    if i < len(steps) - 1:
        txt(s, cx + cw + 0.02, cy + 3.20, 0.18, 0.40, '›', 28, bar,
            align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE, bold=True)
    cx += cw + 0.20

# ────────────────────────────────────────────────────
# 10. 費用・助成金
# ────────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '費用・助成金活用')

# 左上: 費用
card(s, 0.40, cy, 5.90, 2.60, WHITE, MGRAY, TEAL, 0.03)
txt(s, 0.65, cy + 0.12, 1.20, 0.38, '料金', 11, TEAL, bold=True,
    align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.65, cy + 0.58, 5.40, 0.70, '1人あたり 15万円（税別）', 28, BLACK, bold=True)
txt(s, 0.65, cy + 1.35, 5.40, 1.00,
    '研修時間：10〜15時間（3〜5日間に分割）\n対象：エンジニア 25名まで（両本部交代参加可）\n形態：御社オフィスへ講師が訪問',
    12, GRAY)

# 左下: 助成金
card(s, 0.40, cy + 2.80, 5.90, 3.10, GREEN_L, None, GREEN, 0.03)
txt(s, 0.65, cy + 2.95, 2.00, 0.38, '助成金活用', 11, GREEN, bold=True,
    align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.65, cy + 3.42, 5.40, 0.60,
    '実質 4.5〜6万円 / 人', 24, GREEN, bold=True)
txt(s, 0.65, cy + 4.08, 5.40, 1.55,
    '人材開発支援助成金（事業展開等リスキリング支援コース）\n• 還元率：6〜7割（最大75%）\n• 適用条件：平日 9:00〜18:00 での実施\n• 申請サポートはエヌイチがご支援します',
    12, GRAY)

# 右: スケジュール
card(s, 6.65, cy, 6.28, 5.90, WHITE, MGRAY, BLUE, 0.03)
txt(s, 6.90, cy + 0.12, 5.80, 0.38, 'スケジュール目安', 13, BLUE, bold=True)
schedule = [
    ('今週〜来週',     '社長1on1で相談・方向性確認（呉様）'),
    ('承認後 Week1',   'ヒアリングMTG（60〜90分）設定'),
    ('承認後 Week2-3', 'カリキュラム設計・お見積り提示'),
    ('承認後1ヶ月目',  '研修キックオフ'),
]
sh = (5.90 - 0.65) / len(schedule) - 0.12
for i, (phase, act) in enumerate(schedule):
    sy2 = cy + 0.65 + i * (sh + 0.12)
    rect(s, 6.85, sy2, 2.00, sh, BLUE_L)
    txt(s, 6.90, sy2, 1.95, sh, phase, 12, BLUE, bold=True, v_anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 8.88, sy2, 3.85, sh, LGRAY, MGRAY)
    txt(s, 8.98, sy2, 3.70, sh, act, 13, BLACK, v_anchor=MSO_ANCHOR.MIDDLE)

# ────────────────────────────────────────────────────
# 11. ネクストアクション
# ────────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'ネクストアクション')

actions = [
    (TEAL,  '社長1on1での相談（呉様）',
     '来週木曜日の社長1on1で本提案をご相談いただき、方向性をご確認ください'),
    (BLUE,  '社長承認後：ヒアリングMTGの設定',
     '承認が得られたら、エヌイチとのヒアリングMTG（60〜90分）を設定します'),
    (GREEN, 'カリキュラム設計・お見積りの提示',
     'ヒアリング完了後、御社専用カリキュラム案と正式見積りをご提示します'),
    (ORANGE,'BPOパートナーご紹介（別途）',
     'AIプロダクト開発のBPO候補をエヌイチ社内で選定し、ご紹介の準備をします'),
]
ah = (7.25 - cy - 0.10) / len(actions) - 0.12
for i, (bar_c, title, desc) in enumerate(actions):
    ay = cy + 0.10 + i * (ah + 0.12)
    # 番号バッジ
    s2 = rect(s, 0.40, ay, 0.55, ah, bar_c, None, 0.03)
    txt(s, 0.40, ay, 0.55, ah, str(i+1), 22, WHITE, bold=True,
        align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
    # コネクター
    if i < len(actions) - 1:
        rect(s, 0.63, ay + ah, 0.08, 0.12, bar_c)
    # カード
    card(s, 1.10, ay, 11.83, ah, WHITE, MGRAY, bar_c, 0.03)
    txt(s, 1.32, ay + 0.10, 6.00, 0.42, title, 16, BLACK, bold=True,
        v_anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 1.25, ay + ah * 0.55, 10.50, 0.02, MGRAY)
    txt(s, 1.32, ay + ah * 0.55 + 0.05, 11.40, ah * 0.40,
        desc, 13, GRAY, v_anchor=MSO_ANCHOR.MIDDLE)

# ────────────────────────────────────────────────────
# 12. 締め
# ────────────────────────────────────────────────────
s = blank(prs)
rect(s, 0, 0, 13.33, 7.50, NAVY)
rect(s, 0, 0, 0.30, 7.50, TEAL)
rect(s, 0.30, 6.10, 13.03, 0.04, TEAL)
txt(s, 1.20, 1.40, 11.0, 1.60,
    '私たちが提供するのは\n「ツール」や「研修」ではありません。',
    26, WHITE, bold=True, align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
txt(s, 1.20, 3.20, 11.0, 1.20,
    '貴社の未来を創る「AI人材」と、\n組織全体の「成長」です。',
    24, TEAL_L, bold=True, align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
rect(s, 4.00, 4.60, 5.33, 0.04, TEAL)
txt(s, 1.20, 4.80, 11.0, 0.55, '株式会社エヌイチ', 20, WHITE,
    bold=True, align=PP_ALIGN.CENTER)
txt(s, 1.20, 5.45, 11.0, 0.45, 'AI（アイ）ある社会に', 15, TEAL_L,
    align=PP_ALIGN.CENTER)
# 連絡先
rect(s, 3.80, 6.20, 5.73, 0.80, RGBColor(0x28, 0x4E, 0x70), None, 0.03)
txt(s, 3.90, 6.28, 5.53, 0.64,
    '大川 龍之介  TEL: 080-2646-2420  r.okawa@n1-inc.co.jp',
    11, RGBColor(0xB0, 0xC8, 0xD8), align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

# ────────────────────────────────────────────────────
out = '/Users/kyouyuu/cloude/output/proposal_leadinks_fresh.pptx'
prs.save(out)
print(f'Saved: {out}  ({len(prs.slides)} slides)')
