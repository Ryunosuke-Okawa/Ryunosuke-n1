"""
リードインクス株式会社 提案書 v3
1スライド1メッセージ構成（19枚）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── カラーパレット ──────────────────────────────────
TEAL    = RGBColor(0x42, 0xA4, 0xAF)
TEAL_L  = RGBColor(0xE0, 0xF4, 0xF6)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x1E, 0x29, 0x3B)
DARK    = RGBColor(0x11, 0x18, 0x27)
NAVY    = RGBColor(0x1C, 0x35, 0x57)
GRAY    = RGBColor(0x6B, 0x72, 0x80)
LGRAY   = RGBColor(0xF1, 0xF5, 0xF9)
MGRAY   = RGBColor(0xE2, 0xE8, 0xF0)
RED     = RGBColor(0xEF, 0x44, 0x44)
RED_L   = RGBColor(0xFE, 0xF2, 0xF2)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
GREEN_L = RGBColor(0xF0, 0xFD, 0xF4)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
BLUE_L  = RGBColor(0xEF, 0xF6, 0xFF)
ORANGE  = RGBColor(0xF5, 0x9E, 0x0B)
ORANGE_L= RGBColor(0xFE, 0xF3, 0xC7)
DARK_H  = RGBColor(0x1C, 0x35, 0x57)   # dark header / closing

FONT = 'Noto Sans JP'
SW, SH = 13.33, 7.50

# ── 基本ユーティリティ ───────────────────────────────

def prs_new():
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, x, y, w, h, fill, line=None, line_w=0.75):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def txt(slide, x, y, w, h, text, size, color=BLACK, bold=False,
        align=PP_ALIGN.LEFT, va=MSO_ANCHOR.TOP, ml=0.10, mt=0.06):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = va
    tf.margin_left   = Inches(ml)
    tf.margin_right  = Inches(0.06)
    tf.margin_top    = Inches(mt)
    tf.margin_bottom = Inches(0.04)
    for i, line in enumerate(str(text).split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
        r.font.name  = FONT
    return tb

def header(slide, title, subtitle=None, dark=False):
    bg_c = DARK_H if dark else WHITE
    if dark:
        rect(slide, 0, 0, SW, SH, DARK_H)
    rect(slide, 0, 0, SW, 0.06, TEAL)
    tc = WHITE if dark else GRAY
    txt(slide, 9.80, 0.10, 3.30, 0.40, '株式会社エヌイチ', 11, tc,
        align=PP_ALIGN.RIGHT)
    title_c = WHITE if dark else DARK
    txt(slide, 0.45, 0.18, 10.50, 0.75, title, 26, title_c, bold=True,
        va=MSO_ANCHOR.MIDDLE)
    sub_c = RGBColor(0xB0, 0xC4, 0xD8) if dark else GRAY
    if subtitle:
        txt(slide, 0.45, 0.96, 12.43, 0.38, subtitle, 13, sub_c)
        rect(slide, 0.45, 1.38, 12.43, 0.03, TEAL)
        return 1.52
    rect(slide, 0.45, 0.96, 12.43, 0.03, TEAL)
    return 1.10

def card(slide, x, y, w, h, fill=WHITE, border=MGRAY, bar=None):
    r = rect(slide, x, y, w, h, fill, border)
    if bar:
        rect(slide, x, y, 0.07, h, bar)
    return r

def label_rows(slide, pairs, x, y, lw, vw, row_h=0.50, gap=0.08,
               lbg=TEAL_L, lcolor=TEAL):
    for i, (lbl, val) in enumerate(pairs):
        iy = y + i * (row_h + gap)
        rect(slide, x, iy, lw, row_h, lbg)
        txt(slide, x + 0.06, iy, lw - 0.10, row_h, lbl, 10,
            lcolor, bold=True, va=MSO_ANCHOR.MIDDLE)
        rect(slide, x + lw + 0.06, iy, vw, row_h, WHITE, MGRAY)
        txt(slide, x + lw + 0.14, iy, vw - 0.18, row_h, val, 12,
            BLACK, va=MSO_ANCHOR.MIDDLE)

def section_tag(slide, x, y, text, color=TEAL):
    rect(slide, x, y + 0.05, 0.06, 0.28, color)
    txt(slide, x + 0.12, y, 3.0, 0.38, text, 11, color, bold=True)

def num_badge(slide, cx, cy, num, bg=TEAL):
    r = 0.28
    s = slide.shapes.add_shape(9,   # OVAL
        Inches(cx - r), Inches(cy - r), Inches(r*2), Inches(r*2))
    s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = str(num)
    run.font.size = Pt(14); run.font.bold = True
    run.font.color.rgb = WHITE; run.font.name = FONT


# ════════════════════════════════════════════════════
# スライド生成
# ════════════════════════════════════════════════════

prs = prs_new()

# ───────────────────────────────────────────────────
# S1. 表紙
# ───────────────────────────────────────────────────
s = blank(prs)
rect(s, 0, 0, 0.40, SH, TEAL)
rect(s, 0.40, SH - 0.08, SW - 0.40, 0.08, TEAL)

txt(s, 1.00, 1.50, 11.80, 0.55, 'Kawaru Team ご提案書', 16, TEAL, bold=True,
    align=PP_ALIGN.CENTER)
rect(s, 3.80, 2.15, 5.73, 0.04, TEAL)
txt(s, 1.00, 2.30, 11.80, 1.20,
    'リードインクス株式会社 様', 40, DARK, bold=True,
    align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
txt(s, 1.00, 3.65, 11.80, 0.60,
    '「AIで自走できるエンジニア組織」をつくる3〜6ヶ月の伴走提案', 17, GRAY,
    align=PP_ALIGN.CENTER)

txt(s, 1.00, 6.40, 11.80, 0.45,
    '2026年3月24日　　株式会社エヌイチ', 13, GRAY,
    align=PP_ALIGN.RIGHT)

# ───────────────────────────────────────────────────
# S2. 事業内容（リードインクス様）
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '保険業界SaaS × ソフトバンクグループ企業',
            subtitle='事業内容')

info = [
    ('グループ',    'ソフトバンクグループ傘下'),
    ('事業内容',    '保険業界向けSaaSプロダクト提供'),
    ('主要顧客',    '保険会社・保険代理店'),
    ('対象本部①',  'SaaS本部（呉振宇様）：プロダクト開発中心'),
    ('対象本部②',  'SI本部：カスタマイズ開発中心'),
    ('エンジニア数', '両本部合計 約50〜60名'),
]
label_rows(s, info, 0.45, cy + 0.10, 1.70, 8.50, row_h=0.58, gap=0.10)

# 右: アクセントカード
card(s, 10.50, cy + 0.10, 2.40, 5.60, TEAL_L, MGRAY, TEAL)
txt(s, 10.60, cy + 0.30, 2.20, 0.40, 'Key Point', 11, TEAL, bold=True)
txt(s, 10.60, cy + 0.80, 2.10, 4.50,
    '保険業界への\nAI実装が\n競争上の必須要件\nになりつつある',
    14, NAVY, bold=True)

# ───────────────────────────────────────────────────
# S3. 現状① エンジニア組織の構造
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'エンジニアの約半数が中国語話者',
            subtitle='現状①　組織構造')

# 左カラム: 数値バッジ
stats = [
    ('50〜60名', '両本部エンジニア合計'),
    ('〜半数',   '中国語話者（日本在籍）'),
    ('0回',      '外部AI研修の受講歴'),
]
sx = 0.45
for i, (num, label) in enumerate(stats):
    bx = sx + i * 4.16
    card(s, bx, cy, 3.80, 2.20, LGRAY, MGRAY, TEAL)
    txt(s, bx + 0.22, cy + 0.25, 3.40, 1.00, num, 40, TEAL, bold=True,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    txt(s, bx + 0.22, cy + 1.35, 3.40, 0.55, label, 13, GRAY,
        align=PP_ALIGN.CENTER)

# 下段: 詳細
details = [
    ('チーム構成', 'SaaS本部・SI本部に分かれ、合計約50〜60名のエンジニアが在籍'),
    ('言語環境',   '約半数が中国語話者。日本在籍だが日本語非対応のメンバーが多数'),
    ('研修歴',     '外部AI研修の受講歴はゼロ。今回が初めての検討'),
    ('個人利用',   'ChatGPT・Claudeを使うエンジニアは存在するが、全員が我流'),
]
label_rows(s, details, 0.45, cy + 2.40, 1.50, 11.00, row_h=0.52, gap=0.08)

# ───────────────────────────────────────────────────
# S4. 現状② AI活用の現状と制約
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'ChatGPT Enterprise縛り・金融規制・本番運用不可という三重制約',
            subtitle='現状②　AI活用と制約環境')

rows = [
    ('社内AI利用',    '個人レベルでChatGPT・Claudeを使用（我流）',          LGRAY,  BLACK),
    ('今後の方針',    'ChatGPT Enterprise版を近日契約予定',                  BLUE_L, BLUE),
    ('資産利用許可',  'ChatGPT Enterpriseのみ（他ツールへの会社資産アップロードは不可）', ORANGE_L, ORANGE),
    ('グループ方針',  'ソフトバンクグループ方針によりGPT中心の活用が前提',   LGRAY,  BLACK),
    ('AI開発現状',    'POC段階で海外ベンダーと共同実施中',                   LGRAY,  BLACK),
    ('本番運用の壁',  '金融規制上、海外ベンダーとの本番運用は不可',           RED_L,  RED),
    ('3月末の予定',   'ソフトバンク親会社への企画書提出期限',                 ORANGE_L, ORANGE),
]
lw, vw = 1.70, 10.20
for i, (lbl, val, bg, bc) in enumerate(rows):
    iy = cy + 0.10 + i * 0.68
    rect(s, 0.45, iy, lw, 0.56, TEAL_L)
    txt(s, 0.51, iy, lw - 0.10, 0.56, lbl, 10, TEAL, bold=True,
        va=MSO_ANCHOR.MIDDLE)
    rect(s, 0.45 + lw + 0.06, iy, vw, 0.56, bg, MGRAY)
    txt(s, 0.45 + lw + 0.18, iy, vw - 0.24, 0.56, val, 12, bc,
        va=MSO_ANCHOR.MIDDLE)

# ───────────────────────────────────────────────────
# S5. 課題① 全体マップ（4層）
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '課題は「人材・プロセス・制約・組織」の4層構造',
            subtitle='課題①　課題の全体マップ')

layers = [
    ('① 人材の課題',   'AI実装できるエンジニアがいない／手法が属人化・我流',  RED,    RED_L),
    ('② プロセスの課題', 'ビジネス要件をAIタスクに分解する方法が社内にない',    ORANGE, ORANGE_L),
    ('③ 制約の課題',   '海外ベンダー依存・金融規制・ChatGPT Enterprise縛り',   BLUE,   BLUE_L),
    ('④ 組織の課題',   '言語の壁・ベストプラクティスの不在・優先順位不明確',   TEAL,   TEAL_L),
]
gap = 0.14
lh = (SH - cy - 0.30 - gap * 3) / 4
for i, (label, body, bar, bg) in enumerate(layers):
    ly = cy + i * (lh + gap)
    card(s, 0.45, ly, 12.43, lh, bg, MGRAY, bar)
    txt(s, 0.72, ly, 3.20, lh, label, 15, bar, bold=True,
        va=MSO_ANCHOR.MIDDLE)
    rect(s, 3.90, ly + lh * 0.2, 0.03, lh * 0.6, bar)
    txt(s, 4.10, ly, 8.50, lh, body, 14, BLACK,
        va=MSO_ANCHOR.MIDDLE)

# ───────────────────────────────────────────────────
# S6. 課題② 競争リスク
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'AI非対応は、競争力喪失に直結する',
            subtitle='課題②　市場リスク')

# 因果連鎖カード（横並び）
chain = [
    ('保険業界\nAIトレンド加速', BLUE_L,   BLUE),
    ('顧客からの\nPOC要望が急増', ORANGE_L, ORANGE),
    ('社内に実装\nできる人材なし', RED_L,   RED),
    ('POCが\n本番化できない', RED_L,     RED),
]
cw = 2.80
gap = 0.22
cx0 = 0.45
for i, (label, bg, bar) in enumerate(chain):
    bx = cx0 + i * (cw + gap)
    card(s, bx, cy, cw, 2.20, bg, MGRAY, bar)
    txt(s, bx + 0.22, cy, cw - 0.30, 2.20, label, 15, bar, bold=True,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    if i < len(chain) - 1:
        txt(s, bx + cw + 0.04, cy + 0.70, 0.20, 0.80, '→', 22, bar,
            bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# 下段: 詳細説明
details = [
    '保険業界のプロダクトにAI機能を組み込むことは、今や競争上の必須要件になりつつある',
    '顧客（保険会社・代理店）からAI活用のPOC要望が次々と来ているが、社内に実装できる人材がいない',
    '海外ベンダーとPOCを実施しているが、金融規制上、本番運用への移行ができない構造的矛盾がある',
    '結果として：POC案件が「本番化できない状態」のまま滞留している',
]
dy = cy + 2.40
for item in details:
    card(s, 0.45, dy, 12.43, 0.56, LGRAY, MGRAY, TEAL)
    txt(s, 0.72, dy, 11.90, 0.56, item, 12, BLACK, va=MSO_ANCHOR.MIDDLE)
    dy += 0.68

# ───────────────────────────────────────────────────
# S7. 課題③ 技術ギャップ（実際の発言）
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'AI開発の手法が、社内に存在しない',
            subtitle='課題③　技術ギャップ')

# 引用ボックス（3つ）
quotes = [
    '「ビジネス要件をAIタスクに\n分解する方法がわからない」',
    '「要件定義・基本設計はわかるが、\nAI開発の標準手法が不明」',
    '「みんな我流。統一されたベスト\nプラクティスが誰も持っていない」',
]
qw = 3.90
qgap = 0.18
for i, q in enumerate(quotes):
    qx = 0.45 + i * (qw + qgap)
    card(s, qx, cy, qw, 2.00, TEAL_L, MGRAY, TEAL)
    txt(s, qx + 0.22, cy + 0.10, qw - 0.30, 1.80, q, 13, NAVY, bold=True,
        va=MSO_ANCHOR.MIDDLE)

# ラベル
txt(s, 0.45, cy + 2.10, 4.0, 0.38, '▼ この発言が示す本質', 12, TEAL, bold=True)

# 下段: 本質の説明
essences = [
    ('エンジニア基礎力はある',      'システム開発の要件定義・設計プロセスは理解している'),
    ('AI固有の手法が空白',           'ビジネス課題→AIタスク分解→エージェント設計のプロセスがない'),
    ('見積もりすらできない状態',     '新プロダクトのAI部分の工数・コスト試算が不可能'),
    ('優先順位が不明確',             'どのユースケースから着手するかの判断基準がない'),
]
label_rows(s, essences, 0.45, cy + 2.56, 2.20, 9.90, row_h=0.52, gap=0.08)

# ───────────────────────────────────────────────────
# S8. 課題④ 構造的ボトルネック
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '3つの構造的ボトルネック',
            subtitle='課題④　構造的制約')

bottlenecks = [
    ('ChatGPT\nEnterprise縛り', '他ツールで会社資産を扱えない\n活用範囲が狭く限定される\nGPT以外のツール選定が困難', BLUE, BLUE_L),
    ('金融規制の壁',             '海外ベンダーとのPOCを本番化できない\n内製化が必須だが人材がいない\nPOC案件が永久に仮運用のまま',     RED,  RED_L),
    ('言語の壁',                 '約半数の中国語話者に日本語研修だと効果半減\nチーム内でナレッジ共有ができない\n共通言語・標準手法の浸透が阻まれる', ORANGE, ORANGE_L),
]
cw = 3.90
gap = 0.22
for i, (title, body, bar, bg) in enumerate(bottlenecks):
    bx = 0.45 + i * (cw + gap)
    card(s, bx, cy, cw, 5.40, bg, MGRAY, bar)
    txt(s, bx + 0.22, cy + 0.20, cw - 0.30, 0.80, title, 16, bar, bold=True,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    rect(s, bx + 0.30, cy + 1.10, cw - 0.60, 0.03, bar)
    txt(s, bx + 0.22, cy + 1.25, cw - 0.30, 3.90, body, 13, BLACK)

# ───────────────────────────────────────────────────
# S9. 課題⑤ タイムライン危機
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '3月末の企画書提出に間に合わないリスク',
            subtitle='課題⑤　タイムラインの危機')

# 期限バナー
rect(s, 0.45, cy, 12.43, 0.70, RED_L, RED)
txt(s, 0.45, cy, 12.43, 0.70,
    '⚠️   期限：3月末 ─ ソフトバンク親会社への企画書提出（AI推進の実行計画を提示する必要あり）',
    14, RED, bold=True, va=MSO_ANCHOR.MIDDLE, ml=0.20)

# 現状の課題チェックリスト
txt(s, 0.45, cy + 0.90, 6.0, 0.38, '現状（提出に必要なものが揃っていない）', 12, GRAY)
problems = [
    '社内にAI実装できる人材がいない',
    'どのユースケースから着手するか不明',
    'AI開発の標準プロセスが存在しない',
    '新プロダクトのAI部分の見積もりすらできない',
]
py = cy + 1.32
for p in problems:
    card(s, 0.45, py, 5.90, 0.58, RED_L, MGRAY, RED)
    txt(s, 0.72, py, 5.40, 0.58, '✗   ' + p, 12, RED, va=MSO_ANCHOR.MIDDLE)
    py += 0.70

# 右: ネクストステップ示唆
txt(s, 6.80, cy + 0.90, 5.80, 0.38, '今すぐ必要なアクション', 12, TEAL, bold=True)
actions = [
    'AI研修の内容・対象・日程を決定',
    '助成金申請のための実施計画策定',
    'エヌイチとのヒアリング（2時間）を実施',
    '社内承認を取得し着手準備を開始',
]
ay = cy + 1.32
for a in actions:
    card(s, 6.80, ay, 5.90, 0.58, GREEN_L, MGRAY, GREEN)
    txt(s, 7.07, ay, 5.40, 0.58, '✓   ' + a, 12, GREEN, va=MSO_ANCHOR.MIDDLE)
    ay += 0.70

# ───────────────────────────────────────────────────
# S10. 提案ゴール① Before / After
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '3〜6ヶ月で、AIで自走できる組織へ',
            subtitle='提案ゴール①　目指す姿')

col_w = 5.70

# Before
card(s, 0.45, cy, col_w, 5.50, RED_L, MGRAY, RED)
txt(s, 0.72, cy + 0.20, col_w - 0.40, 0.45, 'Before（現状）', 14, RED, bold=True)
befores = [
    '要件分解の方法がわからない',
    'AI開発は海外ベンダー依存',
    '各自が我流でツールを使用',
    'POCを本番化できない',
    '見積もりも優先順位もわからない',
]
by = cy + 0.80
for b in befores:
    txt(s, 0.85, by, col_w - 0.60, 0.52, '✗  ' + b, 13, RED)
    by += 0.60

# Arrow
txt(s, 6.35, cy + 2.20, 0.80, 0.80, '→', 30, TEAL, bold=True,
    align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
txt(s, 6.00, cy + 3.10, 1.50, 0.55,
    'Kawaru Team\n研修', 11, TEAL, bold=True, align=PP_ALIGN.CENTER)

# After
card(s, 7.08, cy, col_w, 5.50, GREEN_L, MGRAY, GREEN)
txt(s, 7.35, cy + 0.20, col_w - 0.40, 0.45, 'After（3〜6ヶ月後）', 14, GREEN, bold=True)
afters = [
    'ビジネス要件→AIタスク分解が自力でできる',
    'ChatGPT Enterpriseの共通プロセスが確立',
    'チーム全員が同じ標準手法を使える',
    'POC案件を社内主導で本番化できる',
    '新プロダクトのAI部分を見積もれる',
]
afy = cy + 0.80
for a in afters:
    txt(s, 7.50, afy, col_w - 0.60, 0.52, '✓  ' + a, 13, GREEN)
    afy += 0.60

# ───────────────────────────────────────────────────
# S11. 提案ゴール② 3つの成果指標
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '研修後に達成する3つの成果指標',
            subtitle='提案ゴール②　成果指標')

goals = [
    ('スキルの定着',
     'ビジネス要件→AIタスク分解→エージェント設計のプロセスを自力で回せるエンジニアが育つ',
     '"わからない"の根本解消',
     BLUE, BLUE_L),
    ('開発の標準化',
     'ChatGPT Enterprise活用を前提とした、社内共通のAI開発プロセスが確立される',
     'チームの共通言語が生まれる',
     TEAL, TEAL_L),
    ('本番化への道筋',
     '現在POC中の案件（保険金請求自動化・AI-OCR等）を社内主導で本番化できる素地ができる',
     '海外ベンダー依存からの脱却',
     GREEN, GREEN_L),
]
gw = 3.90
ggap = 0.28
for i, (title, body, tag, bar, bg) in enumerate(goals):
    gx = 0.45 + i * (gw + ggap)
    card(s, gx, cy, gw, 5.40, bg, MGRAY, bar)
    num_badge(s, gx + 0.45, cy + 0.48, i + 1, bar)
    txt(s, gx + 0.22, cy + 0.90, gw - 0.36, 0.55, title, 15, bar, bold=True)
    rect(s, gx + 0.30, cy + 1.55, gw - 0.60, 0.03, bar)
    txt(s, gx + 0.22, cy + 1.70, gw - 0.36, 2.60, body, 13, BLACK)
    # tag at bottom
    rect(s, gx + 0.22, cy + 4.60, gw - 0.44, 0.50, bar)
    txt(s, gx + 0.22, cy + 4.60, gw - 0.44, 0.50, tag, 11, WHITE,
        bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# ───────────────────────────────────────────────────
# S12. 提案根拠① 比較表
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '一般研修との6つの決定的な違い',
            subtitle='提案根拠①　差別化比較')

# Table header
col0, col1, col2 = 0.45, 3.50, 8.10
cw0, cw1, cw2 = 3.00, 4.55, 4.74
row_h = 0.68
rect(s, col0, cy, cw0, row_h, NAVY)
txt(s, col0 + 0.10, cy, cw0 - 0.20, row_h, '比較軸', 12, WHITE, bold=True,
    va=MSO_ANCHOR.MIDDLE)
rect(s, col1, cy, cw1, row_h, MGRAY)
txt(s, col1 + 0.10, cy, cw1 - 0.20, row_h, '一般的なAI研修', 12, GRAY, bold=True,
    va=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
rect(s, col2, cy, cw2, row_h, TEAL)
txt(s, col2 + 0.10, cy, cw2 - 0.20, row_h, 'Kawaru Team', 12, WHITE, bold=True,
    va=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

comparisons = [
    ('コンテンツ設計', '汎用カリキュラム（既製品）', '御社の課題から0ベースで設計'),
    ('講師',           '教育専門家',                  '自社でAIを日常開発する実務者'),
    ('内容の深さ',     '概念・ツール紹介が中心',       'ChatGPT Enterprise制約下の実践ノウハウ'),
    ('事例',           '他業界の汎用事例',             '保険業界ユースケース（AI-OCR・請求自動化）'),
    ('言語対応',       '日本語のみ',                   '多言語対応可（AIツール活用の翻訳補助）'),
    ('研修後サポート', 'なし',                         '伴走支援が可能（Kawaru Coach）'),
]
for i, (axis, gen, kaw) in enumerate(comparisons):
    ry = cy + row_h + i * (row_h + 0.05)
    bg = LGRAY if i % 2 == 0 else WHITE
    rect(s, col0, ry, cw0, row_h, bg, MGRAY)
    txt(s, col0 + 0.10, ry, cw0 - 0.20, row_h, axis, 12, NAVY, bold=True,
        va=MSO_ANCHOR.MIDDLE)
    rect(s, col1, ry, cw1, row_h, bg, MGRAY)
    txt(s, col1 + 0.14, ry, cw1 - 0.24, row_h, gen, 12, GRAY,
        va=MSO_ANCHOR.MIDDLE)
    rect(s, col2, ry, cw2, row_h, TEAL_L, MGRAY)
    txt(s, col2 + 0.14, ry, cw2 - 0.24, row_h, kaw, 12, NAVY, bold=True,
        va=MSO_ANCHOR.MIDDLE)

# ───────────────────────────────────────────────────
# S13. 提案根拠② なぜエヌイチか
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'エヌイチは「毎日AIを実装している実務者集団」',
            subtitle='提案根拠②　エヌイチが選ばれる理由')

reasons = [
    ('実務者が教える',
     '自社でAIエージェント・ツールを日常業務で開発・運用。「知識を教える」のではなく「現場で機能するノウハウ」を渡す',
     BLUE, BLUE_L),
    ('保険業界ユースケースをそのまま教材に',
     'AI-OCR・保険金請求自動化など、御社に直結するユースケースをケーススタディとして活用可能',
     TEAL, TEAL_L),
    ('ChatGPT制約下のノウハウを保有',
     'Enterprise環境での制約・回避策・活用パターンを実体験として持つ。制約内での最大活用を設計できる',
     ORANGE, ORANGE_L),
    ('中国語話者への対応が可能',
     'AIツールを活用したリアルタイム翻訳補助により、言語の壁を超えた研修設計が可能',
     GREEN, GREEN_L),
]
rw = 5.90
for i, (title, body, bar, bg) in enumerate(reasons):
    rx = 0.45 + (i % 2) * (rw + 0.23)
    ry = cy + (i // 2) * 2.85
    card(s, rx, ry, rw, 2.60, bg, MGRAY, bar)
    txt(s, rx + 0.22, ry + 0.18, rw - 0.34, 0.50, title, 15, bar, bold=True)
    txt(s, rx + 0.22, ry + 0.78, rw - 0.34, 1.60, body, 13, BLACK)

# ───────────────────────────────────────────────────
# S14. 提案詳細① 対象・カリキュラム
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'まず15〜25名でスタート・4モジュール構成',
            subtitle='提案詳細①　研修対象・カリキュラム')

# 対象
card(s, 0.45, cy, 4.20, 1.60, TEAL_L, MGRAY, TEAL)
txt(s, 0.72, cy + 0.15, 3.80, 0.45, '推奨対象者', 12, TEAL, bold=True)
txt(s, 0.72, cy + 0.65, 3.80, 0.80,
    'SaaS本部・SI本部のエンジニア\nまず15〜25名でスタート', 13, NAVY)

card(s, 4.85, cy, 7.80, 1.60, LGRAY, MGRAY, NAVY)
txt(s, 5.12, cy + 0.15, 7.40, 0.45, '横展開の設計', 12, NAVY, bold=True)
txt(s, 5.12, cy + 0.65, 7.40, 0.80,
    '中核メンバー25名を先に育成 → ナレッジ共有 → 全体展開\nいきなり全員よりも定着率が大幅に向上', 13, BLACK)

# カリキュラム（4モジュール）
txt(s, 0.45, cy + 1.80, 6.0, 0.40, '研修カリキュラム（4モジュール）', 12, NAVY, bold=True)
modules = [
    ('Module 1', 'AIエージェント設計の基礎',                   '概念・アーキテクチャ・ユースケース分類'),
    ('Module 2', 'ビジネス課題→AIタスク分解の方法',             '要件定義からAIタスクへの変換手法'),
    ('Module 3', 'ChatGPT Enterpriseを使ったAI開発の標準プロセス', '実務ワークフロー・プロンプト設計・テスト'),
    ('Module 4', '保険業界ユースケースで実践演習',             'AI-OCR・請求自動化をケーススタディに使用'),
]
mh = 1.04
for i, (mnum, mtitle, mbody) in enumerate(modules):
    mx = 0.45 + (i % 2) * 6.30
    my = cy + 2.28 + (i // 2) * (mh + 0.14)
    card(s, mx, my, 6.08, mh, LGRAY, MGRAY, TEAL)
    txt(s, mx + 0.22, my + 0.08, 1.10, 0.38, mnum, 10, TEAL, bold=True)
    txt(s, mx + 1.32, my + 0.08, 4.50, 0.38, mtitle, 13, NAVY, bold=True)
    txt(s, mx + 0.22, my + 0.52, 5.70, 0.42, mbody, 11, GRAY)

# ───────────────────────────────────────────────────
# S15. 提案詳細② 実施形式・日程
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '訪問型・平日開催・多言語対応',
            subtitle='提案詳細②　実施形式・日程')

details = [
    ('研修時間',   '10〜15時間（1日2〜3時間 × 3〜4日間）'),
    ('実施形式',   '御社オフィスへの訪問型（リアル開催）'),
    ('実施日程',   '平日9〜18時（助成金適用の必須条件）'),
    ('言語対応',   'AIツールを活用したリアルタイム翻訳補助（中国語話者対応可）'),
    ('事前準備',   'エヌイチとのヒアリング2時間でコンテンツを共同設計（無料）'),
    ('研修後',     '振り返りセッション実施・アフターフォロー相談可'),
]
label_rows(s, details, 0.45, cy + 0.10, 1.60, 10.60, row_h=0.66, gap=0.10)

# ポイント強調
txt(s, 0.45, cy + 5.40, 12.43, 0.40, '※ 平日・訪問型の実施は助成金適用の条件です。週末・オンライン対応は助成金対象外となります。',
    11, RED, bold=True)

# ───────────────────────────────────────────────────
# S16. 提案詳細③ 費用・助成金
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '助成金活用で実質 4.5〜6万円 ／ 名',
            subtitle='提案詳細③　費用・助成金シミュレーション')

# 定価 → 助成金 → 実質負担 の3カード
cards3 = [
    ('定価', '15万円/名\n× 20名 = 300万円', LGRAY, MGRAY, NAVY),
    ('助成金', '6〜7割還元\n= 180〜210万円', ORANGE_L, MGRAY, ORANGE),
    ('実質負担', '約90〜120万円\n（4.5〜6万円/名）', GREEN_L, MGRAY, GREEN),
]
c3w = 3.80
c3gap = 0.28
for i, (title, body, bg, border, tc) in enumerate(cards3):
    cx3 = 0.45 + i * (c3w + c3gap)
    card(s, cx3, cy, c3w, 2.40, bg, border, tc)
    txt(s, cx3 + 0.22, cy + 0.20, c3w - 0.36, 0.50, title, 14, tc, bold=True,
        align=PP_ALIGN.CENTER)
    txt(s, cx3 + 0.22, cy + 0.88, c3w - 0.36, 1.30, body, 16, tc, bold=True,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    if i < 2:
        txt(s, cx3 + c3w + 0.05, cy + 0.90, 0.28, 0.60, '→', 22, tc,
            bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# 助成金詳細
txt(s, 0.45, cy + 2.65, 12.43, 0.40, '助成金について', 13, NAVY, bold=True)
subsidy_rows = [
    ('種類',       '人材開発支援助成金（経済産業省・厚生労働省系）'),
    ('適用条件',   '平日9〜18時実施・訪問型・受講報告等の書類整備'),
    ('申請サポート', 'エヌイチがフルサポート（書類作成・申請手続きを含む）'),
    ('想定スケジュール', '研修実施→書類提出→2〜4ヶ月後に還付'),
]
label_rows(s, subsidy_rows, 0.45, cy + 3.10, 1.70, 10.60, row_h=0.55, gap=0.08)

# ───────────────────────────────────────────────────
# S17. 提案詳細④ カリキュラム設計フロー
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'ヒアリング2時間でカリキュラムをゼロ設計',
            subtitle='提案詳細④　カリキュラム設計フロー')

steps = [
    ('STEP 1', 'ヒアリング（約2時間）',
     '・研修対象者のスキルレベル確認\n・習得させたいスキル・ユースケースの優先順位確認\n・日程・実施環境の詳細確認',
     BLUE, BLUE_L, '無料'),
    ('STEP 2', 'コンテンツ設計（1〜2週間）',
     '・御社固有の課題に合わせたカリキュラム作成\n・保険業界ユースケースの教材化\n・ChatGPT Enterprise環境での演習設計',
     TEAL, TEAL_L, '1〜2週間'),
    ('STEP 3', '詳細見積もり・スケジュール提示',
     '・人数・時間・日程を確定した正式提案書を提出\n・助成金申請スケジュールの確認\n・ご承認後、研修実施へ',
     GREEN, GREEN_L, '確定後即着手'),
]
sw3 = 3.90
sg3 = 0.22
for i, (step, title, body, bar, bg, tag) in enumerate(steps):
    sx3 = 0.45 + i * (sw3 + sg3)
    card(s, sx3, cy, sw3, 5.40, bg, MGRAY, bar)
    txt(s, sx3 + 0.22, cy + 0.18, sw3 - 0.36, 0.40, step, 11, bar, bold=True)
    txt(s, sx3 + 0.22, cy + 0.66, sw3 - 0.36, 0.55, title, 15, NAVY, bold=True)
    rect(s, sx3 + 0.30, cy + 1.30, sw3 - 0.60, 0.03, bar)
    txt(s, sx3 + 0.22, cy + 1.50, sw3 - 0.36, 3.20, body, 13, BLACK)
    # tag at bottom
    rect(s, sx3 + 0.22, cy + 4.72, sw3 - 0.44, 0.44, bar)
    txt(s, sx3 + 0.22, cy + 4.72, sw3 - 0.44, 0.44, tag, 11, WHITE,
        bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    if i < 2:
        txt(s, sx3 + sw3 + 0.04, cy + 2.40, 0.22, 0.60, '→', 18, bar,
            bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# ───────────────────────────────────────────────────
# S18. クロージング① 4ステップ
# ───────────────────────────────────────────────────
s = blank(prs)
rect(s, 0, 0, SW, SH, DARK_H)
rect(s, 0, 0, SW, 0.06, TEAL)
txt(s, 0.45, 0.18, 10.50, 0.75,
    '4STEPのネクストアクション',
    28, WHITE, bold=True, va=MSO_ANCHOR.MIDDLE)
rect(s, 0.45, 0.96, 12.43, 0.03, TEAL)
txt(s, 9.80, 0.10, 3.30, 0.40, '株式会社エヌイチ', 11, GRAY, align=PP_ALIGN.RIGHT)

steps4 = [
    ('STEP 1', '社内承認',
     '呉様より社長への上申・承認取得',
     '3月上旬', BLUE, BLUE_L),
    ('STEP 2', 'ヒアリング',
     '研修対象者・習得スキル・日程を詳細確認\n（約2時間・無料）',
     '承認後\n1〜2週間', TEAL, TEAL_L),
    ('STEP 3', 'コンテンツ設計',
     '御社専用カリキュラム作成\n詳細費用・スケジュール提示',
     'ヒアリング後\n1〜2週間', GREEN, GREEN_L),
    ('STEP 4', '研修実施',
     '3〜4日間のフルカスタマイズ研修\n振り返り・アフターフォロー',
     'STEP3完了後', ORANGE, ORANGE_L),
]
sw4 = 2.80
sg4 = 0.22
cy4 = 1.25
for i, (step, title, body, period, bar, bg) in enumerate(steps4):
    bx4 = 0.45 + i * (sw4 + sg4)
    card(s, bx4, cy4, sw4, 5.60, bg, MGRAY, bar)
    txt(s, bx4 + 0.22, cy4 + 0.20, sw4 - 0.36, 0.40, step, 11, bar, bold=True)
    txt(s, bx4 + 0.22, cy4 + 0.72, sw4 - 0.36, 0.55, title, 17, NAVY, bold=True)
    rect(s, bx4 + 0.28, cy4 + 1.36, sw4 - 0.56, 0.03, bar)
    txt(s, bx4 + 0.22, cy4 + 1.55, sw4 - 0.36, 2.80, body, 13, BLACK)
    # period badge
    rect(s, bx4 + 0.22, cy4 + 4.85, sw4 - 0.44, 0.55, bar)
    txt(s, bx4 + 0.22, cy4 + 4.85, sw4 - 0.44, 0.55, period, 11, WHITE,
        bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    if i < 3:
        txt(s, bx4 + sw4 + 0.04, cy4 + 2.50, 0.22, 0.60, '→', 18, bar,
            bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# ───────────────────────────────────────────────────
# S19. クロージング② 本日の確認3点
# ───────────────────────────────────────────────────
s = blank(prs)
rect(s, 0, 0, SW, SH, DARK_H)
rect(s, 0, 0, SW, 0.06, TEAL)
txt(s, 0.45, 0.18, 10.50, 0.75,
    '本日確認したい3点',
    28, WHITE, bold=True, va=MSO_ANCHOR.MIDDLE)
rect(s, 0.45, 0.96, 12.43, 0.03, TEAL)
txt(s, 9.80, 0.10, 3.30, 0.40, '株式会社エヌイチ', 11, GRAY, align=PP_ALIGN.RIGHT)

confirms = [
    ('① 社内承認の見込みスケジュール',
     '社長への1on1は来週木曜のご予定でしたが、その後の見通しはいかがでしょうか？\n承認が取れ次第、すぐにヒアリング日程を設定させてください。',
     BLUE, BLUE_L),
    ('② 両本部長との連携タイミング',
     'SaaS本部（呉様）とSI本部との連携は、どの段階で行う予定でしょうか？\n両本部を合わせた対象人数が確定すれば、費用・日程設計がより精緻になります。',
     TEAL, TEAL_L),
    ('③ 中国語対応の優先度',
     '言語サポートは「あれば嬉しい」レベルか、「必須」レベルか、どちらに近いですか？\n必須の場合は、カリキュラム設計の段階から翻訳補助の仕組みを組み込みます。',
     GREEN, GREEN_L),
]
cy19 = 1.25
ch = 1.80
cgap = 0.20
for i, (title, body, bar, bg) in enumerate(confirms):
    cy_i = cy19 + i * (ch + cgap)
    card(s, 0.45, cy_i, 12.43, ch, bg, MGRAY, bar)
    txt(s, 0.72, cy_i + 0.18, 11.50, 0.50, title, 15, bar, bold=True)
    txt(s, 0.72, cy_i + 0.76, 11.50, 0.90, body, 13, BLACK)

# ════════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════════
out = '/Users/kyouyuu/cloude/output/proposal_leadinks_v3.pptx'
prs.save(out)
print(f'Saved: {out}')
