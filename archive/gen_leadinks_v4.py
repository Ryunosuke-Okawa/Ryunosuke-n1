"""
リードインクス株式会社 ご提案書 v4
骨子設計（Step3承認済み）27枚構成
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── カラーパレット ──────────────────────────────────────────────
TEAL    = RGBColor(0x42, 0xA4, 0xAF)
TEAL_L  = RGBColor(0xEA, 0xF6, 0xF7)
SLATE   = RGBColor(0x4A, 0x6C, 0x8C)
SLATE_L = RGBColor(0xEB, 0xF0, 0xF6)
CORAL   = RGBColor(0xC9, 0x54, 0x54)
CORAL_L = RGBColor(0xFD, 0xEF, 0xEF)
SAGE    = RGBColor(0x3A, 0x8F, 0x63)
SAGE_L  = RGBColor(0xED, 0xF7, 0xF2)
NAVY    = RGBColor(0x1C, 0x35, 0x57)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x11, 0x18, 0x27)
GRAY    = RGBColor(0x6B, 0x72, 0x80)
MGRAY   = RGBColor(0xE5, 0xE7, 0xEB)
CARD    = RGBColor(0xF9, 0xFA, 0xFB)

FONT = 'Noto Sans JP'
SW, SH = 13.33, 7.50

# ── ユーティリティ ───────────────────────────────────────────────

def prs_new():
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs

def blank(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill; bg.solid(); bg.fore_color.rgb = WHITE
    return sl

def rect(slide, x, y, w, h, fill, line=None, lw=0.5):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

_FSZ = {9: 11, 10: 12, 11: 13, 12: 14, 13: 15, 14: 16, 15: 17}

def txt(slide, x, y, w, h, text, size, color=DARK, bold=False,
        align=PP_ALIGN.LEFT, va=MSO_ANCHOR.TOP, ml=0.08, mt=0.05):
    size = _FSZ.get(size, size)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    tf.vertical_anchor = va
    tf.margin_left = Inches(ml); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(0.04)
    for i, line in enumerate(str(text).split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = FONT
    return tb

def hdr(slide, title, subtitle=None):
    rect(slide, 0, 0, SW, 0.06, TEAL)
    txt(slide, 9.8, 0.10, 3.3, 0.38, '株式会社エヌイチ', 11, GRAY, align=PP_ALIGN.RIGHT)
    txt(slide, 0.45, 0.18, 10.0, 0.62, title, 22, DARK, bold=True, va=MSO_ANCHOR.MIDDLE)
    if subtitle:
        txt(slide, 0.45, 0.90, 12.0, 0.30, subtitle, 11, GRAY, ml=0.05)
        rect(slide, 0.45, 1.26, SW - 0.9, 0.025, TEAL)
        return 1.42
    rect(slide, 0.45, 0.90, SW - 0.9, 0.025, TEAL)
    return 1.05

def ftr(slide):
    rect(slide, 0, SH - 0.25, SW, 0.02, MGRAY)
    txt(slide, 0.45, SH - 0.23, 4.0, 0.20, '© 株式会社エヌイチ', 9, GRAY)

# ────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ────────────────────────────────────────────────────────────────

def slide_cover(prs, client, date_str):
    sl = blank(prs)
    rect(sl, 0, 0, 0.12, SH, TEAL)
    rect(sl, 0.12, 0, SW - 0.12, 0.70, NAVY)
    txt(sl, 0.5, 0.14, 4.0, 0.42, '株式会社エヌイチ', 14, WHITE, bold=True)
    rect(sl, 0.28, 0.85, SW - 0.28, 6.40, WHITE)
    txt(sl, 1.2, 1.9,  11.0, 1.0, '生成AI支援', 40, TEAL, bold=True)
    txt(sl, 1.2, 2.85, 11.0, 1.0, 'ご提案書',   40, DARK, bold=True)
    txt(sl, 1.2, 4.0,  11.0, 0.5, f'{client} 様向け', 18, GRAY)
    rect(sl, 1.2, 4.65, 7.0, 0.03, TEAL)
    txt(sl, 1.2,  SH - 0.58, 3.5, 0.35, date_str,        12, GRAY)
    txt(sl, 9.0,  SH - 0.58, 4.0, 0.35, '株式会社エヌイチ', 12, GRAY, align=PP_ALIGN.RIGHT)


def toc_slide(prs, sections):
    sl = blank(prs)
    rect(sl, 0, 0, SW, 0.06, TEAL)
    txt(sl, 9.8, 0.10, 3.3, 0.38, '株式会社エヌイチ', 11, GRAY, align=PP_ALIGN.RIGHT)
    txt(sl, 0.45, 0.18, 5.0, 0.55, '目次', 22, DARK, bold=True, va=MSO_ANCHOR.MIDDLE)
    rect(sl, 0.45, 0.82, SW - 0.9, 0.025, TEAL)

    col_w = (SW - 0.95) / 2
    item_h = 0.88
    start_y = 1.02
    colors = [TEAL, TEAL, SLATE, SAGE, NAVY]

    for i, (num, title, desc) in enumerate(sections):
        col = i % 2; row = i // 2
        lx = 0.45 + col * (col_w + 0.05)
        ty = start_y + row * (item_h + 0.12)
        c = colors[i % len(colors)]
        rect(sl, lx, ty, col_w, item_h, CARD, MGRAY, 0.5)
        rect(sl, lx, ty, 0.05, item_h, c)
        txt(sl, lx + 0.14, ty + 0.08, 0.6, 0.28, num, 10, c, bold=True)
        txt(sl, lx + 0.14, ty + 0.34, col_w - 0.25, 0.30, title, 14, DARK, bold=True)
        if desc:
            txt(sl, lx + 0.14, ty + 0.62, col_w - 0.25, 0.22, desc, 9, GRAY)
    ftr(sl)


def section_div(prs, num, title, desc=''):
    sl = blank(prs)
    rect(sl, 0, 0, 0.55, SH, TEAL)
    txt(sl, 0.85, 1.8,  3.5,  1.0, num,   52, TEAL, bold=True)
    rect(sl, 0.85, 3.0, 9.0,  0.03, TEAL)
    txt(sl, 0.85, 3.15, 12.0, 0.8, title, 28, DARK, bold=True)
    if desc:
        txt(sl, 0.85, 4.05, 11.5, 0.4, desc, 13, GRAY)


def ref_div(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 0.55, SH, SLATE)
    txt(sl, 0.85, 2.0, 4.0,  0.8, 'Appendix', 28, SLATE, bold=True)
    rect(sl, 0.85, 2.9, 9.0, 0.03, SLATE)
    txt(sl, 0.85, 3.1, 12.0, 0.8, '参考資料',   28, DARK,  bold=True)


# ── 振り返り ────────────────────────────────────────────────────
def slide_recap(prs):
    sl = blank(prs)
    cy = hdr(sl, '前回お打ち合わせの振り返り', '振り返り')
    ftr(sl)

    info = [
        ('日時',       '2026年2月27日（初回商談）'),
        ('先方参加者', '呉振宇 氏（本部長）'),
        ('自社参加者', '大川龍之介・髙橋悠（COO）'),
        ('きっかけ',   '産業展示会での名刺交換'),
    ]
    lw = 6.2
    row_h = 0.52
    for i, (lbl, val) in enumerate(info):
        bg = TEAL_L if i % 2 == 0 else WHITE
        rect(sl, 0.45, cy + i * row_h, 1.8,      row_h, TEAL,  MGRAY)
        rect(sl, 2.25, cy + i * row_h, lw - 1.8, row_h, bg,    MGRAY)
        txt(sl, 0.52,  cy + i * row_h + 0.10, 1.65,      row_h - 0.14, lbl, 11, WHITE, bold=True)
        txt(sl, 2.35,  cy + i * row_h + 0.10, lw - 1.95, row_h - 0.14, val, 11, DARK)

    ty = cy + len(info) * row_h + 0.1
    rect(sl, 0.45, ty, lw, 1.0, CARD, MGRAY)
    rect(sl, 0.45, ty, 0.05, 1.0, TEAL)
    txt(sl, 0.58, ty + 0.06, lw - 0.2, 0.26, '話した内容', 11, TEAL, bold=True)
    txt(sl, 0.58, ty + 0.36, lw - 0.2, 0.58,
        '• Kawaru 4サービスをご紹介\n• AI活用状況・エンジニア研修ニーズをヒアリング\n• 費用・人数・スケジュール感まで具体的に確認', 10, DARK)

    rx, rw = 6.85, SW - 6.85 - 0.4

    rect(sl, rx, cy, rw, 1.45, CARD, MGRAY)
    rect(sl, rx, cy, 0.05, 1.45, SLATE)
    txt(sl, rx + 0.12, cy + 0.06, rw - 0.18, 0.26, '先方の反応・関心', 11, SLATE, bold=True)
    txt(sl, rx + 0.12, cy + 0.36, rw - 0.22, 1.0,
        '• 保険SaaSにAI機能を組み込みたい\n• エンジニアの我流活用に課題感あり\n• Kawaru Team（フルカスタマイズ研修）に強い関心\n• 社長への上申を検討中', 10, DARK)

    na_y = cy + 1.6
    na_h = SH - na_y - 0.32
    rect(sl, rx, na_y, rw, na_h, SAGE_L, SAGE, 1.0)
    txt(sl, rx + 0.12, na_y + 0.06, rw - 0.22, 0.26, '合意したネクストアクション', 11, SAGE, bold=True)
    txt(sl, rx + 0.12, na_y + 0.36, rw - 0.22, na_h - 0.46,
        '呉さん：今週木曜日の社長1on1でAI研修を上申\n大川・髙橋：提案書を作成してお届け\nBPOパートナー紹介準備（承認後）', 10, DARK)


# ── 現状① ──────────────────────────────────────────────────────
def slide_current1(prs):
    sl = blank(prs)
    cy = hdr(sl, 'エンジニア約50〜60名、半数が中国語話者', '現状①　組織・チーム構成')
    ftr(sl)

    badges = [
        ('50〜60名', '両本部\nエンジニア合計', TEAL),
        ('約半数',   '中国語話者\n（日本語非対応）', CORAL),
        ('2本部',    'SaaS本部\n＋ SI本部', SLATE),
        ('初回',     '外部AI研修\n（今回が初）', NAVY),
    ]
    bw = (SW - 0.95) / 4
    for i, (num, lbl, c) in enumerate(badges):
        lx = 0.45 + i * (bw + 0.02)
        rect(sl, lx, cy,          bw, 1.48, CARD, MGRAY)
        rect(sl, lx, cy,          bw, 0.06, c)
        txt(sl, lx, cy + 0.12,    bw, 0.60, num, 26, c, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, lx + 0.1, cy + 0.74, bw - 0.2, 0.65, lbl, 10, GRAY, align=PP_ALIGN.CENTER)

    ty = cy + 1.62
    details = [
        ('言語構成',   '日本語話者：約半数 ／ 中国語話者：約半数（国内オフィス在籍・プログラマー職）'),
        ('AI活用レベル', '個人でChatGPT・Copilot等を我流使用。組織としての共通スキル・標準はなし'),
        ('制約',       'GPT Enterprise のみ業務利用可。ソフトバンクグループの方針でGPT一択'),
        ('外部研修',   '今回が初の外部AI研修検討。これまで外部トレーニングは未受講'),
    ]
    row_h = 0.52
    for i, (lbl, val) in enumerate(details):
        bg = TEAL_L if i % 2 == 0 else WHITE
        rect(sl, 0.45, ty + i * row_h, 2.1,      row_h, SLATE, MGRAY)
        rect(sl, 2.55, ty + i * row_h, SW - 3.0, row_h, bg,    MGRAY)
        txt(sl, 0.55,  ty + i * row_h + 0.09, 1.95,      row_h - 0.13, lbl, 10, WHITE, bold=True)
        txt(sl, 2.65,  ty + i * row_h + 0.09, SW - 3.12, row_h - 0.13, val, 10, DARK)


# ── 現状② ──────────────────────────────────────────────────────
def slide_current2(prs):
    sl = blank(prs)
    cy = hdr(sl, 'AI活用はChatGPT個人利用が中心', '現状②　AI活用の実態')
    ftr(sl)

    cw = (SW - 0.95) / 4
    ch = SH - cy - 0.33
    cards = [
        (TEAL,  '使用ツール',
         'ChatGPT（個人利用）\n→ Enterprise版を近日契約予定\n個人ではClaude等も使用'),
        (SLATE, '業務利用の制限',
         '会社資産のアップロードは\nChatGPT Enterprise のみ許可\nソフトバンクグループ方針'),
        (CORAL, '組織的課題',
         '「みんな我流でやっている。\nベストプラクティスという\n知識はみんな持っていない」\n（呉さん発言）'),
        (SAGE,  '個人活用の限界',
         'Copilotなど個人コーディング補助は使用\nAIでプロダクト全体を設計・\n開発する手法は未着手'),
    ]
    for i, (c, t, b) in enumerate(cards):
        lx = 0.45 + i * (cw + 0.02)
        rect(sl, lx, cy,          cw, ch,   CARD, MGRAY)
        rect(sl, lx, cy,          cw, 0.05, c)
        txt(sl, lx + 0.1, cy + 0.12, cw - 0.18, 0.30, t, 12, c, bold=True)
        rect(sl, lx + 0.1, cy + 0.46, cw - 0.2, 0.02, MGRAY)
        txt(sl, lx + 0.1, cy + 0.54, cw - 0.18, ch - 0.64, b, 10, DARK)


# ── 現状③ ──────────────────────────────────────────────────────
def slide_current3(prs):
    sl = blank(prs)
    cy = hdr(sl, 'GPT縛り・金融規制・海外外注の三重制約', '現状③　環境・制約')
    ftr(sl)

    constraints = [
        (SLATE, 'ツール制約',
         'ChatGPT Enterprise のみ業務利用可。Claude・Cursorは個人利用のみ（会社資産アップロード不可）'),
        (CORAL, 'セキュリティ規制',
         '金融業界の個人情報保護規制により、海外ベンダーとの本番運用は困難 → 本番化は国内のみ'),
        (NAVY,  'グループ方針',
         'ソフトバンクグループのためGPT中心が会社方針（OpenAI Azure日本DCの活用を検討中）'),
        (CORAL, 'AI開発体制',
         '現在はPOCフェーズで海外ベンダーに外注中。本番運用移行に向けて内製化が急務'),
    ]

    rh = (SH - cy - 0.33) / len(constraints) - 0.05
    for i, (c, lbl, desc) in enumerate(constraints):
        ty = cy + i * (rh + 0.05)
        rect(sl, 0.45, ty, SW - 0.9, rh, CARD, MGRAY)
        rect(sl, 0.45, ty, 0.06, rh, c)
        txt(sl, 0.60, ty + 0.08, 2.2,        0.30, lbl,  12, c,    bold=True)
        txt(sl, 0.60, ty + 0.42, SW - 1.25, rh - 0.55, desc, 10, DARK)


# ── 課題① ロジックツリー ───────────────────────────────────────
def slide_issue_tree(prs):
    sl = blank(prs)
    cy = hdr(sl, '課題は4層構造', '課題①　全体マップ')
    ftr(sl)

    content_h = SH - cy - 0.33

    # Root
    rw, rh = 2.1, 1.05
    rl = 0.4
    rt = cy + (content_h - rh) / 2
    rect(sl, rl, rt, rw, rh, CORAL)
    txt(sl, rl + 0.08, rt + 0.08, rw - 0.16, rh - 0.16,
        'AI内製化が進まず\n顧客への価値提供が\nPOC止まり', 10, WHITE, bold=True)
    # Horizontal connector
    rect(sl, rl + rw, rt + rh / 2 - 0.01, 0.45, 0.02, MGRAY)

    factors = [
        ('A', 'AI開発プロセスが\n社内に存在しない',    CORAL),
        ('B', '我流活用・\nベストプラクティスなし',     SLATE),
        ('C', '多国籍チームへの\n教育が困難',           SLATE),
        ('D', '規制・方針による\nツール制約',           NAVY),
    ]
    fac_l = rl + rw + 0.45
    fac_w = 2.25
    fac_h = 0.82
    gap = (content_h - len(factors) * fac_h) / (len(factors) + 1)

    # Vertical trunk
    trunk_top = cy + gap + fac_h / 2
    trunk_bot = cy + gap + (len(factors) - 1) * (fac_h + gap) + fac_h / 2
    rect(sl, fac_l - 0.22, trunk_top, 0.02, trunk_bot - trunk_top, MGRAY)

    for i, (key, lbl, c) in enumerate(factors):
        ft = cy + gap + i * (fac_h + gap)
        rect(sl, fac_l - 0.22, ft + fac_h / 2 - 0.01, 0.22, 0.02, MGRAY)
        rect(sl, fac_l, ft, fac_w, fac_h, CARD, c, 1.0)
        rect(sl, fac_l, ft, 0.06, fac_h, c)
        txt(sl, fac_l + 0.14, ft + 0.04, fac_w - 0.22, 0.24, key, 11, c, bold=True)
        txt(sl, fac_l + 0.14, ft + 0.30, fac_w - 0.22, fac_h - 0.36, lbl, 10, DARK, bold=True)
        rect(sl, fac_l + fac_w, ft + fac_h / 2 - 0.01, 0.32, 0.02, MGRAY)

    # Sub items
    item_l = fac_l + fac_w + 0.32
    item_w = SW - item_l - 0.3
    items_data = [
        ['要件→AIタスク落とし込みが不明【事実】',       'エージェント設計ノウハウがゼロ【事実】'],
        ['全員が我流でChatGPT活用【事実】',             'ベストプラクティス未確立【事実】'],
        ['エンジニア約半数が日本語非対応【事実】',       '翻訳ツールでは研修品質が低下【推察】'],
        ['GPT Enterprise縛り（SB方針）【事実】',        '金融規制で海外BPO本番利用不可【事実】'],
    ]

    for i, items in enumerate(items_data):
        ft = cy + gap + i * (fac_h + gap)
        sub_h = (fac_h - 0.06) / len(items)
        for j, item in enumerate(items):
            st = ft + j * (sub_h + 0.03)
            tag = '事実' if '事実' in item else '推察'
            tag_c = CORAL if tag == '事実' else SLATE
            disp = item.replace('【事実】', '').replace('【推察】', '')
            rect(sl, item_l, st, item_w, sub_h - 0.02, CARD, MGRAY, 0.3)
            rect(sl, item_l, st, 0.04, sub_h - 0.02, tag_c)
            txt(sl, item_l + 0.1, st + 0.04, item_w - 0.16, sub_h - 0.1, disp, 9, DARK)


# ── 課題② ──────────────────────────────────────────────────────
def slide_issue2(prs):
    sl = blank(prs)
    cy = hdr(sl, 'ビジネス要件をAIタスクに変換できない', '課題②　AI開発プロセスの不在')
    ftr(sl)

    q = ('「一般的なプログラムは、どうやればいいかわかるんですけど。'
         'AIの世界に飛び込むと、AIで人間の代わりに判断するために'
         '何をやらないといけないのか、そういう知識がないんです」\n― 呉振宇 氏（本部長）')
    rect(sl, 0.45, cy, SW - 0.9, 0.98, TEAL_L, TEAL, 0.8)
    txt(sl, 0.60, cy + 0.08, SW - 1.2, 0.83, q, 10, DARK, ml=0.12)

    comp_y = cy + 1.12
    comp_h = SH - comp_y - 0.33
    half_w = (SW - 1.1) / 2

    # LEFT: Traditional
    rect(sl, 0.45, comp_y, half_w, comp_h, CARD, MGRAY)
    txt(sl, 0.45, comp_y + 0.05, half_w, 0.32, '従来の開発手法（経験あり）',
        12, SLATE, bold=True, align=PP_ALIGN.CENTER)
    old_steps = ['要件定義', '基本設計', '詳細設計', 'コーディング', 'テスト・リリース']
    step_h = (comp_h - 0.48) / len(old_steps)
    for j, s in enumerate(old_steps):
        rect(sl, 0.65, comp_y + 0.44 + j * (step_h + 0.04), half_w - 0.4, step_h, SLATE)
        txt(sl, 0.65, comp_y + 0.44 + j * (step_h + 0.04), half_w - 0.4, step_h,
            s, 11, WHITE, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

    # Arrow
    txt(sl, half_w + 0.45 + 0.05, comp_y + comp_h / 2 - 0.25, 0.4, 0.5,
        '→', 22, TEAL, bold=True, align=PP_ALIGN.CENTER)

    # RIGHT: AI development
    rx = 0.5 + half_w
    rect(sl, rx, comp_y, half_w, comp_h, CARD, CORAL, 1.0)
    txt(sl, rx, comp_y + 0.05, half_w, 0.32, 'AI開発手法（ノウハウなし）',
        12, CORAL, bold=True, align=PP_ALIGN.CENTER)
    ai_steps = [
        ('ゴール設定',      CORAL, '← ？'),
        ('タスク分解',      CORAL, '← ？'),
        ('エージェント設計', CORAL, '← ？'),
        ('評価・改善',      SAGE,  ''),
        ('本番化',          SAGE,  ''),
    ]
    for j, (s, c, note) in enumerate(ai_steps):
        rect(sl, rx + 0.2, comp_y + 0.44 + j * (step_h + 0.04), half_w - 0.4, step_h, c)
        txt(sl, rx + 0.2, comp_y + 0.44 + j * (step_h + 0.04), half_w - 0.4, step_h,
            f'{s}  {note}'.strip(), 10, WHITE, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)


# ── 課題③ ──────────────────────────────────────────────────────
def slide_issue3(prs):
    sl = blank(prs)
    cy = hdr(sl, '個人の努力は限界。組織的な底上げが急務', '課題③　我流活用の限界')
    ftr(sl)

    half_w = (SW - 1.1) / 2
    card_h = SH - cy - 0.33

    rect(sl, 0.45, cy, half_w, card_h, CORAL_L, CORAL, 1.0)
    txt(sl, 0.45, cy + 0.05, half_w, 0.34, '現状：個人任せ', 14, CORAL, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate([
        '✗  各自がネットで独学・自分のやり方で使用',
        '✗  ツール・プロンプト設計がバラバラ',
        '✗  成功事例・失敗事例が共有されない',
        '✗  AI開発スキルが組織に積み上がらない',
        '✗  品質・成果に個人差が大きい',
    ]):
        txt(sl, 0.60, cy + 0.50 + j * 0.54, half_w - 0.25, 0.46, item, 10, CORAL)

    txt(sl, half_w + 0.5, cy + card_h / 2 - 0.25, 0.4, 0.5,
        '→', 22, TEAL, bold=True, align=PP_ALIGN.CENTER)

    rx = 0.5 + half_w
    rect(sl, rx, cy, half_w, card_h, SAGE_L, SAGE, 1.0)
    txt(sl, rx, cy + 0.05, half_w, 0.34, '目指す姿：組織化', 14, SAGE, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate([
        '✓  全員が共通のAI活用フレームワークを持つ',
        '✓  ベストプラクティスが組織ナレッジとして共有',
        '✓  ビジネス課題→AIタスク変換を誰でもできる',
        '✓  AI開発スキルが組織として積み上がる',
        '✓  顧客の「AIでできますか？」に答えられる',
    ]):
        txt(sl, rx + 0.15, cy + 0.50 + j * 0.54, half_w - 0.25, 0.46, item, 10, SAGE)


# ── 課題④ ──────────────────────────────────────────────────────
def slide_issue4(prs):
    sl = blank(prs)
    cy = hdr(sl, '本番化できなければ、保険SaaSの競争力を失う', '課題④　タイムライン・緊急性')
    ftr(sl)

    milestones = [
        ('2026\n3月末', 'ソフトバンク（親会社）への\nAIプロダクト企画書提出', CORAL, True),
        ('2026\n4〜6月', '親会社の投資承認・\nAI開発ベンダー選定',           SLATE, False),
        ('2026\n下半期', 'POC→本番化フェーズ移行\n内製エンジニアが必要',       NAVY,  False),
        ('2027〜',       'AI機能搭載の保険SaaSで\n競合との差別化を実現',        SAGE,  False),
    ]
    mw = (SW - 0.95) / len(milestones)
    line_y = cy + 1.08
    rect(sl, 0.45, line_y, SW - 0.9, 0.04, MGRAY)

    card_h = SH - line_y - 0.85
    for i, (t, desc, c, urgent) in enumerate(milestones):
        lx = 0.45 + i * mw
        cx = lx + mw / 2
        txt(sl, lx, cy + 0.06, mw, 0.88, t, 14, c, bold=True, align=PP_ALIGN.CENTER)
        dot_s = 0.22
        rect(sl, cx - dot_s / 2, line_y - dot_s / 2 + 0.02, dot_s, dot_s, c)
        rect(sl, lx + 0.06, line_y + 0.22, mw - 0.12, card_h, CARD, c, 1.2 if urgent else 0.5)
        txt(sl, lx + 0.14, line_y + 0.30, mw - 0.24, card_h - 0.38,
            desc, 10, c if urgent else DARK, bold=urgent)
        if urgent:
            rect(sl, lx + 0.06, line_y + card_h - 0.28, mw - 0.12, 0.26, CORAL)
            txt(sl, lx + 0.06, line_y + card_h - 0.28, mw - 0.12, 0.26,
                '★ 今ここ', 10, WHITE, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

    risk_y = line_y + card_h + 0.28
    rect(sl, 0.45, risk_y, SW - 0.9, 0.50, CORAL_L, CORAL, 0.8)
    txt(sl, 0.58, risk_y + 0.06, SW - 1.1, 0.38,
        '⚠  リスク：AIを搭載した競合SaaSが先行した場合、保険会社・保険代理店からの受注を失うリスクがある',
        10, CORAL, bold=True)


# ── ゴール① ────────────────────────────────────────────────────
def slide_goal1(prs):
    sl = blank(prs)
    cy = hdr(sl, '研修後、自社エンジニアがAIで開発できる状態へ', 'ゴール　Before / After')
    ftr(sl)

    half_w = (SW - 1.1) / 2
    card_h = SH - cy - 0.33

    rect(sl, 0.45, cy, half_w, card_h, CORAL_L, CORAL, 1.0)
    rect(sl, 0.45, cy, half_w, 0.40, CORAL)
    txt(sl, 0.45, cy + 0.04, half_w, 0.32, 'Before', 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate([
        '✗  AIはチャットツール感覚の個人利用にとどまる',
        '✗  ビジネス要件をAI開発に落とし込めない',
        '✗  AI開発はPOCのみ・海外BPOに依存',
        '✗  中国語話者エンジニアが教育から取り残される',
        '✗  顧客の「AIでできますか？」に答えられない',
    ]):
        txt(sl, 0.60, cy + 0.50 + j * 0.51, half_w - 0.25, 0.43, item, 10, CORAL)

    txt(sl, half_w + 0.5, cy + card_h / 2 - 0.22, 0.4, 0.44,
        '→', 22, TEAL, bold=True, align=PP_ALIGN.CENTER)

    rx = 0.5 + half_w
    rect(sl, rx, cy, half_w, card_h, SAGE_L, SAGE, 1.0)
    rect(sl, rx, cy, half_w, 0.40, SAGE)
    txt(sl, rx, cy + 0.04, half_w, 0.32, 'After（研修後）', 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate([
        '✓  エンジニア全員がAI開発プロセスを理解',
        '✓  ビジネス課題→AIタスク分解を自力でできる',
        '✓  GPT Enterprise活用のベストプラクティス確立',
        '✓  中国語話者も含め組織全体で底上げ完了',
        '✓  保険業界ユースケースを自社で実装できる基礎',
    ]):
        txt(sl, rx + 0.15, cy + 0.50 + j * 0.51, half_w - 0.25, 0.43, item, 10, SAGE)


# ── ゴール② ────────────────────────────────────────────────────
def slide_goal2(prs):
    sl = blank(prs)
    cy = hdr(sl, '研修完了後に達成すべき4つの状態', 'ゴール　成果指標')
    ftr(sl)

    goals = [
        ('01', 'AIエージェント基礎の習得',
         'AIエージェントとは何か・どう設計するかを\n全エンジニアが説明・実践できる状態', TEAL),
        ('02', 'AIタスク分解の自走',
         'ビジネス要件を受け取ったとき、\nどのAIタスクに分解すべきか\n自分で判断できる状態', SLATE),
        ('03', 'GPT活用ベストプラクティス確立',
         '社内で統一されたプロンプト設計・\nAPI活用の標準手順が\n組織として存在する状態', SAGE),
        ('04', '保険ユースケースの実装基礎',
         'AI-OCR・保険金自動判断など\n業界固有のユースケースを\n実装に向けて設計できる状態', NAVY),
    ]

    cols = 2
    cw = (SW - 0.95) / cols
    ch = (SH - cy - 0.38) / 2 - 0.08

    for i, (num, title, body, c) in enumerate(goals):
        col = i % cols; row = i // cols
        lx = 0.45 + col * (cw + 0.05)
        ty = cy + row * (ch + 0.10)
        rect(sl, lx, ty, cw, ch, CARD, MGRAY)
        rect(sl, lx, ty, cw, 0.05, c)
        txt(sl, lx + 0.12, ty + 0.10, 0.55, 0.38, num,   22, c,    bold=True)
        txt(sl, lx + 0.75, ty + 0.14, cw - 0.90, 0.30, title, 13, DARK, bold=True)
        txt(sl, lx + 0.12, ty + 0.50, cw - 0.24, ch - 0.60, body, 10, GRAY)


# ── 根拠① ──────────────────────────────────────────────────────
def slide_evidence1(prs):
    sl = blank(prs)
    cy = hdr(sl, 'なぜ一般的な研修では解決しないのか', '提案根拠①　一般研修との違い')
    ftr(sl)

    cols_txt = ['比較軸', '一般的なAI研修', 'Kawaru Team']
    col_w    = [2.3, 4.3, 5.2]
    rows = [
        ['カリキュラム',         '汎用・固定コンテンツ',            '御社専用フルカスタマイズ（0から設計）'],
        ['講師のバックグラウンド', '研修専門家・講師',               '現役でAI開発を行う実践者'],
        ['内容の焦点',           'ツール紹介・一般的な使い方',        '保険業界ユースケース＋GPT活用法'],
        ['対象者への配慮',        '標準的な受講者を想定',             '中国語話者も含めた設計対応'],
        ['研修後のゴール',        'AIの概要・知識習得',              '翌日から業務で使える実践力'],
    ]

    row_h = (SH - cy - 0.33) / (len(rows) + 1)
    lx = 0.45

    for j, (ct, cw) in enumerate(zip(cols_txt, col_w)):
        cx = lx + sum(col_w[:j])
        hc = [SLATE, SLATE, TEAL][j]
        rect(sl, cx, cy, cw, row_h, hc, MGRAY)
        txt(sl, cx + 0.08, cy + 0.05, cw - 0.12, row_h - 0.08,
            ct, 12, WHITE, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

    for i, row in enumerate(rows):
        bg = TEAL_L if i % 2 == 0 else WHITE
        for j, (cell, cw) in enumerate(zip(row, col_w)):
            cx = lx + sum(col_w[:j])
            if j == 0:
                fb = SLATE_L
            elif j == 2:
                fb = RGBColor(0xE8, 0xF8, 0xEC) if i % 2 == 0 else RGBColor(0xF0, 0xFD, 0xF6)
            else:
                fb = bg
            rect(sl, cx, cy + (i + 1) * row_h, cw, row_h, fb, MGRAY, 0.3)
            c = TEAL if j == 2 else (SLATE if j == 0 else DARK)
            txt(sl, cx + 0.10, cy + (i + 1) * row_h + 0.05, cw - 0.15, row_h - 0.08,
                cell, 10, c, bold=(j == 2), va=MSO_ANCHOR.MIDDLE)


# ── 根拠② ──────────────────────────────────────────────────────
def slide_evidence2(prs):
    sl = blank(prs)
    cy = hdr(sl, 'AI開発の現場にいるからこそ伝えられる', '提案根拠②　自社の強み')
    ftr(sl)

    sw3 = (SW - 0.95) / 3
    sh3 = SH - cy - 0.33
    strengths = [
        ('01', '現役AI開発者が講師',
         '私たち自身が日常的にAIエージェント・\nツールを開発しています。\n「理論」ではなく、実際に動く現場の\n知見を直接お伝えできます。\nChatGPT活用のリアルなベストプラクティスも\n実体験ベースです。', TEAL),
        ('02', '0から作るフルカスタマイズ',
         '「汎用コンテンツを流す研修」ではなく、\n御社のAI課題・業界特性・チーム構成を\n2時間のヒアリングで徹底把握した上で\nカリキュラムをゼロから設計します。\n中国語話者エンジニアへの配慮も含みます。', SLATE),
        ('03', '豊富な導入実績',
         'キリングループ・株式会社ミズカラなど\n多様な業種・規模に導入済。\n「AI未経験でも全員が使いたいと回答\n満足度8.5点」という実績があります。\n保険・金融業界への知見も対応可能です。', SAGE),
    ]
    for i, (num, title, body, c) in enumerate(strengths):
        lx = 0.45 + i * (sw3 + 0.02)
        rect(sl, lx, cy, sw3, sh3, CARD, MGRAY)
        rect(sl, lx, cy, sw3, 0.05, c)
        txt(sl, lx + 0.14, cy + 0.10, 0.55, 0.42, num,   22, c,    bold=True)
        txt(sl, lx + 0.14, cy + 0.54, sw3 - 0.25, 0.32, title, 13, DARK, bold=True)
        rect(sl, lx + 0.14, cy + 0.92, sw3 - 0.30, 0.02, MGRAY)
        txt(sl, lx + 0.14, cy + 1.00, sw3 - 0.25, sh3 - 1.10, body, 10, GRAY)


# ── 詳細① ──────────────────────────────────────────────────────
def slide_detail1(prs):
    sl = blank(prs)
    cy = hdr(sl, 'ヒアリング2時間でカリキュラムをゼロ設計', '提案詳細①　内容・カスタマイズフロー')
    ftr(sl)

    steps = [
        ('STEP 1', 'ヒアリング\n（2時間）',  '御社の課題・目標・\nチーム構成を徹底把握', TEAL),
        ('STEP 2', 'カリキュラム\n設計',     '御社専用の研修内容を\nゼロから設計',          TEAL),
        ('STEP 3', '資料作成・準備',         '保険業界ユースケース含む\n演習資料を作成',      SLATE),
        ('STEP 4', '研修実施\n（3〜4日間）', '御社オフィスへ訪問\n1日2〜3時間×3〜4日',     SAGE),
        ('STEP 5', '振り返り\nフォロー',     '研修後のフォロー\nベストプラクティス定着',     NAVY),
    ]
    sw5 = (SW - 0.95) / len(steps)
    step_h = 2.32
    sy = cy + 0.08

    for i, (lbl, title, body, c) in enumerate(steps):
        lx = 0.45 + i * (sw5 + 0.02)
        rect(sl, lx, sy, sw5 - 0.02, step_h, CARD, MGRAY)
        rect(sl, lx, sy, sw5 - 0.02, 0.05, c)
        txt(sl, lx, sy + 0.08,  sw5 - 0.02, 0.22, lbl,   9,  c,    bold=True, align=PP_ALIGN.CENTER)
        txt(sl, lx, sy + 0.34,  sw5 - 0.02, 0.52, title, 11, DARK, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, lx + 0.10, sy + 0.94, sw5 - 0.22, step_h - 1.04, body, 9, GRAY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            txt(sl, lx + sw5 - 0.08, sy + step_h / 2 - 0.18, 0.18, 0.36,
                '→', 14, TEAL, bold=True, align=PP_ALIGN.CENTER)

    cy2 = sy + step_h + 0.16
    txt(sl, 0.45, cy2, 4.0, 0.26, 'カリキュラム内容例（御社向け）', 12, TEAL, bold=True)

    curric = [
        'AIエージェントとは何か・設計の基礎',
        'ビジネス課題 → AIタスク分解フレームワーク',
        'ChatGPT Enterprise 実践的活用法（GPT縛り環境対応）',
        '保険業界ユースケース演習（AI-OCR・保険金自動判断）',
    ]
    cw2 = (SW - 0.95) / 2
    for i, item in enumerate(curric):
        col = i % 2; row = i // 2
        lx = 0.45 + col * (cw2 + 0.05)
        ty = cy2 + 0.33 + row * 0.40
        rect(sl, lx, ty, cw2, 0.35, TEAL_L, TEAL, 0.4)
        txt(sl, lx + 0.12, ty + 0.05, cw2 - 0.2, 0.26, f'• {item}', 10, DARK)


# ── 詳細② ──────────────────────────────────────────────────────
def slide_detail2(prs):
    sl = blank(prs)
    cy = hdr(sl, '25名×2〜3セッション、平日9〜18時実施', '提案詳細②　実施形式・スケジュール')
    ftr(sl)

    lw = 5.8
    info = [
        ('対象',          '両本部エンジニア　50〜60名'),
        ('1セッション定員', '15〜25名（業務への影響を最小化するため分割）'),
        ('セッション数',   '25名 × 2〜3セッション'),
        ('実施形式',       'オフライン訪問型（御社オフィスへ弊社講師が訪問）'),
        ('1セッション期間', '1日2〜3時間 × 3〜4日間（合計10〜15時間）'),
        ('助成金条件',     '平日9:00〜18:00 の実施が適用条件（時間外は対象外）'),
    ]
    row_h = 0.53
    for i, (lbl, val) in enumerate(info):
        bg = TEAL_L if i % 2 == 0 else WHITE
        rect(sl, 0.45, cy + i * row_h, 1.9,      row_h, TEAL,  MGRAY)
        rect(sl, 2.35, cy + i * row_h, lw - 1.9, row_h, bg,    MGRAY)
        txt(sl, 0.52, cy + i * row_h + 0.09, 1.76,     row_h - 0.13, lbl, 10, WHITE, bold=True)
        txt(sl, 2.44, cy + i * row_h + 0.09, lw - 2.05, row_h - 0.13, val, 10, DARK)

    rx = 6.45
    rw = SW - rx - 0.3
    txt(sl, rx, cy, rw, 0.26, 'セッション分割（60名の場合）', 11, TEAL, bold=True)

    sessions = [
        ('Session 1', '25名', '1日目〜4日目',   TEAL),
        ('Session 2', '25名', '5日目〜8日目',   SLATE),
        ('Session 3', '10名', '9日目〜12日目',  SAGE),
    ]
    total_h = len(info) * row_h - 0.36
    sess_h = total_h / len(sessions) - 0.08
    for i, (lbl, count, period, c) in enumerate(sessions):
        ty = cy + 0.34 + i * (sess_h + 0.08)
        rect(sl, rx, ty, rw, sess_h, CARD, MGRAY)
        rect(sl, rx, ty, 0.05, sess_h, c)
        txt(sl, rx + 0.12, ty + 0.06, 1.6,        0.28, lbl,    12, c,    bold=True)
        txt(sl, rx + 1.8,  ty + 0.04, 0.9, sess_h - 0.08, count, 20, c,    bold=True, va=MSO_ANCHOR.MIDDLE)
        txt(sl, rx + 2.85, ty + 0.10, rw - 3.0,   0.26, period, 10, GRAY)


# ── 詳細③ ──────────────────────────────────────────────────────
def slide_detail3(prs):
    sl = blank(prs)
    cy = hdr(sl, '1名15万円、助成金活用で実質負担30〜40%', '提案詳細③　費用・助成金')
    ftr(sl)

    cw2 = (SW - 0.95) / 2
    card_h = 2.2
    cases = [
        ('50名の場合', '750万円', '実質負担：225〜300万円', TEAL),
        ('60名の場合', '900万円', '実質負担：270〜360万円', SLATE),
    ]
    for i, (lbl, total, real, c) in enumerate(cases):
        lx = 0.45 + i * (cw2 + 0.05)
        rect(sl, lx, cy,           cw2, card_h, CARD, MGRAY)
        rect(sl, lx, cy,           cw2, 0.05,   c)
        txt(sl, lx, cy + 0.10,     cw2, 0.28, lbl,   12, c,    bold=True, align=PP_ALIGN.CENTER)
        txt(sl, lx, cy + 0.44,     cw2, 0.64, total, 28, DARK, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, lx, cy + 1.10,     cw2, 0.26, '（1名15万円ベース）', 10, GRAY, align=PP_ALIGN.CENTER)
        rect(sl, lx + 0.2, cy + 1.40, cw2 - 0.4, 0.02, MGRAY)
        txt(sl, lx, cy + 1.50,     cw2, 0.55, real, 14, SAGE, bold=True, align=PP_ALIGN.CENTER)

    sub_y = cy + card_h + 0.16
    rect(sl, 0.45, sub_y, SW - 0.9, 1.08, TEAL_L, TEAL, 0.8)
    txt(sl, 0.58, sub_y + 0.08, 5.0, 0.26, '助成金について', 12, TEAL, bold=True)
    txt(sl, 0.58, sub_y + 0.38, SW - 1.1, 0.65,
        '• 人材開発支援助成金（訓練等）が適用可能 → 研修費用の60〜75%を助成\n'
        '• 適用条件：平日9:00〜18:00 の実施（時間外・休日は対象外）\n'
        '• 申請手続きのサポートも弊社にてご支援可能です',
        10, DARK)


# ── ネクストアクション ──────────────────────────────────────────
def slide_na(prs):
    sl = blank(prs)
    cy = hdr(sl, 'ネクストアクション', '次のステップ')
    ftr(sl)

    steps = [
        ('STEP 1', '社内承認',
         '今週木曜日の社長との1on1にて\nAI研修のご承認をいただく', TEAL),
        ('STEP 2', 'ヒアリング日程調整',
         '承認後、カリキュラム設計のための\n2時間ヒアリングを設定', SLATE),
        ('STEP 3', 'カリキュラム設計・研修実施',
         'フルカスタマイズカリキュラムを\n確認いただき研修を実施', SAGE),
    ]
    sw3 = (SW - 0.95) / len(steps)
    sy = cy + 0.28
    sh3 = SH - sy - 0.36

    for i, (lbl, title, body, c) in enumerate(steps):
        lx = 0.45 + i * (sw3 + 0.02)
        rect(sl, lx, sy, sw3 - 0.02, sh3, CARD, MGRAY)
        rect(sl, lx, sy, sw3 - 0.02, 0.05, c)
        dot_s = 0.60
        dot_l = lx + (sw3 - 0.02 - dot_s) / 2
        rect(sl, dot_l, sy + 0.14, dot_s, dot_s, c)
        txt(sl, dot_l, sy + 0.14, dot_s, dot_s,
            str(i + 1), 24, WHITE, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
        txt(sl, lx, sy + 0.90, sw3 - 0.02, 0.34,
            title, 14, DARK, bold=True, align=PP_ALIGN.CENTER)
        txt(sl, lx + 0.15, sy + 1.38, sw3 - 0.32, sh3 - 1.48,
            body, 10, GRAY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            txt(sl, lx + sw3 - 0.06, sy + sh3 / 2 - 0.22, 0.18, 0.44,
                '→', 18, TEAL, bold=True, align=PP_ALIGN.CENTER)


# ── 締め ────────────────────────────────────────────────────────
def slide_closing(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 0.12, SH, TEAL)
    rect(sl, 0.12, 0, SW - 0.12, 0.70, NAVY)
    txt(sl, 0.5, 0.14, 4.0, 0.42, '株式会社エヌイチ', 14, WHITE, bold=True)
    txt(sl, 1.2, 1.8, SW - 1.6, 3.2,
        '私たちが提供するのは\n「ツール」や「研修」ではありません。\n\n'
        '貴社の未来を創る「AI人材」と、\n組織全体の「成長」です。',
        20, DARK, bold=True)
    rect(sl, 1.2, 5.25, 6.5, 0.03, TEAL)
    txt(sl, 1.2, SH - 0.56, 4.5, 0.36, '株式会社エヌイチ', 14, GRAY)


# ── 会社概要 ────────────────────────────────────────────────────
def slide_company(prs):
    sl = blank(prs)
    cy = hdr(sl, '会社概要')
    ftr(sl)

    info = [
        ('会社名',        '株式会社エヌイチ（n1 inc.）'),
        ('代表取締役CEO', '奥山幸生'),
        ('取締役COO',     '髙橋悠'),
        ('取締役CHRO',    '米澤浩明'),
        ('資本金',        '1,000万円'),
        ('設立',          '2023年10月'),
        ('所在地',        '東京都新宿区西新宿3-7-30\nフロンティアグラン西新宿12F'),
        ('Mission',       '新しい働く「カタチ」を創る'),
        ('Vision',        'AI（アイ）ある社会に'),
    ]
    lw = 5.6
    row_h = 0.46
    for i, (lbl, val) in enumerate(info):
        ty = cy + i * row_h
        if ty + row_h > SH - 0.3:
            break
        bg = TEAL_L if i % 2 == 0 else WHITE
        rect(sl, 0.45, ty, 1.75, row_h, TEAL,  MGRAY)
        rect(sl, 2.20, ty, lw - 1.75, row_h, bg, MGRAY)
        txt(sl, 0.53, ty + 0.07, 1.58,      row_h - 0.11, lbl, 9, WHITE, bold=True)
        txt(sl, 2.30, ty + 0.05, lw - 1.90, row_h - 0.09, val, 9, DARK)

    rx = 6.2
    rw = SW - rx - 0.3
    txt(sl, rx, cy, rw, 0.26, 'サービスラインアップ', 12, TEAL, bold=True)

    svcs = [
        ('Kawaru',       '業務自動化SaaS',         TEAL),
        ('Kawaru Team',  'フルカスタマイズAI研修',  SLATE),
        ('Kawaru BPO',   'AI業務効率化代行',        SAGE),
        ('Kawaru Coach', 'AI顧問サービス',          NAVY),
    ]
    sh_svc = (SH - cy - 0.48) / len(svcs) - 0.06
    for i, (name, desc, c) in enumerate(svcs):
        ty = cy + 0.36 + i * (sh_svc + 0.06)
        rect(sl, rx, ty, rw, sh_svc, CARD, MGRAY)
        rect(sl, rx, ty, 0.05, sh_svc, c)
        txt(sl, rx + 0.12, ty + 0.04, rw - 0.20, 0.28, name, 13, c,    bold=True)
        txt(sl, rx + 0.12, ty + 0.34, rw - 0.20, sh_svc - 0.42, desc, 10, GRAY)


# ── 導入実績 ────────────────────────────────────────────────────
def slide_cases(prs):
    sl = blank(prs)
    cy = hdr(sl, '導入実績', 'あらゆる規模の企業様に導入されています')
    ftr(sl)

    companies = [
        '株式会社ブロードリンク', '株式会社中西製作所',
        '株式会社ミズカラ',       'エル・ティー・エス リンク',
        'FRUOR株式会社',          '株式会社BIRDY',
        '株式会社オンシナジー',   'キリングループ',
    ]
    bw = (SW - 0.95) / 4
    bh = 0.44
    for i, name in enumerate(companies):
        col = i % 4; row = i // 4
        lx = 0.45 + col * (bw + 0.02)
        ty = cy + row * (bh + 0.08)
        rect(sl, lx, ty, bw, bh, CARD, MGRAY)
        txt(sl, lx + 0.10, ty + 0.06, bw - 0.18, bh - 0.10,
            name, 9, DARK, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

    case_y = cy + 2 * (bh + 0.08) + 0.16
    txt(sl, 0.45, case_y, 3.0, 0.26, '導入事例', 13, TEAL, bold=True)

    cases = [
        ('キリングループ様',      'AI未経験、活用イメージなし',
         '全員が「業務でAIを使いたい」と回答\n満足度 8.5点（10点満点）', TEAL),
        ('株式会社ミズカラ様',    '問い合わせ対応に月283時間',
         '月85時間へ削減\n約200時間＝25営業日分の削減', SLATE),
        ('株式会社BIRDY様',       '業務の属人化・営業効率の壁',
         'バックオフィス効率化\n経営資源を成長領域にシフト', SAGE),
    ]
    csw = (SW - 0.95) / 3
    cch = SH - case_y - 0.65
    for i, (co, issue, result, c) in enumerate(cases):
        lx = 0.45 + i * (csw + 0.02)
        ty = case_y + 0.33
        rect(sl, lx, ty, csw, cch, CARD, MGRAY)
        rect(sl, lx, ty, 0.05, cch, c)
        txt(sl, lx + 0.12, ty + 0.08, csw - 0.20, 0.28, co, 12, c, bold=True)
        rect(sl, lx + 0.12, ty + 0.40, csw - 0.24, 0.60, CORAL_L, CORAL, 0.5)
        txt(sl, lx + 0.20, ty + 0.46, csw - 0.35, 0.48, f'課題：{issue}', 9, CORAL)
        rect(sl, lx + 0.12, ty + 1.06, csw - 0.24, 0.02, MGRAY)
        rect(sl, lx + 0.12, ty + 1.12, csw - 0.24, cch - 1.22, SAGE_L, SAGE, 0.5)
        txt(sl, lx + 0.20, ty + 1.18, csw - 0.35, cch - 1.32, f'成果：{result}', 9, SAGE, bold=True)


# ── MAIN ────────────────────────────────────────────────────────
def main():
    prs = prs_new()
    client   = 'リードインクス株式会社'
    date_str = '2026年3月25日'

    sections = [
        ('01', '振り返り',   '前回商談の認識合わせ'),
        ('02', '現状',       '組織・AI活用・制約の整理'),
        ('03', '課題',       '根本課題と4つの要因'),
        ('04', 'ゴール',     '研修後の目指す姿'),
        ('05', '提案詳細',   'Kawaru Team 詳細・費用'),
    ]

    slide_cover(prs, client, date_str)           # 1
    toc_slide(prs, sections)                     # 2
    section_div(prs, '01', '振り返り',   '前回商談の内容と合意事項を確認します')  # 3
    slide_recap(prs)                             # 4
    section_div(prs, '02', '現状',       '貴社のAI活用・組織の現状を整理します')  # 5
    slide_current1(prs)                          # 6
    slide_current2(prs)                          # 7
    slide_current3(prs)                          # 8
    section_div(prs, '03', '課題',       '現状から導き出される根本課題を整理します')  # 9
    slide_issue_tree(prs)                        # 10
    slide_issue2(prs)                            # 11
    slide_issue3(prs)                            # 12
    slide_issue4(prs)                            # 13
    section_div(prs, '04', 'ゴール',     '研修後に実現したい状態を定義します')    # 14
    slide_goal1(prs)                             # 15
    slide_goal2(prs)                             # 16
    section_div(prs, '05', '提案詳細',   'Kawaru Team のご提案内容・費用をお伝えします')  # 17
    slide_evidence1(prs)                         # 18
    slide_evidence2(prs)                         # 19
    slide_detail1(prs)                           # 20
    slide_detail2(prs)                           # 21
    slide_detail3(prs)                           # 22
    slide_na(prs)                                # 23
    slide_closing(prs)                           # 24
    ref_div(prs)                                 # 25
    slide_company(prs)                           # 26
    slide_cases(prs)                             # 27

    out = '/Users/kyouyuu/Downloads/proposal_leadinks_v4_20260325.pptx'
    prs.save(out)
    print(f'Saved: {out}  ({len(prs.slides)} slides)')


if __name__ == '__main__':
    main()
