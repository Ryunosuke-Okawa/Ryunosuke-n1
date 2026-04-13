"""
リードインクス株式会社 提案書 v2
全スライド白背景・内容展開版（14枚）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── カラーパレット ──────────────────────────────────
TEAL    = RGBColor(0x42, 0xA4, 0xAF)
TEAL_L  = RGBColor(0xE0, 0xF4, 0xF6)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x1E, 0x29, 0x3B)
DARK    = RGBColor(0x11, 0x18, 0x27)
GRAY    = RGBColor(0x6B, 0x72, 0x80)
LGRAY   = RGBColor(0xF1, 0xF5, 0xF9)
MGRAY   = RGBColor(0xE2, 0xE8, 0xF0)
RED     = RGBColor(0xEF, 0x44, 0x44)
RED_L   = RGBColor(0xFE, 0xF2, 0xF2)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
GREEN_L = RGBColor(0xF0, 0xFD, 0xF4)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
BLUE_L  = RGBColor(0xEF, 0xF6, 0xFF)
ORANGE  = RGBColor(0xF5, 0x9E, 0x0B)
ORANGE_L= RGBColor(0xFE, 0xF3, 0xC7)

FONT = 'Noto Sans JP'
SW, SH = 13.33, 7.50   # slide width / height in inches

# ── 基本ユーティリティ ───────────────────────────────

def prs_new():
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def txt(slide, x, y, w, h, text, size, color=BLACK, bold=False,
        align=PP_ALIGN.LEFT, va=MSO_ANCHOR.TOP, ml=0.10, mt=0.06):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = va
    tf.margin_left   = Inches(ml)
    tf.margin_right  = Inches(0.06)
    tf.margin_top    = Inches(mt)
    tf.margin_bottom = Inches(0.04)
    for i, line in enumerate(str(text).split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
        r.font.name  = FONT
    return tb

# ── 共通ヘッダー ─────────────────────────────────────

def header(slide, title, subtitle=None):
    """
    白背景前提。タイトルエリア（上部）を構成し、コンテンツ開始Yを返す。
    """
    # 上端アクセントライン
    rect(slide, 0, 0, SW, 0.06, TEAL)
    # 右上社名
    txt(slide, 9.80, 0.10, 3.30, 0.40, '株式会社エヌイチ', 11, GRAY,
        align=PP_ALIGN.RIGHT)
    # タイトル
    txt(slide, 0.45, 0.18, 10.50, 0.75, title, 26, DARK, bold=True,
        va=MSO_ANCHOR.MIDDLE)
    if subtitle:
        txt(slide, 0.45, 0.96, 12.43, 0.38, subtitle, 13, GRAY)
        # セパレーター
        rect(slide, 0.45, 1.38, 12.43, 0.03, TEAL)
        return 1.52
    rect(slide, 0.45, 0.96, 12.43, 0.03, TEAL)
    return 1.10

# ── カードユーティリティ ─────────────────────────────

def card(slide, x, y, w, h, fill=WHITE, border=MGRAY, bar=None):
    r = rect(slide, x, y, w, h, fill, border)
    if bar:
        rect(slide, x, y, 0.07, h, bar)
    return r

def label_rows(slide, pairs, x, y, lw, vw, row_h=0.50, gap=0.08,
               lbg=TEAL_L, lcolor=TEAL):
    for i, (lbl, val) in enumerate(pairs):
        iy = y + i * (row_h + gap)
        rect(slide, x, iy, lw, row_h, lbg)
        txt(slide, x + 0.06, iy, lw - 0.10, row_h, lbl, 10,
            lcolor, bold=True, va=MSO_ANCHOR.MIDDLE)
        rect(slide, x + lw + 0.06, iy, vw, row_h, WHITE, MGRAY)
        txt(slide, x + lw + 0.14, iy, vw - 0.18, row_h, val, 12,
            BLACK, va=MSO_ANCHOR.MIDDLE)

def bullet_rows(slide, items, x, y, w, row_h, gap=0.12,
                size=13, bg=LGRAY, bar=TEAL):
    for i, text in enumerate(items):
        iy = y + i * (row_h + gap)
        card(slide, x, iy, w, row_h, bg, MGRAY, bar)
        txt(slide, x + 0.22, iy, w - 0.32, row_h, text, size,
            BLACK, va=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════
# スライド生成
# ════════════════════════════════════════════════════

prs = prs_new()

# ───────────────────────────────────────────────────
# 1. 表紙（白背景）
# ───────────────────────────────────────────────────
s = blank(prs)
# 白背景のまま。左アクセントバー + 底ライン
rect(s, 0, 0, 0.40, SH, TEAL)
rect(s, 0.40, SH - 0.08, SW - 0.40, 0.08, TEAL)

# 中央コンテンツ
txt(s, 1.00, 1.80, 11.80, 0.55, '生成AI支援 ご提案書', 16, TEAL, bold=True,
    align=PP_ALIGN.CENTER)
rect(s, 3.50, 2.42, 6.33, 0.04, TEAL)
txt(s, 1.00, 2.60, 11.80, 1.40,
    'リードインクス株式会社 様', 44, DARK, bold=True,
    align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
txt(s, 1.00, 4.20, 11.80, 0.55,
    '「AI開発の型」を持って自走できる組織をつくる提案', 18, GRAY,
    align=PP_ALIGN.CENTER)

# 右下: 日付・社名
txt(s, 1.00, 6.40, 11.80, 0.45,
    '2026年3月24日　　株式会社エヌイチ', 13, GRAY,
    align=PP_ALIGN.RIGHT)

# ───────────────────────────────────────────────────
# 2. 会社概要
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '会社概要')

# 左: 基本情報
card(s, 0.45, cy, 6.00, SH - cy - 0.30, WHITE, MGRAY, TEAL)
txt(s, 0.70, cy + 0.15, 5.60, 0.38, '基本情報', 12, TEAL, bold=True)
info = [
    ('会社名',  '株式会社エヌイチ'),
    ('経営陣',  'CEO 奥山幸生 / COO 髙橋悠 / CHRO 米澤浩明'),
    ('資本金',  '1,000万円'),
    ('設立',    '2023年10月'),
    ('所在地',  '東京都新宿区西新宿3-7-30'),
    ('事業',    'AI活用DX支援 / AI研修 / AIスクール'),
    ('Mission', '新しい「働くカタチ」を創る'),
    ('Vision',  'AI（アイ）ある社会に'),
]
label_rows(s, info, 0.62, cy + 0.62, 1.20, 4.48,
           row_h=0.56, gap=0.07)

# 右: サービスカード
svcs = [
    ('Kawaru',       '業務自動化SaaS',          TEAL,   TEAL_L),
    ('Kawaru Team',  'フルカスタマイズAI研修',   BLUE,   BLUE_L),
    ('Kawaru BPO',   'AI業務効率化代行',         GREEN,  GREEN_L),
    ('Kawaru Coach', 'AI顧問サービス',           ORANGE, ORANGE_L),
]
sx = 6.80
sy = cy
scard_h = (SH - cy - 0.30 - 0.15 * 3) / 4
for name, desc, bcolor, bg in svcs:
    card(s, sx, sy, 6.08, scard_h, bg, MGRAY, bcolor)
    txt(s, sx + 0.22, sy + 0.12, 5.70, 0.45,
        name, 16, DARK, bold=True)
    txt(s, sx + 0.22, sy + scard_h * 0.50, 5.70, 0.40,
        desc, 13, GRAY)
    sy += scard_h + 0.15

# ───────────────────────────────────────────────────
# 3. 導入実績
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '導入実績', 'あらゆる規模・業種の企業様に導入されています')

# 企業バッジ
companies = [
    '株式会社ブロードリンク', '株式会社中西製作所',
    '株式会社ミズカラ',       'エル・ティー・エス リンク',
    'FRUOR株式会社',          '株式会社BIRDY',
    '株式会社オンシナジー',   'キリングループ',
]
bw, bh, bgap = 2.93, 0.42, 0.20
for i, name in enumerate(companies):
    bx = 0.45 + (i % 4) * (bw + bgap)
    by = cy + 0.10 + (i // 4) * (bh + 0.14)
    rect(s, bx, by, bw, bh, LGRAY, MGRAY)
    txt(s, bx + 0.10, by, bw - 0.10, bh, name, 11, GRAY,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# 事例カード（3列）
cases = [
    ('キリングループ様',
     'AI未経験でツール選定もできない状態',
     '全員が「業務でAIを使いたい」と回答\n満足度 8.5点'),
    ('株式会社ミズカラ様',
     '問い合わせ対応に月283時間を費やしていた',
     '月85時間に削減\n約200時間（25営業日分）の工数削減'),
    ('株式会社BIRDY様',
     '業務の属人化・営業効率改善の停滞',
     'バックオフィス効率化を実現\n経営資源を成長領域にシフト'),
]
ct = cy + 1.18
ch = SH - ct - 0.25
cw = (SW - 0.90 - 0.24) / 3
for i, (co, issue, result) in enumerate(cases):
    cx = 0.45 + i * (cw + 0.12)
    card(s, cx, ct, cw, ch, WHITE, MGRAY, TEAL)
    # 会社名
    txt(s, cx + 0.22, ct + 0.16, cw - 0.35, 0.45,
        co, 14, DARK, bold=True)
    # 課題
    rect(s, cx + 0.15, ct + 0.72, cw - 0.30, 0.03, MGRAY)
    txt(s, cx + 0.15, ct + 0.80, 0.60, 0.26, '課題', 9,
        RED, bold=True, va=MSO_ANCHOR.MIDDLE)
    rect(s, cx + 0.15, ct + 0.80, cw - 0.30, ch * 0.30, RED_L, None)
    rect(s, cx + 0.15, ct + 0.80, 0.07, ch * 0.30, RED)
    txt(s, cx + 0.30, ct + 0.80, cw - 0.45, ch * 0.30,
        issue, 12, BLACK, va=MSO_ANCHOR.MIDDLE)
    # 成果
    txt(s, cx + 0.15, ct + 0.80 + ch*0.30 + 0.06,
        0.60, 0.26, '成果', 9, GREEN, bold=True, va=MSO_ANCHOR.MIDDLE)
    rect(s, cx + 0.15, ct + 0.80 + ch*0.30 + 0.06,
         cw - 0.30, ch - (0.80 + ch*0.30 + 0.12), GREEN_L, None)
    rect(s, cx + 0.15, ct + 0.80 + ch*0.30 + 0.06,
         0.07, ch - (0.80 + ch*0.30 + 0.12), GREEN)
    txt(s, cx + 0.30, ct + 0.80 + ch*0.30 + 0.06,
        cw - 0.45, ch - (0.80 + ch*0.30 + 0.12),
        result, 12, DARK, bold=True, va=MSO_ANCHOR.MIDDLE)

# ───────────────────────────────────────────────────
# 4. 前回商談の振り返り
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '前回お打合せの振り返り', '2026年2月27日（初回商談）')

# 上段: 商談サマリー
card(s, 0.45, cy, SW - 0.90, 2.80, WHITE, MGRAY, TEAL)
txt(s, 0.70, cy + 0.18, 12.00, 0.38, '商談サマリー', 12, TEAL, bold=True)
overview = [
    ('参加者（先方）', 'リードインクス株式会社　呉振宇氏（本部長）'),
    ('参加者（自社）', '大川龍之介・髙橋悠（COO）'),
    ('商談内容',       'エヌイチのサービス紹介 → 御社の現状・ニーズのヒアリング'),
    ('先方の反応',     'Kawaru Team（フルカスタマイズ研修）に強い関心を示していただいた'),
    ('合意事項',       '研修（Kawaru Team）を軸に進める方向性で合意。社長稟議を経て詳細設計へ'),
]
label_rows(s, overview, 0.60, cy + 0.65, 1.50, 10.50,
           row_h=0.36, gap=0.06)

# 下段: 合意したネクストアクション
card(s, 0.45, cy + 2.95, SW - 0.90, 4.25, WHITE, MGRAY, BLUE)
txt(s, 0.70, cy + 3.12, 12.00, 0.38, '前回合意したネクストアクション', 12, BLUE, bold=True)
nas = [
    ('呉振宇氏',
     '来週木曜日の社長1on1でAI研修提案（本提案書の内容）についてご相談いただく'),
    ('エヌイチ',
     'AI開発BPOパートナーを社内で選定し、承認後にご紹介できる準備を進める'),
]
nah = (4.25 - 0.65) / 2 - 0.15
ny = cy + 3.62
for who, na in nas:
    rect(s, 0.60, ny, 1.60, nah, BLUE_L)
    txt(s, 0.64, ny, 1.52, nah, who, 13, BLUE, bold=True, va=MSO_ANCHOR.MIDDLE)
    rect(s, 2.25, ny, SW - 2.75, nah, LGRAY, MGRAY)
    txt(s, 2.40, ny, SW - 2.95, nah, na, 13, BLACK, va=MSO_ANCHOR.MIDDLE)
    ny += nah + 0.15

# ───────────────────────────────────────────────────
# 5. ヒアリングで把握した現状（事実）
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'ヒアリングで把握した御社の現状', '※ いずれも商談での発言に基づく事実です')

facts = [
    ('事業・組織',
     'ソフトバンクグループ会社として保険会社・保険代理店向けSaaSプロダクトを開発・提供。'
     'エンジニアは両本部合計で50〜60名。多国籍チーム（日本語話者・中国語話者が混在）。'),
    ('AI活用の現状',
     'AI活用はメンバー個人レベルのChatGPT利用が中心。組織としてのベストプラクティスは未確立。'
     'ソフトバンクの方針によりChatGPT Enterpriseを中心とした活用が前提となっている。'),
    ('法規制上の制約',
     '金融規制上、個人情報を含む本番運用は海外ベンダーへの委託が不可。'
     'AIのPOCは海外ベンダーと実施中だが、本番運用は国内ベンダー限定。'),
    ('直近の締め切り',
     '2026年3月末に親会社（ソフトバンク）へのAIプロダクト企画書提出を予定。'
     'この締め切りに向けて、AI開発・活用の方向性を早急に固める必要がある。'),
]
avail_h = SH - cy - 0.20
row_h = (avail_h - 0.10 * (len(facts) - 1)) / len(facts)
for i, (label, detail) in enumerate(facts):
    fy = cy + 0.10 + i * (row_h + 0.10)
    card(s, 0.45, fy, SW - 0.90, row_h, WHITE, MGRAY, TEAL)
    txt(s, 0.72, fy + 0.10, 2.10, row_h - 0.20, label,
        13, TEAL, bold=True, va=MSO_ANCHOR.MIDDLE)
    rect(s, 2.85, fy + row_h * 0.15, 0.03, row_h * 0.70, MGRAY)
    txt(s, 3.00, fy, SW - 3.60, row_h, detail,
        13, BLACK, va=MSO_ANCHOR.MIDDLE)

# ───────────────────────────────────────────────────
# 6. 課題のロジックツリー
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '御社が直面している「3つの壁」', '事実ベースのヒアリング内容から特定した構造的な課題')

# 根本課題ブロック（左）
root_x, root_w = 0.45, 2.40
root_y, root_h = cy, SH - cy - 0.20
card(s, root_x, root_y, root_w, root_h, RED_L, None, RED)
txt(s, root_x + 0.18, root_y, root_w - 0.25, root_h,
    '組織として\nAIを活用・開発\nする力がなく\n顧客への\nAI価値提供が\n遅れている',
    15, DARK, bold=True,
    align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# コネクター（根本→中間）
rect(s, root_x + root_w, cy + root_h / 2 - 0.03, 0.30, 0.06, MGRAY)

# 3つの中間ノード・葉ノード
branches = [
    ('AI開発の\n方法論・知識不足', BLUE_L, BLUE, [
        ('AIエージェント開発のプロセスが社内に存在しない', True),
        ('ビジネス要件をAIタスクに分解する方法が不明', True),
        ('個人の我流利用にとどまり組織知が蓄積されない', True),
    ]),
    ('組織環境の制約', ORANGE_L, ORANGE, [
        ('ツール利用がChatGPT Enterprise限定（ソフトバンク方針）', True),
        ('多言語チーム（日中混在）による研修実施の障壁', True),
        ('金融規制により本番運用は国内ベンダー限定', True),
    ]),
    ('POCから本番実装への壁', RED_L, RED, [
        ('海外ベンダーは本番運用に使えない（法規制）', True),
        ('自社エンジニアのAI開発スキルが不足', True),
        ('保険業界特有ユースケースの実現方法が不明', False),
    ]),
]

mid_x = root_x + root_w + 0.30
mid_w = 2.30
leaf_x = mid_x + mid_w + 0.18
leaf_w = SW - leaf_x - 0.35

total_h = SH - cy - 0.20
branch_h = total_h / 3

for i, (blabel, bbg, bbar, leaves) in enumerate(branches):
    by = cy + i * branch_h
    bmy = by + branch_h / 2

    # コネクター（根本→中）
    rect(s, mid_x - 0.18, bmy - 0.03, 0.18, 0.06, MGRAY)
    # 中間ノード
    card(s, mid_x, by + 0.10, mid_w, branch_h - 0.20,
         bbg, None, bbar)
    txt(s, mid_x + 0.18, by + 0.10, mid_w - 0.28, branch_h - 0.20,
        blabel, 13, DARK, bold=True,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    # コネクター（中→葉）
    rect(s, mid_x + mid_w, bmy - 0.03, 0.18, 0.06, MGRAY)

    # 葉ノード
    leaf_h = (branch_h - 0.28) / 3
    for j, (ltext, is_fact) in enumerate(leaves):
        ly = by + 0.14 + j * leaf_h
        lbg = LGRAY
        lbar = TEAL if is_fact else ORANGE
        label = '事実' if is_fact else '推察'
        lc = TEAL if is_fact else ORANGE
        card(s, leaf_x, ly, leaf_w, leaf_h - 0.06, lbg, MGRAY, lbar)
        # ラベルバッジ
        rect(s, leaf_x + 0.16, ly + 0.06,
             0.60, leaf_h - 0.18, TEAL_L if is_fact else ORANGE_L)
        txt(s, leaf_x + 0.16, ly + 0.06,
            0.60, leaf_h - 0.18, label, 9,
            lc, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
        txt(s, leaf_x + 0.86, ly, leaf_w - 0.96, leaf_h - 0.06,
            ltext, 12, DARK, va=MSO_ANCHOR.MIDDLE)

# ───────────────────────────────────────────────────
# 7. ご提案のゴール
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'ご提案のゴール')

# メインメッセージ
card(s, 0.45, cy, SW - 0.90, 1.60, TEAL_L, None, TEAL)
txt(s, 0.70, cy, SW - 1.20, 1.60,
    '3ヶ月で、エンジニアが「AI開発の型」を持って\n自走できる組織をつくります',
    28, DARK, bold=True,
    align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# 4つのポイント
points = [
    ('AI開発の方法論を習得',
     'AIエージェント開発の標準プロセスと、ビジネス要件をAIタスクに分解する思考法を体系的に習得します。'),
    ('保険業界特化の実践演習',
     '御社のユースケース（AI-OCR・保険金自動判断）を題材にした実践演習で、研修後すぐに業務へ応用できます。'),
    ('ChatGPT Enterprise前提の設計',
     'ソフトバンクの方針に合わせ、ChatGPT Enterpriseを前提に設計。現場にそのままつながる内容にします。'),
    ('多言語チームへの対応',
     '日中バイリンガル資料オプションで、中国語話者も含む全メンバーが同じレベルで参加可能にします。'),
]
avail_h = SH - cy - 1.75
ph = (avail_h - 0.10 * (len(points) - 1)) / len(points)
for i, (pt, pd) in enumerate(points):
    py = cy + 1.70 + i * (ph + 0.10)
    card(s, 0.45, py, SW - 0.90, ph, WHITE, MGRAY, TEAL)
    # 番号
    rect(s, 0.45, py, 0.55, ph, TEAL)
    txt(s, 0.45, py, 0.55, ph, str(i + 1), 22, WHITE, bold=True,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    txt(s, 1.15, py + 0.08, 3.00, ph - 0.16,
        pt, 14, TEAL, bold=True, va=MSO_ANCHOR.MIDDLE)
    rect(s, 4.20, py + ph * 0.15, 0.03, ph * 0.70, MGRAY)
    txt(s, 4.40, py, SW - 5.00, ph,
        pd, 13, BLACK, va=MSO_ANCHOR.MIDDLE)

# ───────────────────────────────────────────────────
# 8. 助成金活用で実質負担を大幅削減
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '助成金活用で、実質負担を大幅削減')

# メイン数字
rect(s, 0.45, cy, SW - 0.90, 1.80, ORANGE_L, None)
txt(s, 0.70, cy, SW - 1.20, 1.80,
    '通常 1人15万円 → 助成金適用で実質 4.5〜6万円 / 人',
    26, DARK, bold=True,
    align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# 詳細カード（3列）
details = [
    ('制度名',
     '人材開発支援助成金\n事業展開等リスキリング支援コース',
     TEAL, TEAL_L),
    ('還元率・条件',
     '最大75%（6〜7割程度）\n平日9:00〜18:00での実施が条件',
     ORANGE, ORANGE_L),
    ('申請サポート',
     'エヌイチが申請サポートを提供\n適用条件の確認からご支援します',
     GREEN, GREEN_L),
]
dw = (SW - 0.90 - 0.20) / 3
for i, (dlbl, dval, dc, dbg) in enumerate(details):
    dx = 0.45 + i * (dw + 0.10)
    dy = cy + 2.00
    dh = 2.40
    card(s, dx, dy, dw, dh, dbg, None, dc)
    txt(s, dx + 0.18, dy + 0.15, dw - 0.30, 0.45, dlbl, 13, dc, bold=True)
    txt(s, dx + 0.18, dy + 0.68, dw - 0.30, dh - 0.80, dval, 14, DARK)

# 注釈
txt(s, 0.45, cy + 4.60, SW - 0.90, 0.60,
    '※ 助成金の適用条件・金額は個社の状況により異なります。詳細はエヌイチがヒアリング後にご案内します。',
    11, GRAY)

# ───────────────────────────────────────────────────
# 9. なぜKawaru Teamか
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'なぜKawaru Teamか', '一般的なAI研修との違い')

cols = [
    ('一般的なAI研修', LGRAY,  GRAY,  BLACK, [
        '既成カリキュラムをそのまま提供',
        '自社業務・業界への応用は自己解決',
        '研修後の実践サポートなし',
        '多言語・多国籍チームへの対応が困難',
    ]),
    ('Kawaru Team', TEAL_L, TEAL,  DARK, [
        '課題ヒアリング後にカリキュラムを0から設計',
        '保険業界のユースケースを題材に構成',
        'ChatGPT Enterprise前提で実務直結',
        '日中バイリンガル対応オプションあり',
    ]),
    ('期待できる成果', GREEN_L, GREEN, DARK, [
        'AI開発の型が組織に定着する',
        '海外BPO依存を段階的に脱却できる',
        'エンジニアが自走してAIを実装できる',
        '3月末の企画書提出に間に合う',
    ]),
]
avail_h = SH - cy - 0.20
cw = (SW - 0.90 - 0.20) / 3
cx = 0.45
for ctitle, cbg, cbar, ctc, items in cols:
    card(s, cx, cy, cw, avail_h, cbg, None, cbar)
    txt(s, cx + 0.18, cy + 0.15, cw - 0.30, 0.55,
        ctitle, 17, ctc, bold=True, align=PP_ALIGN.CENTER)
    rect(s, cx + 0.15, cy + 0.78, cw - 0.30, 0.04, cbar)
    ih = (avail_h - 0.98 - 0.12 * (len(items) - 1)) / len(items)
    for j, item in enumerate(items):
        iy = cy + 0.96 + j * (ih + 0.12)
        card(s, cx + 0.15, iy, cw - 0.30, ih, WHITE, MGRAY, cbar)
        txt(s, cx + 0.32, iy, cw - 0.55, ih,
            item, 13, ctc, va=MSO_ANCHOR.MIDDLE)
    cx += cw + 0.10

# ───────────────────────────────────────────────────
# 10. Kawaru Teamの進め方
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'Kawaru Teamの進め方', 'ヒアリングから研修後フォローまで、すべて伴走します')

steps = [
    ('STEP 1', 'ヒアリング\n課題整理', TEAL, TEAL_L, [
        '現状・ゴールの深掘りヒアリング（60〜90分）',
        '学習対象者のスキルレベル確認',
        '研修テーマ・到達目標の設定',
        'ChatGPT Enterprise環境・前提条件の確認',
    ]),
    ('STEP 2', 'カリキュラム\n設計', BLUE, BLUE_L, [
        '御社専用カリキュラムを0から設計',
        '保険業界ユースケースを題材に構成',
        '日中バイリンガル資料の作成（オプション）',
        '平日9〜18時での実施スケジュール調整',
    ]),
    ('STEP 3', '研修実施\nフォロー', GREEN, GREEN_L, [
        '1日2〜3時間 × 3〜5日、25名以内',
        '実施中の質疑応答・進捗確認',
        '実施後フォローアップMTG',
        '助成金申請サポート（対象条件確認）',
    ]),
]
avail_h = SH - cy - 0.20
cw = (SW - 0.90 - 0.20) / 3
cx = 0.45
badge_h = 0.50
card_h = avail_h - badge_h - 0.08

for i, (badge, title, bcolor, bg, items) in enumerate(steps):
    # バッジ
    rect(s, cx + cw / 2 - 0.80, cy, 1.60, badge_h, bcolor)
    txt(s, cx + cw / 2 - 0.80, cy, 1.60, badge_h, badge,
        13, WHITE, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    # 矢印
    if i < len(steps) - 1:
        txt(s, cx + cw + 0.02, cy + avail_h / 2, 0.18, 0.40,
            '›', 30, bcolor, align=PP_ALIGN.CENTER, bold=True)
    # カード
    card(s, cx, cy + badge_h + 0.08, cw, card_h, bg, None, bcolor)
    txt(s, cx + 0.18, cy + badge_h + 0.18, cw - 0.28, 1.00,
        title, 19, DARK, bold=True, align=PP_ALIGN.CENTER)
    rect(s, cx + 0.15, cy + badge_h + 1.22, cw - 0.30, 0.04, bcolor)
    ih = (card_h - 1.38 - 0.10 * (len(items) - 1)) / len(items)
    for j, item in enumerate(items):
        iy = cy + badge_h + 1.36 + j * (ih + 0.10)
        card(s, cx + 0.15, iy, cw - 0.30, ih, WHITE, MGRAY, bcolor)
        txt(s, cx + 0.32, iy, cw - 0.55, ih, item, 13, DARK,
            va=MSO_ANCHOR.MIDDLE)
    cx += cw + 0.10

# ───────────────────────────────────────────────────
# 11. 費用
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, '費用')

# メイン料金
card(s, 0.45, cy, SW - 0.90, 2.00, TEAL_L, None, TEAL)
txt(s, 0.70, cy, SW - 1.20, 2.00,
    '1人あたり 15万円（税別）',
    44, DARK, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# 詳細
details2 = [
    ('研修時間',  '10〜15時間（3〜5日間に分割して実施）'),
    ('対象人数',  'エンジニア 25名まで（両本部交代参加も可）'),
    ('実施形態',  '御社オフィスへ講師が訪問して実施'),
    ('カスタマイズ', '御社の課題・業界・ツール環境に合わせてカリキュラムを完全設計'),
]
dh2 = (SH - cy - 2.20 - 0.10 * (len(details2) - 1)) / len(details2)
for i, (dlbl, dval) in enumerate(details2):
    dy = cy + 2.20 + i * (dh2 + 0.10)
    card(s, 0.45, dy, SW - 0.90, dh2, WHITE, MGRAY, TEAL)
    txt(s, 0.70, dy, 2.20, dh2, dlbl, 13, TEAL, bold=True,
        va=MSO_ANCHOR.MIDDLE)
    rect(s, 2.92, dy + dh2 * 0.15, 0.03, dh2 * 0.70, MGRAY)
    txt(s, 3.10, dy, SW - 3.70, dh2, dval, 14, DARK,
        va=MSO_ANCHOR.MIDDLE)

# ───────────────────────────────────────────────────
# 12. スケジュール
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'スケジュール目安', '社長承認から研修キックオフまでの流れ')

schedule = [
    ('今週〜来週',       TEAL,   TEAL_L,
     '社長1on1での相談（呉様）',
     '来週木曜日の社長1on1でAI研修の方向性をご確認いただく。'
     'ご承認いただければ、翌週よりヒアリングMTGを設定します。'),
    ('承認後 Week1',     BLUE,   BLUE_L,
     'ヒアリングMTG（60〜90分）',
     '御社の現状・課題・ゴールを深掘りヒアリング。'
     '対象者のスキルレベルや業務内容も確認し、カリキュラム設計の前提を揃えます。'),
    ('承認後 Week2〜3',  GREEN,  GREEN_L,
     'カリキュラム設計・お見積り提示',
     '御社専用カリキュラム案を作成し、正式見積りとともにご提示します。'
     '内容確認・調整後、研修スケジュールを確定します。'),
    ('承認後 1ヶ月目',   ORANGE, ORANGE_L,
     '研修キックオフ',
     'エンジニアチームへの研修をスタート。'
     '実施期間中は進捗確認とQ&A対応を行い、終了後にフォローアップMTGを実施します。'),
]
avail_h = SH - cy - 0.20
sh2 = (avail_h - 0.12 * (len(schedule) - 1)) / len(schedule)
for i, (phase, sc, sbg, title, desc) in enumerate(schedule):
    sy2 = cy + 0.10 + i * (sh2 + 0.12)
    # フェーズバッジ
    rect(s, 0.45, sy2, 2.20, sh2, sbg)
    rect(s, 0.45, sy2, 0.08, sh2, sc)
    txt(s, 0.58, sy2, 2.04, sh2, phase, 13, sc, bold=True,
        va=MSO_ANCHOR.MIDDLE)
    # コンテンツ
    card(s, 2.72, sy2, SW - 3.17, sh2, WHITE, MGRAY, sc)
    txt(s, 2.96, sy2 + 0.06, 6.50, 0.40, title, 14, DARK, bold=True)
    txt(s, 2.96, sy2 + 0.50, SW - 3.70, sh2 - 0.56, desc, 12, GRAY)

# ───────────────────────────────────────────────────
# 13. ネクストアクション
# ───────────────────────────────────────────────────
s = blank(prs)
cy = header(s, 'ネクストアクション')

actions = [
    (TEAL,
     '① 社長1on1での相談（呉様）',
     '来週木曜日の社長1on1で本提案書の内容についてご相談いただき、'
     'AI研修実施の方向性をご確認ください。'),
    (BLUE,
     '② 社長承認後：ヒアリングMTGの設定',
     '社長承認が得られたら、エヌイチとのヒアリングMTG（60〜90分）を設定します。'
     '日程はご承認直後に調整可能です。'),
    (GREEN,
     '③ カリキュラム設計・お見積り提示',
     'ヒアリング完了後、御社専用カリキュラム案と正式見積りをご提示します。'
     '内容確認後、研修スケジュールを確定します。'),
    (ORANGE,
     '④ BPOパートナーご紹介（別途）',
     'AIプロダクト開発のBPO候補をエヌイチ社内で選定し、'
     'ご紹介の準備を並行して進めます（承認後にご連絡します）。'),
]
avail_h = SH - cy - 0.20
ah = (avail_h - 0.12 * (len(actions) - 1)) / len(actions)
for i, (acolor, atitle, adesc) in enumerate(actions):
    ay = cy + 0.10 + i * (ah + 0.12)
    # 左バー（カラー）
    rect(s, 0.45, ay, 0.55, ah, acolor)
    txt(s, 0.45, ay, 0.55, ah, str(i + 1), 22, WHITE, bold=True,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    # コネクター
    if i < len(actions) - 1:
        rect(s, 0.66, ay + ah, 0.12, 0.12, acolor)
    # カード本体
    card(s, 1.08, ay, SW - 1.53, ah, WHITE, MGRAY, acolor)
    txt(s, 1.32, ay + 0.10, SW - 2.20, 0.44,
        atitle, 16, DARK, bold=True, va=MSO_ANCHOR.MIDDLE)
    rect(s, 1.20, ay + ah * 0.52, SW - 2.00, 0.03, MGRAY)
    txt(s, 1.32, ay + ah * 0.55, SW - 2.20, ah * 0.40,
        adesc, 13, GRAY, va=MSO_ANCHOR.MIDDLE)

# ───────────────────────────────────────────────────
# 14. 締め（白背景）
# ───────────────────────────────────────────────────
s = blank(prs)
rect(s, 0, 0, 0.40, SH, TEAL)
rect(s, 0.40, SH - 0.08, SW - 0.40, 0.08, TEAL)
txt(s, 10.20, 0.18, 2.90, 0.40, '株式会社エヌイチ', 11, GRAY,
    align=PP_ALIGN.RIGHT)

txt(s, 0.90, 1.80, SW - 1.30, 0.60,
    '私たちが提供するのは', 20, GRAY, align=PP_ALIGN.CENTER)
txt(s, 0.90, 2.48, SW - 1.30, 1.20,
    '「ツール」や「研修」ではありません。',
    34, DARK, bold=True, align=PP_ALIGN.CENTER)

rect(s, 4.00, 3.84, 5.33, 0.05, TEAL)

txt(s, 0.90, 4.02, SW - 1.30, 1.20,
    '貴社の未来を創る「AI人材」と、\n組織全体の「成長」です。',
    26, TEAL, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

txt(s, 0.90, 5.50, SW - 1.30, 0.55,
    '株式会社エヌイチ', 18, DARK, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0.90, 6.08, SW - 1.30, 0.40,
    '大川 龍之介　080-2646-2420　r.okawa@n1-inc.co.jp',
    13, GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
out = '/Users/kyouyuu/cloude/output/proposal_leadinks_v2.pptx'
prs.save(out)
print(f'Saved: {out}  ({len(prs.slides)} slides)')
