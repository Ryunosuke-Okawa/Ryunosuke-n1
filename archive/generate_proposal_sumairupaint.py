"""
提案書生成スクリプト: 株式会社住まいるペイント 様
生成日: 2026年3月30日
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

# ============================================================
# カラー定義
# ============================================================
TEAL      = RGBColor(0x42, 0xA4, 0xAF)
SLATE     = RGBColor(0x4A, 0x6C, 0x8C)
CORAL     = RGBColor(0xC9, 0x54, 0x54)
SAGE      = RGBColor(0x3A, 0x8F, 0x63)
NAVY      = RGBColor(0x1C, 0x35, 0x57)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL= RGBColor(0xEA, 0xF6, 0xF7)
CARD_BG   = RGBColor(0xF9, 0xFA, 0xFB)
BORDER    = RGBColor(0xE5, 0xE7, 0xEB)
TITLE_CLR = RGBColor(0x11, 0x18, 0x27)
BODY_CLR  = RGBColor(0x1F, 0x29, 0x37)
SUB_CLR   = RGBColor(0x4B, 0x55, 0x63)
LABEL_CLR = RGBColor(0x6B, 0x72, 0x80)
RED       = RGBColor(0xEF, 0x44, 0x44)
LIGHT_RED = RGBColor(0xFF, 0xF5, 0xF5)
LIGHT_GRN = RGBColor(0xF0, 0xFA, 0xF5)

FONT = 'Noto Sans JP'
W = Inches(13.33)
H = Inches(7.50)

# ============================================================
# ユーティリティ関数
# ============================================================
def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # 白背景
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide

def rect(slide, left, top, width, height, fill=None, line=None, lw=Pt(0.5)):
    """矩形シェイプを追加"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = lw
    else:
        shape.line.width = Pt(0)
        shape.line.fill.background()
    return shape

def txt(slide, text, left, top, width, height,
        size=13, bold=False, color=BODY_CLR,
        align=PP_ALIGN.LEFT, wrap=True):
    """テキストボックスを追加"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def shape_txt(slide, text, left, top, width, height,
              fill=CARD_BG, size=13, bold=False, color=BODY_CLR,
              align=PP_ALIGN.LEFT, line=None, lw=Pt(0.5),
              v_anchor=MSO_ANCHOR.MIDDLE):
    """テキスト入りシェイプ"""
    shape = rect(slide, left, top, width, height, fill=fill, line=line, lw=lw)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape

def header(slide, title, subtitle=None):
    """標準ヘッダー: タイトル + TEALアクセントライン + 社名"""
    txt(slide, '株式会社エヌイチ',
        Inches(9.5), Inches(0.12), Inches(3.7), Inches(0.4),
        size=11, color=SUB_CLR, align=PP_ALIGN.RIGHT)
    txt(slide, title,
        Inches(0.5), Inches(0.18), Inches(9.0), Inches(0.72),
        size=26, bold=True, color=TITLE_CLR)
    rect(slide, Inches(0.5), Inches(0.93), Inches(12.33), Pt(2), fill=TEAL)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.35),
            size=11, color=SUB_CLR)

# ============================================================
# スライド1: 表紙
# ============================================================
def make_cover(prs):
    slide = blank_slide(prs)
    rect(slide, Inches(0), Inches(0), Inches(0.5), H, fill=TEAL)
    rect(slide, Inches(0.5), Inches(5.5), Inches(12.83), Pt(1), fill=BORDER)
    txt(slide, '生成AI支援ご提案書',
        Inches(1.0), Inches(1.8), Inches(11.0), Inches(1.4),
        size=38, bold=True, color=TITLE_CLR)
    txt(slide, '株式会社住まいるペイント 様向け',
        Inches(1.0), Inches(3.1), Inches(11.0), Inches(0.8),
        size=20, color=SLATE)
    rect(slide, Inches(1.0), Inches(3.9), Inches(4.5), Pt(2), fill=TEAL)
    txt(slide, '2026年3月30日',
        Inches(1.0), Inches(6.7), Inches(4.0), Inches(0.5),
        size=12, color=LABEL_CLR)
    txt(slide, '株式会社エヌイチ',
        Inches(9.0), Inches(6.7), Inches(4.0), Inches(0.5),
        size=12, color=SUB_CLR, align=PP_ALIGN.RIGHT)

# ============================================================
# スライド2: 目次
# ============================================================
def make_toc(prs):
    slide = blank_slide(prs)
    header(slide, '目次')
    sections = [
        ('01', '振り返り', '本日のお打合せ概要'),
        ('02', '現状', '住まいるペイント様の現状'),
        ('03', '課題', 'AI活用を阻む3つの壁'),
        ('04', 'ゴール', '目指す未来の姿'),
        ('05', '提案詳細', 'エヌイチからのご提案'),
    ]
    cw = Inches(2.3)
    ch = Inches(4.5)
    gap = Inches(0.18)
    sy = Inches(1.5)
    for i, (num, sec, desc) in enumerate(sections):
        x = Inches(0.5) + i * (cw + gap)
        rect(slide, x, sy, cw, ch, fill=CARD_BG)
        rect(slide, x, sy, cw, Inches(0.07), fill=TEAL)
        txt(slide, num, x + Inches(0.15), sy + Inches(0.12),
            Inches(0.8), Inches(0.55), size=22, bold=True, color=TEAL)
        txt(slide, sec, x + Inches(0.15), sy + Inches(0.65),
            cw - Inches(0.3), Inches(0.5), size=16, bold=True, color=TITLE_CLR)
        txt(slide, desc, x + Inches(0.15), sy + Inches(1.25),
            cw - Inches(0.3), Inches(0.8), size=11, color=SUB_CLR)

# ============================================================
# セクション区切り
# ============================================================
def make_section(prs, num, title, desc=''):
    slide = blank_slide(prs)
    rect(slide, Inches(0), Inches(0), Inches(0.5), H, fill=TEAL)
    txt(slide, num,
        Inches(1.0), Inches(1.2), Inches(5.0), Inches(2.5),
        size=96, bold=True, color=LIGHT_TEAL)
    txt(slide, title,
        Inches(1.2), Inches(3.5), Inches(11.5), Inches(1.2),
        size=38, bold=True, color=TITLE_CLR)
    if desc:
        txt(slide, desc,
            Inches(1.2), Inches(4.65), Inches(11.5), Inches(0.6),
            size=15, color=SUB_CLR)
    txt(slide, '株式会社エヌイチ',
        Inches(9.5), Inches(6.85), Inches(3.5), Inches(0.45),
        size=11, color=LABEL_CLR, align=PP_ALIGN.RIGHT)

# ============================================================
# 参考資料区切り
# ============================================================
def make_ref_divider(prs):
    slide = blank_slide(prs)
    rect(slide, Inches(0), Inches(0), Inches(0.5), H, fill=SLATE)
    txt(slide, '参考資料',
        Inches(1.2), Inches(3.2), Inches(11.5), Inches(1.2),
        size=38, bold=True, color=TITLE_CLR)
    txt(slide, '株式会社エヌイチ',
        Inches(9.5), Inches(6.85), Inches(3.5), Inches(0.45),
        size=11, color=LABEL_CLR, align=PP_ALIGN.RIGHT)

# ============================================================
# スライド4: 振り返り
# ============================================================
def make_review(prs):
    slide = blank_slide(prs)
    header(slide, '本日のお打合せ概要', '振り返り')

    rows = [
        ('参加者', '堺 様（株式会社住まいるペイント） / 大川 龍之介（株式会社エヌイチ）'),
        ('実施日時', '2026年3月30日　Zoom'),
        ('主なトピック',
         '・Claude と Claude Code の違い（定常業務自動化 vs 対話型AI）\n'
         '・ブログ・定型業務の自動化可能性（WordPress投稿・メール送信なども実現可能）\n'
         '・ChatGPT GPT作成の基礎（プロンプト設計の重要性）\n'
         '・Kawaru Teamの体験型研修スタイル（ゲーム要素×AI活用の伏線回収）←「面白いですね」\n'
         '・エヌイチのフルカスタマイズアプローチ（事業ゴール→現状→課題特定→設計）'),
        ('合意したネクストアクション',
         '・エヌイチより本日の議事録＋資料を堺様へ送付\n'
         '・堺様が内容をご確認後、ご興味があればご連絡　→　詳細ヒアリングへ'),
    ]

    row_heights = [Inches(0.6), Inches(0.6), Inches(1.9), Inches(1.05)]
    y = Inches(1.38)
    label_w = Inches(3.0)
    val_w = Inches(9.7)

    for (label, value), rh in zip(rows, row_heights):
        shape_txt(slide, label,
                  Inches(0.5), y, label_w, rh,
                  fill=TEAL, size=13, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
        shape_txt(slide, value,
                  Inches(3.6), y, val_w, rh,
                  fill=CARD_BG, size=12, color=BODY_CLR,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.TOP)
        # 薄いグリッド線
        rect(slide, Inches(0.5), y + rh - Pt(0.5), label_w + val_w + Inches(0.1), Pt(1), fill=BORDER)
        y += rh

# ============================================================
# スライド6: 現状①
# ============================================================
def make_current1(prs):
    slide = blank_slide(prs)
    header(slide, 'すでに3ツールを使いこなし、SEOでも上位表示', '現状①　AI活用実績')

    cards = [
        (TEAL,  '3ツール併用',    'ChatGPT・Gemini・Claude を\n日常的に使いこなし、\n用途に応じて使い分けを実践'),
        (SAGE,  'SEO・MEO 上位',  '自社ホームページで\n検索1〜3位・マップ検索でも\n上位表示を自力で達成'),
        (SLATE, '5つのナレッジDB', '性格・業務ナレッジなど\n5つのDBを整備し\nAIにデータを渡して活用'),
        (NAVY,  '独学で継続学習',  '七里さんなど複数の情報源から\nAI活用を継続学習中\n（助成金知識も保有）'),
    ]

    cw = Inches(2.95)
    ch = Inches(4.5)
    gap = Inches(0.2)
    y = Inches(1.4)

    for i, (color, title, desc) in enumerate(cards):
        x = Inches(0.5) + i * (cw + gap)
        rect(slide, x, y, cw, ch, fill=CARD_BG)
        rect(slide, x, y, cw, Inches(0.07), fill=color)
        txt(slide, title,
            x + Inches(0.18), y + Inches(0.15), cw - Inches(0.35), Inches(0.5),
            size=14, bold=True, color=color)
        txt(slide, desc,
            x + Inches(0.18), y + Inches(0.75), cw - Inches(0.35), Inches(3.5),
            size=13, color=BODY_CLR)

    txt(slide, '※ AI活用への意欲・リテラシーは高い状態からのスタートです',
        Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4),
        size=11, color=LABEL_CLR)

# ============================================================
# スライド7: 現状②
# ============================================================
def make_current2(prs):
    slide = blank_slide(prs)
    header(slide, '業務自動化への意欲は高い、仕組みはこれから', '現状②　自動化の現状')

    col_w = Inches(5.9)
    col_h = Inches(5.2)
    y = Inches(1.35)

    # Left: できていること
    rect(slide, Inches(0.5), y, col_w, col_h, fill=LIGHT_TEAL)
    rect(slide, Inches(0.5), y, col_w, Inches(0.07), fill=SAGE)
    txt(slide, '✓  できていること',
        Inches(0.7), y + Inches(0.12), col_w - Inches(0.3), Inches(0.45),
        size=15, bold=True, color=SAGE)
    done = [
        '3つのAIツールを使いこなしている',
        'SEO・MEO で上位表示を自力で達成',
        '5つのナレッジDBを整備しAIに渡して活用',
        'ブログ記事の執筆にAIを活用中',
        'AIの基礎知識は独学で習得済み',
    ]
    for j, item in enumerate(done):
        txt(slide, f'・  {item}',
            Inches(0.7), y + Inches(0.68) + j * Inches(0.87),
            col_w - Inches(0.3), Inches(0.75), size=13, color=BODY_CLR)

    # Right: まだ手動
    rx = Inches(6.9)
    rect(slide, rx, y, col_w, col_h, fill=LIGHT_RED)
    rect(slide, rx, y, col_w, Inches(0.07), fill=CORAL)
    txt(slide, '✗  まだ手動・仕組み化できていないこと',
        rx + Inches(0.2), y + Inches(0.12), col_w - Inches(0.3), Inches(0.45),
        size=15, bold=True, color=CORAL)
    todo = [
        'ブログ・記事の自動投稿フローがない',
        'Claude Code / AIエージェントの作り方がわからない',
        '定型業務（見積・顧客フォロー）の自動化がない',
        '自社業種特化のAI活用ロードマップがない',
        'AI学習の「正解ルート」が定まっていない',
    ]
    for j, item in enumerate(todo):
        txt(slide, f'・  {item}',
            rx + Inches(0.2), y + Inches(0.68) + j * Inches(0.87),
            col_w - Inches(0.3), Inches(0.75), size=13, color=BODY_CLR)

# ============================================================
# スライド9: 課題① ロジックツリー
# ============================================================
def make_issue_tree(prs):
    slide = blank_slide(prs)
    header(slide, '課題は3層構造', '課題①　全体マップ')

    # Root
    shape_txt(slide,
              'リフォーム会社として\nAIで収益を上げる\n仕組みが作れていない',
              Inches(0.4), Inches(2.6), Inches(2.8), Inches(1.4),
              fill=NAVY, size=12, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER)

    # Connector line (root → middle)
    rect(slide, Inches(3.2), Inches(3.25), Inches(0.3), Pt(1), fill=NAVY)

    # Middle nodes
    mid_data = [
        (CORAL, '① AI活用の\n「正解」が見えない', Inches(1.35)),
        (SLATE, '② ルーティン業務が\nまだ手動で動いている', Inches(2.95)),
        (TEAL,  '③ AI活用戦略を\n体系化する伴走者がいない', Inches(4.55)),
    ]
    detail_data = [
        ['情報源が多すぎ優先順位不明【事実】',
         'Claude/Code/GPT使い分け曖昧【事実】',
         '業種特化のロードマップがない【推察】'],
        ['ブログ自動化フローがない【事実】',
         'AIエージェントの作り方不明【事実】',
         '見積・顧客フォローが手動【推察】'],
        ['独学で学んでも実装で止まる【推察】',
         '助成金活用が未検討【推察】',
         ''],
    ]

    for (color, label, my), dets in zip(mid_data, detail_data):
        # Vertical connector
        rect(slide, Inches(3.5), my + Inches(0.45), Pt(1), Inches(0.35), fill=color)
        # Middle node
        shape_txt(slide, label,
                  Inches(3.55), my, Inches(2.5), Inches(0.9),
                  fill=color, size=12, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
        # Detail items
        for k, det in enumerate(dets):
            if det:
                dy = my + k * Inches(0.5)
                shape_txt(slide, f'  {det}',
                          Inches(6.2), dy, Inches(6.8), Inches(0.45),
                          fill=CARD_BG, size=10, color=BODY_CLR,
                          line=BORDER, lw=Pt(0.5),
                          v_anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# スライド10: 課題②
# ============================================================
def make_issue2(prs):
    slide = blank_slide(prs)
    header(slide, 'AI情報の多さが、行動を止めている', '課題②　正解不明問題')

    cards = [
        (CORAL, '情報の氾濫',
         '「AIを学ぶ」コンテンツが多すぎて\n何が自分に合っているか判断できない\n\n'
         'ChatGPT・Gemini・Claude・七里さんなど\n複数の学習源を並列で追いかけている状態'),
        (SLATE, '判断の停止',
         '「何が正解か」がわからず\n行動の優先順位が決まらない\n\n'
         '「一番結果が出るのがどれか\nわからない」（商談での発言）'),
        (NAVY,  '機会損失',
         'AI活用による業務効率化・\n収益向上のスタートが遅れている\n\n'
         'リフォーム業界でAIを先行活用する\nライバル業者との差が広がるリスク'),
    ]

    cw = Inches(3.8)
    ch = Inches(5.0)
    gap = Inches(0.35)
    y = Inches(1.4)

    for i, (color, title, desc) in enumerate(cards):
        x = Inches(0.5) + i * (cw + gap)
        rect(slide, x, y, cw, ch, fill=CARD_BG)
        rect(slide, x, y, cw, Inches(0.07), fill=color)
        txt(slide, title,
            x + Inches(0.2), y + Inches(0.14), cw - Inches(0.4), Inches(0.5),
            size=16, bold=True, color=color)
        txt(slide, desc,
            x + Inches(0.2), y + Inches(0.75), cw - Inches(0.4), Inches(4.0),
            size=12, color=BODY_CLR)
        if i < 2:
            txt(slide, '→',
                x + cw + Inches(0.05), y + ch / 2 - Inches(0.25),
                Inches(0.28), Inches(0.5),
                size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ============================================================
# スライド11: 課題③
# ============================================================
def make_issue3(prs):
    slide = blank_slide(prs)
    header(slide, '業種特化の自動化がまだ手つかず', '課題③　業務自動化の壁')

    col_w = Inches(5.8)
    col_h = Inches(5.1)
    y = Inches(1.38)

    # Before
    rect(slide, Inches(0.5), y, col_w, col_h, fill=LIGHT_RED)
    rect(slide, Inches(0.5), y, col_w, Inches(0.07), fill=CORAL)
    txt(slide, 'Before（現状）',
        Inches(0.7), y + Inches(0.12), col_w - Inches(0.3), Inches(0.45),
        size=15, bold=True, color=CORAL)
    before = [
        'ブログ記事は毎回手動でAIに依頼→コピペ',
        '顧客への返信はその都度手打ち',
        '見積書作成は感覚・経験ベース',
        '施工事例の記事化ができていない',
        'AI活用の「型」がなく属人的なまま',
    ]
    for j, item in enumerate(before):
        txt(slide, f'✗   {item}',
            Inches(0.7), y + Inches(0.68) + j * Inches(0.88),
            col_w - Inches(0.3), Inches(0.75), size=13, color=BODY_CLR)

    # Arrow
    txt(slide, '→',
        Inches(6.4), y + col_h / 2 - Inches(0.3),
        Inches(0.5), Inches(0.6),
        size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # After
    ax = Inches(7.1)
    rect(slide, ax, y, col_w, col_h, fill=LIGHT_GRN)
    rect(slide, ax, y, col_w, Inches(0.07), fill=SAGE)
    txt(slide, 'After（目指す姿）',
        ax + Inches(0.2), y + Inches(0.12), col_w - Inches(0.3), Inches(0.45),
        size=15, bold=True, color=SAGE)
    after = [
        '写真を撮るだけでSEO記事が自動生成・投稿',
        '問い合わせ内容に応じた返信を自動ドラフト',
        '顧客情報から見積書を自動作成',
        '施工事例がコンテンツ資産として蓄積され続ける',
        'AI活用の「型」が社内に定着',
    ]
    for j, item in enumerate(after):
        txt(slide, f'✓   {item}',
            ax + Inches(0.2), y + Inches(0.68) + j * Inches(0.88),
            col_w - Inches(0.3), Inches(0.75), size=13, color=BODY_CLR)

# ============================================================
# スライド13: ゴール①
# ============================================================
def make_goal1(prs):
    slide = blank_slide(prs)
    header(slide, '住まいるペイントがAIで実現する未来', 'ゴール　目指す姿')

    goals = [
        (TEAL, '01', '集客',
         '施工事例・ブログが自動で蓄積されSEO資産が\n増え続ける仕組みへ\n→ 問い合わせが自動で増える構造を実現'),
        (SAGE, '02', '業務効率',
         '見積・顧客フォロー・記事作成の定型業務が\n半自動化される\n→ 堺さんが本質的な仕事に集中できる時間を創出'),
        (NAVY, '03', '収益向上',
         'AI活用の型が定まり、業務品質が\n属人から標準化へ\n→ スタッフへの横展開で組織全体の生産性向上'),
    ]

    cw = Inches(3.95)
    ch = Inches(4.5)
    gap = Inches(0.2)
    y = Inches(1.45)

    for i, (color, num, title, desc) in enumerate(goals):
        x = Inches(0.5) + i * (cw + gap)
        rect(slide, x, y, cw, ch, fill=CARD_BG)
        rect(slide, x, y, cw, Inches(0.07), fill=color)
        txt(slide, num,
            x + Inches(0.18), y + Inches(0.13), Inches(0.5), Inches(0.5),
            size=20, bold=True, color=color)
        txt(slide, title,
            x + Inches(0.65), y + Inches(0.13), cw - Inches(0.8), Inches(0.5),
            size=18, bold=True, color=color)
        txt(slide, desc,
            x + Inches(0.18), y + Inches(0.75), cw - Inches(0.35), Inches(3.5),
            size=13, color=BODY_CLR)

# ============================================================
# スライド14: ゴール②
# ============================================================
def make_goal2(prs):
    slide = blank_slide(prs)
    header(slide, 'リフォーム会社×AI：3つの成果指標', 'ゴール　KPI')

    kpis = [
        (TEAL, '集客KPI',  '月間問い合わせ数',   '現状比 +30%',       'SEO記事の自動蓄積による\n検索流入増'),
        (SAGE, '業務KPI',  '定型業務の月間工数',  '月20時間以上削減',   '見積・記事・返信の\n半自動化による時間創出'),
        (NAVY, '収益KPI',  'ROI（投資対効果）',   '初年度 3〜4倍',      '業務効率化×集客強化の\n複合効果'),
    ]

    cw = Inches(3.95)
    ch = Inches(4.8)
    gap = Inches(0.2)
    y = Inches(1.45)

    for i, (color, cat, metric, target, reason) in enumerate(kpis):
        x = Inches(0.5) + i * (cw + gap)
        rect(slide, x, y, cw, ch, fill=CARD_BG)
        rect(slide, x, y, cw, Inches(0.07), fill=color)
        txt(slide, cat,
            x + Inches(0.2), y + Inches(0.14), cw - Inches(0.4), Inches(0.4),
            size=13, bold=True, color=color)
        txt(slide, metric,
            x + Inches(0.2), y + Inches(0.62), cw - Inches(0.4), Inches(0.45),
            size=13, color=SUB_CLR)
        txt(slide, target,
            x + Inches(0.2), y + Inches(1.12), cw - Inches(0.4), Inches(0.9),
            size=26, bold=True, color=color)
        rect(slide, x + Inches(0.2), y + Inches(2.1), cw - Inches(0.4), Pt(1), fill=BORDER)
        txt(slide, reason,
            x + Inches(0.2), y + Inches(2.25), cw - Inches(0.4), Inches(2.3),
            size=12, color=BODY_CLR)

# ============================================================
# スライド16: 提案根拠①
# ============================================================
def make_reason1(prs):
    slide = blank_slide(prs)
    header(slide, '型化研修では現場に落とせない', '提案根拠①　業界の課題')

    headers_row = ['比較項目', '一般的なAI研修', 'エヌイチのアプローチ']
    rows = [
        ('カリキュラム設計', 'あらかじめ決まった型を提供', '事業・組織・業種からゼロ設計'),
        ('現場への落とし込み', '研修で終了。実装支援なし', '「明日から使える」まで伴走'),
        ('カスタマイズ深度', '業種・役職でざっくり分類', '担当者の業務内容まで個別設計'),
        ('研修形式', '動画・eラーニングが中心', '体験型・ゲーム要素入りの双方向型'),
        ('成果の確認', '受講完了で終わり', '活用定着・改善まで継続支援'),
    ]

    col_w = [Inches(2.5), Inches(4.6), Inches(4.6)]
    row_h = Inches(0.72)
    sx = Inches(0.5)
    sy = Inches(1.38)
    header_colors = [NAVY, SLATE, TEAL]

    for j, (h, w, hc) in enumerate(zip(headers_row, col_w, header_colors)):
        x = sx + sum(col_w[:j])
        shape_txt(slide, h, x, sy, w, row_h,
                  fill=hc, size=13, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        bg = CARD_BG if i % 2 == 0 else WHITE
        for j, (cell, w) in enumerate(zip(row, col_w)):
            x = sx + sum(col_w[:j])
            y = sy + (i + 1) * row_h
            c = TEAL if j == 2 else BODY_CLR
            b = True if j == 0 else False
            align = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
            shape_txt(slide, f'  {cell}' if j > 0 else cell,
                      x, y, w, row_h,
                      fill=bg, size=12, bold=b, color=c,
                      align=align, line=BORDER, lw=Pt(0.5))

# ============================================================
# スライド17: 提案根拠②
# ============================================================
def make_reason2(prs):
    slide = blank_slide(prs)
    header(slide, 'リフォーム業界への支援実績と強み', '提案根拠②　エヌイチの強み')

    cards = [
        (TEAL,  '業種問わずフルカスタマイズ',
         'エンジニア向け・現場スタッフ向け・管理職向けを問わず、\n業種業界に関係なく対応。\nリフォーム業の実務に合わせた研修設計が可能。'),
        (SAGE,  'キリングループ様：満足度8.5点',
         'AI未経験スタッフへの体験型研修を実施。\n「全員が業務でAIを使いたい」と回答。\n和気あいあいとしたゲーム要素入りの研修設計が好評。'),
        (SLATE, 'ミズカラ様：月283h→85hへ削減',
         '問い合わせ対応業務を月283時間から85時間へ削減。\n約200時間（25営業日分）の工数削減を実現。\n現場への落とし込みまで伴走した成果。'),
    ]

    cw = Inches(3.95)
    ch = Inches(5.0)
    gap = Inches(0.2)
    y = Inches(1.45)

    for i, (color, title, desc) in enumerate(cards):
        x = Inches(0.5) + i * (cw + gap)
        rect(slide, x, y, cw, ch, fill=CARD_BG)
        rect(slide, x, y, cw, Inches(0.08), fill=color)
        txt(slide, title,
            x + Inches(0.2), y + Inches(0.15), cw - Inches(0.4), Inches(0.6),
            size=15, bold=True, color=color)
        txt(slide, desc,
            x + Inches(0.2), y + Inches(0.9), cw - Inches(0.4), Inches(3.8),
            size=13, color=BODY_CLR)

# ============================================================
# スライド18: Kawaru Team
# ============================================================
def make_kawaru_team(prs):
    slide = blank_slide(prs)
    header(slide, 'フルカスタマイズのAI研修', '提案詳細①　Kawaru Team')

    y = Inches(1.38)

    # 差別化ラベル
    shape_txt(slide, '他との差別化ポイント',
              Inches(0.5), y, Inches(3.8), Inches(0.42),
              fill=TEAL, size=13, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER)
    diffs = [
        '型化された研修ではなく、御社の業務・目標・メンバー構成からゼロ設計',
        '「事業ゴール → 現状把握 → 課題特定 → カリキュラム設計」の順で構築',
        'ゲーム要素・体験型でAI活用の伏線回収をする設計（堺様が「面白いですね」と反応）',
    ]
    for j, d in enumerate(diffs):
        txt(slide, f'・  {d}',
            Inches(0.5), y + Inches(0.52) + j * Inches(0.58),
            Inches(12.3), Inches(0.52), size=13, color=BODY_CLR)

    rect(slide, Inches(0.5), y + Inches(2.3), Inches(12.3), Pt(1), fill=BORDER)

    # 成果ラベル
    oy = y + Inches(2.5)
    shape_txt(slide, 'どうなるか（成果）',
              Inches(0.5), oy, Inches(3.8), Inches(0.42),
              fill=SAGE, size=13, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER)
    outcomes = [
        '「明日から自分の業務でAIが使える」状態になる',
        'キリングループ様：満足度8.5点 / 「全員が業務でAIを使いたい」と回答',
        '参加スタッフ全員が自走できるAIリテラシーを獲得',
    ]
    for j, o in enumerate(outcomes):
        txt(slide, f'✓   {o}',
            Inches(0.5), oy + Inches(0.52) + j * Inches(0.58),
            Inches(12.3), Inches(0.52), size=13, color=BODY_CLR)

    # 3ステップフロー
    steps = [
        ('01', 'ヒアリング',  '事業ゴール・現状・\n課題を特定'),
        ('02', 'カリキュラム設計', 'フルカスタマイズで\n完全ゼロ設計'),
        ('03', '研修実施',   '体験型研修を実施\n→ 定着まで支援'),
    ]
    step_colors = [TEAL, SLATE, SAGE]
    sw = Inches(3.7)
    fy = y + Inches(4.55)
    for i, (num, title, desc) in enumerate(steps):
        sx = Inches(0.5) + i * (sw + Inches(0.5))
        rect(slide, sx, fy, sw, Inches(1.15), fill=step_colors[i])
        txt(slide, f'{num}　{title}',
            sx + Inches(0.15), fy + Inches(0.1), sw - Inches(0.3), Inches(0.5),
            size=15, bold=True, color=WHITE)
        txt(slide, desc,
            sx + Inches(0.15), fy + Inches(0.6), sw - Inches(0.3), Inches(0.5),
            size=11, color=WHITE)
        if i < 2:
            txt(slide, '→',
                sx + sw + Inches(0.08), fy + Inches(0.38),
                Inches(0.4), Inches(0.4),
                size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ============================================================
# スライド19: Kawaru Coach
# ============================================================
def make_kawaru_coach(prs):
    slide = blank_slide(prs)
    header(slide, '堺さん自身がAIプロになる', '提案詳細②　Kawaru Coach')

    y = Inches(1.38)

    # 差別化ラベル
    shape_txt(slide, '他との差別化ポイント',
              Inches(0.5), y, Inches(3.8), Inches(0.42),
              fill=TEAL, size=13, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER)
    diffs = [
        '一般的なコンサルは「提案だけ」。Kawaru Coachは実装・定着まで1on1で伴走',
        '経営・事業・組織の視点でAI活用を設計。技術支援だけでなくビジネス視点のサポート',
        '住まいるペイント専用のAI活用ロードマップをゼロから構築する',
    ]
    for j, d in enumerate(diffs):
        txt(slide, f'・  {d}',
            Inches(0.5), y + Inches(0.52) + j * Inches(0.58),
            Inches(12.3), Inches(0.52), size=13, color=BODY_CLR)

    rect(slide, Inches(0.5), y + Inches(2.3), Inches(12.3), Pt(1), fill=BORDER)

    # 成果ラベル
    oy = y + Inches(2.5)
    shape_txt(slide, 'どうなるか（成果）',
              Inches(0.5), oy, Inches(3.8), Inches(0.42),
              fill=SAGE, size=13, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER)
    outcomes = [
        '堺さん自身がAIプロになり、自社業務の自動化を自走できる状態になる',
        'ブログ自動化・見積効率化など具体業務への実装まで支援',
        '属人業務70%削減・ROI約4倍の実績（Kawaru Coach 支援実績より）',
    ]
    for j, o in enumerate(outcomes):
        txt(slide, f'✓   {o}',
            Inches(0.5), oy + Inches(0.52) + j * Inches(0.58),
            Inches(12.3), Inches(0.52), size=13, color=BODY_CLR)

    # 引用カード
    qy = oy + Inches(2.3)
    rect(slide, Inches(0.5), qy, Inches(12.3), Inches(1.05), fill=LIGHT_TEAL)
    txt(slide,
        '「AIが"よくわからない"という経営者にこそ使ってほしい。'
        'ITリテラシーが低い会社でも、ちゃんと導入・運用まで伴走してくれます」\n'
        '— 株式会社BIRDY 代表取締役 鳥屋直弘様',
        Inches(0.8), qy + Inches(0.12), Inches(11.8), Inches(0.82),
        size=12, color=SLATE)

# ============================================================
# スライド20: ネクストアクション
# ============================================================
def make_next_action(prs):
    slide = blank_slide(prs)
    header(slide, '次のステップ', 'ネクストアクション')

    steps = [
        (TEAL,  '01', '資料のご確認',
         '本日共有の議事録・本資料を\nご確認ください\n\nご質問はいつでもお気軽にご連絡ください'),
        (SLATE, '02', '詳細ヒアリング（30〜60分）',
         '住まいるペイント様の\n事業ゴール・現状・課題を\nお聞かせいただきます\n（無料でのご対応です）'),
        (SAGE,  '03', 'AI活用ロードマップのご提示',
         'ヒアリング内容をもとに\n住まいるペイント様専用の\nAI活用ロードマップと\n最適なご提案をお届けします'),
    ]

    cw = Inches(3.95)
    ch = Inches(5.0)
    gap = Inches(0.2)
    y = Inches(1.45)

    for i, (color, num, title, desc) in enumerate(steps):
        x = Inches(0.5) + i * (cw + gap)
        rect(slide, x, y, cw, ch, fill=CARD_BG)
        rect(slide, x, y, cw, Inches(0.08), fill=color)
        # バッジ（右上）
        shape_txt(slide, num,
                  x + cw - Inches(0.75), y + Inches(0.12),
                  Inches(0.58), Inches(0.58),
                  fill=color, size=14, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
        txt(slide, title,
            x + Inches(0.2), y + Inches(0.85), cw - Inches(0.4), Inches(0.75),
            size=15, bold=True, color=color)
        txt(slide, desc,
            x + Inches(0.2), y + Inches(1.7), cw - Inches(0.4), Inches(3.0),
            size=13, color=BODY_CLR)
        if i < 2:
            txt(slide, '→',
                x + cw + Inches(0.0), y + ch / 2 - Inches(0.25),
                Inches(0.22), Inches(0.5),
                size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ============================================================
# 締めスライド
# ============================================================
def make_closing(prs):
    slide = blank_slide(prs)
    rect(slide, Inches(0), Inches(0), Inches(0.5), H, fill=TEAL)
    txt(slide,
        '私たちが提供するのは「ツール」や「研修」ではありません。',
        Inches(1.0), Inches(2.0), Inches(11.5), Inches(0.85),
        size=21, bold=True, color=TITLE_CLR)
    txt(slide,
        '貴社の未来を創る「AI人材」と、組織全体の「成長」です。',
        Inches(1.0), Inches(2.85), Inches(11.5), Inches(0.85),
        size=21, bold=True, color=TEAL)
    rect(slide, Inches(1.0), Inches(3.85), Inches(5.0), Pt(2), fill=TEAL)
    txt(slide, '株式会社エヌイチ',
        Inches(1.0), Inches(4.1), Inches(5.0), Inches(0.5),
        size=15, color=SUB_CLR)

# ============================================================
# 会社概要
# ============================================================
def make_company(prs):
    slide = blank_slide(prs)
    header(slide, '会社概要')

    table_data = [
        ('会社名',   '株式会社エヌイチ'),
        ('経営陣',   '代表取締役CEO 奥山幸生 / 取締役COO 髙橋悠 / 取締役CHRO 米澤浩明'),
        ('資本金',   '1,000万円'),
        ('設立',     '2023年10月'),
        ('所在地',   '東京都新宿区西新宿3-7-30 フロンティアグラン西新宿12F'),
        ('Mission', '新しい働く『カタチ』を創る'),
        ('Vision',  'AI（アイ）ある社会に'),
    ]

    lw = Inches(2.0)
    vw = Inches(5.5)
    rh = Inches(0.63)
    sy = Inches(1.35)

    for i, (label, value) in enumerate(table_data):
        bg = CARD_BG if i % 2 == 0 else WHITE
        y = sy + i * rh
        shape_txt(slide, label, Inches(0.5), y, lw, rh,
                  fill=TEAL, size=12, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
        shape_txt(slide, f'  {value}', Inches(2.6), y, vw, rh,
                  fill=bg, size=11, color=BODY_CLR,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE)

    services = [
        (TEAL,  'Kawaru',       '業務自動化SaaS'),
        (SLATE, 'Kawaru Team',  'フルカスタマイズAI研修'),
        (SAGE,  'Kawaru BPO',   'AI業務効率化代行'),
        (NAVY,  'Kawaru Coach', 'AI顧問サービス'),
    ]
    svx = Inches(8.5)
    svw = Inches(4.5)
    svh = Inches(1.1)
    svgap = Inches(0.15)
    svy = Inches(1.5)

    for i, (color, name, desc) in enumerate(services):
        y = svy + i * (svh + svgap)
        rect(slide, svx, y, svw, svh, fill=CARD_BG)
        rect(slide, svx, y, Inches(0.07), svh, fill=color)
        txt(slide, name,
            svx + Inches(0.2), y + Inches(0.08), svw - Inches(0.3), Inches(0.5),
            size=14, bold=True, color=color)
        txt(slide, desc,
            svx + Inches(0.2), y + Inches(0.58), svw - Inches(0.3), Inches(0.45),
            size=12, color=BODY_CLR)

# ============================================================
# 導入実績
# ============================================================
def make_achievements(prs):
    slide = blank_slide(prs)
    header(slide, '導入実績', 'あらゆる規模の企業様に導入されています')

    companies = [
        '株式会社ブロードリンク', '株式会社中西製作所', '株式会社ミズカラ',
        '株式会社エル・ティー・エス リンク', 'FRUOR株式会社',
        '株式会社BIRDY', '株式会社オンシナジー', 'キリングループ',
    ]
    cmpw = Inches(1.45)
    cmph = Inches(0.5)
    cmpgap = Inches(0.1)
    cy = Inches(1.35)
    for i, company in enumerate(companies):
        x = Inches(0.5) + i * (cmpw + cmpgap)
        shape_txt(slide, company, x, cy, cmpw, cmph,
                  fill=CARD_BG, size=9, color=BODY_CLR,
                  align=PP_ALIGN.CENTER, line=BORDER, lw=Pt(0.5))

    cases = [
        (TEAL,  'キリングループ様',
         'AI未経験・活用イメージなし',
         '全員が「業務でAIを使いたい」と回答\n満足度8.5点'),
        (SAGE,  '株式会社ミズカラ様',
         '問い合わせ対応に月283時間',
         '月85時間へ削減\n（約200時間＝25営業日分の削減）'),
        (SLATE, '株式会社BIRDY様',
         '業務の属人化・営業効率の壁',
         'バックオフィス効率化\n経営資源を成長領域にシフト'),
    ]

    casew = Inches(4.0)
    caseh = Inches(4.8)
    casegap = Inches(0.25)
    casey = Inches(2.1)

    for i, (color, company, issue, result) in enumerate(cases):
        x = Inches(0.5) + i * (casew + casegap)
        rect(slide, x, casey, casew, caseh, fill=CARD_BG)
        rect(slide, x, casey, casew, Inches(0.07), fill=color)
        txt(slide, company,
            x + Inches(0.2), casey + Inches(0.14), casew - Inches(0.4), Inches(0.5),
            size=14, bold=True, color=color)
        shape_txt(slide, '課題',
                  x + Inches(0.2), casey + Inches(0.75), Inches(1.1), Inches(0.32),
                  fill=CORAL, size=10, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
        txt(slide, issue,
            x + Inches(0.2), casey + Inches(1.15), casew - Inches(0.4), Inches(0.7),
            size=12, color=BODY_CLR)
        shape_txt(slide, '成果',
                  x + Inches(0.2), casey + Inches(1.95), Inches(1.1), Inches(0.32),
                  fill=SAGE, size=10, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
        txt(slide, result,
            x + Inches(0.2), casey + Inches(2.38), casew - Inches(0.4), Inches(2.2),
            size=13, bold=True, color=color)

# ============================================================
# メイン実行
# ============================================================
prs = new_prs()

make_cover(prs)          # 1: 表紙
make_toc(prs)            # 2: 目次
make_section(prs, '01', '振り返り', '本日のお打合せ概要')  # 3
make_review(prs)         # 4: 振り返り
make_section(prs, '02', '現状', '住まいるペイント様の現状')  # 5
make_current1(prs)       # 6: 現状①
make_current2(prs)       # 7: 現状②
make_section(prs, '03', '課題', 'AI活用を阻む3つの壁')  # 8
make_issue_tree(prs)     # 9: 課題① ロジックツリー
make_issue2(prs)         # 10: 課題②
make_issue3(prs)         # 11: 課題③
make_section(prs, '04', 'ゴール', '目指す未来の姿')  # 12
make_goal1(prs)          # 13: ゴール①
make_goal2(prs)          # 14: ゴール②
make_section(prs, '05', '提案詳細', 'エヌイチからのご提案')  # 15
make_reason1(prs)        # 16: 提案根拠①
make_reason2(prs)        # 17: 提案根拠②
make_kawaru_team(prs)    # 18: Kawaru Team
make_kawaru_coach(prs)   # 19: Kawaru Coach
make_next_action(prs)    # 20: ネクストアクション
make_closing(prs)        # 21: 締めスライド
make_ref_divider(prs)    # 22: 参考資料区切り
make_company(prs)        # 23: 会社概要
make_achievements(prs)   # 24: 導入実績

output = '/Users/kyouyuu/Downloads/proposal_sumairupaint_20260330.pptx'
prs.save(output)
print(f'✅ 保存完了: {output}')
print(f'📊 総スライド数: {len(prs.slides)}枚')
