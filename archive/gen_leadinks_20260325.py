"""
リードインクス株式会社 ご提案書 2026-03-25
骨子設計（Step3）に基づく19枚構成
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── カラーパレット（統一トーン: 寒色系4系統） ─────────────
# ブランド: TEAL（メインアクセント）
TEAL    = RGBColor(0x42, 0xA4, 0xAF)   # ティール
TEAL_L  = RGBColor(0xE0, 0xF4, 0xF6)   # ティール薄
# サブアクセント: SLATE（寒色グレーブルー）→ 旧BLUE/ORANGE代替
SLATE   = RGBColor(0x4A, 0x6C, 0x8C)   # スレート
SLATE_L = RGBColor(0xEB, 0xF0, 0xF6)   # スレート薄
# 機能色①: CORAL（課題・問題）→ 旧RED代替
CORAL   = RGBColor(0xC9, 0x54, 0x54)   # コーラル
CORAL_L = RGBColor(0xFD, 0xEF, 0xEF)   # コーラル薄
# 機能色②: SAGE（ポジティブ・解決）→ 旧GREEN代替
SAGE    = RGBColor(0x3A, 0x8F, 0x63)   # セージ
SAGE_L  = RGBColor(0xED, 0xF7, 0xF2)   # セージ薄
# ニュートラル
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x1E, 0x29, 0x3B)
DARK    = RGBColor(0x11, 0x18, 0x27)
NAVY    = RGBColor(0x1C, 0x35, 0x57)   # 濃紺（ステップ最終・強調）
GRAY    = RGBColor(0x6B, 0x72, 0x80)
LGRAY   = RGBColor(0xF1, 0xF5, 0xF9)
MGRAY   = RGBColor(0xCB, 0xD5, 0xE1)
# 旧名エイリアス（既存コードとの互換のため）
RED     = CORAL;  RED_L    = CORAL_L
GREEN   = SAGE;   GREEN_L  = SAGE_L
BLUE    = SLATE;  BLUE_L   = SLATE_L
ORANGE  = SLATE;  ORANGE_L = SLATE_L
DARK_H  = NAVY

FONT = 'Noto Sans JP'
SW, SH = 13.33, 7.50

# ── 基本ユーティリティ ───────────────────────────────

def prs_new():
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, x, y, w, h, fill, line=None, line_w=0.75):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

# フォントサイズ スケールマップ（小〜中サイズを+2pt拡大、大サイズは維持）
_FSZ = {8: 10, 9: 11, 10: 12, 11: 13, 12: 14, 13: 15, 14: 16, 15: 17}

def txt(slide, x, y, w, h, text, size, color=BLACK, bold=False,
        align=PP_ALIGN.LEFT, va=MSO_ANCHOR.TOP, ml=0.10, mt=0.06):
    size = _FSZ.get(size, size)  # スケール適用
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = va
    tf.margin_left   = Inches(ml)
    tf.margin_right  = Inches(0.06)
    tf.margin_top    = Inches(mt)
    tf.margin_bottom = Inches(0.04)
    for i, line in enumerate(str(text).split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(size)
        r.font.color.rgb = color
        r.font.bold  = bold
        r.font.name  = FONT
    return tb

def header(slide, title, subtitle=None, dark=False):
    if dark:
        rect(slide, 0, 0, SW, SH, DARK_H)
    rect(slide, 0, 0, SW, 0.06, TEAL)
    tc = WHITE if dark else GRAY
    txt(slide, 9.80, 0.10, 3.30, 0.40, '株式会社エヌイチ', 11, tc,
        align=PP_ALIGN.RIGHT)
    title_c = WHITE if dark else DARK
    txt(slide, 0.45, 0.18, 10.50, 0.75, title, 25, title_c, bold=True,
        va=MSO_ANCHOR.MIDDLE)
    sub_c = RGBColor(0xB0, 0xC4, 0xD8) if dark else GRAY
    if subtitle:
        txt(slide, 0.45, 0.96, 12.43, 0.38, subtitle, 13, sub_c)
        rect(slide, 0.45, 1.38, 12.43, 0.03, TEAL)
        return 1.52
    rect(slide, 0.45, 0.96, 12.43, 0.03, TEAL)
    return 1.10

def card(slide, x, y, w, h, fill=WHITE, border=MGRAY, bar=None):
    r = rect(slide, x, y, w, h, fill, border)
    if bar:
        rect(slide, x, y, 0.07, h, bar)
    return r

def label_rows(slide, pairs, x, y, lw, vw, row_h=0.50, gap=0.08,
               lbg=TEAL_L, lcolor=TEAL):
    for i, (lbl, val) in enumerate(pairs):
        iy = y + i * (row_h + gap)
        rect(slide, x, iy, lw, row_h, lbg)
        txt(slide, x + 0.06, iy, lw - 0.10, row_h, lbl, 10,
            lcolor, bold=True, va=MSO_ANCHOR.MIDDLE)
        rect(slide, x + lw + 0.06, iy, vw, row_h, WHITE, MGRAY)
        txt(slide, x + lw + 0.14, iy, vw - 0.18, row_h, val, 12,
            BLACK, va=MSO_ANCHOR.MIDDLE)

def num_badge(slide, cx, cy, num, bg=TEAL):
    r = 0.28
    s = slide.shapes.add_shape(9,
        Inches(cx - r), Inches(cy - r), Inches(r*2), Inches(r*2))
    s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = str(num)
    run.font.size = Pt(16); run.font.bold = True
    run.font.color.rgb = WHITE; run.font.name = FONT

# ════════════════════════════════════════════════════
# 目次・セクション区切りスライド ヘルパー
# ════════════════════════════════════════════════════

def toc_slide(prs, sections):
    """目次スライド（全体構成 5セクション一覧）"""
    s = blank(prs)
    rect(s, 0, 0, 0.40, SH, TEAL)
    rect(s, 0.40, 0, SW - 0.40, 0.08, TEAL)
    txt(s, 0.60, 0.18, 10.50, 0.70, '目次', 28, DARK, bold=True, va=MSO_ANCHOR.MIDDLE)
    txt(s, 10.20, 0.10, 3.00, 0.40, '株式会社エヌイチ', 11, GRAY, align=PP_ALIGN.RIGHT)
    rect(s, 0.60, 0.90, 12.73, 0.03, TEAL)
    for i, (num, title, desc) in enumerate(sections):
        ty = 1.10 + i * 1.18
        bg = TEAL_L if i % 2 == 0 else LGRAY
        card(s, 0.60, ty, 12.73, 1.04, bg, MGRAY, TEAL)
        # セクション番号バッジ
        rect(s, 0.60, ty, 1.10, 1.04, TEAL)
        txt(s, 0.60, ty, 1.10, 1.04, num, 26, WHITE, bold=True,
            align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
        # 仕切り
        rect(s, 1.74, ty + 0.18, 0.04, 0.68, TEAL)
        # タイトルと説明
        txt(s, 1.92, ty + 0.10, 5.00, 0.46, title, 16, NAVY, bold=True)
        txt(s, 1.92, ty + 0.60, 11.00, 0.38, desc, 12, GRAY)
    return s


def section_divider(prs, section_num, section_title, section_desc=''):
    """セクション区切りスライド（白背景・左TEALバー・大きな番号）"""
    s = blank(prs)
    # 左アクセントバー
    rect(s, 0, 0, 0.55, SH, TEAL)
    # 上部細ライン
    rect(s, 0.55, 0, SW - 0.55, 0.08, TEAL)
    # 会社名
    txt(s, 10.20, 0.10, 3.00, 0.40, '株式会社エヌイチ', 11, GRAY, align=PP_ALIGN.RIGHT)
    # 大きな背景数字（装飾）
    txt(s, 0.70, 1.20, 5.50, 4.50, section_num, 130, LGRAY, bold=True, va=MSO_ANCHOR.MIDDLE)
    # セクションタイトル
    txt(s, 5.60, 2.20, 7.30, 1.60, section_title, 38, NAVY, bold=True, va=MSO_ANCHOR.MIDDLE)
    # TEAL区切り線
    rect(s, 5.60, 4.00, 7.30, 0.05, TEAL)
    # セクション説明
    if section_desc:
        txt(s, 5.60, 4.20, 7.30, 0.60, section_desc, 14, GRAY)
    return s


def ref_divider(prs):
    """参考資料 区切りスライド"""
    s = blank(prs)
    rect(s, 0, 0, 0.55, SH, MGRAY)
    rect(s, 0.55, 0, SW - 0.55, 0.08, MGRAY)
    txt(s, 10.20, 0.10, 3.00, 0.40, '株式会社エヌイチ', 11, GRAY, align=PP_ALIGN.RIGHT)
    txt(s, 0.70, 1.20, 5.50, 4.50, 'REF', 110, LGRAY, bold=True, va=MSO_ANCHOR.MIDDLE)
    txt(s, 5.60, 2.40, 7.30, 1.40, '参考資料', 38, GRAY, bold=True, va=MSO_ANCHOR.MIDDLE)
    rect(s, 5.60, 4.00, 7.30, 0.05, MGRAY)
    txt(s, 5.60, 4.20, 7.30, 0.60, '会社概要・導入実績', 14, GRAY)
    return s


# ════════════════════════════════════════════════════
# スライド生成開始
# ════════════════════════════════════════════════════

prs = prs_new()

# ═══════════════════════════════════════════════════════
# 固定スライド1: 表紙
# ═══════════════════════════════════════════════════════
s = blank(prs)
rect(s, 0, 0, 0.40, SH, TEAL)
rect(s, 0.40, SH - 0.08, SW - 0.40, 0.08, TEAL)

txt(s, 1.00, 1.50, 11.80, 0.55, '生成AI支援 ご提案書', 16, TEAL, bold=True,
    align=PP_ALIGN.CENTER)
rect(s, 3.80, 2.15, 5.73, 0.04, TEAL)
txt(s, 1.00, 2.30, 11.80, 1.20,
    'リードインクス株式会社 様', 38, DARK, bold=True,
    align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
txt(s, 1.00, 3.65, 11.80, 0.60,
    '「AIで自走できるエンジニア組織」を実現するための提案', 17, GRAY,
    align=PP_ALIGN.CENTER)
txt(s, 1.00, 6.40, 11.80, 0.45,
    '2026年3月25日　　株式会社エヌイチ', 13, GRAY,
    align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════
# Slide 2: 目次
# ═══════════════════════════════════════════════════════
toc_sections = [
    ('01', '振り返り',  '前回お打合せの振り返り・ヒアリング内容の確認'),
    ('02', '現状',      'リードインクス様のビジネス環境・AI活用の現状分析'),
    ('03', '課題',      '3つの課題と市場リスク・開発能力ギャップの整理'),
    ('04', 'ゴール',    '研修後に目指す姿・3〜6ヶ月後の成果指標'),
    ('05', '提案詳細',  'カリキュラム設計・実施形式・費用・助成金活用'),
]
toc_slide(prs, toc_sections)

# ═══════════════════════════════════════════════════════
# Section Divider: 01 振り返り
# ═══════════════════════════════════════════════════════
section_divider(prs, '01', '振り返り', '前回お打合せの内容確認')

# ═══════════════════════════════════════════════════════
# P1: 振り返り
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, '前回のお打合せ振り返り', subtitle='振り返り')

recap_rows = [
    ('日時',             '2026年2月27日'),
    ('参加者（先方）',   '呉振宇 様（リードインクス株式会社 SaaS本部長）'),
    ('参加者（自社）',   '大川龍之介・髙橋悠（株式会社エヌイチ）'),
    ('お話しした内容',   'エヌイチ4サービス紹介 → リードインクス様のAI活用課題ヒアリング'),
    ('先方の関心',       'Kawaru Team（フルカスタマイズ研修）に最も強い関心\nAI開発BPOは将来的な連携として検討'),
    ('合意したNX',       '呉様 → 社長への1on1で研修提案を相談\nエヌイチ → 本提案書を作成・お送り'),
]
label_rows(s, recap_rows, 0.45, cy + 0.10, 2.20, 10.00, row_h=0.72, gap=0.10)

# ═══════════════════════════════════════════════════════
# P2-1: 現状① 事業・組織概要
# ═══════════════════════════════════════════════════════
# Section Divider: 02 現状
# ═══════════════════════════════════════════════════════
section_divider(prs, '02', '現状', 'ビジネス環境・AI活用の現状分析')

# ═══════════════════════════════════════════════════════
# P2-1: 現状① 事業・組織概要
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, 'ソフトバンクグループ・保険SaaS特化の開発会社',
            subtitle='現状①　事業・組織概要')

info_rows = [
    ('グループ',        'ソフトバンクグループ子会社'),
    ('主力事業',        '保険会社・保険代理店向けSaaSプロダクトの開発・提供'),
    ('SaaS本部',        '保険業界向けSaaS製品の開発・運営（呉振宇 様ご担当）'),
    ('SI本部',          'フルカスタマイズ開発中心（もう一方の本部）'),
    ('エンジニア規模',  '両本部合計 約50〜60名'),
]
label_rows(s, info_rows, 0.45, cy + 0.10, 1.90, 10.10, row_h=0.76, gap=0.10)

# 強調カード: 中国語話者の割合
cy2 = cy + 0.10 + 5 * (0.76 + 0.10)
card(s, 0.45, cy2, 12.43, 1.50, RED_L, MGRAY, RED)
txt(s, 0.72, cy2 + 0.15, 5.0, 0.50,
    '【重要】エンジニアの約半数が中国語話者', 14, RED, bold=True)
txt(s, 0.72, cy2 + 0.72, 11.50, 0.60,
    '日本在籍だが日本語コミュニケーションが困難なメンバーが多数。研修設計において言語対応が必須条件となる。', 12, BLACK)

# ═══════════════════════════════════════════════════════
# P2-2: 現状② AI活用の実態
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, '個人の我流、ベストプラクティスなし',
            subtitle='現状②　AI活用の実態')

# 比較表: 現状 vs あるべき姿
col0, col1, col2 = 0.45, 3.80, 8.45
cw0, cw1, cw2 = 3.30, 4.60, 4.45
rh = 0.76
# ヘッダー
rect(s, col0, cy, cw0, rh, NAVY)
txt(s, col0 + 0.10, cy, cw0, rh, '項目', 12, WHITE, bold=True, va=MSO_ANCHOR.MIDDLE)
rect(s, col1, cy, cw1, rh, RED_L, RED)
txt(s, col1 + 0.10, cy, cw1, rh, '現状（事実）', 12, RED, bold=True,
    va=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
rect(s, col2, cy, cw2, rh, GREEN_L, GREEN)
txt(s, col2 + 0.10, cy, cw2, rh, '目指す姿（研修後）', 12, GREEN, bold=True,
    va=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

rows = [
    ('AI活用状況',    '個人が我流でChatGPT/Claude利用',             'チームで標準プロセスを共有'),
    ('ベストプラクティス', '未確立・誰も持っていない',                'AI開発の共通手法が社内に存在'),
    ('スキルレベル',  'IT知識は高い。AI開発手法は未習得',          '要件→タスク分解→エージェント設計ができる'),
    ('組織的取り組み', '外部研修の受講歴ゼロ',                        'フルカスタマイズ研修で一気に習得'),
    ('AI開発経験',    '海外ベンダーにBPO委託（POCのみ）',           '社内主導で本番化できる体制へ'),
]
for i, (axis, before, after) in enumerate(rows):
    ry = cy + rh + i * (rh + 0.06)
    bg = LGRAY if i % 2 == 0 else WHITE
    rect(s, col0, ry, cw0, rh, bg, MGRAY)
    txt(s, col0 + 0.10, ry, cw0 - 0.20, rh, axis, 11, NAVY, bold=True, va=MSO_ANCHOR.MIDDLE)
    rect(s, col1, ry, cw1, rh, bg, MGRAY)
    txt(s, col1 + 0.10, ry, cw1 - 0.20, rh, before, 11, RED, va=MSO_ANCHOR.MIDDLE)
    rect(s, col2, ry, cw2, rh, bg, MGRAY)
    txt(s, col2 + 0.10, ry, cw2 - 0.20, rh, after, 11, GREEN, va=MSO_ANCHOR.MIDDLE)

# 引用
qy = cy + rh + 5 * (rh + 0.06) + 0.10
card(s, 0.45, qy, 12.43, 0.90, TEAL_L, MGRAY, TEAL)
txt(s, 0.72, qy + 0.10, 11.80, 0.70,
    '「みんな自分でネットで勉強して、自分のやり方でやってる。決まったベストプラクティスという知識は誰も持っていない」\n─ 呉振宇 様（商談中の発言）',
    11, NAVY, bold=False)

# ═══════════════════════════════════════════════════════
# P2-3: 現状③ 三重制約
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, 'GPT縛り・金融規制・言語バリアの三重制約',
            subtitle='現状③　制約・環境')

constraints = [
    ('ツール制約',
     'ChatGPT Enterprise版のみ会社資産のアップロードを許可\nソフトバンクグループ方針によりGPT中心の活用が前提\nCursor / Claude等は個人利用のみ可',
     SLATE, SLATE_L),
    ('規制制約',
     '保険業界の個人情報保護規制により海外ベンダーとの本番運用は法律上困難\nPOCは海外ベンダーと実施中だが、本番化するには国内完結が必須',
     CORAL, CORAL_L),
    ('言語制約',
     'エンジニアの約半数が中国語話者で日本語コミュニケーション困難\n翻訳ツール活用での研修に懸念あり（意味は伝わるが熱量が薄れる可能性）',
     NAVY, LGRAY),
]
cw3 = 3.90
cgap = 0.28
for i, (title, body, bar, bg) in enumerate(constraints):
    cx3 = 0.45 + i * (cw3 + cgap)
    card(s, cx3, cy, cw3, 5.50, bg, MGRAY, bar)
    txt(s, cx3 + 0.20, cy + 0.20, cw3 - 0.30, 0.50, title, 15, bar, bold=True)
    rect(s, cx3 + 0.20, cy + 0.80, cw3 - 0.40, 0.03, bar)
    txt(s, cx3 + 0.20, cy + 1.00, cw3 - 0.30, 4.20, body, 13, BLACK)

# ═══════════════════════════════════════════════════════
# P3-1: 課題① 全体マップ（ロジックツリー）
# ═══════════════════════════════════════════════════════
# Section Divider: 03 課題
# ═══════════════════════════════════════════════════════
section_divider(prs, '03', '課題', '3つの課題と市場リスクの整理')

# ═══════════════════════════════════════════════════════
# P3-1: 課題① ロジックツリー
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, '課題は3層構造', subtitle='課題①　全体マップ')

# ルートノード
root_h = 0.70
root_y = cy + 0.10
rect(s, 0.45, root_y, 12.43, root_h, NAVY)
txt(s, 0.45, root_y, 12.43, root_h,
    'AI開発・活用能力の組織的欠如', 17, WHITE, bold=True,
    align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# 縦線（ルート → 3ブランチ）
line_top_y = root_y + root_h
line_bot_y = line_top_y + 0.40
mid_x = 0.45 + 12.43 / 2
rect(s, mid_x - 0.015, line_top_y, 0.03, line_bot_y - line_top_y, MGRAY)

# 横線（3ブランチを繋ぐ）
branch_xs = [0.45 + 0.20 + 0.65, 0.45 + 4.30 + 0.65, 0.45 + 8.20 + 0.65]
bw = 3.80
branch_y = line_bot_y
rect(s, branch_xs[0] - 0.015, branch_y, branch_xs[2] - branch_xs[0] + 0.03, 0.03, MGRAY)

# 3つのブランチブロック
branches = [
    ('① プロダクトへの\nAI組み込み能力不足', CORAL, CORAL_L),
    ('② エンジニアのAIスキルの\n属人化・標準化不足', SLATE, SLATE_L),
    ('③ 多国籍チームへの\n展開障壁', NAVY,  LGRAY),
]
branch_x_starts = [0.45, 4.57, 8.69]
branch_w = 4.00
branch_bh = 0.90
for i, (label, bar, bg) in enumerate(branches):
    bx = branch_x_starts[i]
    # 縦線（横線 → ブランチ）
    rect(s, bx + branch_w / 2 - 0.015, branch_y, 0.03, 0.30, MGRAY)
    card(s, bx, branch_y + 0.30, branch_w, branch_bh, bg, MGRAY, bar)
    txt(s, bx + 0.16, branch_y + 0.30, branch_w - 0.20, branch_bh, label, 13, bar,
        bold=True, va=MSO_ANCHOR.MIDDLE)

# リーフ（各ブランチの詳細）
leaves = [
    [   # ① プロダクト
        ('ビジネス要件→AI実装への変換方法が未確立', CORAL, CORAL_L, '【事実】'),
        ('エージェント設計・タスク分解のノウハウがない', CORAL, CORAL_L, '【事実】'),
        ('AI非対応プロダクトは市場競争力を失うリスク', CORAL, CORAL_L, '【推察】'),
    ],
    [   # ② スキル属人化
        ('個人の我流でAI活用（ベストプラクティス未共有）', SLATE, SLATE_L, '【事実】'),
        ('AI開発の標準プロセスを誰も知らない', SLATE, SLATE_L, '【事実】'),
        ('ChatGPT Enterprise縛り・GPT中心方針が前提制約', SLATE, SLATE_L, '【事実】'),
    ],
    [   # ③ 言語バリア
        ('エンジニアの約半数が中国語話者（日本語不可）', NAVY, LGRAY, '【事実】'),
        ('翻訳ツール活用での研修効果に懸念', NAVY, LGRAY, '【事実】'),
    ],
]
leaf_top_y = branch_y + 0.30 + branch_bh + 0.15
for col, (col_leaves, bx_start) in enumerate(zip(leaves, branch_x_starts)):
    lw_each = branch_w
    lx = bx_start
    for j, (text, bar, bg, label) in enumerate(col_leaves):
        ly = leaf_top_y + j * 0.70
        card(s, lx, ly, lw_each - 0.08, 0.60, bg, MGRAY)
        txt(s, lx + 0.14, ly, lw_each - 0.22, 0.60, text, 10, bar, va=MSO_ANCHOR.MIDDLE)
        # ラベル
        rect(s, lx + lw_each - 0.68, ly + 0.10, 0.58, 0.38, bar)
        txt(s, lx + lw_each - 0.68, ly + 0.10, 0.58, 0.38, label, 8, WHITE,
            bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════
# P3-2: 課題② 市場リスク
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, 'AI非対応は、競争力喪失に直結する', subtitle='課題②　市場リスク')

# 因果連鎖カード
chain = [
    ('保険業界\nAIトレンド加速', TEAL_L,  TEAL),
    ('顧客からの\nPOC要望が急増', SLATE_L, SLATE),
    ('社内に実装\nできる人材なし', CORAL_L, CORAL),
    ('POCが\n本番化できない',   LGRAY,   NAVY),
]
cw_ch = 2.80
cgap_ch = 0.22
cx0_ch = 0.45
for i, (label, bg, bar) in enumerate(chain):
    bx_ch = cx0_ch + i * (cw_ch + cgap_ch)
    card(s, bx_ch, cy, cw_ch, 2.20, bg, MGRAY, bar)
    txt(s, bx_ch + 0.20, cy, cw_ch - 0.30, 2.20, label, 15, bar, bold=True,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    if i < len(chain) - 1:
        txt(s, bx_ch + cw_ch + 0.04, cy + 0.70, 0.20, 0.80, '→', 22, bar,
            bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# 根拠リスト
dy = cy + 2.45
items = [
    ('【事実】', '保険業界のSaaSプロダクトへのAI機能組み込みは今や競争上の必須要件になりつつある', RED),
    ('【事実】', '「従来のフォーム形式だとAIトレンドから見て古くなるリスクがある」（呉振宇 様）', RED),
    ('【事実】', '顧客（保険会社）からAI活用のPOC要望が来ているが、社内に実装できる人材がいない', ORANGE),
    ('【推察】', '海外ベンダーに依存した状態では、本番運用移行後に自社開発力が不足する可能性がある', ORANGE),
]
for tag, text, color in items:
    rect(s, 0.45, dy, 12.43, 0.58, LGRAY, MGRAY)
    rect(s, 0.45, dy, 0.70, 0.58, color)
    txt(s, 0.45, dy, 0.70, 0.58, tag, 8, WHITE, bold=True,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    txt(s, 1.22, dy, 11.50, 0.58, text, 12, BLACK, va=MSO_ANCHOR.MIDDLE)
    dy += 0.70

# ═══════════════════════════════════════════════════════
# P3-3: 課題③ 開発能力ギャップ
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, 'ビジネス要件をAIに落とせない', subtitle='課題③　開発能力ギャップ')

# 比較表
col0 = 0.45; col1 = 3.00; col2 = 8.10
cw0 = 2.50; cw1 = 5.05; cw2 = 4.74
rh2 = 0.84

rect(s, col0, cy, cw0, rh2, NAVY)
txt(s, col0 + 0.10, cy, cw0, rh2, '開発フェーズ', 12, WHITE, bold=True, va=MSO_ANCHOR.MIDDLE)
rect(s, col1, cy, cw1, rh2, LGRAY, MGRAY)
txt(s, col1 + 0.10, cy, cw1, rh2, '従来のシステム開発（習熟済み）', 12, NAVY, bold=True,
    va=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
rect(s, col2, cy, cw2, rh2, RED_L, RED)
txt(s, col2 + 0.10, cy, cw2, rh2, 'AI開発（手法が不明）', 12, RED, bold=True,
    va=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

table_rows = [
    ('起点',        '顧客要件',                          'ビジネス課題・解決ゴール'),
    ('設計',        '要件定義 → 基本設計 → 詳細設計',   'ゴール設定 → タスク分解 → エージェント設計'),
    ('実装',        'コーディング（経験豊富）',           'AIモデル/API活用（ノウハウなし）'),
    ('検証',        'テスト・QA（標準化済み）',           '精度評価・プロンプト調整（方法不明）'),
    ('習熟レベル',  '✅  標準手法を全員が理解',           '❌  誰も手法を知らない'),
]
for i, (phase, trad, ai) in enumerate(table_rows):
    ry = cy + rh2 + i * (rh2 + 0.06)
    bg2 = LGRAY if i % 2 == 0 else WHITE
    rect(s, col0, ry, cw0, rh2, bg2, MGRAY)
    txt(s, col0 + 0.10, ry, cw0 - 0.20, rh2, phase, 11, NAVY, bold=True, va=MSO_ANCHOR.MIDDLE)
    rect(s, col1, ry, cw1, rh2, bg2, MGRAY)
    txt(s, col1 + 0.10, ry, cw1 - 0.20, rh2, trad, 11, BLACK, va=MSO_ANCHOR.MIDDLE)
    rect(s, col2, ry, cw2, rh2, bg2, MGRAY)
    txt(s, col2 + 0.10, ry, cw2 - 0.20, rh2, ai, 11, RED, va=MSO_ANCHOR.MIDDLE)

# 引用
qy2 = cy + rh2 + 5 * (rh2 + 0.06) + 0.10
card(s, 0.45, qy2, 12.43, 0.85, TEAL_L, MGRAY, TEAL)
txt(s, 0.72, qy2 + 0.10, 11.80, 0.65,
    '「こういうものを作ってほしいと言われたら、どう形にしたらいいかが分からない」\n「AIの世界で、こういう目標を実現するためにタスクをどう分配するか、そのノウハウは今ない」（呉振宇 様）',
    11, NAVY)

# ═══════════════════════════════════════════════════════
# Section Divider: 04 ゴール
# ═══════════════════════════════════════════════════════
section_divider(prs, '04', 'ゴール', '研修後に目指す姿・成果指標')

# ═══════════════════════════════════════════════════════
# P4-1: ゴール① 目指す姿
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, 'エンジニア全員がAIを武器に、プロダクトを進化させる',
            subtitle='ゴール①　目指す姿')

col_w = 5.60
# Before
card(s, 0.45, cy, col_w, 5.50, RED_L, MGRAY, RED)
txt(s, 0.72, cy + 0.20, col_w - 0.40, 0.45, 'Before（現在）', 14, RED, bold=True)
befores = [
    'AI活用: 個人の我流・ベストプラクティスなし',
    '開発力: ビジネス要件→AI実装の変換ができない',
    'プロダクト: フォーム中心の従来型SaaS',
    '依存状況: 海外BPO依存（本番移行に課題）',
    '言語環境: 約半数が中国語話者で研修が難しい',
]
by = cy + 0.80
for b in befores:
    txt(s, 0.72, by, col_w - 0.50, 0.58, '✗  ' + b, 12, RED)
    by += 0.70

# 矢印
txt(s, 6.30, cy + 2.10, 0.80, 0.80, '→', 30, TEAL, bold=True,
    align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
txt(s, 6.00, cy + 3.00, 1.50, 0.65,
    'Kawaru\nTeam研修', 11, TEAL, bold=True, align=PP_ALIGN.CENTER)

# After
card(s, 7.16, cy, col_w, 5.50, GREEN_L, MGRAY, GREEN)
txt(s, 7.43, cy + 0.20, col_w - 0.40, 0.45, 'After（3〜6ヶ月後）', 14, GREEN, bold=True)
afters = [
    'AI活用: チームで標準プロセスを共有',
    '開発力: 要件→タスク分解→エージェント設計が自力で可能',
    'プロダクト: AI機能を組み込んだ競争力あるSaaSへ',
    '依存状況: 自社エンジニアが主体的にAI開発を推進',
    '言語環境: 多国籍チームへの研修設計を実現',
]
afy = cy + 0.80
for a in afters:
    txt(s, 7.43, afy, col_w - 0.50, 0.58, '✓  ' + a, 12, GREEN)
    afy += 0.70

# ═══════════════════════════════════════════════════════
# P4-2: ゴール② 成果指標
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, '3ヶ月後・6ヶ月後に何が変わるか', subtitle='ゴール②　成果指標')

timeline = [
    ('研修直後',   'AI開発の標準プロセスを全員が理解\nChatGPT Enterpriseを業務で活用できる',       TEAL,  TEAL_L),
    ('1ヶ月後',   '自社サービスへのAI機能組み込みアイデアを\nチームで出せるようになる',              SLATE, SLATE_L),
    ('3ヶ月後',   '社内でAIエージェントのプロトタイプを\n自力で設計・実装できる',                    SAGE,  SAGE_L),
    ('6ヶ月後',   '海外BPO依存から脱却し\n自社主導でAI開発を推進する体制が整う',                    NAVY,  LGRAY),
]
tw = 5.90
tgap = 0.22
for i, (period, result, bar, bg) in enumerate(timeline):
    tx = 0.45 + (i % 2) * (tw + tgap)
    ty = cy + (i // 2) * 2.70
    card(s, tx, ty, tw, 2.50, bg, MGRAY, bar)
    # バッジを右上に配置してテキストと重ならないようにする
    num_badge(s, tx + tw - 0.42, ty + 0.38, i + 1, bar)
    txt(s, tx + 0.20, ty + 0.12, tw - 1.00, 0.52, period, 14, bar, bold=True)
    rect(s, tx + 0.20, ty + 0.78, tw - 0.40, 0.03, bar)
    txt(s, tx + 0.20, ty + 0.92, tw - 0.30, 1.36, result, 13, BLACK)

# ═══════════════════════════════════════════════════════
# Section Divider: 05 提案詳細（根拠・詳細を含む）
# ═══════════════════════════════════════════════════════
section_divider(prs, '05', '提案詳細', 'カリキュラム・実施形式・費用・助成金')

# ═══════════════════════════════════════════════════════
# P5-1: 根拠① 比較表
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, '汎用研修では御社の課題は解けない', subtitle='根拠①　他社との違い')

col0 = 0.45; col1 = 3.50; col2 = 8.10
cw0 = 3.00; cw1 = 4.55; cw2 = 4.74
rh3 = 0.80

rect(s, col0, cy, cw0, rh3, NAVY)
txt(s, col0 + 0.10, cy, cw0 - 0.20, rh3, '比較軸', 12, WHITE, bold=True, va=MSO_ANCHOR.MIDDLE)
rect(s, col1, cy, cw1, rh3, MGRAY)
txt(s, col1 + 0.10, cy, cw1 - 0.20, rh3, '一般的なAI研修', 12, GRAY, bold=True,
    va=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
rect(s, col2, cy, cw2, rh3, TEAL)
txt(s, col2 + 0.10, cy, cw2 - 0.20, rh3, 'Kawaru Team', 12, WHITE, bold=True,
    va=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

comparisons = [
    ('カリキュラム',     '汎用コンテンツ（固定）',              'ヒアリング2時間でゼロから設計'),
    ('講師',             '教育専門家',                           '自社でAIを日常開発する実務者'),
    ('対象者',           '非エンジニア向けが多い',              '御社エンジニアに特化した内容'),
    ('ツール対応',       '汎用ツール前提',                       'ChatGPT Enterprise縛り前提で設計'),
    ('業界特化',         '汎用事例',                             '保険業界ユースケースを教材化'),
    ('言語対応',         '日本語のみ',                           '多国籍チームへの対応を考慮した設計'),
]
for i, (axis, gen, kaw) in enumerate(comparisons):
    ry = cy + rh3 + i * (rh3 + 0.05)
    bg3 = LGRAY if i % 2 == 0 else WHITE
    rect(s, col0, ry, cw0, rh3, bg3, MGRAY)
    txt(s, col0 + 0.10, ry, cw0 - 0.20, rh3, axis, 12, NAVY, bold=True, va=MSO_ANCHOR.MIDDLE)
    rect(s, col1, ry, cw1, rh3, bg3, MGRAY)
    txt(s, col1 + 0.14, ry, cw1 - 0.24, rh3, gen, 12, GRAY, va=MSO_ANCHOR.MIDDLE)
    rect(s, col2, ry, cw2, rh3, TEAL_L, MGRAY)
    txt(s, col2 + 0.14, ry, cw2 - 0.24, rh3, kaw, 12, NAVY, bold=True, va=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════
# P5-2: 根拠② エヌイチの強み
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, 'エヌイチが選ばれる2つの理由', subtitle='根拠②　強み')

reasons = [
    ('理由①  実務者が教える',
     '自社でAIエージェント・ツールを日常業務で開発・運用している実務者が直接指導\n「知識を教える」のではなく「現場で機能するノウハウ」を渡す\n\n「私たちも普段からAIを使って開発を進めている。ゴール設定とプロセスをどう進めているかを実践目線でお伝えできる」（髙橋悠）',
     SLATE, SLATE_L),
    ('理由②  完全フルカスタマイズ — 御社専用をゼロ設計',
     '動画提供・汎用コンテンツではなく、御社の課題・ゴール・ツール環境・言語条件をすべて踏まえた上でカリキュラムを一から構築\n\n「御社が欲しい情報だけを整理して提供できる」（髙橋悠）\nChatGPT Enterprise制約・保険業界ユースケース・中国語対応まで設計に組み込む',
     TEAL, TEAL_L),
]
rw = 12.43
rh4 = 2.70
rgap = 0.30
for i, (title, body, bar, bg) in enumerate(reasons):
    ry = cy + i * (rh4 + rgap)
    card(s, 0.45, ry, rw, rh4, bg, MGRAY, bar)
    txt(s, 0.72, ry + 0.20, rw - 0.40, 0.50, title, 15, bar, bold=True)
    rect(s, 0.72, ry + 0.78, rw - 0.50, 0.03, bar)
    txt(s, 0.72, ry + 0.95, rw - 0.40, 1.55, body, 12, BLACK)

# ═══════════════════════════════════════════════════════
# P6-1: 詳細① カリキュラム概要
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, '研修対象・カリキュラム概要', subtitle='詳細①　内容')

# 対象
card(s, 0.45, cy, 5.90, 1.40, TEAL_L, MGRAY, TEAL)
txt(s, 0.72, cy + 0.15, 5.50, 0.40, '対象', 12, TEAL, bold=True)
txt(s, 0.72, cy + 0.62, 5.50, 0.65,
    'リードインクス SaaS本部・SI本部のエンジニア\nまず1セッション（15〜25名）からスタートを推奨', 12, NAVY)

card(s, 6.55, cy, 6.33, 1.40, LGRAY, MGRAY, NAVY)
txt(s, 6.82, cy + 0.15, 5.90, 0.40, '実施規模・セッション設計', 12, NAVY, bold=True)
txt(s, 6.82, cy + 0.62, 5.90, 0.65,
    '1セッション = 15〜25名（Kawaru Teamの定員）\n両本部60名を受講させる場合 → 25名 × 3セッション', 12, BLACK)

# モジュール
txt(s, 0.45, cy + 1.60, 6.0, 0.40, '研修カリキュラム（4モジュール / 合計10〜15時間）', 13, NAVY, bold=True)
modules = [
    ('Module 1', 'AI開発の思考法',
     'ビジネス課題をAIタスクに分解する考え方\nエージェントとは何か・どう機能するか',
     TEAL, TEAL_L),
    ('Module 2', 'エージェント設計の実践',
     'ChatGPT Enterprise前提でエージェントを構築\n実際に手を動かす演習形式',
     BLUE, BLUE_L),
    ('Module 3', '保険業界ユースケース演習',
     'AI-OCR自動化・小額短期保険の自動判断を題材\n呉様が言及した具体事例をそのまま教材化',
     ORANGE, ORANGE_L),
    ('Module 4', 'AI開発プロセスの標準化',
     'チームでのAI開発の進め方・ドキュメント化\n共通言語・ベストプラクティスの確立',
     GREEN, GREEN_L),
]
mw = 5.90; mgap = 0.22
for i, (mnum, mtitle, mbody, bar, bg) in enumerate(modules):
    mx = 0.45 + (i % 2) * (mw + mgap)
    my = cy + 2.10 + (i // 2) * 1.85  # 2.30→1.85: 下段が枠内に収まるよう縮小
    card(s, mx, my, mw, 1.70, bg, MGRAY, bar)  # 2.10→1.70
    txt(s, mx + 0.18, my + 0.10, 1.20, 0.32, mnum, 10, bar, bold=True)
    txt(s, mx + 1.40, my + 0.10, mw - 1.55, 0.38, mtitle, 14, NAVY, bold=True)
    rect(s, mx + 0.18, my + 0.58, mw - 0.36, 0.03, bar)
    txt(s, mx + 0.18, my + 0.70, mw - 0.30, 0.85, mbody, 11, BLACK)

# ═══════════════════════════════════════════════════════
# P6-2: 詳細② 実施形式・スケジュール
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, '3〜4日間、交代制で実施', subtitle='詳細②　実施形式・スケジュール')

impl_rows = [
    ('実施形式',   '御社オフィスへの訪問型（リアル開催）'),
    ('1日の時間',  '2〜3時間（業務への影響を最小化）'),
    ('実施日数',   '3〜4日間 / グループあたり'),
    ('グループ',   '15〜25名ずつ交代制で全員受講'),
    ('実施時間帯', '平日 9:00〜18:00 ★助成金適用の必須条件★'),
    ('言語対応',   '日本語＋翻訳ツール活用。資料は図解中心で設計（中国語話者への配慮）'),
]
label_rows(s, impl_rows, 0.45, cy + 0.10, 1.70, 10.50, row_h=0.62, gap=0.06)

# ステップ図（テーブル端 = cy+0.10+6×0.68=5.78" → ラベルy=cy+4.33, カードy=cy+4.73）
txt(s, 0.45, cy + 4.33, 5.0, 0.35, '実施ロードマップ（目安）', 13, NAVY, bold=True)
steps_r = [
    ('STEP 1', 'ヒアリング\nカリキュラム設計', '2〜3週間', TEAL,  TEAL_L),
    ('STEP 2', '第1グループ実施\n（SaaS本部 前半）', '3〜4日', SLATE, SLATE_L),
    ('STEP 3', '第2グループ実施\n（後半・SI本部）', '3〜4日', SAGE,  SAGE_L),
]
sw3 = 3.80; sg3 = 0.24
sy3 = cy + 4.73  # 1.52+4.73=6.25", カード端=6.25+1.22=7.47" ✓
for i, (step, title, period, bar, bg) in enumerate(steps_r):
    sx3 = 0.45 + i * (sw3 + sg3)
    card(s, sx3, sy3, sw3, 1.22, bg, MGRAY, bar)
    txt(s, sx3 + 0.18, sy3 + 0.08, sw3 - 0.26, 0.30, step, 10, bar, bold=True)
    txt(s, sx3 + 0.18, sy3 + 0.42, sw3 - 0.26, 0.58, title, 11, NAVY, bold=True)
    txt(s, sx3 + 0.18, sy3 + 1.03, sw3 - 0.26, 0.18, '目安: ' + period, 9, bar)
    if i < 2:
        txt(s, sx3 + sw3 + 0.04, sy3 + 0.42, 0.22, 0.50, '→', 16, bar,
            bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════
# P6-3: 詳細③ 費用・助成金
# ═══════════════════════════════════════════════════════
s = blank(prs)
cy = header(s, '助成金活用で実質負担6〜7割減', subtitle='詳細③　費用・投資対効果')

# 3カード
cards3 = [
    ('定価（1セッション）', '15万円 / 名\n× 25名 / セッション\n= 375万円 / 回', LGRAY,   MGRAY, NAVY),
    ('助成金還元', '6〜7割還元\n= 約225〜262万円 / 回', SLATE_L, MGRAY, SLATE),
    ('実質負担',  '約112〜150万円 / 回\n（4.5〜6万円 / 名）', SAGE_L, MGRAY, SAGE),
]
c3w = 3.80; c3gap = 0.28
for i, (title, body, bg, border, tc) in enumerate(cards3):
    cx3 = 0.45 + i * (c3w + c3gap)
    card(s, cx3, cy, c3w, 2.40, bg, border, tc)
    txt(s, cx3 + 0.22, cy + 0.20, c3w - 0.36, 0.50, title, 14, tc, bold=True,
        align=PP_ALIGN.CENTER)
    txt(s, cx3 + 0.22, cy + 0.88, c3w - 0.36, 1.30, body, 15, tc, bold=True,
        align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    if i < 2:
        txt(s, cx3 + c3w + 0.05, cy + 0.90, 0.28, 0.60, '→', 22, tc,
            bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# 助成金詳細
txt(s, 0.45, cy + 2.65, 12.43, 0.40, '助成金について', 13, NAVY, bold=True)
subsidy_rows = [
    ('種類',             '人材開発支援助成金（厚生労働省）'),
    ('適用条件',         '平日9:00〜18:00での実施・訪問型・受講記録等の書類整備'),
    ('申請サポート',     'エヌイチがフルサポート（書類作成・申請手続きを含む）'),
    ('スケジュール目安', '研修実施 → 書類提出 → 2〜4ヶ月後に還付'),
]
label_rows(s, subsidy_rows, 0.45, cy + 3.10, 1.70, 10.50, row_h=0.58, gap=0.08)

# 注記
txt(s, 0.45, cy + 5.55, 12.43, 0.38,
    '※ 1セッション25名が基本設計。50名受講の場合は25名×2セッション（750万円→実質225〜300万円）。正式見積りはヒアリング後に提示します。',
    10, GRAY)

# ═══════════════════════════════════════════════════════
# P7: ネクストアクション
# ═══════════════════════════════════════════════════════
s = blank(prs)
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, SW, 0.06, TEAL)
txt(s, 0.45, 0.18, 10.50, 0.75, '次のステップ', 28, DARK, bold=True, va=MSO_ANCHOR.MIDDLE)
txt(s, 9.80, 0.10, 3.30, 0.40, '株式会社エヌイチ', 11, GRAY, align=PP_ALIGN.RIGHT)
rect(s, 0.45, 0.96, 12.43, 0.03, TEAL)

steps_na = [
    ('STEP 1', '社内承認',
     '呉様より社長への上申・承認取得',
     '（承認済の場合はSTEP 2へ）',
     SLATE, SLATE_L),
    ('STEP 2', 'ヒアリング',
     '研修対象者・習得スキル・日程を詳細確認\n（約2時間・無料）',
     '承認後 即対応',
     TEAL, TEAL_L),
    ('STEP 3', 'カリキュラム設計',
     '御社専用カリキュラム作成\n詳細費用・スケジュール提示',
     'ヒアリング後\n1〜2週間',
     SAGE, SAGE_L),
    ('STEP 4', '研修実施',
     '3〜4日間のフルカスタマイズ研修\n振り返り・アフターフォロー',
     'STEP 3完了後',
     NAVY, LGRAY),
]
sw4 = 2.80; sg4 = 0.22; cy4 = 1.25
for i, (step, title, body, period, bar, bg) in enumerate(steps_na):
    bx4 = 0.45 + i * (sw4 + sg4)
    card(s, bx4, cy4, sw4, 5.70, bg, MGRAY, bar)
    txt(s, bx4 + 0.22, cy4 + 0.18, sw4 - 0.36, 0.40, step, 11, bar, bold=True)
    txt(s, bx4 + 0.22, cy4 + 0.70, sw4 - 0.36, 0.55, title, 17, NAVY, bold=True)
    rect(s, bx4 + 0.28, cy4 + 1.35, sw4 - 0.56, 0.03, bar)
    txt(s, bx4 + 0.22, cy4 + 1.55, sw4 - 0.36, 2.60, body, 13, BLACK)
    rect(s, bx4 + 0.22, cy4 + 4.94, sw4 - 0.44, 0.55, bar)
    txt(s, bx4 + 0.22, cy4 + 4.94, sw4 - 0.44, 0.55, period, 11, WHITE,
        bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
    if i < 3:
        txt(s, bx4 + sw4 + 0.04, cy4 + 2.50, 0.22, 0.60, '→', 18, bar,
            bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

# 本日の確認事項（y=1.25+5.75=7.00", end=7.00+0.48=7.48" ✓）
rect(s, 0.45, cy4 + 5.75, 12.43, 0.48, TEAL)
txt(s, 0.72, cy4 + 5.75, 12.00, 0.48,
    '本日確認させていただきたいこと  ① カリキュラム・実施形式のご認識合わせ  ② 言語対応の方針  ③ ヒアリングの日程調整',
    12, WHITE, bold=True, va=MSO_ANCHOR.MIDDLE)

# ═══════════════════════════════════════════════════════
# 固定スライド: 締め
# ═══════════════════════════════════════════════════════
s = blank(prs)
rect(s, 0, 0, SW, SH, WHITE)
rect(s, 0, 0, SW, 0.10, TEAL)
rect(s, 0, SH - 0.10, SW, 0.10, TEAL)

txt(s, 0.80, 1.80, 11.73, 2.80,
    '私たちが提供するのは\n「ツール」や「研修」ではありません。\n貴社の未来を創る「AI人材」と、\n組織全体の「成長」です。',
    30, DARK, bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

txt(s, 0.80, 6.00, 11.73, 0.55,
    '株式会社エヌイチ', 18, TEAL, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════
# 参考資料: 会社概要・導入実績（最後に配置）
# ════════════════════════════════════════════════════

# 参考資料 区切り
ref_divider(prs)

# 参考資料①: 会社概要
s = blank(prs)
cy = header(s, '会社概要')

info = [
    ('会社名',   '株式会社エヌイチ'),
    ('経営陣',   '代表取締役CEO 奥山幸生\n取締役COO 髙橋悠　取締役CHRO 米澤浩明'),
    ('資本金',   '1,000万円'),
    ('設立',     '2023年10月'),
    ('所在地',   '東京都新宿区西新宿3-7-30 フロンティアグラン西新宿12F'),
    ('Mission',  '新しい働く「カタチ」を創る'),
    ('Vision',   'AI（アイ）ある社会に'),
]
lw, vw = 1.30, 7.30
for i, (lbl, val) in enumerate(info):
    iy = cy + i * 0.72
    rect(s, 0.45, iy, lw, 0.65, TEAL_L)
    txt(s, 0.51, iy, lw - 0.10, 0.65, lbl, 10, TEAL, bold=True, va=MSO_ANCHOR.MIDDLE)
    rect(s, 0.45 + lw + 0.06, iy, vw, 0.65, WHITE, MGRAY)
    txt(s, 0.45 + lw + 0.14, iy, vw - 0.18, 0.65, val, 11, BLACK, va=MSO_ANCHOR.MIDDLE)

services = [
    ('Kawaru',       '業務自動化SaaS',           TEAL,  TEAL_L),
    ('Kawaru Team',  'フルカスタマイズAI研修',    SLATE, SLATE_L),
    ('Kawaru BPO',   'AI業務効率化代行',          SAGE,  SAGE_L),
    ('Kawaru Coach', 'AI顧問サービス',            NAVY,  LGRAY),
]
sc_x = 9.30; sc_w = 3.60; sc_h = 1.40; sc_gap = 0.18
for i, (name, desc, bar, bg) in enumerate(services):
    sy = cy + i * (sc_h + sc_gap)
    card(s, sc_x, sy, sc_w, sc_h, bg, MGRAY, bar)
    txt(s, sc_x + 0.20, sy + 0.15, sc_w - 0.30, 0.45, name, 14, bar, bold=True)
    txt(s, sc_x + 0.20, sy + 0.65, sc_w - 0.30, 0.55, desc, 12, BLACK)

# 参考資料②: 導入実績
s = blank(prs)
cy = header(s, '導入実績')
txt(s, 0.45, cy, 12.43, 0.38, 'あらゆる規模の企業様に導入されています', 13, GRAY)

companies = ['株式会社ブロードリンク', '株式会社中西製作所', '株式会社ミズカラ',
             '株式会社エル・ティー・エス リンク', 'FRUOR株式会社',
             '株式会社BIRDY', '株式会社オンシナジー', 'キリングループ']
cw_each = (SW - 0.90) / len(companies)
for i, name in enumerate(companies):
    cx_c = 0.45 + i * cw_each
    rect(s, cx_c, cy + 0.48, cw_each - 0.05, 0.55, LGRAY, MGRAY)
    txt(s, cx_c, cy + 0.48, cw_each - 0.05, 0.55, name, 9, NAVY,
        bold=True, align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)

cases = [
    ('キリングループ様',
     '課題：AI未経験、活用イメージなし',
     '成果：全員が「業務でAIを使いたい」と回答\n満足度 8.5点',
     TEAL, TEAL_L),
    ('株式会社ミズカラ様',
     '課題：問い合わせ対応に月283時間',
     '成果：月85時間へ削減\n（約200時間＝25営業日分の削減）',
     SLATE, SLATE_L),
    ('株式会社BIRDY様',
     '課題：業務の属人化、営業効率の壁',
     '成果：バックオフィス効率化\n経営資源を成長領域にシフト',
     SAGE, SAGE_L),
]
case_w = (SW - 0.90 - 0.40) / 3
case_y = cy + 1.20
for i, (name, challenge, result, bar, bg) in enumerate(cases):
    cx_c = 0.45 + i * (case_w + 0.20)
    card(s, cx_c, case_y, case_w, 5.10, bg, MGRAY, bar)
    txt(s, cx_c + 0.20, case_y + 0.20, case_w - 0.30, 0.45, name, 14, bar, bold=True)
    rect(s, cx_c + 0.20, case_y + 0.72, case_w - 0.40, 0.03, bar)
    txt(s, cx_c + 0.20, case_y + 0.90, case_w - 0.30, 0.80, challenge, 12, BLACK)
    txt(s, cx_c + 0.20, case_y + 1.80, case_w - 0.30, 1.20, result, 13, bar, bold=True)

# ════════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════════
out = '/Users/kyouyuu/cloude/output/proposal_leadinks_20260325.pptx'
prs.save(out)
print(f'Saved: {out}  ({len(prs.slides)} slides)')
