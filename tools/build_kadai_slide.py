"""
課題スライド（差し込み用・単独1枚PPTX）
3.現状 と 5.テーマ の間に挿入する想定
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

C = {
    'white':         RGBColor(0xFF, 0xFF, 0xFF),
    'navy':          RGBColor(0x1C, 0x35, 0x57),
    'kawaru':        RGBColor(0x25, 0x63, 0xEB),
    'kawaru_l':      RGBColor(0xDB, 0xEA, 0xFE),
    'gray':          RGBColor(0x6B, 0x72, 0x80),
    'light_gray':    RGBColor(0xF3, 0xF4, 0xF6),
    'coral':         RGBColor(0xC9, 0x54, 0x54),
    'coral_l':       RGBColor(0xFD, 0xEF, 0xEF),
}

FONT = 'Noto Sans JP'
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)


def _block(slide, x, y, w, h, text, size, color=None, bg=None,
           bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, rounded=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    if bg:
        s.fill.solid()
        s.fill.fore_color.rgb = bg
    else:
        s.fill.background()
    s.line.fill.background()
    if rounded and hasattr(s, 'adjustments') and len(s.adjustments) > 0:
        s.adjustments[0] = 0.05

    tf = s.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)

    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color or C['navy']
    return s


def _line(slide, x, y, w, color, thick=0.05):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(thick))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _arrow(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景（白）
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C['white']
    bg.line.fill.background()

    # タイトル
    _block(s, 0.5, 0.4, 11.5, 0.7, 'このままでは事業が回らない', 24, bold=True, color=C['navy'])
    _line(s, 0.5, 1.15, 1.8, C['kawaru'], 0.07)

    # 3つの課題（横並び）
    issues = [
        ('業務量が\n増え続ける',
         'ローンチ・展示会・デリバリーで\n業務は増える一方'),
        ('人を増やすには\n時間がかかる',
         '採用・育成にコストと\nスピードが見合わない'),
        ('業務が\n毎月変わる',
         '既存の業務効率化フレームでは\n追いつかない'),
    ]
    cw = 4.0
    gap = 0.165
    cx_start = 0.5
    for i, (title, desc) in enumerate(issues):
        x = cx_start + i*(cw + gap)
        # 課題タイトル（コーラル背景）
        _block(s, x, 1.7, cw, 1.6, title, 22, bold=True,
               color=C['white'], bg=C['coral'], rounded=True,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 課題説明（薄い背景）
        _block(s, x, 3.35, cw, 1.4, desc, 14,
               color=C['navy'], bg=C['coral_l'], rounded=True,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 矢印
    _arrow(s, 6.42, 5.0, 0.5, 0.6, C['kawaru'])

    # 解決策
    _block(s, 0.5, 5.8, 12.33, 1.2,
           'AI社員 5人 で解決する',
           36, bold=True, color=C['kawaru'],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
           bg=C['kawaru_l'], rounded=True)

    # ノート
    s.notes_slide.notes_text_frame.text = """【課題・15秒】
ここで、Kawaru事業部の課題を3つお話しします。
1つ目、業務量が増え続けます。ローンチ、展示会、デリバリーで業務は増える一方です。
2つ目、人を増やすには時間がかかります。採用・育成にコストとスピードが見合いません。
3つ目、業務が毎月変わります。既存の業務効率化フレームでは追いつかない。
だから、AI社員5人で解決する、というのが私の答えです。
"""

    output = '/Users/kyouyuu/claude/strategy/AI社員5人_発表_課題スライド.pptx'
    prs.save(output)
    print(f"Saved: {output}")


if __name__ == '__main__':
    build()
