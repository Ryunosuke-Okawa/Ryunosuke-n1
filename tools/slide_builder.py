"""
Reusable PPTX slide builder for proposal documents.
Key design: text is placed INSIDE colored AutoShapes (not floating TextBoxes).
This creates visually rich, high-readability slides matching the Copilot reference.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === COLORS ===
C = {
    'bg1':       RGBColor(0xF3, 0xF4, 0xF6),   # slide base
    'bg2':       RGBColor(0xFF, 0xFF, 0xFF),     # content area
    'accent':    RGBColor(0x42, 0xA4, 0xAF),     # teal accent
    'card':      RGBColor(0xF9, 0xFA, 0xFB),     # card bg
    'highlight': RGBColor(0xEA, 0xF6, 0xF7),     # teal highlight
    'info_bg':   RGBColor(0xEF, 0xF6, 0xFF),     # light blue section
    'info_bar':  RGBColor(0x3B, 0x82, 0xF6),     # blue bar
    'ok_bg':     RGBColor(0xF0, 0xFD, 0xF4),     # light green section
    'ok_bar':    RGBColor(0x16, 0xA3, 0x4A),     # green bar
    'ok_text':   RGBColor(0x16, 0x6D, 0x37),     # dark green text
    'warn_bg':   RGBColor(0xFE, 0xF3, 0xC7),     # light yellow
    'warn_bar':  RGBColor(0xD9, 0x77, 0x06),     # orange bar
    'warn_text': RGBColor(0xD9, 0x77, 0x06),     # orange text
    'issue_bg':  RGBColor(0xFE, 0xF2, 0xF2),     # light red
    'issue_bar': RGBColor(0xEF, 0x44, 0x44),     # red bar
    'divider':   RGBColor(0xE5, 0xE7, 0xEB),
    'title':     RGBColor(0x11, 0x18, 0x27),     # near-black
    'body':      RGBColor(0x1F, 0x29, 0x37),
    'sub':       RGBColor(0x4B, 0x55, 0x63),
    'label':     RGBColor(0x6B, 0x72, 0x80),
    'white':     RGBColor(0xFF, 0xFF, 0xFF),
    'black':     RGBColor(0x00, 0x00, 0x00),
    'light_acc': RGBColor(0xE0, 0xF2, 0xF4),
    'dark_head': RGBColor(0x1E, 0x40, 0x5F),     # dark header bar
    # === Kawaru Brand Colors（展示会フライヤー準拠） ===
    'kawaru_blue':   RGBColor(0x1E, 0x90, 0xFF),  # Kawaru（本体）= 青
    'kawaru_orange': RGBColor(0xFF, 0x66, 0x00),  # Kawaru Team   = オレンジ
    'kawaru_green':  RGBColor(0x10, 0xBF, 0x16),  # Kawaru BPO    = 緑
    'kawaru_purple': RGBColor(0x7C, 0x3A, 0xED),  # Kawaru Coach  = 紫
}

FONT = 'Noto Sans JP'
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
M = 0.50     # margin
CW = 12.33   # content width


class ProposalBuilder:
    def __init__(self, client_name, date_str):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.client = client_name
        self.date = date_str

    def _slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _bg(self, slide):
        for c in [C['bg1'], C['bg2']]:
            s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
            s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

    def _rect(self, slide, x, y, w, h, color, rounded=False):
        st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        s = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
        if rounded and hasattr(s, 'adjustments') and len(s.adjustments) > 0:
            s.adjustments[0] = 0.03
        return s

    # === CORE: Text inside shape (NOT floating textbox) ===

    def _block(self, slide, x, y, w, h, text, size,
               bg=None, tc=None, bold=False, align=PP_ALIGN.LEFT,
               anchor=MSO_ANCHOR.TOP, bar_color=None, rounded=True,
               ml=0.12, mr=0.08, mt=0.06, mb=0.06):
        """AutoShape with fill + text inside. The core visual element."""
        st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        s = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
        if bg:
            s.fill.solid(); s.fill.fore_color.rgb = bg
        else:
            s.fill.background()
        s.line.fill.background()
        if rounded and hasattr(s, 'adjustments') and len(s.adjustments) > 0:
            s.adjustments[0] = 0.02

        tf = s.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.vertical_anchor = anchor
        tf.margin_left = Inches(ml)
        tf.margin_right = Inches(mr)
        tf.margin_top = Inches(mt)
        tf.margin_bottom = Inches(mb)

        for i, line in enumerate(text.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_before = Pt(1); p.space_after = Pt(1)
            r = p.add_run(); r.text = line
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = tc or C['body']; r.font.name = FONT

        if bar_color:
            self._rect(slide, x, y, 0.05, h, bar_color)

        return s

    def _text(self, slide, x, y, w, h, text, size, bold=False, color=None,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        """Transparent text (no background) — use sparingly."""
        return self._block(slide, x, y, w, h, text, size,
                           bg=None, tc=color, bold=bold, align=align,
                           anchor=anchor, rounded=False, ml=0.04, mr=0.04, mt=0.02, mb=0.02)

    def _section_head(self, slide, x, y, w, text, size=13):
        """Section header: accent bar + bold text"""
        self._rect(slide, x, y + 0.02, 0.05, 0.26, C['accent'])
        self._block(slide, x + 0.15, y, w - 0.15, 0.30, text, size,
                    tc=C['title'], bold=True, rounded=False)

    def _title_bar(self, slide, title, subtitle=None):
        """Top bar with title, company label, accent line. Returns content start Y."""
        # Header band
        self._rect(slide, 0, 0, 13.33, 0.06, C['accent'])
        # Title
        self._block(slide, M, 0.35, 8.5, 0.55, title, 27,
                    tc=C['title'], bold=True, rounded=False, ml=0.08)
        # Company
        self._block(slide, 10.50, 0.45, 2.60, 0.35, '株式会社エヌイチ', 13,
                    tc=C['sub'], bold=True, align=PP_ALIGN.RIGHT, rounded=False)
        if subtitle:
            self._block(slide, M, 0.95, 10.0, 0.35, subtitle, 13,
                        tc=C['sub'], rounded=False, ml=0.08)
            self._rect(slide, M, 1.38, CW, 0.025, C['accent'])
            return 1.55
        else:
            self._rect(slide, M, 0.98, CW, 0.025, C['accent'])
            return 1.15

    def _badge(self, slide, cx, cy, r, text, bg=None, tc=None):
        s = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(r * 2), Inches(r * 2))
        s.fill.solid(); s.fill.fore_color.rgb = bg or C['accent']; s.line.fill.background()
        tf = s.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = text
        run.font.size = Pt(max(8, int(r * 26))); run.font.bold = True
        run.font.color.rgb = tc or C['white']; run.font.name = FONT

    def _arrow_h(self, slide, x, y):
        self._block(slide, x, y, 0.40, 0.40, '→', 22,
                    tc=C['accent'], bold=True, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE, rounded=False)

    # ========== SLIDE TYPES ==========

    def cover(self):
        slide = self._slide(); self._bg(slide)
        # Left accent bar
        self._rect(slide, 0, 0, 0.15, 7.50, C['accent'])
        # Top header band
        self._rect(slide, 0.15, 0, 13.18, 0.65, C['dark_head'])
        self._block(slide, 0.50, 0.12, 5.0, 0.42, '株式会社エヌイチ', 14,
                    tc=C['white'], rounded=False)
        # White content area
        self._rect(slide, 0.30, 0.80, 12.73, 6.50, C['white'], rounded=True)
        # Main title
        self._block(slide, 1.5, 2.10, 10.33, 1.10, '生成AI支援ご提案書', 48,
                    tc=C['black'], bold=True, align=PP_ALIGN.CENTER, rounded=False)
        # Subtitle
        self._block(slide, 1.5, 3.40, 10.33, 0.55, f'{self.client} 様向け', 24,
                    tc=C['sub'], align=PP_ALIGN.CENTER, rounded=False)
        # Accent line
        self._rect(slide, 5.17, 4.15, 3.0, 0.035, C['accent'])
        # Date
        self._block(slide, 0.67, 6.50, 2.5, 0.30, self.date, 13,
                    tc=C['sub'], rounded=False)
        # Company
        self._block(slide, 9.00, 6.45, 3.80, 0.35, '株式会社エヌイチ', 15,
                    tc=C['body'], bold=True, align=PP_ALIGN.RIGHT, rounded=False)

    def company_overview(self):
        slide = self._slide(); self._bg(slide)
        y = self._title_bar(slide, '会社概要')
        # Left: info in colored block
        self._rect(slide, M, y, 6.50, 5.70, C['card'], rounded=True)
        info = [
            ('会社名', '株式会社エヌイチ（n1 inc.）'),
            ('設立', '2023年11月'), ('代表', '奥山 幸生'),
            ('所在地', '東京都新宿区'), ('事業内容', 'AI導入支援・研修・SaaS開発'),
            ('従業員数', '約150名（業務委託・インターン生含む）'),
        ]
        ty = y + 0.20
        for label, val in info:
            # Label in accent colored block
            self._block(slide, 0.70, ty, 1.60, 0.32, label, 11,
                        bg=C['light_acc'], tc=C['accent'], bold=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ml=0.05)
            self._block(slide, 2.40, ty, 4.40, 0.32, val, 13,
                        tc=C['body'], rounded=False, ml=0.10)
            ty += 0.45
        # Mission block
        ty += 0.15
        self._block(slide, 0.70, ty, 6.10, 0.80,
                    'ミッション\n「AI（アイ）ある社会に」', 14,
                    bg=C['highlight'], tc=C['body'], bold=True,
                    bar_color=C['accent'], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        # Right: service cards with Kawaru brand colors（展示会フライヤー準拠）
        svcs = [
            ('Kawaru', '業務自動化SaaS', C['kawaru_blue']),
            ('Kawaru Team', 'フルカスタマイズAI研修', C['kawaru_orange']),
            ('Kawaru BPO', 'AI業務効率化代行', C['kawaru_green']),
            ('Kawaru Coach', 'AI顧問サービス', C['kawaru_purple']),
        ]
        sx, sw = 7.30, 5.53
        sy = y
        for name, desc, bar_c in svcs:
            self._rect(slide, sx, sy, sw, 1.20, C['card'], rounded=True)
            self._rect(slide, sx, sy, 0.06, 1.20, bar_c)
            self._block(slide, sx + 0.20, sy + 0.15, sw - 0.40, 0.35, name, 15,
                        tc=C['body'], bold=True, rounded=False, ml=0.08)
            self._block(slide, sx + 0.20, sy + 0.55, sw - 0.40, 0.40, desc, 12,
                        tc=C['sub'], rounded=False, ml=0.08)
            sy += 1.40

    def case_studies(self):
        slide = self._slide(); self._bg(slide)
        y = self._title_bar(slide, '導入実績', 'あらゆる規模の企業様に導入されています')
        # Badges
        companies = [
            '株式会社ブロードリンク', '株式会社中西製作所', '株式会社ミズカラ', 'エル・ティー・エス リンク',
            'FRUOR株式会社', '株式会社BIRDY', '株式会社オンシナジー', 'キリングループ'
        ]
        bw, bh, gap = 2.85, 0.33, 0.23
        for i, name in enumerate(companies):
            bx = M + (i % 4) * (bw + gap)
            by = y + 0.05 + (i // 4) * (bh + 0.10)
            self._block(slide, bx, by, bw, bh, name, 10,
                        bg=C['card'], tc=C['sub'], align=PP_ALIGN.CENTER,
                        anchor=MSO_ANCHOR.MIDDLE)
        # Case cards — top at y+0.90, height 4.60 → bottom ≤ y+5.50 ≤ 7.10 (safe)
        cases = [
            ('キリングループ様', 'AI未経験、活用イメージなし',
             '全員が「業務でAIを使いたい」と回答\n満足度 8.5点'),
            ('株式会社ミズカラ様', '問い合わせ対応に月283時間',
             '月85時間へ削減\n約200時間＝25営業日分の削減'),
            ('株式会社BIRDY様', '業務の属人化、営業効率の壁',
             'バックオフィス効率化\n経営資源を成長領域にシフト'),
        ]
        cw, cg = 3.85, 0.39
        ct = y + 0.90
        ch = 4.60
        for i, (co, issue, result) in enumerate(cases):
            cx = M + i * (cw + cg)
            # Card bg
            self._rect(slide, cx, ct, cw, ch, C['card'], rounded=True)
            # Badge + name
            self._badge(slide, cx + 0.35, ct + 0.30, 0.20, str(i + 1))
            self._block(slide, cx + 0.65, ct + 0.12, cw - 0.85, 0.35, co, 14,
                        tc=C['body'], bold=True, rounded=False)
            # Issue block (red-tinted)
            iy = ct + 0.60
            self._block(slide, cx + 0.12, iy, cw - 0.24, 0.85,
                        f'課題\n{issue}', 11,
                        bg=C['issue_bg'], tc=C['body'], bar_color=C['issue_bar'])
            # Arrow
            self._block(slide, cx + cw / 2 - 0.20, iy + 0.90, 0.40, 0.28,
                        '▼', 14, tc=C['accent'], align=PP_ALIGN.CENTER,
                        anchor=MSO_ANCHOR.MIDDLE, rounded=False)
            # Result block (green-tinted)
            ry = iy + 1.25
            self._block(slide, cx + 0.12, ry, cw - 0.24, 1.60,
                        f'成果\n{result}', 12,
                        bg=C['ok_bg'], tc=C['body'], bold=True, bar_color=C['ok_bar'])

    def intent_slide(self, main_msg, sub_msg, bullets, accent_msg):
        slide = self._slide(); self._bg(slide)
        # Left accent bar
        self._rect(slide, 0, 0, 0.15, 7.50, C['accent'])
        self._rect(slide, 0.15, 0, 13.18, 0.65, C['dark_head'])
        self._block(slide, 10.00, 0.12, 3.10, 0.42, '株式会社エヌイチ', 13,
                    tc=C['white'], align=PP_ALIGN.RIGHT, rounded=False)
        # White content area
        self._rect(slide, 0.30, 0.80, 12.73, 6.50, C['white'], rounded=True)
        # Main message in blue info block
        self._block(slide, 0.65, 1.10, 12.00, 1.20, main_msg, 24,
                    bg=C['info_bg'], tc=C['title'], bold=True, bar_color=C['accent'])
        # Sub message
        self._block(slide, 0.65, 2.45, 12.00, 0.50, sub_msg, 14,
                    tc=C['sub'], rounded=False, ml=0.20)
        # Section header
        self._section_head(slide, 0.65, 3.15, 6.0, 'ご提案のポイント')
        # Bullet blocks
        by = 3.55
        for b in bullets:
            self._block(slide, 0.80, by, 11.60, 0.42, f'●  {b}', 13,
                        bg=C['card'], tc=C['body'], bar_color=C['accent'])
            by += 0.52
        # Bottom accent block
        self._block(slide, 0.65, 5.80, 12.00, 0.60, accent_msg, 15,
                    bg=C['highlight'], tc=C['accent'], bold=True,
                    bar_color=C['accent'], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        # Note
        self._block(slide, 0.65, 6.50, 12.00, 0.30,
                    '※ 本ご提案は、初回ヒアリングの内容に基づいて設計しています。',
                    9, tc=C['label'], rounded=False, ml=0.10)

    def status_issues(self, review_items, conclusion, issues,
                      right_title='特定された課題',
                      right_bar_color=None, right_item_bg=None, right_item_bar=None):
        """Two-panel slide: review on left, issues on right.
        right_bar_color defaults to issue_bar (red). Pass C['accent'] for neutral contexts."""
        rbc  = right_bar_color  or C['issue_bar']
        ribg = right_item_bg    or C['issue_bg']
        ribar= right_item_bar   or C['issue_bar']
        slide = self._slide(); self._bg(slide)
        y = self._title_bar(slide, '貴社の現状と課題', '初回ヒアリングを踏まえて')
        panel_h = min(5.60, 7.35 - y)
        lw = 5.80
        # === Left panel ===
        self._rect(slide, M, y, lw, panel_h, C['white'], rounded=True)
        self._rect(slide, M, y, 0.06, panel_h, C['accent'])
        self._section_head(slide, M + 0.20, y + 0.15, lw - 0.40, '前回のお打合せ 振り返り', 14)
        iy = y + 0.60
        for label, val in review_items:
            self._block(slide, M + 0.20, iy, 0.90, 0.30, label, 10,
                        bg=C['light_acc'], tc=C['accent'], bold=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
            self._block(slide, M + 1.20, iy, lw - 1.55, 0.30, val, 12,
                        bg=C['card'], tc=C['body'], rounded=False, ml=0.10)
            iy += 0.42
        iy += 0.15
        self._block(slide, M + 0.20, iy, lw - 0.40, 0.90,
                    f'結論\n{conclusion}', 13,
                    bg=C['ok_bg'], tc=C['body'], bold=True, bar_color=C['ok_bar'])
        # === Right panel ===
        rx = M + lw + 0.25
        rw = CW - lw - 0.25
        self._rect(slide, rx, y, rw, panel_h, C['white'], rounded=True)
        self._rect(slide, rx, y, 0.06, panel_h, rbc)
        self._section_head(slide, rx + 0.20, y + 0.15, rw - 0.40, right_title, 14)
        iy = y + 0.60
        for i, (t, d) in enumerate(issues):
            self._badge(slide, rx + 0.45, iy + 0.35, 0.20, str(i + 1))
            self._block(slide, rx + 0.80, iy, rw - 1.10, 0.28, t, 13,
                        tc=C['body'], bold=True, rounded=False)
            self._block(slide, rx + 0.80, iy + 0.32, rw - 1.10, 0.55, d, 11,
                        bg=ribg, tc=C['sub'], bar_color=ribar)
            iy += 1.10

    def proposal_overview(self, items):
        """items = [(label, title, [bullets]), ...]"""
        slide = self._slide(); self._bg(slide)
        y = self._title_bar(slide, 'ご提案の全体像',
                           'Kawaru Team のフルカスタマイズAI研修で、AIエージェント活用力を高めます')
        n = len(items)
        cw = 3.50
        total = n * cw + (n - 1) * 0.70
        start_x = M + (CW - total) / 2
        for i, (label, title, bullets) in enumerate(items):
            cx = start_x + i * (cw + 0.70)
            is_last = (i == n - 1)
            # Label tag
            lbg = C['ok_bg'] if is_last else C['accent']
            lc = C['ok_text'] if is_last else C['white']
            bar = C['ok_bar'] if is_last else None
            self._block(slide, cx, y + 0.05, cw, 0.35, label, 12,
                        bg=lbg, tc=lc, bold=True, align=PP_ALIGN.CENTER,
                        anchor=MSO_ANCHOR.MIDDLE, bar_color=bar)
            # Card
            ct = y + 0.50
            ch = 5.00
            self._rect(slide, cx, ct, cw, ch, C['white'], rounded=True)
            self._rect(slide, cx, ct, 0.06, ch, C['accent'] if not is_last else C['ok_bar'])
            # Title inside blue/green block
            tbg = C['info_bg'] if not is_last else C['ok_bg']
            self._block(slide, cx + 0.15, ct + 0.15, cw - 0.30, 0.70, title, 15,
                        bg=tbg, tc=C['body'], bold=True, anchor=MSO_ANCHOR.MIDDLE)
            # Bullets
            by = ct + 1.00
            for b in bullets:
                self._block(slide, cx + 0.15, by, cw - 0.30, 0.35, f'●  {b}', 11,
                            bg=C['card'], tc=C['sub'])
                by += 0.42
            # Arrow
            if not is_last:
                ax = cx + cw + 0.15
                self._arrow_h(slide, ax, y + 2.80)

    def detail_3col(self, title, subtitle, columns):
        """3-column detail slide"""
        slide = self._slide(); self._bg(slide)
        y = self._title_bar(slide, title, subtitle)
        cw = 3.50; gap = 0.67
        total = 3 * cw + 2 * gap
        sx = M + (CW - total) / 2
        for i, (label, col_title, bullets) in enumerate(columns):
            cx = sx + i * (cw + gap)
            # Label
            self._block(slide, cx, y + 0.05, cw, 0.30, label, 11,
                        bg=C['light_acc'], tc=C['accent'], bold=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            ct = y + 0.45
            ch = 5.30
            self._rect(slide, cx, ct, cw, ch, C['white'], rounded=True)
            self._rect(slide, cx, ct, 0.06, ch, C['accent'])
            # Title block
            self._block(slide, cx + 0.15, ct + 0.12, cw - 0.30, 0.65, col_title, 14,
                        bg=C['info_bg'], tc=C['body'], bold=True,
                        anchor=MSO_ANCHOR.MIDDLE, bar_color=C['info_bar'])
            by = ct + 0.90
            for b in bullets:
                self._block(slide, cx + 0.15, by, cw - 0.30, 0.32, f'●  {b}', 11,
                            bg=C['card'], tc=C['sub'])
                by += 0.38

    def steps_3(self, title, subtitle, steps):
        """3-step flow slide"""
        slide = self._slide(); self._bg(slide)
        y = self._title_bar(slide, title, subtitle)
        cw = 3.50
        total = 3 * cw + 2 * 0.70
        sx = M + (CW - total) / 2
        for i, (step_label, step_title, bullets) in enumerate(steps):
            cx = sx + i * (cw + 0.70)
            # Step label
            self._block(slide, cx + cw / 2 - 0.60, y + 0.05, 1.20, 0.32, step_label, 11,
                        bg=C['accent'], tc=C['white'], bold=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            ct = y + 0.48
            ch = 5.20
            self._rect(slide, cx, ct, cw, ch, C['white'], rounded=True)
            self._rect(slide, cx, ct, 0.06, ch, C['accent'])
            # Title block
            self._block(slide, cx + 0.15, ct + 0.12, cw - 0.30, 0.65, step_title, 14,
                        bg=C['highlight'], tc=C['body'], bold=True,
                        anchor=MSO_ANCHOR.MIDDLE, bar_color=C['accent'])
            by = ct + 0.90
            for b in bullets:
                self._block(slide, cx + 0.15, by, cw - 0.30, 0.32, f'●  {b}', 11,
                            bg=C['card'], tc=C['sub'])
                by += 0.38
            if i < len(steps) - 1:
                ax = cx + cw + 0.15
                self._arrow_h(slide, ax, y + 2.80)

    def curriculum(self, sessions):
        """sessions = [(num_str, title, description), ...]"""
        slide = self._slide(); self._bg(slide)
        y = self._title_bar(slide, '研修カリキュラム案（全5回）')
        rh = 0.88
        rg = 0.10
        for i, (num, title, desc) in enumerate(sessions):
            ry = y + 0.10 + i * (rh + rg)
            # Row block with accent bar
            self._rect(slide, M, ry, CW, rh, C['white'], rounded=True)
            self._rect(slide, M, ry, 0.06, rh, C['accent'])
            # Badge
            self._badge(slide, M + 0.50, ry + rh / 2, 0.24, num)
            # Title (in blue block)
            self._block(slide, M + 0.95, ry + 0.10, 3.80, 0.35, title, 14,
                        bg=C['info_bg'], tc=C['body'], bold=True,
                        anchor=MSO_ANCHOR.MIDDLE, bar_color=C['info_bar'])
            # Desc
            self._block(slide, M + 0.95, ry + 0.50, CW - 1.50, 0.32, desc, 11,
                        tc=C['sub'], rounded=False, ml=0.10)
        # Note
        ny = y + 0.10 + len(sessions) * (rh + rg) + 0.08
        self._block(slide, M, ny, CW, 0.28,
                    '※ カリキュラムは貴社の課題・ゴールに応じてフルカスタマイズいたします',
                    10, bg=C['warn_bg'], tc=C['warn_text'], bar_color=C['warn_bar'])

    def pricing(self, price_main, price_detail, subsidy_price, subsidy_detail, schedule):
        slide = self._slide(); self._bg(slide)
        y = self._title_bar(slide, '料金・スケジュール')
        lw = 5.80
        # === Price card ===
        self._rect(slide, M, y, lw, 2.20, C['white'], rounded=True)
        self._rect(slide, M, y, 0.06, 2.20, C['accent'])
        self._block(slide, M + 0.20, y + 0.12, 1.00, 0.28, '料金', 10,
                    bg=C['light_acc'], tc=C['accent'], bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
        self._block(slide, M + 0.20, y + 0.50, lw - 0.40, 0.55, price_main, 26,
                    tc=C['body'], bold=True, rounded=False, ml=0.10)
        self._block(slide, M + 0.20, y + 1.15, lw - 0.40, 0.80, price_detail, 12,
                    tc=C['sub'], rounded=False, ml=0.10)
        # === Subsidy card ===
        sy = y + 2.35
        self._rect(slide, M, sy, lw, 2.85, C['ok_bg'], rounded=True)
        self._rect(slide, M, sy, 0.06, 2.85, C['ok_bar'])
        self._block(slide, M + 0.20, sy + 0.12, 1.40, 0.28, '助成金活用', 10,
                    bg=C['ok_bar'], tc=C['white'], bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
        self._block(slide, M + 0.20, sy + 0.50, lw - 0.40, 0.50, subsidy_price, 20,
                    tc=C['ok_text'], bold=True, rounded=False, ml=0.10)
        self._block(slide, M + 0.20, sy + 1.10, lw - 0.40, 1.50, subsidy_detail, 12,
                    tc=C['sub'], rounded=False, ml=0.10)
        # === Schedule ===
        rx = M + lw + 0.25
        rw = CW - lw - 0.25
        self._rect(slide, rx, y, rw, 5.20, C['white'], rounded=True)
        self._rect(slide, rx, y, 0.06, 5.20, C['info_bar'])
        self._section_head(slide, rx + 0.18, y + 0.15, rw - 0.40, 'スケジュール（目安）')
        sy2 = y + 0.60
        for j, (label, val) in enumerate(schedule):
            self._block(slide, rx + 0.18, sy2, rw - 0.35, 0.30, label, 13,
                        bg=C['info_bg'], tc=C['body'], bold=True,
                        bar_color=C['info_bar'], anchor=MSO_ANCHOR.MIDDLE)
            self._block(slide, rx + 0.18, sy2 + 0.35, rw - 0.35, 0.30, val, 11,
                        tc=C['sub'], rounded=False, ml=0.20)
            sy2 += 0.80

    def next_actions(self, actions):
        """actions = [(title, desc), ...]  — auto-sizes row height to fill slide."""
        slide = self._slide(); self._bg(slide)
        y = self._title_bar(slide, 'ネクストアクション')
        n = len(actions)
        available_h = 7.30 - y - 0.10
        rg = 0.18
        rh = (available_h - rg * (n - 1)) / n
        for i, (title, desc) in enumerate(actions):
            ry = y + 0.10 + i * (rh + rg)
            self._rect(slide, M + 0.80, ry, CW - 0.80, rh, C['white'], rounded=True)
            self._rect(slide, M + 0.80, ry, 0.06, rh, C['accent'])
            self._badge(slide, M + 0.35, ry + rh / 2, 0.26, str(i + 1))
            if i < n - 1:
                self._rect(slide, M + 0.33, ry + rh, 0.04, rg, C['accent'])
            self._block(slide, M + 1.10, ry + 0.10, 3.80, 0.34, title, 14,
                        bg=C['info_bg'], tc=C['body'], bold=True,
                        anchor=MSO_ANCHOR.MIDDLE, bar_color=C['info_bar'])
            self._block(slide, M + 1.10, ry + 0.52, CW - 1.50, rh - 0.60, desc, 12,
                        tc=C['sub'], rounded=False, ml=0.10, anchor=MSO_ANCHOR.TOP)

    def closing(self):
        slide = self._slide(); self._bg(slide)
        self._rect(slide, 0, 0, 0.15, 7.50, C['accent'])
        self._rect(slide, 0.15, 0, 13.18, 0.65, C['dark_head'])
        self._block(slide, 0.50, 0.12, 5.0, 0.42, '株式会社エヌイチ', 14,
                    tc=C['white'], rounded=False)
        self._rect(slide, 0.30, 0.80, 12.73, 6.50, C['white'], rounded=True)
        # Main message in highlight block
        self._block(slide, 1.20, 1.60, 10.93, 2.20,
                    '私たちが提供するのは\n「ツール」や「研修」ではありません。\n\n'
                    '貴社の未来を創る「AI人材」と、\n組織全体の「成長」です。',
                    22, bg=C['highlight'], tc=C['title'], bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    bar_color=C['accent'])
        # Sub
        self._block(slide, 2.0, 4.10, 9.33, 0.45,
                    '貴社のAIパートナーとして、共に未来を創ることをお約束します。',
                    14, tc=C['sub'], align=PP_ALIGN.CENTER, rounded=False)
        # Company
        self._block(slide, 3.0, 4.80, 7.33, 0.50, '株式会社エヌイチ', 20,
                    tc=C['body'], bold=True, align=PP_ALIGN.CENTER, rounded=False)
        self._block(slide, 3.0, 5.30, 7.33, 0.35, 'AI（アイ）ある社会に', 14,
                    tc=C['accent'], align=PP_ALIGN.CENTER, rounded=False)
        # Contact in card
        self._block(slide, 3.50, 5.85, 6.33, 0.85,
                    '大川 龍之介（Okawa Ryunosuke）\n'
                    '法人事業部 セールスマネージャー\n'
                    'TEL: 080-2646-2420  Mail: r.okawa@n1-inc.co.jp',
                    10, bg=C['card'], tc=C['label'], align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)

    def logic_tree(self, title, root_text, branches):
        """
        Horizontal logic tree: root → branches → leaves
        branches: [(branch_label, [(leaf_text, is_fact), ...]), ...]
        """
        slide = self._slide(); self._bg(slide)
        y = self._title_bar(slide, title)
        content_h = 6.80 - y
        n_branches = len(branches)
        branch_h = content_h / n_branches
        root_x, root_w = 0.50, 2.10
        self._rect(slide, root_x, y, root_w, content_h, C['issue_bg'], rounded=True)
        self._rect(slide, root_x, y, 0.06, content_h, C['issue_bar'])
        self._block(slide, root_x + 0.12, y + 0.10, root_w - 0.22, content_h - 0.20,
                    root_text, 13, tc=C['body'], bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, rounded=False)
        mid_x = root_x + root_w + 0.38
        mid_w = 2.65
        leaf_x = mid_x + mid_w + 0.38
        leaf_w = 13.33 - leaf_x - 0.33
        for i, (branch_label, leaves) in enumerate(branches):
            by = y + i * branch_h + 0.10
            bh = branch_h - 0.20
            conn_y = by + bh / 2 - 0.02
            self._rect(slide, root_x + root_w, conn_y, 0.38, 0.04, C['label'])
            self._rect(slide, mid_x, by, mid_w, bh, C['info_bg'], rounded=True)
            self._rect(slide, mid_x, by, 0.06, bh, C['info_bar'])
            self._block(slide, mid_x + 0.12, by + 0.08, mid_w - 0.22, bh - 0.16,
                        branch_label, 12, tc=C['body'], bold=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, rounded=False)
            n_leaves = len(leaves)
            leaf_item_h = bh / n_leaves
            for j, (leaf_text, is_fact) in enumerate(leaves):
                ly = by + j * leaf_item_h + 0.04
                lh = leaf_item_h - 0.08
                conn_ly = ly + lh / 2 - 0.02
                self._rect(slide, mid_x + mid_w, conn_ly, 0.38, 0.04, C['label'])
                label_text = '事実' if is_fact else '推察'
                label_bg = C['light_acc'] if is_fact else C['warn_bg']
                label_tc = C['accent'] if is_fact else C['warn_bar']
                bar_c = C['accent'] if is_fact else C['warn_bar']
                self._rect(slide, leaf_x, ly, leaf_w, lh, C['card'], rounded=True)
                self._rect(slide, leaf_x, ly, 0.06, lh, bar_c)
                self._block(slide, leaf_x + 0.12, ly + 0.04, 0.58, lh - 0.08,
                            label_text, 9, bg=label_bg, tc=label_tc, bold=True,
                            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ml=0.03)
                self._block(slide, leaf_x + 0.78, ly, leaf_w - 0.88, lh,
                            leaf_text, 11, tc=C['body'], rounded=False, ml=0.05,
                            anchor=MSO_ANCHOR.MIDDLE)

    def save(self, path):
        self.prs.save(path)
        return path
