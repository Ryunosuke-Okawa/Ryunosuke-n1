"""
AI社員5人 発表スライド生成スクリプト（4月末発表用・13枚構成）
白×紺×Kawaru ブルー基調 / 4サービスカラーは一覧時のみ
ノート（台本）入り
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE

C = {
    'white':         RGBColor(0xFF, 0xFF, 0xFF),
    'navy':          RGBColor(0x1C, 0x35, 0x57),
    'navy_l':        RGBColor(0x4A, 0x6C, 0x8C),
    'kawaru':        RGBColor(0x25, 0x63, 0xEB),
    'kawaru_l':      RGBColor(0xDB, 0xEA, 0xFE),
    'team':          RGBColor(0xFF, 0x66, 0x00),
    'coach':         RGBColor(0x7C, 0x3A, 0xED),
    'bpo':           RGBColor(0x10, 0xBF, 0x16),
    'gray':          RGBColor(0x6B, 0x72, 0x80),
    'light_gray':    RGBColor(0xF3, 0xF4, 0xF6),
    'border':        RGBColor(0xE5, 0xE7, 0xEB),
}

FONT = 'Noto Sans JP'
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
TOTAL = 12


class Builder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

    def _slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _bg(self, slide):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        s.fill.solid()
        s.fill.fore_color.rgb = C['white']
        s.line.fill.background()

    def _block(self, slide, x, y, w, h, text, size, color=None, bg=None,
               bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, rounded=False):
        st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        s = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
        if bg:
            s.fill.solid()
            s.fill.fore_color.rgb = bg
        else:
            s.fill.background()
        s.line.fill.background()
        if rounded and hasattr(s, 'adjustments') and len(s.adjustments) > 0:
            s.adjustments[0] = 0.05

        tf = s.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.vertical_anchor = anchor
        tf.margin_left = Inches(0.10)
        tf.margin_right = Inches(0.10)
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.06)

        for i, line in enumerate(text.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color or C['navy']
        return s

    def _line(self, slide, x, y, w, color, thick=0.05):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(thick))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def _connector(self, slide, x1, y1, x2, y2, color, weight=2.5):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        line.line.color.rgb = color
        line.line.width = Pt(weight)
        return line

    def _arrow(self, slide, x, y, w, h, color):
        s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def _arrow_right(self, slide, x, y, w, h, color):
        s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def _circle(self, slide, x, y, w, h, color):
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def _title(self, slide, text, page_num):
        self._block(slide, 12.0, 0.3, 1.1, 0.3, f"{page_num} / {TOTAL}", 11,
                    color=C['gray'], align=PP_ALIGN.RIGHT)
        self._block(slide, 0.5, 0.4, 11.5, 0.7, text, 24, bold=True, color=C['navy'])
        self._line(slide, 0.5, 1.15, 1.8, C['kawaru'], 0.07)

    def _notes(self, slide, text):
        slide.notes_slide.notes_text_frame.text = text

    # ========== Slides ==========

    def s01_cover(self):
        s = self._slide()
        self._bg(s)
        self._line(s, 0, 0, 13.33, C['kawaru'], 0.15)
        self._circle(s, 10.5, 0.5, 2.3, 2.3, C['kawaru_l'])
        self._circle(s, 0.5, 5.5, 1.8, 1.8, C['kawaru_l'])

        self._block(s, 0.5, 2.5, 12.33, 1.6, 'AIハッカソン',
                    60, bold=True, color=C['kawaru'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._line(s, 5.5, 4.3, 2.33, C['navy'], 0.06)
        self._block(s, 0.5, 4.5, 12.33, 0.8, 'Kawaru事業部',
                    28, bold=True, color=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 0.5, 5.4, 12.33, 0.7, '大川 龍之介',
                    22, color=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 0.5, 6.8, 12.33, 0.4, '2026.04',
                    14, color=C['gray'], align=PP_ALIGN.CENTER)

        self._notes(s, """【表紙・5秒】
Kawaru事業部の大川です。AIハッカソンの発表を始めます。
タイトルは「AI社員を動かして、Kawaru事業部の事業を推進する」です。
""")

    def s02_premise(self):
        s = self._slide()
        self._bg(s)
        self._block(s, 0.5, 0.4, 12.33, 0.5, '前提', 14,
                    color=C['gray'], align=PP_ALIGN.CENTER, bold=True)
        self._line(s, 6.0, 1.0, 1.33, C['kawaru'], 0.06)

        self._block(s, 0.5, 1.6, 12.33, 1.4, 'Kawaru事業部は、',
                    28, color=C['navy'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 0.5, 2.7, 12.33, 1.2, '決まった業務が少ない',
                    48, bold=True, color=C['kawaru'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        self._block(s, 0.7, 4.6, 5.5, 2.0, '既存事業',
                    18, bold=True, color=C['gray'], bg=C['light_gray'], rounded=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        self._block(s, 0.7, 5.2, 5.5, 1.4,
                    '業務が決まっている\n改善対象が見える',
                    16, color=C['gray'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._arrow_right(s, 6.4, 5.3, 0.8, 0.7, C['kawaru'])
        self._block(s, 7.4, 4.6, 5.3, 2.0, 'Kawaru事業部',
                    18, bold=True, color=C['white'], bg=C['kawaru'], rounded=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        self._block(s, 7.4, 5.2, 5.3, 1.4,
                    '業務が毎月変わる\n改善より、業務を作る段階',
                    16, color=C['white'], bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        self._block(s, 12.0, 7.1, 1.1, 0.3, "2 / 12", 11, color=C['gray'], align=PP_ALIGN.RIGHT)

        self._notes(s, """【前提・10秒】
最初に前提を共有します。
Kawaru事業部は、決まった業務が少ない事業部です。
既存事業のように業務が固まっておらず、毎月内容が変わります。
業務改善ではなく、業務を作っていく段階。
だから今日は「業務削減」の話ではなく、「AI社員を動かして事業を推進する」という話をします。
""")

    def s03_status(self):
        s = self._slide()
        self._bg(s)
        self._title(s, 'Kawaru事業部の現状：4サービスを2人体制で動かしている', 3)

        services = [
            ('Kawaru', C['kawaru']),
            ('Kawaru Team', C['team']),
            ('Kawaru Coach', C['coach']),
            ('Kawaru BPO', C['bpo']),
        ]
        sw = 3.0
        gap = 0.05
        for i, (name, col) in enumerate(services):
            x = 0.5 + i*(sw + gap)
            self._block(s, x, 1.6, sw, 1.6, name, 24, bold=True,
                        color=C['white'], bg=col, rounded=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        self._arrow(s, 6.42, 3.4, 0.5, 0.6, C['kawaru'])

        self._block(s, 3.5, 4.3, 2.5, 1.4, '髙橋', 28, bold=True,
                    color=C['white'], bg=C['navy'], rounded=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 6.1, 4.3, 1.0, 1.4, '+', 36, bold=True,
                    color=C['gray'], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 7.2, 4.3, 2.5, 1.4, '大川', 28, bold=True,
                    color=C['white'], bg=C['kawaru'], rounded=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        self._block(s, 0.5, 6.0, 12.33, 1.0,
                    '現状、Kawaru事業部の重要な4サービスを2人体制でやっている',
                    22, bold=True, color=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    bg=C['kawaru_l'], rounded=True)

        self._notes(s, """【現状・15秒】
Kawaru事業部の現状をお見せします。
Kawaru、Kawaru Team、Kawaru Coach、Kawaru BPOの4サービスがあります。
これを髙橋と私の2人体制で動かしています。
2人で4サービスを動かすので、当然手が回りません。
だから次のような体制を作りました。
""")

    def s04_theme(self):
        """今回のテーマ：9月末までにAI社員5人を作る（シンプル版）"""
        s = self._slide()
        self._bg(s)
        self._title(s, '今回のテーマ', 4)

        # サブ
        self._block(s, 0.5, 1.7, 12.33, 0.6, '今回のハッカソンで取り組むこと', 16,
                    color=C['gray'], align=PP_ALIGN.CENTER, bold=True)

        # メイン（中央に大きく）
        self._block(s, 0.5, 2.7, 12.33, 1.2, '9月末までに、',
                    36, color=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 0.5, 3.9, 12.33, 1.8, 'AI社員 5人 を作る',
                    72, bold=True, color=C['kawaru'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # アクセント
        self._line(s, 5.5, 5.9, 2.33, C['kawaru'], 0.07)

        # フッター
        self._block(s, 0.5, 6.3, 12.33, 0.5,
                    '月35万円 × 5人 = 月175万円分の業務をAI社員にやってもらう',
                    16, color=C['navy'], align=PP_ALIGN.CENTER, bold=True)

        self._notes(s, """【テーマ・10秒】
今回のハッカソンで私が取り組むテーマはこれです。
9月末までに、AI社員5人を作る。
基準は1人あたり月35万円分の業務回避、5人で月175万円分です。
ではどう作るか、次のスライドで。
""")

    def s05_concept(self):
        """私はこう考えた + なんでこれに（統合）"""
        s = self._slide()
        self._bg(s)
        self._title(s, '9月末に向けて、私はこう考えた', 5)

        # 上メッセージ
        self._block(s, 0.5, 1.7, 12.33, 0.6,
                    '事業に必要な機能を、5人のAI社員に分担する',
                    18, color=C['navy'], align=PP_ALIGN.CENTER, bold=True)

        # 5機能 + 役割ラベル
        funcs = [
            ('マーケ', '認知獲得'),
            ('営業', '商談・受注'),
            ('デリバリー', '納品'),
            ('事務・法務', '管理・記録'),
            ('戦略', 'リサーチ・分析'),
        ]
        cw = 2.4
        gap = 0.07
        for i, (name, role) in enumerate(funcs):
            x = 0.5 + i*(cw + gap)
            self._block(s, x, 2.6, cw, 1.5, name, 24, bold=True,
                        color=C['white'], bg=C['kawaru'], rounded=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._block(s, x, 4.15, cw, 0.8, role, 14,
                        color=C['navy'], bg=C['light_gray'], rounded=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 下メッセージ（強調）
        self._block(s, 0.5, 5.4, 12.33, 1.5,
                    '事業の全部の仕事を 5人で網羅している',
                    26, bold=True, color=C['kawaru'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    bg=C['kawaru_l'], rounded=True)

        self._notes(s, """【私はこう考えた + なんでこれに・25秒】
9月末までにAI社員5人を作るために、私はこう考えました。
事業に必要な機能を5つに分けて、5人のAI社員に分担してもらう。
マーケは認知獲得、営業は商談・受注、デリバリーは納品、事務法務は管理・記録、戦略はリサーチ・分析。
事業の全部の仕事を、この5人で網羅しています。
だから案件が変わっても、この5人で受け止められる、というのが私の考えです。
""")

    def s06_directing(self):
        s = self._slide()
        self._bg(s)
        self._title(s, '実は、もう動かしている', 6)

        self._block(s, 5.67, 1.4, 2.0, 0.7, '髙橋', 18, bold=True,
                    color=C['white'], bg=C['navy'], rounded=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._connector(s, 6.67, 2.1, 6.67, 2.5, C['gray'], 2.0)

        ok_x = 5.17; ok_y = 2.5; ok_w = 3.0; ok_h = 1.0
        self._block(s, ok_x, ok_y, ok_w, ok_h, '大川', 28, bold=True,
                    color=C['white'], bg=C['kawaru'], rounded=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        ok_cx = ok_x + ok_w/2
        ok_by = ok_y + ok_h

        labels = ['マーケAI', '営業AI', 'デリバリーAI', '事務法務AI', '戦略AI']
        ai_y = 4.6
        ai_w = 2.3
        ai_h = 1.5
        for i, l in enumerate(labels):
            x = 0.5 + i*2.55
            cx = x + ai_w/2
            self._connector(s, ok_cx, ok_by, cx, ai_y, C['kawaru'], 2.5)
            self._block(s, x, ai_y, ai_w, ai_h, l, 16, bold=True,
                        color=C['white'], bg=C['kawaru'], rounded=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        self._block(s, 0.5, 6.4, 12.33, 0.8,
                    '大川主導で AI社員5人 を動かして事業部を回している',
                    20, bold=True, color=C['kawaru'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        self._notes(s, """【実はもう動かしている・15秒】
ここがポイントなんですが、これは構想ではなく、もう動いています。
私はすでにAI社員5人を動かしています。
マーケ、営業、デリバリー、事務法務、戦略の5人。
私が指示を出して、AI社員に動いてもらう。
これによって、2人だけのKawaru事業部でも4サービスを回せています。
""")

    def s07_5agents(self):
        s = self._slide()
        self._bg(s)
        self._title(s, 'AI社員5人の構成', 7)

        roles = [
            ('マーケAI', 'LP / SNS / 展示会 / LINE配信 / 事例記事'),
            ('営業AI', '商談前準備 / 要件ヒアリング / 提案書 / フォロー'),
            ('デリバリーAI', '研修運営 / Dify・n8n・GAS構築 / CSコンサル / PM'),
            ('事務・法務AI', '契約書 / 請求書 / 利用契約書 / プライバシーポリシー / PL'),
            ('戦略・プロダクトAI', '競合調査 / 市場調査 / 料金設計 / データ分析'),
        ]
        y = 1.5
        h = 0.95
        for name, desc in roles:
            self._block(s, 0.5, y, 3.5, h, name, 18, bold=True,
                        color=C['white'], bg=C['kawaru'], rounded=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._block(s, 4.1, y, 8.7, h, desc, 14, color=C['navy'],
                        bg=C['light_gray'], rounded=True,
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            y += 1.05
        self._block(s, 0.5, 6.85, 12.33, 0.4,
                    '4人 = 各領域の戦略 + 実行  /  戦略・プロダクトAI = リサーチ・分析担当',
                    11, color=C['gray'], align=PP_ALIGN.CENTER, bold=True)

        self._notes(s, """【5人の構成・25秒】
5人の構成を詳しくお見せします。
マーケAIは、LP・SNS・展示会・LINE配信・事例記事を担当。
営業AIは、商談前準備・要件ヒアリング・提案書・フォロー。
デリバリーAIは、研修運営・Dify、n8n、GAS構築・CSコンサル・PM。
事務法務AIは、契約書・請求書・利用契約書・プライバシーポリシー・PL。
戦略プロダクトAIは、競合調査・市場調査・料金設計・データ分析。
各領域の戦略と実行を5人で持っています。戦略AIだけは横断的にリサーチと分析を担当します。
""")

    def s08_why(self):
        """なんでこれにしたのか"""
        s = self._slide()
        self._bg(s)
        self._title(s, 'なんでこれにしたのか', 8)

        funcs = [
            ('マーケ', '認知獲得'),
            ('営業', '商談・受注'),
            ('デリバリー', '納品'),
            ('事務・法務', '管理・記録'),
            ('戦略', 'リサーチ・分析'),
        ]
        cw = 2.4
        cx = 0.5
        gap = 0.07
        for i, (name, desc) in enumerate(funcs):
            x = cx + i*(cw + gap)
            self._block(s, x, 2.0, cw, 1.6, name, 22, bold=True,
                        color=C['white'], bg=C['kawaru'], rounded=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._block(s, x, 3.65, cw, 0.8, desc, 14,
                        color=C['navy'], bg=C['light_gray'], rounded=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        self._block(s, 0.5, 5.0, 12.33, 1.6,
                    '事業の全部の仕事を 5人で網羅している',
                    26, bold=True, color=C['kawaru'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    bg=C['kawaru_l'], rounded=True)

        self._notes(s, """【なんでこれにしたのか・15秒】
なんでこの5人にしたのか。
マーケ、営業、デリバリー、事務法務、戦略。これは事業に必要な機能をMECEで切ったものです。
認知獲得から商談・受注、納品、管理・記録、リサーチ・分析まで。
事業の全部の仕事を、この5人で網羅しています。
だから案件が変わっても、5人で受け止められます。
""")

    def s09_how_use(self):
        s = self._slide()
        self._bg(s)
        self._title(s, '今こう使ってます', 7)

        items = [
            ('マーケAI', [
                '展示会クリエイティブ',
                'LINE配信シナリオ',
                'SNS投稿生成',
                'LP構成案・コピー',
                'YouTube → SNS転記',
            ]),
            ('営業AI', [
                '商談前ブリーフィング',
                '提案書骨子',
                'フォローメール',
                '要件ヒアリング整理',
                '見積書作成',
            ]),
            ('デリバリーAI', [
                'NTT 研修200枚',
                '中西 Dify構築 (5h)',
                'キリン 研修PM',
                'ブロードリンク 省エネ',
                'カリキュラム設計',
            ]),
            ('事務・法務AI', [
                '契約書チェック',
                'PL自動更新',
                '受信トレイ整理',
                '月次レポート',
                'ツールコスト最適化',
            ]),
            ('戦略AI', [
                '競合調査',
                '料金設計試算',
                'データ分析',
                'プロダクト改善提案',
                '朝の工程表',
            ]),
        ]
        cw = 2.4
        gap = 0.1
        ch_y = 1.5
        ch_h = 4.7
        for i, (name, examples) in enumerate(items):
            x = 0.5 + i * (cw + gap)
            self._block(s, x, ch_y, cw, ch_h, '', 11,
                        bg=C['light_gray'], rounded=True)
            self._block(s, x+0.1, ch_y+0.1, cw-0.2, 0.7, name, 14, bold=True,
                        color=C['white'], bg=C['kawaru'], rounded=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            ex_text = '\n'.join([f'✓ {e}' for e in examples])
            self._block(s, x+0.15, ch_y+0.95, cw-0.3, ch_h-1.05, ex_text, 11,
                        color=C['navy'],
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

        self._block(s, 0.5, 6.4, 12.33, 0.8,
                    'いろんな領域で AI社員 を使い倒している',
                    20, bold=True, color=C['kawaru'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        self._notes(s, """【実際の使い方・25秒】
実際に今こう使っています。
マーケAIは、展示会クリエイティブやLP制作。
営業AIは、商談前準備や提案書の骨子作成。
デリバリーAIは、NTT研修200枚を1人2週間で、中西のDify構築を5時間で完納しました。
事務法務AIは、契約書チェックやPL自動更新。
戦略AIは、競合調査や料金設計、データ分析。
細かく読まなくて大丈夫ですが、いろんな領域でAI社員を使い倒しているということが伝わればOKです。
""")

    def s10_to_results(self):
        s = self._slide()
        self._bg(s)
        self._block(s, 0.5, 0.4, 12.33, 0.5, 'そして', 14,
                    color=C['gray'], align=PP_ALIGN.CENTER, bold=True)
        self._line(s, 6.0, 1.0, 1.33, C['kawaru'], 0.06)

        self._block(s, 0.5, 2.2, 12.33, 1.5, 'これが',
                    32, color=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 0.5, 3.5, 12.33, 1.6, '売上に直結している',
                    52, bold=True, color=C['kawaru'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._arrow(s, 6.42, 5.5, 0.5, 0.6, C['kawaru'])
        self._block(s, 12.0, 7.1, 1.1, 0.3, "8 / 12", 11, color=C['gray'], align=PP_ALIGN.RIGHT)

        self._notes(s, """【橋渡し・8秒】
そして、これが売上に直結しています。
具体的な数字をお見せします。
""")

    def s11_results(self):
        s = self._slide()
        self._bg(s)
        self._title(s, '半期 ¥6,710,000、4案件すべてにAI社員が動いている', 9)

        cases = [
            ('NTTネクシア', '¥5,000,000'),
            ('中西製作所', '¥1,000,000'),
            ('中部キリンビバレッジ', '¥450,000'),
            ('ブロードリンク', '¥260,000'),
        ]
        y = 1.7
        h = 0.7
        self._block(s, 0.5, y, 5.0, h, '案件', 16, bold=True,
                    color=C['white'], bg=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 5.55, y, 3.0, h, '売上 (税抜)', 16, bold=True,
                    color=C['white'], bg=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.78
        for case, amount in cases:
            self._block(s, 0.5, y, 5.0, h, case, 16, color=C['navy'],
                        bg=C['light_gray'],
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            self._block(s, 5.55, y, 3.0, h, amount, 16, color=C['navy'],
                        bg=C['light_gray'],
                        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            y += 0.78
        self._block(s, 0.5, y, 5.0, h, '半期合計', 18, bold=True,
                    color=C['white'], bg=C['kawaru'],
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 5.55, y, 3.0, h, '¥6,710,000', 18, bold=True,
                    color=C['white'], bg=C['kawaru'],
                    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

        self._block(s, 9.0, 1.7, 3.83, 4.6,
                    '月換算\n\n¥464,000\n\n=\n\nAI社員\n1.3人分\n稼働中',
                    22, bold=True, color=C['white'], bg=C['kawaru'], rounded=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        self._notes(s, """【売上実績・20秒】
10月から3月の半期で、売上6,710,000円。
NTTネクシア、中西製作所、キリン、ブロードリンク、4案件すべてにAI社員が関与しています。
NTT研修スライド200枚は2週間で完納、中西のDify構築は5時間で設計しました。
月換算では464,000円、AI社員1.3人分がすでに稼働している計算になります。
これは机上の空論ではなく、実際にお客さんに納品している事実です。
""")

    def s12_now_future(self):
        s = self._slide()
        self._bg(s)
        self._line(s, 5.5, 1.5, 2.33, C['kawaru'], 0.08)
        self._block(s, 0.5, 2.0, 12.33, 1.2, 'ここまでが現状の話',
                    34, bold=True, color=C['gray'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._arrow(s, 6.17, 3.4, 1.0, 0.8, C['kawaru'])
        self._block(s, 0.5, 4.4, 12.33, 1.5, 'ここから今後の話',
                    50, bold=True, color=C['kawaru'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 0.5, 6.0, 12.33, 0.6,
                    '9月末までにAI社員5人を完成させる、その具体ロードマップ',
                    18, color=C['navy'], bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 12.0, 7.1, 1.1, 0.3, "10 / 12", 11,
                    color=C['gray'], align=PP_ALIGN.RIGHT)

        self._notes(s, """【現状→今後・7秒】
ここまでが現状の話です。
ここから9月末までの今後の話、具体的なロードマップを共有します。
""")

    def s13_roadmap(self):
        """ロードマップ（具体タスク + 理由）"""
        s = self._slide()
        self._bg(s)
        self._title(s, '9月末までのロードマップ', 11)

        rows = [
            ('4月\n(現在)', '1.3人分',
             '展示会・ローンチに向けた制作物作成（クリエイティブ・LP・LINE）',
             '5月の本番に向けた助走期間', False),
            ('5月', '1.6人分',
             'β版ローンチ実施(5/2) ／ 6月8EXPOの戦略・戦術設計\n（集客リクエスト配信／当日オペ／営業トークスクリプト）',
             'マーケ・営業の比重が大きい時期。AI社員を育てながら走らせる', False),
            ('6月末\n(中間)', '2人分',
             '8EXPO出展(6/3-4) ／ ローンチ後の運用フィードバック反映 ／ 5人連携テスト',
             'デリバリー内容が変化する時期。顧客の声を聞きながらブラッシュアップ', True),
            ('7-8月', '3〜4人分',
             'β版経由の新規案件対応 ／ プロダクト改善ループ確立 ／ 事務・法務AI追加',
             '顧客が増えて事務処理が増える時期。事務AIを稼働させて対応', False),
            ('9月末\n(最終)', '5人分',
             '5人フル稼働 ／ プロダクト × デリバリー両輪完成',
             'Agent Teamsで5人協働を実装し、事業を動かす完成形へ', True),
        ]
        y = 1.5
        h = 0.85
        # ヘッダー
        self._block(s, 0.5, y, 1.5, h, '月', 13, bold=True,
                    color=C['white'], bg=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 2.05, y, 1.4, h, '達成水準', 13, bold=True,
                    color=C['white'], bg=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 3.50, y, 5.5, h, '具体タスク', 13, bold=True,
                    color=C['white'], bg=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 9.05, y, 3.78, h, 'この時期に重点を置く理由', 13, bold=True,
                    color=C['white'], bg=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.95

        for month, level, task, reason, emp in rows:
            if emp:
                bg_l = C['kawaru']; bg_r = C['kawaru_l']
                fc_l = C['white']; fc_r = C['navy']
            else:
                bg_l = C['light_gray']; bg_r = C['light_gray']
                fc_l = C['navy']; fc_r = C['navy']
            self._block(s, 0.5, y, 1.5, h, month, 11, bold=emp, color=fc_l, bg=bg_l,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._block(s, 2.05, y, 1.4, h, level, 14, bold=True, color=fc_l, bg=bg_l,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._block(s, 3.50, y, 5.5, h, task, 11, color=fc_r, bg=bg_r,
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            self._block(s, 9.05, y, 3.78, h, reason, 11, color=fc_r, bg=bg_r,
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            y += 0.85

        self._block(s, 0.5, 6.85, 12.33, 0.45,
                    '新規事業なので進行で変わることも多いが、5人完成は確実に進める',
                    13, bold=True, color=C['kawaru'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        self._notes(s, """【ロードマップ・25秒】
9月末までのロードマップです。
4月の現在は、5月のβ版ローンチと6月の8EXPO展示会に向けた制作物作成が中心です。
5月はローンチ実施と、6月の展示会に向けた戦略・戦術設計にシフト。集客リクエスト配信、当日オペ、営業トークスクリプトなど、マーケと営業の比重が大きい時期です。AI社員を育てながら走らせます。
6月末で2人分。8EXPO出展とローンチ後の運用フィードバック反映、5人の連携テスト。デリバリー内容が変化する時期で、顧客の声を聞きながらブラッシュアップします。
7-8月は3〜4人分。顧客が増えて事務処理も増えるため、事務・法務AIを稼働させて対応します。
9月末で5人分。Agent Teamsで5人協働を実装し、事業を動かす完成形を目指します。
新規事業なので進行で変わることも多いですが、5人完成というゴールは確実に進めます。
""")

    def s14_close(self):
        s = self._slide()
        self._bg(s)
        self._line(s, 0, 0, 13.33, C['kawaru'], 0.15)
        self._circle(s, 10.5, 0.8, 2.0, 2.0, C['kawaru_l'])
        self._circle(s, 0.7, 5.5, 1.6, 1.6, C['kawaru_l'])

        self._block(s, 0.5, 1.5, 12.33, 1.5, 'AI社員を動かして、',
                    40, bold=True, color=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 0.5, 2.9, 12.33, 1.6,
                    'Kawaru事業部の事業を推進する',
                    50, bold=True, color=C['kawaru'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._line(s, 5.5, 4.8, 2.33, C['kawaru'], 0.07)
        self._block(s, 0.5, 5.0, 12.33, 0.6,
                    '2人 + AI社員5人 + Agent Teams で動かす',
                    20, bold=True, color=C['navy'],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        self._block(s, 1.0, 6.0, 3.5, 0.8, 'Kawaru', 15, bold=True,
                    color=C['white'], bg=C['kawaru'], rounded=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 4.7, 6.0, 2.6, 0.8, 'Kawaru Team', 14, bold=True,
                    color=C['white'], bg=C['team'], rounded=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 7.4, 6.0, 2.6, 0.8, 'Kawaru Coach', 14, bold=True,
                    color=C['white'], bg=C['coach'], rounded=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 10.1, 6.0, 2.4, 0.8, 'Kawaru BPO', 14, bold=True,
                    color=C['white'], bg=C['bpo'], rounded=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self._block(s, 12.0, 7.1, 1.1, 0.3, '12 / 12', 11,
                    color=C['gray'], align=PP_ALIGN.RIGHT)

        self._notes(s, """【締め・15秒】
AI社員を動かして、Kawaru事業部の事業を推進する。
2人＋AI社員5人＋Agent Teamsで、Kawaru事業部を動かしていきます。
以上です。ありがとうございました。
""")

    def build(self, output):
        self.s01_cover()
        self.s02_premise()
        self.s03_status()
        self.s04_theme()
        self.s05_concept()
        self.s06_directing()
        self.s09_how_use()      # 7枚目：今こう使ってます
        self.s10_to_results()   # 8枚目：橋渡し
        self.s11_results()      # 9枚目：半期¥6,710,000
        self.s12_now_future()   # 10枚目：ここから今後
        self.s13_roadmap()      # 11枚目：ロードマップ
        self.s14_close()        # 12枚目：締め
        self.prs.save(output)
        print(f"Saved: {output}")


if __name__ == '__main__':
    output_path = '/Users/kyouyuu/claude/strategy/AI社員5人_発表_20260426.pptx'
    Builder().build(output_path)
